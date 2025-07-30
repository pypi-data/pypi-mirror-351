from doctomarkdown.base import BaseConverter, PageResult, ConversionResult
from pptx import Presentation

class PptxToMarkdown(BaseConverter):
    """Converter for PPTX files to Markdown format."""
    
    def extract_content(self):
        doc = Presentation(self.filepath)
        pages = []
        markdown_lines = []

        for page_number, slide in enumerate(doc.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            page_content = "\n".join(slide_content)
            pages.append(PageResult(page_number, page_content))
            markdown_lines.append(f"## Slide {page_number}\n\n{page_content}\n")

        self._markdown = "\n".join(markdown_lines)
        return pages