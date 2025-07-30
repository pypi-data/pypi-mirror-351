from tkinter import FALSE
import numpy as np
import pandas as pd
import sys # built-in
import os # built-in
from IPython.display import display
import shutil
import logging
from pathlib import Path
from datetime import datetime, date, time 
import re # built-in
import stat
import platform
import subprocess

from typing import Dict, List, Optional, Union, Any, Tuple, Literal,Callable
from regex import X

try:
    get_ipython().run_line_magic("load_ext", "autoreload")
    get_ipython().run_line_magic("autoreload", "2")
except NameError:
    pass

import warnings

warnings.simplefilter("ignore", category=pd.errors.SettingWithCopyWarning)
warnings.filterwarnings("ignore", category=pd.errors.PerformanceWarning)
warnings.filterwarnings("ignore")
try:
    import pkg_resources
except ImportError:
    pkg_resources = None
import glob  # built-in
import pkg_resources # built-in
import importlib
import inspect
import pkgutil
import pytest
try:
    import importlib.metadata as metadata  # Python 3.8+
except ImportError:
    import importlib_metadata as metadata  # For older versions via backport

class PkgManager:
    """
    PkgManager.uninstall("py2ls")
    PkgManager.uninstall("py2ls", mode="startswith")
    PkgManager.uninstall("py2ls", mode="endswith")
    PkgManager.uninstall("py2ls", mode="contains")
    PkgManager.uninstall("py2ls", mode="regex")

    PkgManager.timemachine()
    """

    @staticmethod
    def uninstall(
        kw: Union[str, List[str]],
        mode: str = "exact",
        dry_run: bool = False,
        make_backup: bool = True,
        make_log: bool = True,
        station: Optional[str] = None,
    ) -> None:
        if station is None:
            station = os.path.dirname(os.path.dirname(sys.executable)) 
        os.makedirs(station, exist_ok=True)
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

        if isinstance(kw, str):
            kw = [kw]
        kw = [k.lower() for k in kw] if mode != "regex" else kw
        mode = mode.lower()
        valid_modes = {"exact", "startswith", "endswith", "contains", "regex"}
        if mode not in valid_modes:
            raise ValueError(f"Mode must be one of {valid_modes}")

        installed_packages = {pkg.key: pkg.version for pkg in pkg_resources.working_set}
        matched: Set[str] = set()

        for name in installed_packages:
            for key in kw:
                if (
                    (mode == "exact" and name == key)
                    or (mode == "startswith" and name.startswith(key))
                    or (mode == "endswith" and name.endswith(key))
                    or (mode == "contains" and key in name)
                    or (mode == "regex" and re.search(key, name))
                ):
                    matched.add(name)
                    break

        if not matched:
            print("No packages matched the criteria.")
            return

        if make_backup and not dry_run:
            backup_path = os.path.join(station, f"requirements_backup_{timestamp}.txt")
            with open(backup_path, "w") as f:
                subprocess.run(["pip", "freeze"], stdout=f, check=True)
            print(f"Backup created at: '{backup_path}'")

        if dry_run:
            print("[DRY RUN] The following packages would be uninstalled:")
            for pkg in sorted(matched):
                print(f"  - {pkg}=={installed_packages[pkg]}")
            return

        print(f"[UNINSTALLING] {len(matched)} packages:")
        for pkg in sorted(matched):
            print(f"  - {pkg}=={installed_packages[pkg]}")
            subprocess.run(["pip", "uninstall", "-y", pkg], check=True)

        if make_log:
            log_path = os.path.join(station, f"log_uninstall.txt")
            with open(log_path, "w") as f:
                f.write(f"# Uninstallation log created at {timestamp}\n")
                f.write(f"# Mode: {mode}, Keywords: {kw}\n\n")
                for pkg in sorted(matched):
                    f.write(f"{pkg}=={installed_packages[pkg]}\n")
            print(f"Log written to '{log_path}'")

    @staticmethod
    def list_backups(station: Optional[str] = None) -> List[str]:
        if station is None:
            station = os.path.dirname(sys.executable)
            if os.name == "nt":
                station = os.path.dirname(station)
        return sorted(glob.glob(os.path.join(station, "requirements_backup_*.txt")))

    @staticmethod
    def list_logs(station: Optional[str] = None) -> List[str]:
        if station is None:
            station = os.path.dirname(sys.executable)
            if os.name == "nt":
                station = os.path.dirname(station)
        return sorted(glob.glob(os.path.join(station, "uninstall_*.txt")))

    @staticmethod
    def restore(
        timestamp: Optional[str] = None,
        station: Optional[str] = None,
        dry_run: bool = False,
    ) -> None:
        if station is None:
            station = os.path.dirname(sys.executable)
            if os.name == "nt":
                station = os.path.dirname(station)

        backups = PkgManager.list_backups(station)
        logs = PkgManager.list_logs(station)

        if not timestamp:
            print("Available restore points:\n\nBackups:")
            for i, backup in enumerate(backups, 1):
                ts = os.path.basename(backup)[18:-4]
                print(f"  {i}. {ts} (backup)")
            print("\nUninstall logs:")
            for i, log in enumerate(logs, len(backups) + 1):
                ts = os.path.basename(log)[10:-4]
                print(f"  {i}. {ts} (log)")
            print("\nSpecify timestamp or selection number to restore.")
            return

        try:
            selection = int(timestamp)
            all_files = backups + logs
            if 1 <= selection <= len(all_files):
                file_path = all_files[selection - 1]
                is_log = selection > len(backups)
            else:
                raise ValueError("Invalid selection number")
        except ValueError:
            backup_pattern = os.path.join(
                station, f"requirements_backup_{timestamp}.txt"
            )
            log_pattern = os.path.join(station, f"uninstall_{timestamp}.txt")
            matching_backups = glob.glob(backup_pattern)
            matching_logs = glob.glob(log_pattern)

            if matching_backups:
                file_path = matching_backups[0]
                is_log = False
            elif matching_logs:
                file_path = matching_logs[0]
                is_log = True
            else:
                print(f"No backup or log found for timestamp: {timestamp}")
                return

        with open(file_path, "r") as f:
            packages = [
                line.strip() for line in f if line.strip() and not line.startswith("#")
            ]

        if dry_run:
            print(
                f"[DRY RUN] Would restore {len(packages)} packages from:\n  {file_path}"
            )
            for pkg in packages:
                print(f"  - {pkg}")
            return

        print(f"[RESTORING] {len(packages)} packages from:\n  {file_path}")
        for pkg in packages:
            print(f"  - Installing {pkg}")
            subprocess.run(["pip", "install", pkg], check=True)

    @staticmethod
    def timemachine(station: Optional[str] = None) -> None:
        if station is None:
            station = os.path.dirname(sys.executable)
            if os.name == "nt":
                station = os.path.dirname(station)

        backups = PkgManager.list_backups(station)
        logs = PkgManager.list_logs(station)

        if not backups and not logs:
            print("No backup or log files found.")
            return

        print("\nTime Machine - Available Restore Points:")
        print("--------------------------------------")
        print("\nBackups (complete environment snapshots):")
        for i, backup in enumerate(backups, 1):
            ts = os.path.basename(backup)[18:-4]
            print(f"  {i}. {ts}")
        print("\nUninstall Logs (specific package lists):")
        for i, log in enumerate(logs, len(backups) + 1):
            ts = os.path.basename(log)[10:-4]
            print(f"  {i}. {ts}")
        print("\n0. Exit Time Machine")

        while True:
            try:
                choice = input("\nSelect a restore point (number) or '0' to exit: ")
                if choice == "0":
                    return
                selection = int(choice)
                all_files = backups + logs
                if 1 <= selection <= len(all_files):
                    file_path = all_files[selection - 1]
                    timestamp = os.path.basename(file_path).split("_")[-1][:-4]
                    PkgManager.restore(timestamp, station)
                    return
                else:
                    print("Invalid selection. Please try again.")
            except ValueError:
                print("Please enter a valid number.")  
    @staticmethod
    def listfunc(
        where: Union[str, Any]= None,
        query: Optional[str] = None,
        return_output:bool=False,
        show_all: bool = False,
        include_dunder: bool = False,
        only_defined: bool = True,
        verbose: bool = False
    ) -> Dict[str, Any]:
        """
        Recursively list functions defined in a package/module and its submodules.
        If `where=None`, returns the installed pip packages instead.

        Args:
            where (str or module or None): Module/package to inspect, or None for full pip list.
            query (str): Optional search string for fuzzy match.
            show_all (bool): Show all callables including those starting with '_'.
            include_dunder (bool): Include dunder (__) methods like __init__.
            only_defined (bool): Only show functions actually defined in the module.
            verbose (bool): Show detailed skip/load error messages.

        Returns:
            dict: Nested dictionary with module names and their function lists or pip list.
        """
        if where is None:
            # Return pip list instead
            print("📦 Installed pip packages:")
            pip_packages = {dist.metadata['Name']: dist.version for dist in metadata.distributions()}
            if query:
                func_OI = strcmp(query, list(pip_packages.keys()))[0]
                print(f"  - {func_OI}=={pip_packages[func_OI]}")
                return {func_OI: pip_packages[func_OI]}
            for name, version in sorted(pip_packages.items()):
                print(f"  - {name}=={version}")
            return pip_packages
        if isinstance(where, str):
            try:
                mod = importlib.import_module(where)
            except ModuleNotFoundError:
                print(f"Module '{where}' not found.")
                return {}
        else:
            mod = where

        root_name = mod.__name__
        results = {}

        def list_functions_in_module(module, module_name) -> List[str]:
            funcs = []
            for name in dir(module):
                attr = getattr(module, name)
                if callable(attr):
                    if not show_all:
                        if name.startswith("__") and not include_dunder:
                            continue
                        if name.startswith("_") and not name.startswith("__"):
                            continue
                    if only_defined and getattr(attr, "__module__", "") != module_name:
                        continue
                    funcs.append(name)
            if query:
                from difflib import get_close_matches
                funcs = get_close_matches(query, funcs, n=10, cutoff=0.3)
            return sorted(set(funcs))

        def walk_package(mod) -> Dict[str, Any]:
            subresults = {}
            modname = mod.__name__
            funcs = list_functions_in_module(mod, modname)
            if funcs:
                print(f"\n🎒: {modname}")
                for f in funcs:
                    print(f"  - {f}")
                subresults[modname] = funcs

            if hasattr(mod, '__path__'):  # If it's a package
                for finder, name, ispkg in pkgutil.walk_packages(mod.__path__, prefix=mod.__name__ + "."):
                    try:
                        submod = importlib.import_module(name)
                        submod_result = walk_package(submod)
                        subresults.update(submod_result)
                    except pytest.skip.Exception as e:
                        if verbose:
                            print(f"正常跳过: Skipped test module {name}: {e}")
                    except Exception as e:
                        if verbose:
                            print(f"因错跳过 {name}: {e}")
            return subresults

        results[root_name] = walk_package(mod)
        return results if return_output else None

def _yaoshi_fernet(mima="mimashigudingde",yan=b"mimashigudingde",verbose=True):
    import base64
    from cryptography.hazmat.primitives.kdf.pbkdf2 import PBKDF2HMAC
    from cryptography.hazmat.primitives import hashes
    from cryptography.hazmat.backends import default_backend 
    kdf = PBKDF2HMAC(
        algorithm=hashes.SHA256(),
        length=32,  
        salt=yan,
        iterations=100000,
        backend=default_backend()
    ) 
    return base64.urlsafe_b64encode(kdf.derive(mima.encode())) 
def fenpass(fpath: str, password: str, copy: bool = True):
    """对文件或整个文件夹进行加密
    Example:
        fpath =  r'path'
        pwd = '****'
        copy_mode =   "yes" # True/False
        fenpass(fpath, pwd, copy_mode)
        """

    def encrypt_xlsx(file_path: str, password: str, copy: bool = True):
        """使用 msoffcrypto 进行 Excel 加密，使其需要密码才能打开"""
        import msoffcrypto
        encrypted_path = file_path.replace(".xlsx", "_en.xlsx") if copy else file_path

        try:
            with open(file_path, "rb") as f:
                office_file = msoffcrypto.OfficeFile(f)

                if office_file.is_encrypted():
                    print(f"Excel 文件已加密，跳过: {file_path}")
                    return

                with open(encrypted_path, "wb") as ef:
                    office_file.encrypt(password, ef)
            
            print(f"Excel 文件已加密: {encrypted_path}")
        except Exception as e:
            print(f"加密 Excel 文件失败: {file_path}, 错误信息: {str(e)}")

    def encrypt_pdf(file_path: str, password: str, copy: bool = True):
        """加密 PDF 文件，使其需要密码才能打开"""
        from PyPDF2 import PdfReader, PdfWriter 
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            print(f"PDF 文件已加密，跳过: {file_path}")
            return
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        encrypted_path = file_path.replace(".pdf", "_en.pdf") if copy else file_path
        with open(encrypted_path, "wb") as f:
            writer.write(f)
        print(f"PDF 文件已加密: {encrypted_path}") 
    def encrypt_txt(file_path: str, password: str, copy: bool = True):
        """加密文本文件，使用对称加密"""
        from cryptography.fernet import Fernet
        key = _yaoshi_fernet()
        cipher = Fernet(key)
        with open(file_path, "rb") as f:
            encrypted_data = cipher.encrypt(f.read())
        encrypted_path = file_path.replace(".txt", "_en.txt") if copy else file_path
        with open(encrypted_path, "wb") as f:
            f.write(encrypted_data)
        print(f"文本文件已加密: {encrypted_path}")

    def encrypt_zip(file_path: str, password: str, copy: bool = True):
        """加密 ZIP 文件，使其需要密码才能解压"""
        import zipfile
        encrypted_path = file_path.replace(".zip", "_en.zip") if copy else file_path
        with zipfile.ZipFile(encrypted_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            zf.write(file_path, os.path.basename(file_path))
            zf.setpassword(password.encode())
        print(f"ZIP 文件已加密: {encrypted_path}")

    def encrypt_file(file_path: str, password: str, copy: bool = True):
        """根据文件类型选择合适的加密方法"""
        ext = os.path.splitext(file_path)[1].lower()
        try:
            if ext == ".xlsx":
                encrypt_xlsx(file_path, password, copy)
            elif ext == ".pdf":
                encrypt_pdf(file_path, password, copy)
            elif ext == ".docx":
                print(f"不支持的文件类型: {ext}")
            elif ext == ".txt":
                encrypt_txt(file_path, password, copy)
            elif ext == ".zip":
                encrypt_zip(file_path, password, copy)
            else:
                print(f"不支持的文件类型: {ext}")
        except Exception as e:
            print(e)
    # ===== Main Function =======
    if os.path.isdir(fpath):
        for root, _, files in os.walk(fpath):
            for file in files:
                encrypt_file(os.path.join(root, file), password, copy)
    else:
        encrypt_file(fpath, password, copy)
   
def fdepass(fpath: str, password: str, copy: bool = True):
    """
    对文件或整个文件夹进行解密
    Example:
        fpath =  r'path'
        pwd = '****'
        copy_mode =   "yes" # True/False
        fdepass(fpath, pwd, copy_mode)
    """
    def decrypt_xlsx(file_path: str, password: str, copy: bool = True):
        """使用 msoffcrypto 解密 Excel 文件"""
        import msoffcrypto
        decrypted_path = file_path.replace(".xlsx", "_de.xlsx") if copy else file_path

        try:
            original_mod_time = os.path.getmtime(file_path)
            original_access_time = os.path.getatime(file_path)
            with open(file_path, "rb") as f:
                office_file = msoffcrypto.OfficeFile(f)
                
                # Check if the file is already decrypted or unsupported format
                if not office_file.is_encrypted():
                    print(f"文件已经是解密状态，跳过解密: {file_path}")
                    return

                # Load the password and decrypt the file
                office_file.load_key(password=password)

                # Attempt decryption
                with open(decrypted_path, "wb") as df:
                    office_file.decrypt(df) 
            print(f"Excel 文件已解密: {decrypted_path}")
            os.utime(decrypted_path, (original_access_time, original_mod_time))
        except msoffcrypto.exceptions.DecryptionError:
            print(f"无法解密 Excel 文件，密码可能错误: {file_path}")
        except msoffcrypto.exceptions.FileFormatError:
            print(f"不支持的文件格式，无法处理文件: {file_path}")
        except Exception as e:
            print(f"解密过程中发生了错误: {str(e)}")

    def decrypt_pdf(file_path: str, password: str, copy: bool = True):
        """解密 PDF 文件"""
        from PyPDF2 import PdfReader, PdfWriter
        reader = PdfReader(file_path)
        
        try:
            reader.decrypt(password)
        except Exception as e:
            print(f"PDF 文件解密失败: {file_path}, 错误信息: {str(e)}")
            return

        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        
        decrypted_path = file_path.replace(".pdf", "_de.pdf") if copy else file_path
        with open(decrypted_path, "wb") as f:
            writer.write(f)

        print(f"PDF 文件已解密: {decrypted_path}")

    def decrypt_txt(file_path: str, password: str, copy: bool = True):
        """解密文本文件"""
        from cryptography.fernet import Fernet 
        key = _yaoshi_fernet() 

        try:
            cipher = Fernet(key)  # Use the derived key to create the Fernet object
            with open(file_path, "rb") as f:
                encrypted_data = f.read()
                decrypted_data = cipher.decrypt(encrypted_data)

            decrypted_path = file_path.replace(".txt", "_de.txt") if copy else file_path
            with open(decrypted_path, "wb") as f:
                f.write(decrypted_data)

            print(f"文本文件已解密: {decrypted_path}")

        except Exception as e:
            print(f"文本文件解密失败: {file_path}, 错误信息: {str(e)}")

    def decrypt_zip(file_path: str, password: str, copy: bool = True):
        """解密 ZIP 文件"""
        import zipfile
        decrypted_path = file_path.replace(".zip", "_.zip") if copy else file_path

        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                zf.setpassword(password.encode())
                zf.extractall(decrypted_path)
            print(f"ZIP 文件已解密: {decrypted_path}")

        except Exception as e:
            print(f"ZIP 文件解密失败: {file_path}, 错误信息: {str(e)}")

    def decrypt_file(file_path: str, password: str, copy: bool = True):
        """根据文件类型选择合适的解密方法"""
        try:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".xlsx":
                decrypt_xlsx(file_path, password, copy)
            elif ext == ".pdf":
                decrypt_pdf(file_path, password, copy)
            elif ext == ".docx":
                print(f"不支持的文件类型: {ext}")
            elif ext == ".txt":
                decrypt_txt(file_path, password, copy)
            elif ext == ".zip":
                decrypt_zip(file_path, password, copy)
            else:
                print(f"不支持的文件类型: {ext}")
        except Exception as e:
            print(e)
    
    # ===== Main Function =======
    if os.path.isdir(fpath):
        for root, _, files in os.walk(fpath):
            for file in files:
                decrypt_file(os.path.join(root, file), password, copy)
    else:
        decrypt_file(fpath, password, copy)


def fbackup(
    fpath: str,
    backup_dir: str,
    backup_keep_days: int = 30,
    max_backups: Optional[int] = 5,
    interval: Optional[int] = None,  # 新增参数：时间间隔（秒），例如 3600*12 表示每12小时算一次备份
    compress: bool = False,
    zip_mode: bool = False,
    tar_mode: bool = False,
    checksum: bool = False,
    metadata: bool = False,
    verbose: bool = True,
    preserve_extension: bool = True,
    timestamp_format: str = "%Y%m%d_%H%M%S",
    exclude_patterns: Optional[List[str]] = None,
    include_hidden: bool = True
) -> Optional[str]:
    """
    ULTIMATE all-in-one backup function for files AND folders.
    
    Parameters:
    -----------
    fpath : str
        File or folder to back up
    backup_dir : str
        Where to store backups
    backup_keep_days : int
        Delete backups older than X days (default: 30)
    max_backups : int (optional)
        Maximum number of backups to keep
    compress : bool
        Use gzip compression (single files only)
    zip_mode : bool
        Use ZIP compression (files or folders)
    tar_mode : bool
        Use TAR compression (files or folders)
    checksum : bool
        Generate SHA256 checksum file
    metadata : bool
        Save backup metadata as JSON
    verbose : bool
        Print progress messages
    preserve_extension : bool
        Keep original file extension
    timestamp_format : str
        Custom timestamp format
    exclude_patterns : List[str]
        File patterns to exclude (e.g., ["*.tmp", "temp_*"])
    include_hidden : bool
        Include hidden files/folders (starting with .)
    
    Returns:
    --------
    str or None
        Path to the new backup, or None if failed

    # Backup a file with gzip compression
    fbackup(
        fpath="important_document.pdf",
        backup_dir="backups",
        compress=True,
        checksum=True,
        verbose=True
    )
    
    # Backup a folder with ZIP compression (excluding temp files)
    fbackup(
        fpath="project_folder",
        backup_dir="backups",
        zip_mode=True,
        exclude_patterns=["*.tmp", "temp_*"],
        metadata=True,
        backup_keep_days=60,
        max_backups=5,
        verbose=True
    )
    
    # Backup a folder with tar.gz (including hidden files)
    fbackup(
        fpath=".config",
        backup_dir="backups",
        tar_mode=True,
        include_hidden=True,
        verbose=True
    )

    """

    import os
    import shutil
    import gzip
    import hashlib
    import time
    from datetime import datetime
    from typing import Optional, List, Tuple
    import zipfile
    import tarfile
    import json
    import logging
    from pathlib import Path
    # --- Setup logging ---
    def _log(msg: str, level: str = "info"):
        if verbose:
            now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
            print(f"[{now}] [{level.upper()}] {msg}")

    # --- Helper functions ---
    def _generate_checksum(file_path: str) -> str:
        sha256 = hashlib.sha256()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                sha256.update(chunk)
        return sha256.hexdigest()

    def _should_exclude(file_path: str) -> bool:
        if not exclude_patterns:
            return False
        filename = os.path.basename(file_path) 
        return (
                any(filename.endswith(pattern.lstrip('*')) for pattern in exclude_patterns) or
                any(filename.startswith(pattern.rstrip('*')) for pattern in exclude_patterns)
                )
 
    def _cleanup_old_backups(backup_prefix: str) -> None:
        try:
            now = time.time()
            cutoff = now - (backup_keep_days * 86400)
            
            backups = []
            prefix_len = len(backup_prefix)
            for f in os.listdir(backup_dir):
                if not f.startswith(backup_prefix):
                    continue
                full_path = os.path.join(backup_dir, f)
                if not os.path.isfile(full_path):
                    continue
                
                # 提取时间戳部分
                suffix = f[prefix_len:]
                try:
                    # 动态计算时间戳长度
                    example_ts = datetime(2000, 1, 1, 0, 0, 0).strftime(timestamp_format)
                    ts_length = len(example_ts)
                    if len(suffix) < ts_length:
                        continue
                    timestamp_str = suffix[:ts_length]
                    file_time = datetime.strptime(timestamp_str, timestamp_format).timestamp()
                    backups.append((full_path, file_time))
                except ValueError:
                    continue  # 时间戳解析失败则跳过
            
            deleted = set()
            # 按保留天数删除
            for path, file_time in backups:
                if file_time < cutoff:
                    try:
                        os.remove(path)
                        deleted.add(path)
                        _log(f"Deleted old backup: {path}", "debug")
                        # 删除关联文件
                        for ext in ['.sha256', '.meta']:
                            if os.path.exists(path + ext):
                                os.remove(path + ext)
                    except Exception as e:
                        _log(f"Error deleting {path}: {e}", "warning")
            
            # 按备份数量删除（新增 interval 逻辑）
            if max_backups and max_backups > 0:
                remaining = [b for b in backups if b[0] not in deleted]
                remaining.sort(key=lambda x: x[1])  # 按时间升序排列（最旧在前）
                
                if interval:
                    # 按时间间隔分组
                    grouped_backups = []
                    current_group = []
                    last_time = None
                    for path, file_time in remaining:
                        if last_time is None or (file_time - last_time) >= interval:
                            if current_group:
                                grouped_backups.append(current_group)
                            current_group = [path]
                            last_time = file_time
                        else:
                            current_group.append(path)
                    if current_group:
                        grouped_backups.append(current_group)
                    
                    # 计算需要删除的组数
                    num_to_delete = len(grouped_backups) - max_backups
                    if num_to_delete > 0:
                        for i in range(num_to_delete):
                            for path in grouped_backups[i]:
                                try:
                                    os.remove(path)
                                    _log(f"Deleted excess backup (interval-based): {path}", "debug")
                                    for ext in ['.sha256', '.meta']:
                                        if os.path.exists(path + ext):
                                            os.remove(path + ext)
                                except Exception as e:
                                    _log(f"Error deleting {path}: {e}", "warning")
                else:
                    # 原来的逻辑（直接按文件数删除）
                    num_to_delete = len(remaining) - max_backups
                    if num_to_delete > 0:
                        for i in range(num_to_delete):
                            path = remaining[i][0]
                            try:
                                os.remove(path)
                                _log(f"Deleted excess backup: {path}", "debug")
                                for ext in ['.sha256', '.meta']:
                                    if os.path.exists(path + ext):
                                        os.remove(path + ext)
                            except Exception as e:
                                _log(f"Error deleting {path}: {e}", "warning")
        except Exception as e:
            _log(f"Cleanup error: {e}", "error")

    # --- Main backup logic ---
    try:
        # Validate source
        if not os.path.exists(fpath):
            _log(f"Source path not found: {fpath}", "error")
            return None

        is_dir = os.path.isdir(fpath)
        is_file = os.path.isfile(fpath)
        
        if not (is_dir or is_file):
            _log(f"Source is neither file nor directory: {fpath}", "error")
            return None

        # Create backup directory
        os.makedirs(backup_dir, exist_ok=True)

        # Generate backup filename
        source_name = os.path.basename(fpath.rstrip('/'))
        name, ext = os.path.splitext(source_name)
        timestamp = datetime.now().strftime(timestamp_format)
        
        if preserve_extension and is_file:
            backup_prefix = f"{name}_"
            backup_name = f"{name}_{timestamp}{ext}"
        else:
            backup_prefix = f"{source_name}_"
            backup_name = f"{source_name}_{timestamp}"

        # Add compression extension
        if tar_mode:
            backup_name += ".tar.gz"
            compress_type = "tar"
        elif zip_mode:
            backup_name += ".zip"
            compress_type = "zip"
        elif compress and is_file:
            backup_name += ".gz"
            compress_type = "gzip"
        else:
            compress_type = "none"

        backup_path = os.path.join(backup_dir, backup_name)

        # --- Perform backup ---
        _log(f"Starting backup of {fpath}...")
        
        if is_file:
            if compress_type == "gzip":
                with open(fpath, 'rb') as f_in:
                    with gzip.open(backup_path, 'wb') as f_out:
                        shutil.copyfileobj(f_in, f_out)
            elif compress_type == "zip":
                with zipfile.ZipFile(backup_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    zipf.write(fpath, os.path.basename(fpath))
            elif compress_type == "tar":
                with tarfile.open(backup_path, "w:gz") as tar:
                    tar.add(fpath, arcname=os.path.basename(fpath))
            else:
                shutil.copy2(fpath, backup_path)
        else:  # Directory
            if compress_type == "zip":
                with zipfile.ZipFile(backup_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for root, dirs, files in os.walk(fpath):
                        if not include_hidden:
                            dirs[:] = [d for d in dirs if not d.startswith('.')]
                            files = [f for f in files if not f.startswith('.')]
                        
                        for file in files:
                            full_path = os.path.join(root, file)
                            if not _should_exclude(full_path):
                                arcname = os.path.relpath(full_path, start=fpath)
                                zipf.write(full_path, arcname)
            elif compress_type == "tar":
                with tarfile.open(backup_path, "w:gz") as tar:
                    for root, dirs, files in os.walk(fpath):
                        if not include_hidden:
                            dirs[:] = [d for d in dirs if not d.startswith('.')]
                            files = [f for f in files if not f.startswith('.')]
                        
                        for file in files:
                            full_path = os.path.join(root, file)
                            if not _should_exclude(full_path):
                                arcname = os.path.relpath(full_path, start=fpath)
                                tar.add(full_path, arcname=arcname)
            else:
                # Simple folder copy
                shutil.copytree(
                    fpath,
                    backup_path,
                    ignore=shutil.ignore_patterns(*exclude_patterns) if exclude_patterns else None,
                    dirs_exist_ok=False
                )

        _log(f"Backup created: {backup_path}")

        # Generate checksum
        if checksum and os.path.isfile(backup_path):
            checksum_value = _generate_checksum(backup_path)
            with open(backup_path + '.sha256', 'w') as f:
                f.write(checksum_value)
            _log(f"Checksum saved: {backup_path}.sha256")

        # Save metadata
        if metadata:
            meta = {
                "source": fpath,
                "backup_time": datetime.now().isoformat(),
                "type": "directory" if is_dir else "file",
                "size": os.path.getsize(backup_path) if os.path.isfile(backup_path) else 
                      sum(f.stat().st_size for f in Path(backup_path).rglob('*') if f.is_file()),
                "compression": compress_type,
                "excluded": exclude_patterns if exclude_patterns else None,
                "system": {
                    "platform": os.name,
                    "user": os.getlogin() if hasattr(os, 'getlogin') else None
                }
            }
            with open(backup_path + '.meta', 'w') as f:
                json.dump(meta, f, indent=2)
            _log(f"Metadata saved: {backup_path}.meta")

        # Clean up old backups
        if backup_keep_days > 0 or max_backups:
            _cleanup_old_backups(backup_prefix)

        return backup_path

    except Exception as e:
        _log(f"Backup FAILED: {str(e)}", "error")
        return None
 
 

# ------ Excel文化查看------
import threading
MONTH_TRANSLATIONS = {
    "January": "Januar",
    "February": "Februar",
    "March": "März",
    "April": "April",
    "May": "Mai",
    "June": "Juni",
    "July": "Juli",
    "August": "August",
    "September": "September",
    "October": "Oktober",
    "November": "November",
    "December": "Dezember"
}

def watchdog(fpath, dir_save=None, dir_backup=None, check_interval=30, verbose=True, backup_interval=3600*8, backup_keep_days=30, max_retries=100, retry_delay=30):
    """用来查看Excel文件的变化"""
    import os
    import time
    import csv
    import shutil
    from datetime import datetime
    from openpyxl import load_workbook

    def init_dir_save(dir_save, verbose=False):
        """初始化日志文件，如果不存在则创建"""
        if not os.path.exists(dir_save):
            with open(dir_save, mode="w", newline="") as file:
                writer = csv.writer(file)
                writer.writerow(["when", "who", "file", "sheetname", "change_type", "details", "old_value", "new_value"])
            if verbose:
                print(f"📄 Log file created: {dir_save}")

    def delete_old_backups(dir_backup, file_path=None, days=30, verbose=False):
        """删除超过指定天数的备份文件"""
        if not os.path.exists(dir_backup):
            return
        cutoff_time = time.time() - (days * 86400)
        for filename in os.listdir(dir_backup):
            fpath_ = os.path.join(dir_backup, filename)
            if os.path.isfile(fpath_) and (file_path is None or os.path.basename(file_path) in filename):
                file_mtime = os.path.getmtime(fpath_)
                if file_mtime < cutoff_time:
                    os.remove(fpath_)
                    if verbose:
                        print(f"🗑 Deleted old backup: {fpath_}")

    def load_excel_snapshot(file_path, max_retries=100, retry_delay=30, verbose=False):
        """加载Excel文件并返回其数据和最后修改用户"""
        snapshot = {}
        last_modified_by = "Unknown"
        retries = 0
        
        while retries < max_retries:
            if not os.path.exists(file_path):
                return (snapshot, last_modified_by)
            
            try:
                wb = load_workbook(file_path, data_only=True)
                last_modified_by = wb.properties.last_modified_by or "Unknown"
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_data = []
                    for row in ws.iter_rows(values_only=True):
                        sheet_data.append(list(row))
                    snapshot[sheet_name] = sheet_data
                return (snapshot, last_modified_by)
            except Exception as e:
                if verbose:
                    print(f"⚠ Could not read file: {e} (Retry {retries + 1}/{max_retries})")
                time.sleep(retry_delay)
                retries += 1
        
        return (snapshot, last_modified_by)

    def log_event(dir_save, file, sheet, change_type, details, old_value, new_value, user, verbose=True):
        """记录变化到日志文件"""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(dir_save, mode="a", newline="") as file_:
            writer = csv.writer(file_)
            writer.writerow([timestamp, user, file, sheet, change_type, details, old_value, new_value])
            try:
                old_value=old_value.replace("\n", " ")
                new_value=new_value.replace("\n", " ")
            except:
                pass
        if verbose:
            print(f"(User: {user}) [{timestamp}] {os.path.basename(file)} \n\t@{sheet} {change_type} {details}: \n\t{old_value} \n\t→{new_value} ")

    def create_backup(file_path, dir_backup, backup_keep_days=30, verbose=False):
        """创建Excel文件的备份"""
        if not os.path.exists(dir_backup):
            os.makedirs(dir_backup)
        if os.path.exists(file_path):
            backup_name = os.path.join(dir_backup, f"{os.path.basename(file_path)}.backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx")
            shutil.copy(file_path, backup_name)
            if verbose:
                print(f"📂 Backup created: {backup_name}") 
        delete_old_backups(dir_backup, file_path=file_path, days=backup_keep_days, verbose=verbose)

    def detect_row_changes(prev_sheet_data, new_sheet_data, file, sheet, user, dir_save, verbose):
        """检测并记录行插入或删除"""
        if len(new_sheet_data) > len(prev_sheet_data):  # 行插入
            for i in range(len(new_sheet_data)):
                if i >= len(prev_sheet_data) or new_sheet_data[i] != prev_sheet_data[i]:
                    inserted_row = new_sheet_data[i]
                    log_event(dir_save, file, sheet, "row_insert", f"at row {i + 1}", None, inserted_row, user, verbose)
                    return True  # 跳过单元格变化检测
        elif len(new_sheet_data) < len(prev_sheet_data):  # 行删除
            for i in range(len(prev_sheet_data)):
                if i >= len(new_sheet_data) or prev_sheet_data[i] != new_sheet_data[i]:
                    deleted_row = prev_sheet_data[i]
                    log_event(dir_save, file, sheet, "row_delete", f"at row {i + 1}", deleted_row, None, user, verbose)
                    return True  # 跳过单元格变化检测
        return False

    def detect_column_changes(prev_sheet_data, new_sheet_data, file, sheet, user, dir_save, verbose):
        """检测并记录列插入或删除"""
        prev_cols = max(len(row) for row in prev_sheet_data) if prev_sheet_data else 0
        new_cols = max(len(row) for row in new_sheet_data) if new_sheet_data else 0
        
        if new_cols > prev_cols:  # 列插入
            for col_idx in range(new_cols):
                if col_idx >= prev_cols or any(
                    (row[col_idx] if col_idx < len(row) else None) != (prev_row[col_idx] if col_idx < len(prev_row) else None)
                    for row, prev_row in zip(new_sheet_data, prev_sheet_data)
                ):
                    inserted_col = [row[col_idx] if col_idx < len(row) else None for row in new_sheet_data]
                    log_event(dir_save, file, sheet, "col_instert", f"at column {chr(65 + col_idx)}", None, inserted_col, user, verbose)
                    return True  # 跳过单元格变化检测
        elif new_cols < prev_cols:  # 列删除
            for col_idx in range(prev_cols):
                if col_idx >= new_cols or any(
                    (row[col_idx] if col_idx < len(row) else None) != (new_row[col_idx] if col_idx < len(new_row) else None)
                    for row, new_row in zip(prev_sheet_data, new_sheet_data)
                ):
                    deleted_col = [row[col_idx] if col_idx < len(row) else None for row in prev_sheet_data]
                    log_event(dir_save, file, sheet, "col_delete", f"at column {chr(65 + col_idx)}", deleted_col, None, user, verbose)
                    return True  # 跳过单元格变化检测
        return False

    def detect_cell_changes(prev_sheet_data, new_sheet_data, file, sheet, user, dir_save, verbose):
        """检测并记录单元格变化"""
        for row_idx in range(min(len(prev_sheet_data), len(new_sheet_data))):
            prev_row = prev_sheet_data[row_idx]
            new_row = new_sheet_data[row_idx]
            for col_idx in range(min(len(prev_row), len(new_row))):
                prev_val = prev_row[col_idx]
                new_val = new_row[col_idx]
                if prev_val == new_val:
                    continue
                # Ignore changes if they are just month name translations
                if (prev_val in MONTH_TRANSLATIONS and MONTH_TRANSLATIONS[prev_val] == new_val) or \
                   (new_val in MONTH_TRANSLATIONS and MONTH_TRANSLATIONS[new_val] == prev_val):
                    continue  # Ignore this change 
                cell = f"{chr(65 + col_idx)}{row_idx + 1}"
                log_event(dir_save, file, sheet, " ", cell, prev_val, new_val, user, verbose)

    def monitor_file(fpath, dir_save, dir_backup, check_interval, verbose, backup_interval, backup_keep_days, max_retries, retry_delay):
        """监控单个Excel文件的变化"""
        if not os.path.exists(fpath):
            if verbose:
                print(f"❌ Error: The file '{fpath}' does not exist!")
            return

        # 初始化日志文件
        log_file = os.path.join(dir_save, f"log_{os.path.basename(fpath)}.csv")
        init_dir_save(log_file, verbose=verbose)

        # 加载初始快照
        previous_snapshot, previous_user = load_excel_snapshot(fpath, max_retries=max_retries, retry_delay=retry_delay, verbose=verbose)

        start_time = time.time()
        while True:
            time.sleep(check_interval)
            new_snapshot, new_user = load_excel_snapshot(fpath, max_retries=max_retries, retry_delay=retry_delay, verbose=verbose)
            
            if not new_snapshot:
                continue
            
            for sheet in new_snapshot:
                new_sheet_data = new_snapshot[sheet]
                prev_sheet_data = previous_snapshot.get(sheet, [])
                
                # 检测行变化
                if detect_row_changes(prev_sheet_data, new_sheet_data, fpath, sheet, new_user, log_file, verbose):
                    continue  # 跳过单元格变化检测
                
                # 检测列变化
                if detect_column_changes(prev_sheet_data, new_sheet_data, fpath, sheet, new_user, log_file, verbose):
                    continue  # 跳过单元格变化检测
                
                # 检测单元格变化
                detect_cell_changes(prev_sheet_data, new_sheet_data, fpath, sheet, new_user, log_file, verbose)
            
            previous_snapshot = new_snapshot
            previous_user = new_user
            
            # 备份逻辑
            if dir_backup is not None and (time.time() - start_time) >= backup_interval:
                create_backup(fpath, dir_backup, backup_keep_days=backup_keep_days, verbose=verbose)
                start_time = time.time()
            
            # 停止监控
            if os.path.exists("stop_monitoring.txt"):
                if verbose:
                    print(f"🛑 Stopping monitoring for {fpath} due to user request.")
                break

    # 确保fpath是列表
    if isinstance(fpath, str):
        fpath = [fpath]

    # 创建日志目录
    if dir_save is None:
        dir_save = "excel_monitor_logs"
    if not os.path.exists(dir_save):
        os.makedirs(dir_save)

    # 创建线程监控每个文件
    threads = []
    for fpath_ in fpath:
        if verbose:
            print(f"\tChecking: \n{fpath_}\n")
        thread = threading.Thread(
            target=monitor_file,
            args=(fpath_, dir_save, dir_backup, check_interval, verbose, backup_interval, backup_keep_days, max_retries, retry_delay)
        )
        thread.start()
        threads.append(thread)

    # 等待所有线程完成
    for thread in threads:
        thread.join()


def fmonitor(fpath, dir_save=None, dir_backup=None, check_interval=2, verbose=True, backup_interval=3600*8, backup_keep_days=30, max_retries=5, retry_delay=30):
    import os
    import time
    import csv
    import shutil
    from datetime import datetime
    from openpyxl import load_workbook

    def init_dir_save(dir_save, verbose=False):
        """Initialize the log file with headers if it doesn't exist."""
        if not os.path.exists(dir_save):
            with open(dir_save, mode="w", newline="") as file:
                writer = csv.writer(file)
                writer.writerow(["when", "who", "sheetname", "change_type", "details", "old_value", "new_value"])
            if verbose:
                print(f"📄 Log file created: {dir_save}")

    def delete_old_backups(dir_backup, file_path=None, days=30, verbose=False):
        """Delete backups older than the specified number of days."""
        if not os.path.exists(dir_backup):
            return
        cutoff_time = time.time() - (days * 86400)
        for filename in os.listdir(dir_backup):
            fpath_ = os.path.join(dir_backup, filename)
            if os.path.isfile(fpath_) and (file_path is None or os.path.basename(file_path) in filename):
                file_mtime = os.path.getmtime(fpath_)
                if file_mtime < cutoff_time:
                    os.remove(fpath_)
                    if verbose:
                        print(f"🗑 Deleted old backup: {fpath_}")

    def load_excel_snapshot(file_path, max_retries=5, retry_delay=30, verbose=False):
        """Load an Excel file and return a snapshot of its data and last modified user."""
        snapshot = {}
        last_modified_by = "Unknown"
        retries = 0
        
        while retries < max_retries:
            if not os.path.exists(file_path):
                return (snapshot, last_modified_by)
            
            try:
                wb = load_workbook(file_path, data_only=True)
                last_modified_by = wb.properties.last_modified_by or "Unknown"
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    sheet_data = []
                    for row in ws.iter_rows(values_only=True):
                        sheet_data.append(list(row))
                    snapshot[sheet_name] = sheet_data
                return (snapshot, last_modified_by)
            except Exception as e:
                if verbose:
                    print(f"⚠ Could not read file: {e} (Retry {retries + 1}/{max_retries})")
                time.sleep(retry_delay)
                retries += 1
        
        return (snapshot, last_modified_by)

    def log_event(dir_save, sheet, change_type, details, old_value, new_value, user, verbose=True):
        """Log changes with the specified user."""
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        with open(dir_save, mode="a", newline="") as file:
            writer = csv.writer(file)
            writer.writerow([timestamp, user, sheet, change_type, details, old_value, new_value])
        if verbose:
            print(f"(User: {user}) {sheet} {change_type} {details}: {old_value} → {new_value} [{timestamp}]")

    def create_backup(file_path, dir_backup, backup_keep_days=30, verbose=False):
        """Create a backup of the Excel file."""
        if not os.path.exists(dir_backup):
            os.makedirs(dir_backup)
        if os.path.exists(file_path):
            backup_name = os.path.join(dir_backup, f"{os.path.basename(file_path)}.backup_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx")
            shutil.copy(file_path, backup_name)
            if verbose:
                print(f"📂 Backup created: {backup_name}") 
        delete_old_backups(dir_backup, file_path=file_path, days=backup_keep_days, verbose=verbose)

    def detect_row_changes(prev_sheet_data, new_sheet_data, sheet, user, dir_save, verbose):
        """Detect and log row insertions or deletions."""
        if len(new_sheet_data) > len(prev_sheet_data):  # Row inserted
            for i in range(len(new_sheet_data)):
                if i >= len(prev_sheet_data) or new_sheet_data[i] != prev_sheet_data[i]:
                    inserted_row = new_sheet_data[i]
                    log_event(dir_save, sheet, "ROW_INSERTED", f"at row {i + 1}", None, inserted_row, user, verbose)
                    return True  # Skip cell changes for this sheet
        elif len(new_sheet_data) < len(prev_sheet_data):  # Row deleted
            for i in range(len(prev_sheet_data)):
                if i >= len(new_sheet_data) or prev_sheet_data[i] != new_sheet_data[i]:
                    deleted_row = prev_sheet_data[i]
                    log_event(dir_save, sheet, "ROW_DELETED", f"at row {i + 1}", deleted_row, None, user, verbose)
                    return True  # Skip cell changes for this sheet
        return False

    def detect_cell_changes(prev_sheet_data, new_sheet_data, sheet, user, dir_save, verbose):
        """Detect and log individual cell changes."""
        for row_idx in range(min(len(prev_sheet_data), len(new_sheet_data))):
            prev_row = prev_sheet_data[row_idx]
            new_row = new_sheet_data[row_idx]
            for col_idx in range(min(len(prev_row), len(new_row))):
                prev_val = prev_row[col_idx]
                new_val = new_row[col_idx]
                if prev_val != new_val:
                    cell = f"{chr(65 + col_idx)}{row_idx + 1}"
                    log_event(dir_save, sheet, "", cell, prev_val, new_val, user, verbose)

    if not os.path.exists(fpath):
        if verbose:
            print(f"❌ Error: The file '{fpath}' does not exist!")
        return
    if dir_save is None:
        dir_save = f"log_{os.path.basename(fpath)}.csv"
    if verbose:
        print(f"📊 Monitoring Excel file: \n\t{fpath}")
    init_dir_save(dir_save, verbose=verbose)
    previous_snapshot = load_excel_snapshot(fpath, max_retries=max_retries, retry_delay=retry_delay, verbose=verbose)
    previous_data, previous_user = previous_snapshot

    start_time = time.time()
    while True:
        time.sleep(check_interval)
        new_snapshot = load_excel_snapshot(fpath, max_retries=max_retries, retry_delay=retry_delay, verbose=verbose)
        new_data, new_user = new_snapshot
        
        if not new_data:
            continue
        
        for sheet in new_data:
            new_sheet_data = new_data[sheet]
            prev_sheet_data = previous_data.get(sheet, [])
            
            # Detect row changes first
            if detect_row_changes(prev_sheet_data, new_sheet_data, sheet, new_user, dir_save, verbose):
                continue  # Skip cell changes if row changes were detected
            
            # Detect cell changes
            detect_cell_changes(prev_sheet_data, new_sheet_data, sheet, new_user, dir_save, verbose)
        
        previous_snapshot = new_snapshot
        previous_data, previous_user = new_data, new_user
        
        if dir_backup is not None and (time.time() - start_time) >= backup_interval:
            create_backup(fpath, dir_backup, backup_keep_days=backup_keep_days, verbose=verbose)
            start_time = time.time()
        
        if os.path.exists("stop_monitoring.txt"):
            if verbose:
                print("🛑 Stopping monitoring due to user request.")
            break 


# only for backup these scripts
def backup(
    src="/Users/macjianfeng/Dropbox/github/python/py2ls/.venv/lib/python3.12/site-packages/py2ls/",
    tar="/Users/macjianfeng/Dropbox/github/python/py2ls/py2ls/",
    kind="py",
    overwrite=True,
    reverse=False,
    verbose=False,
):
    if reverse:
        src, tar = tar, src
        print(f"reversed")
    f = listdir(src, kind=kind, verbose=verbose)
    [fcopy(i, tar, overwrite=overwrite, verbose=verbose) for i in f.path]

    print(f"backup '*.{kind}'...\nfrom {src} \nto {tar}")


def run_once_within(duration=60, reverse=False):  # default 60s
    import time

    """
    如果reverse is True, 则在第一次运行时并不运行.但是在第二次运行时则运行
    usage:
    if run_once_within():
        print("This code runs once per minute.")
    else:
        print("The code has already been run in the last minute.")
    
    """
    if not hasattr(run_once_within, "time_last"):
        run_once_within.time_last = None
    time_curr = time.time()

    if (run_once_within.time_last is None) or (
        time_curr - run_once_within.time_last >= duration
    ):
        run_once_within.time_last = time_curr  # Update the last execution time
        return False if reverse else True
    else:
        return True if reverse else False


def plt_font(dir_font: str = "/System/Library/Fonts/Hiragino Sans GB.ttc"):
    """
    Add the Chinese (default) font to the font manager
    show chinese
    Args:
        dir_font (str, optional): _description_. Defaults to "/System/Library/Fonts/Hiragino Sans GB.ttc".
    """
    import matplotlib.pyplot as plt
    from matplotlib import font_manager

    slashtype = "/" if "mac" in get_os() else "\\"
    if slashtype in dir_font:
        font_manager.fontManager.addfont(dir_font)
        fontname = os.path.basename(dir_font).split(".")[0]
    else:
        if "cn" in dir_font.lower() or "ch" in dir_font.lower():
            fontname = "Hiragino Sans GB"  # default Chinese font
        else:
            fontname = dir_font

    plt.rcParams["font.sans-serif"] = [fontname]
    # plt.rcParams["font.family"] = "sans-serif"
    plt.rcParams["axes.unicode_minus"] = False
    fonts_in_system = font_manager.findSystemFonts(fontpaths=None, fontext="ttf")
    fontname_in_system = [os.path.basename(i).split(".")[0] for i in fonts_in_system]
    if fontname not in fontname_in_system:
        print(f"Font '{fontname}' not found. Falling back to default.")
        plt.rcParams["font.sans-serif"] = ["Arial"]
    return fontname


# set 'dir_save'
if "dar" in sys.platform:
    dir_save = "/Users/macjianfeng/Dropbox/Downloads/"
else:
    if "win" in sys.platform:
        dir_save = "Z:\\Jianfeng\\temp\\"
    elif "lin" in sys.platform:
        dir_data = "/Users/macjianfeng/Dropbox/github/python/py2ls/confidential_data/gmail_login.json"


def unique(lst, ascending=None):
    """
    移除列表中的重复元素，同时可以选择按升序或降序排序。

    参数:
    lst (list): 输入的列表，其中可能包含重复的元素。
    ascending (bool, 可选): 如果为 True，则按升序排序；如果为 False，则按降序排序；如果为 None，则不排序，只移除重复元素。

    返回:
    list: 一个列表，其中的元素是唯一的，顺序根据参数 `ascending` 进行排序。
    """
    if not lst:
        return []
    if ascending is not None:
        # 移除重复项
        unique_items = list(set(lst))
        # 统一元素类型（例如，转换为字符串）
        try:
            unique_items = sorted(
                unique_items, key=lambda x: str(x), reverse=not ascending
            )
        except TypeError:
            # 如果排序失败（例如，元素类型不一致），返回原始去重列表
            return unique_items
        return unique_items
    else:
        # 移除重复项同时保持原始顺序
        seen = set()  # 用于记录已见的元素
        result = []  # 用于存储结果
        for item in lst:
            if item not in seen:
                seen.add(item)  # 记录该元素
                result.append(item)  # 添加到结果列表中
        return result


# ************* below section: run_when *************
def run_when(when: str = "every 2 min", job=None, wait: int = 60):
    if "every" in when.lower():
        when = when.replace("every", "")
        run_every(when=when, job=job, wait=wait)
    elif any([i in when.lower() for i in ["at", "@", ":", "am", "pm"]]):
        time_words = ["at", "@", ":", "am", "pm"]
        # 判断'时间词'是否存在
        time_words_bool = [i in when.lower() for i in time_words]
        # 找到'时间词'的位置
        true_indices = [index for index, value in enumerate(time_words_bool) if value]
        time_word = time_words[true_indices[0]]  # 找到第一个'时间词'
        when = when.replace(time_word, "")  # 去除 时间词
        run_at(when=when, job=job, wait=wait)


def run_every(when: str = None, job=None, wait: int = 60):
    """
    Schedules a job to run at the given interval.

    :param when: String specifying the interval, e.g. '2 minutes', '4 hours', '1 day'.
    :param job: The function to be scheduled.

    # usage:
        def job():
            print("1 sec")
        run_every(when="1 sec", job=job)
    """
    import schedule
    import time

    if job is None:
        print("No job provided!")
        return

    interval, unit = (
        str2num(when),
        strcmp(when.replace("every", ""), ["seconds", "minutes", "hours", "days"])[0],
    )
    print(interval, unit)
    # Mapping the scheduling based on the unit1
    if unit == "seconds":
        schedule.every(interval).seconds.do(job)
    elif unit == "minutes":
        schedule.every(interval).minutes.do(job)
    elif unit == "hours":
        schedule.every(interval).hours.do(job)
    elif unit == "days":
        schedule.every(interval).days.do(job)
    else:
        print(f"Invalid time unit: {unit}")
        return

    print(f"Scheduled job when {interval} {unit}.")

    # Keep the script running to execute the schedule
    while True:
        schedule.run_pending()
        time.sleep(wait)  # in seconds
    time.sleep(wait)  # in seconds 
def run_at(when: str, job=None, wait: int = 60):
    """
    Schedules a job to run at an exact time of the day.

    # Example usage:
    def my_job():
        print("Job executed at the exact time!")
    # Schedule the job at 14:30 when day
    run_at(when="1.30 pm", job=my_job)

    :param when: String specifying the time, e.g. '1:30 pm','1.30 am','14:30', '1:30 pm', '8:45 am'.
    :param job: The function to be scheduled.
    :param wait: The sleep interval between checks in seconds.
    """
    from datetime import datetime
    import time

    if job is None:
        print("No job provided!")
        return
    when = when.replace("A.M.", "AM").replace("P.M.", "PM")
    when = when.replace(".", ":")
    when = when.strip()

    try:
        # Attempt to parse the time in both 24-hour and 12-hour format
        if "am" in when.lower() or "pm" in when.lower():
            scheduled_time = datetime.strptime(
                when, "%I:%M %p"
            ).time()  # 12-hour format with AM/PM
        else:
            scheduled_time = datetime.strptime(when, "%H:%M").time()  # 24-hour format
    except ValueError:
        print(
            f"Invalid time format: {when}. Use 'HH:MM' (24-hour) or 'H:MM AM/PM' format."
        )
        return
    print(f"Job scheduled to run at {scheduled_time}.")
    # Keep checking the current time
    while True:
        now = datetime.now()
        # Check if current time matches the scheduled time
        if (
            now.time().hour == scheduled_time.hour
            and now.time().minute == scheduled_time.minute
        ):
            job()  # Run the job
            time.sleep(
                wait
            )  # Sleep for a minute to avoid running the job multiple times in the same minute

        time.sleep(wait)  # wait to avoid excessive CPU usage
 
# ************* above section: run_when *************
def get_timezone(timezone: str | list = None):
    if timezone is None:
        usage = """
        usage:
        datetime.now().astimezone(get_timezone("shanghai")).strftime("%H:%M")
        """
        print(usage)
        return None
    from pytz import all_timezones
    import pytz

    if isinstance(timezone, str):
        timezone = [timezone]

    # Extract the part after the "/" in time zones (if exists)
    timezones = [ssplit(i, "/")[1] if "/" in i else i for i in all_timezones]

    # Print all available time zones for debugging purposes
    # print(timezones)

    # Find and return matched time zones using strcmp
    matched_timezones = [all_timezones[strcmp(i, timezones)[1]] for i in timezone]
    if len(matched_timezones) == 1:
        return pytz.timezone(matched_timezones[0])
    else:
        return matched_timezones

def upgrade(module="py2ls", uninstall=False):
    """
    Installs or upgrades a specified Python module.

    Parameters:
    module (str): The name of the module to install/upgrade.
    uninstall (bool): If True, uninstalls the webdriver-manager before upgrading.
    """
    import subprocess
    def is_package_installed(package_name):
        """Check if a package is installed."""
        import importlib.util

        package_spec = importlib.util.find_spec(package_name)
        return package_spec is not None

    if not is_package_installed(module):
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", module])
        except subprocess.CalledProcessError as e:
            print(f"An error occurred while installing {module}: {e}")
    if uninstall:
        subprocess.check_call(["pip", "uninstall", "-y", "webdriver-manager"])
    try:
        subprocess.check_call(
            [sys.executable, "-m", "pip", "install", "--upgrade", module]
        )
    except subprocess.CalledProcessError as e:
        print(f"An error occurred while upgrading py2ls: {e}")


def get_version(pkg):
    import importlib.metadata

    def get_v(pkg_name):
        try:
            version = importlib.metadata.version(pkg_name)
            print(f"version {pkg_name} == {version}")
        except importlib.metadata.PackageNotFoundError:
            print(f"Package '{pkg_name}' not found")

    if isinstance(pkg, str):
        get_v(pkg)
    elif isinstance(pkg, list):
        [get_v(pkg_) for pkg_ in pkg] 

def rm_folder(folder_path, verbose=True):
    import shutil

    try:
        shutil.rmtree(folder_path)
        if verbose:
            print(f"Successfully deleted {folder_path}")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {folder_path}. Reason: {e}")


def fremove(path, verbose=True):
    """
    Remove a folder and all its contents or a single file.
    Parameters:
    path (str): The path to the folder or file to remove.
    verbose (bool): If True, print success or failure messages. Default is True.
    """
    try:
        if os.path.isdir(path):
            import shutil

            shutil.rmtree(path)
            if verbose:
                print(f"Successfully deleted folder {path}")
        elif os.path.isfile(path):
            os.remove(path)
            if verbose:
                print(f"Successfully deleted file {path}")
        else:
            if verbose:
                print(f"Path {path} does not exist")
    except Exception as e:
        if verbose:
            print(f"Failed to delete {path}. Reason: {e}")
 
def get_cwd():
    from pathlib import Path
    # Get the current script's directory as a Path object
    current_directory = Path(__file__).resolve().parent
    return current_directory


def search(
    query,
    limit=5,
    kind="text",
    output="df",
    verbose=False,
    download=False,
    dir_save=None,
    **kwargs,
):
    from duckduckgo_search import DDGS

    if "te" in kind.lower():
        results = DDGS().text(query, max_results=limit)
        res = pd.DataFrame(results)
        res.rename(columns={"href": "links"}, inplace=True)
    if verbose:
        print(f'searching "{query}": got the results below\n{res}')
    if download:
        try:
            downloader(
                url=res.links.tolist(), dir_save=dir_save, verbose=verbose, **kwargs
            )
        except:
            if verbose:
                print(f"failed link")
    return res


def echo(*args, **kwargs):
    """
    query, model="gpt", verbose=True, log=True, dir_save=dir_save
    a ai chat tool
    Args:
        query (str): _description_
        model (str, optional): _description_. Defaults to "gpt".
        verbose (bool, optional): _description_. Defaults to True.
        log (bool, optional): _description_. Defaults to True.
        dir_save (str, path, optional): _description_. Defaults to dir_save.

    Returns:
        str: the answer from ai
    """
    global dir_save
    from duckduckgo_search import DDGS

    query = None
    model = kwargs.get("model", "gpt")
    verbose = kwargs.get("verbose", True)
    log = kwargs.get("log", True)
    dir_save = kwargs.get("dir_save", dir_save)
    for arg in args:
        if isinstance(arg, str):
            if os.path.isdir(arg):
                dir_save = arg
            # elif os.path.isfile(arg):
            #     dir_save = dirname(arg)
            elif len(arg) <= 5:
                model = arg
            else:
                query = arg
        elif isinstance(arg, dict):
            verbose = arg.get("verbose", verbose)
            log = arg.get("log", log)

    def is_in_any(str_candi_short, str_full, ignore_case=True):
        if isinstance(str_candi_short, str):
            str_candi_short = [str_candi_short]
        res_bool = []
        if ignore_case:
            [res_bool.append(i in str_full.lower()) for i in str_candi_short]
        else:
            [res_bool.append(i in str_full) for i in str_candi_short]
        return any(res_bool)

    def valid_mod_name(str_fly):
        if is_in_any(str_fly, "claude-3-haiku"):
            return "claude-3-haiku"
        elif is_in_any(str_fly, "gpt-3.5"):
            return "gpt-3.5"
        elif is_in_any(str_fly, "llama-3-70b"):
            return "llama-3-70b"
        elif is_in_any(str_fly, "mixtral-8x7b"):
            return "mixtral-8x7b"
        else:
            print(
                f"not support your model{model}, supported models: 'claude','gpt(default)', 'llama','mixtral'"
            )
            return "gpt-3.5"  # default model

    model_valid = valid_mod_name(model)
    res = DDGS().chat(query, model=model_valid)
    if verbose:
        from pprint import pp

        pp(res)
    if log:
        from datetime import datetime
        import time

        dt_str = datetime.fromtimestamp(time.time()).strftime("%Y-%m-%d_%H:%M:%S")
        res_ = f"\n\n####Q:{query}\n\n#####Ans:{dt_str}\n\n>{res}\n"
        if bool(os.path.basename(dir_save)):
            fpath = dir_save
        else:
            os.makedirs(dir_save, exist_ok=True)
            fpath = os.path.join(dir_save, f"log_ai.md")
        fupdate(fpath=fpath, content=res_)
        print(f"log file:{fpath}")
    return res


def chat(*args, **kwargs):
    return echo(*args, **kwargs)
def ai(*args, **kwargs):
    return echo(*args, **kwargs)

def detect_lang(text, output="lang", verbose=True):
    from langdetect import detect

    dir_curr_script = os.path.dirname(os.path.abspath(__file__))
    dir_lang_code = dir_curr_script + "/data/lang_code_iso639.json"
    print(dir_curr_script, os.getcwd(), dir_lang_code)
    lang_code_iso639 = fload(dir_lang_code)
    l_lang, l_code = [], []
    [[l_lang.append(v), l_code.append(k)] for v, k in lang_code_iso639.items()]
    try:
        if is_text(text):
            code_detect = detect(text)
            if "c" in output.lower():  # return code
                return l_code[strcmp(code_detect, l_code, verbose=verbose)[1]]
            else:
                return l_lang[strcmp(code_detect, l_code, verbose=verbose)[1]]
        else:
            print(f"{text} is not supported")
            return "no"
    except:
        return "no"


def is_text(s):
    has_alpha = any(char.isalpha() for char in s)
    has_non_alpha = any(not char.isalpha() for char in s)
    # no_special = not re.search(r'[^A-Za-z0-9\s]', s)
    return has_alpha and has_non_alpha

def share(*args, strict=True, n_shared=2, verbose=True):
    """
    check the shared elelements in two list.
    usage:
        list1 = [1, 2, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        list3 = [5, 6, 9, 10]
        a = shared(list1, list2,list3)
    """
    if verbose:
        print("\n********* checking shared elements *********")

    if len(args) == 1 and isinstance(args[0], list):
        lists = args[0]  # Unpack the single list
    else:
        lists = args  # Use the provided arguments as lists
    flattened_lists = [flatten(lst, verbose=verbose) for lst in lists]
    # Ensure all arguments are lists
    if any(not isinstance(lst, list) for lst in flattened_lists):
        print(f"{' ' * 2}All inputs must be lists.")
        return []
    first_list = flattened_lists[0]
    shared_elements = [
        item for item in first_list if all(item in lst for lst in flattened_lists)
    ]
    if strict:
        # Strict mode: require elements to be in all lists
        shared_elements = set(flattened_lists[0])
        for lst in flattened_lists[1:]:
            shared_elements.intersection_update(lst)
    else:
        from collections import Counter

        all_elements = [item for sublist in flattened_lists for item in sublist]
        element_count = Counter(all_elements)
        # Get elements that appear in at least n_shared lists
        shared_elements = [
            item for item, count in element_count.items() if count >= n_shared
        ]

    shared_elements = flatten(shared_elements, verbose=verbose)
    if verbose:
        elements2show = (
            shared_elements if len(shared_elements) < 10 else shared_elements[:5]
        )
        tail = "" if len(shared_elements) < 10 else "......"
        elements2show.append(tail)
        print(f"{' '*2}{len(shared_elements)} elements shared: {' '*2}{elements2show}")
        print("********* checking shared elements *********")
    return shared_elements


def shared(*args, n_shared=None, verbose=True, **kwargs):
    """
    check the shared elelements in two list.
    usage:
        list1 = [1, 2, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        list3 = [5, 6, 9, 10]
        a = shared(list1, list2,list3)
    """
    if verbose:
        print("\n********* checking shared elements *********")

    if len(args) == 1 and isinstance(args[0], list):
        lists = args[0]  # Unpack the single list
    else:
        lists = args  # Use the provided arguments as lists
    flattened_lists = [flatten(lst, verbose=verbose) for lst in lists]

    if n_shared is None:
        n_shared = len(flattened_lists)
        strict = True
    else:
        strict = False
    # Ensure all arguments are lists
    if any(not isinstance(lst, list) for lst in flattened_lists):
        print(f"{' ' * 2}All inputs must be lists.")
        return []
    first_list = flattened_lists[0]
    shared_elements = [
        item for item in first_list if all(item in lst for lst in flattened_lists)
    ]
    if strict:
        # Strict mode: require elements to be in all lists
        shared_elements = set(flattened_lists[0])
        for lst in flattened_lists[1:]:
            shared_elements.intersection_update(lst)
    else:
        from collections import Counter

        all_elements = [item for sublist in flattened_lists for item in sublist]
        element_count = Counter(all_elements)
        # Get elements that appear in at least n_shared lists
        shared_elements = [
            item for item, count in element_count.items() if count >= n_shared
        ]

    shared_elements = flatten(shared_elements, verbose=verbose)
    if verbose:
        elements2show = (
            shared_elements if len(shared_elements) < 10 else shared_elements[:5]
        )
        print(f"{' '*2}{len(shared_elements)} elements shared: {' '*2}{elements2show}")
        print("********* checking shared elements *********")
    return shared_elements


def share_not(*args, n_shared=None, verbose=False):
    """
    To find the elements in list1 that are not shared with list2 while maintaining the original order of list1
    usage:
        list1 = [1, 8, 3, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        not_shared(list1,list2)# output [1,3]
    """
    _common = shared(*args, n_shared=n_shared, verbose=verbose)
    list1 = flatten(args[0], verbose=verbose)
    _not_shared = [item for item in list1 if item not in _common]
    return _not_shared


def not_shared(*args, n_shared=None, verbose=False):
    """
    To find the elements in list1 that are not shared with list2 while maintaining the original order of list1
    usage:
        list1 = [1, 8, 3, 3, 4, 5]
        list2 = [4, 5, 6, 7, 8]
        not_shared(list1,list2)# output [1,3]
    """
    _common = shared(*args, n_shared=n_shared, verbose=verbose)
    list1 = flatten(args[0], verbose=verbose)
    _not_shared = [item for item in list1 if item not in _common]
    return _not_shared


def flatten(nested: Any, unique_list=True, verbose=False):
    """
    Recursively flattens a nested structure (lists, tuples, dictionaries, sets) into a single list.
    Parameters:
        nested : Any, Can be a list, tuple, dictionary, or set.
    Returns: list, A flattened list.
    """
    flattened_list = []
    stack = [nested]
    while stack:
        current = stack.pop()
        if isinstance(current, dict):
            stack.extend(current.values())
        elif isinstance(current, (list, tuple, set)):
            stack.extend(current)
        elif isinstance(current, pd.Series):
            stack.extend(current)
        elif isinstance(
            current, (pd.Index, np.ndarray)
        ):  # df.columns df.index are object of type pd.Index
            stack.extend(current.tolist())
        else:
            flattened_list.append(current)
    if verbose:
        print(
            f"{' '*2}<in info: {len(unique(flattened_list))} elements after flattened>"
        )
    if unique_list:
        return unique(flattened_list)[::-1]
    else:
        return flattened_list


#! ===========extract_text===========
def extract_text(
    text: Union[str, List[str]],
    patterns: Union[str, List[str]],
    *,
    mode: Literal["between", "split", "extract"] = "between",
    keep: Literal["none", "left", "right", "both", "markers"] = "none",
    case: Literal["sensitive", "insensitive"] = "insensitive",
    all_matches: bool = False,
    positions: bool = False,
    regex: bool = False,
    delimiter: Optional[str] = None,
    trim: bool = True,
    as_dict: bool = False,
    verbose: bool = False,
    **kwargs,
) -> Union[List[str], Tuple[int, str], Dict[str, Any], List[Dict[str, Any]], None]:
    """
    Ultimate text extraction tool with enhanced reliability and features.

    Key improvements:
    - Robust split mode with proper delimiter handling
    - Consistent return types across all modes
    - Improved pattern matching logic
    - Better edge case handling


    print(extract_text("A,B,C", ",", mode="split", keep="none", all_matches=True))
    # Correctly returns: ['A', 'B', 'C']

    print(extract_text("A,B,C", ",", mode="split", keep="left"))
    # Returns: ['A,', 'B,', 'C']

    print(extract_text("A,B,C", ",", mode="split", keep="right"))
    # Returns: [',B', ',C']

    print(extract_text("A,B,C", ",", mode="split", keep="both"))
    # Returns: ['A', ',', 'B', ',', 'C']
    """
    if verbose:
        print("""
extract_text(
    text: Union[str, List[str]],
    patterns: Union[str, List[str]],
    *,
    mode: Literal["between", "split", "extract"] = "between",
    keep: Literal["none", "left", "right", "both", "markers"] = "none",
    case: Literal["sensitive", "insensitive"] = "insensitive",
    all_matches: bool = False,
    positions: bool = False,
    regex: bool = False,
    delimiter: Optional[str] = None,
    trim: bool = True,
    as_dict: bool = False,
    verbose: bool = False,
    **kwargs,
) 
              """)
    # Normalization and validation
    text = _normalize_text(text, delimiter)
    patterns = _validate_patterns(patterns)
    flags = re.IGNORECASE if case == "insensitive" else 0

    # Find all matches with enhanced validation
    matches = _find_matches(text, patterns, regex, flags)
    if not matches:
        return None

    # Mode-specific processing
    if mode == "extract":
        return _handle_extract(matches, all_matches, as_dict, positions, trim)
    elif mode == "split":
        return _handle_split(text, matches, keep, all_matches, as_dict, positions, trim)
    elif mode == "between":
        return _handle_between(text, matches, patterns, keep, as_dict, positions, trim)
    else:
        raise ValueError(f"Invalid mode: {mode}")


def _normalize_text(text: Union[str, List[str]], delimiter: Optional[str]) -> str:
    """Normalize text input to single string"""
    if isinstance(text, list):
        return delimiter.join(text) if delimiter else " ".join(text)
    return text


def _validate_patterns(patterns: Union[str, List[str]]) -> List[str]:
    """Validate and normalize patterns"""
    if isinstance(patterns, str):
        return [patterns]
    if not patterns:
        raise ValueError("At least one pattern required")
    return patterns


def _find_matches(
    text: str, patterns: List[str], regex: bool, flags: int
) -> List[dict]:
    """Find all pattern matches with enhanced regex handling"""
    matches = []
    for pattern in patterns:
        try:
            search_pattern = pattern if regex else re.escape(pattern)
            for match in re.finditer(search_pattern, text, flags=flags):
                matches.append(
                    {
                        "text": match.group(),
                        "start": match.start(),
                        "end": match.end(),
                        "pattern": pattern,
                        "full_match": match,
                    }
                )
        except re.error as e:
            raise ValueError(f"Invalid pattern '{pattern}': {e}")
    return sorted(matches, key=lambda x: x["start"])


def _handle_extract(
    matches: List[dict], all_matches: bool, as_dict: bool, positions: bool, trim: bool
) -> Union[List, dict]:
    """Handle text extraction of matched patterns"""
    results = []
    for match in matches if all_matches else [matches[0]]:
        content = match["text"].strip() if trim else match["text"]
        result = (
            {
                "text": content,
                "start": match["start"],
                "end": match["end"],
                "pattern": match["pattern"],
            }
            if as_dict
            else content
        )
        if positions and as_dict:
            result["positions"] = [(match["start"], match["end"])]
        results.append(result)

    return results[0] if not all_matches else results


def _create_part(
    content: str,
    start: int,
    end: int,
    match: Optional[dict],
    as_dict: bool,
    positions: bool,
    trim: bool,
) -> Union[str, dict]:
    """Create a standardized result part"""
    content = content.strip() if trim else content
    if not as_dict:
        return content

    part = {
        "text": content,
        "start": start,
        "end": end,
        "pattern": match["pattern"] if match else None,
    }
    if positions and match:
        part["positions"] = [(match["start"], match["end"])]
    return part


def _handle_between(
    text: str,
    matches: List[dict],
    patterns: List[str],
    keep: str,
    as_dict: bool,
    positions: bool,
    trim: bool,
) -> Union[Tuple, dict]:
    """Reliable between-mode implementation with boundary checks"""
    first_pattern, last_pattern = patterns[0], patterns[-1]
    first_matches = [m for m in matches if m["pattern"] == first_pattern]
    last_matches = [m for m in matches if m["pattern"] == last_pattern]

    if not first_matches or not last_matches:
        return None

    first = first_matches[0]
    last = last_matches[-1]

    if first["start"] > last["start"]:
        return None

    # Calculate extraction window
    start, end = first["start"], last["end"]
    if keep == "none":
        start, end = first["end"], last["start"]
    elif keep == "left":
        end = last["start"]
    elif keep == "right":
        start = first["end"]

    extracted = text[start:end].strip() if trim else text[start:end]

    if as_dict:
        result = {
            "text": extracted,
            "start": start,
            "end": end,
            "patterns": patterns,
            "match_positions": [(m["start"], m["end"]) for m in matches],
        }
        return result

    return (
        (start, extracted)
        if not positions
        else (start, extracted, [(m["start"], m["end"]) for m in matches])
    )


def _handle_split(
    text: str,
    matches: List[dict],
    keep: str,
    all_matches: bool,
    as_dict: bool,
    positions: bool,
    trim: bool,
) -> Union[List, dict]:
    """Split text with proper handling of keep='both' to include delimiters on both sides"""
    if not matches:
        return (
            [text]
            if not as_dict
            else [{"text": text, "start": 0, "end": len(text), "pattern": None}]
        )

    parts = []
    prev_end = 0
    process_matches = matches if all_matches else [matches[0]]

    # Special handling for keep="both"
    if keep == "both":
        for i, match in enumerate(process_matches):
            start, end = match["start"], match["end"]
            matched_text = text[start:end]

            # First segment (text before first delimiter + first delimiter)
            if i == 0:
                segment = text[prev_end:end]  # From start to end of first delimiter
                if trim:
                    segment = segment.strip()
                if segment or not trim:
                    if as_dict:
                        parts.append(
                            {
                                "text": segment,
                                "start": prev_end,
                                "end": end,
                                "pattern": match["pattern"],
                                **({"positions": [(start, end)]} if positions else {}),
                            }
                        )
                    else:
                        parts.append(segment)
                prev_end = end

            # Middle segments (delimiter + text + next delimiter)
            if i > 0 and i < len(process_matches):
                next_match = process_matches[i]
                next_start, next_end = next_match["start"], next_match["end"]
                segment = text[
                    prev_end:next_end
                ]  # From prev_end to end of next delimiter
                if trim:
                    segment = segment.strip()
                if segment or not trim:
                    if as_dict:
                        parts.append(
                            {
                                "text": segment,
                                "start": prev_end,
                                "end": next_end,
                                "pattern": next_match["pattern"],
                                **(
                                    {"positions": [(next_start, next_end)]}
                                    if positions
                                    else {}
                                ),
                            }
                        )
                    else:
                        parts.append(segment)
                prev_end = next_end

        # Last segment (last delimiter + remaining text)
        if process_matches and prev_end < len(text):
            last_match = process_matches[-1]
            segment = text[
                last_match["start"] : len(text)
            ]  # From last delimiter to end
            if trim:
                segment = segment.strip()
            if segment or not trim:
                if as_dict:
                    parts.append(
                        {
                            "text": segment,
                            "start": last_match["start"],
                            "end": len(text),
                            "pattern": last_match["pattern"],
                            **(
                                {
                                    "positions": [
                                        (last_match["start"], last_match["end"])
                                    ]
                                }
                                if positions
                                else {}
                            ),
                        }
                    )
                else:
                    parts.append(segment)

        return parts

    # Original handling for other keep modes
    for i, match in enumerate(process_matches):
        start, end = match["start"], match["end"]
        matched_text = text[start:end]

        # Handle text before the match
        if prev_end < start:
            before = text[prev_end:start]
            if trim:
                before = before.strip()
            if before or not trim:
                if as_dict:
                    parts.append(
                        {
                            "text": before,
                            "start": prev_end,
                            "end": start,
                            "pattern": None,
                            **({"positions": []} if positions else {}),
                        }
                    )
                else:
                    parts.append(before)

        # Handle the match based on keep mode
        if keep == "none":
            pass  # Skip the delimiter
        elif keep == "left":
            if parts:
                if as_dict:
                    parts[-1]["text"] += matched_text
                    parts[-1]["end"] = end
                else:
                    parts[-1] += matched_text
            else:
                if as_dict:
                    parts.append(
                        {
                            "text": matched_text,
                            "start": start,
                            "end": end,
                            "pattern": match["pattern"],
                            **({"positions": [(start, end)]} if positions else {}),
                        }
                    )
                else:
                    parts.append(matched_text)
        elif keep == "right":
            if i < len(process_matches) - 1:
                next_start = process_matches[i + 1]["start"]
                if end < next_start:
                    between = text[end:next_start]
                    if as_dict:
                        parts.append(
                            {
                                "text": matched_text + between,
                                "start": start,
                                "end": next_start,
                                "pattern": match["pattern"],
                                **({"positions": [(start, end)]} if positions else {}),
                            }
                        )
                    else:
                        parts.append(matched_text + between)
                    prev_end = next_start
                    continue

        prev_end = end

    # Handle remaining text after last match
    if prev_end < len(text):
        remaining = text[prev_end:]
        if trim:
            remaining = remaining.strip()
        if remaining or not trim:
            if keep == "right" and parts and process_matches:
                last_match = process_matches[-1]
                matched_text = text[last_match["start"] : last_match["end"]]
                if as_dict:
                    parts.append(
                        {
                            "text": matched_text + remaining,
                            "start": last_match["start"],
                            "end": len(text),
                            "pattern": last_match["pattern"],
                            **(
                                {
                                    "positions": [
                                        (last_match["start"], last_match["end"])
                                    ]
                                }
                                if positions
                                else {}
                            ),
                        }
                    )
                else:
                    parts.append(matched_text + remaining)
            else:
                if as_dict:
                    parts.append(
                        {
                            "text": remaining,
                            "start": prev_end,
                            "end": len(text),
                            "pattern": None,
                            **({"positions": []} if positions else {}),
                        }
                    )
                else:
                    parts.append(remaining)

    # Filter empty parts if trimming
    if trim:
        parts = [p for p in parts if (p["text"].strip() if as_dict else p.strip())]

    return parts


def _merge_parts(
    parts: List[Union[str, dict]], text: str, as_dict: bool, trim: bool
) -> Union[str, dict]:
    """Merge adjacent parts for keep=left mode"""
    if as_dict:
        merged_text = "".join(p["text"] for p in parts)
        return {
            "text": merged_text.strip() if trim else merged_text,
            "start": parts[0]["start"],
            "end": parts[-1]["end"],
            "patterns": list(set(p["pattern"] for p in parts if p["pattern"])),
        }
    return "".join(parts).strip() if trim else "".join(parts)
#! ===========extract_text===========

def strcmp(
    search_term: str,
    candidates: List[str],
    ignore_case: bool = True,
    get_rank: bool = True,
    return_scores: bool = False,
    verbose: bool = False,
    scorer: str = "auto",
    method: Optional[str] = None,
    top_n: Optional[int] = 1,
    exact_match_first: bool = False
) -> Union[Tuple[str, int], List[str], List[Tuple[str, int, int]]]:
    """
    使用多种模糊匹配算法对 search_term 与候选字符串进行相似性比较，并返回最佳匹配。

    参数说明：
    - search_term: 要搜索的目标字符串
    - candidates: 候选字符串列表
    - ignore_case: 是否忽略大小写
    - get_rank: 是否返回所有候选的相似度排序结果
    - return_scores: 是否同时返回相似度分数
    - verbose: 是否打印详细信息
    - scorer: 默认评分方法 ("WR", "part", "ratio" 等)
    - method: 明确指定使用的评分方法（覆盖 scorer）
    - top_n: 返回前 top_n 个结果（仅在 get_rank 为 True 时生效）
    - exact_match_first: 若发现完全匹配，则直接返回

    支持的匹配方法说明：
    - ratio: 字符级逐位比较，最严格，适合字符变动不大的情况。
    - partial_ratio: 适用于短字符串在长字符串中子集匹配，例如 “apple” 和 “green apple”。
    - token_sort_ratio: 对词语排序后再比较，适合词序不同但内容相似的句子。
    - token_set_ratio: 比较词集合，适用于一方包含另一方的情况（重复词不影响）。
    - partial_token_sort_ratio: 词排序 + 部分匹配，适合短词嵌套在长句中。
    - WRatio: 综合使用多种方法的加权结果，适合泛用场景。
    """
    from fuzzywuzzy import fuzz, process
    if verbose:
        method_str="""
        - ratio: 字符级逐位比较，最严格，适合字符变动不大的情况。
        - partial_ratio: 适用于短字符串在长字符串中子集匹配，例如 “apple” 和 “green apple”。
        - token_sort_ratio: 对词语排序后再比较，适合词序不同但内容相似的句子。
        - token_set_ratio: 比较词集合，适用于一方包含另一方的情况（重复词不影响）。
        - partial_token_sort_ratio: 词排序 + 部分匹配，适合短词嵌套在长句中。
        - WRatio: 综合使用多种方法的加权结果，适合泛用场景。
        """
        print(method_str)
    def to_lower(s):
        if ignore_case:
            if isinstance(s, str):
                return s.lower()
            elif isinstance(s, list):
                return [str(i).lower() for i in s]
        return s

    if not candidates or not isinstance(candidates, list):
        raise ValueError("Candidates must be a non-empty list of strings.")

    # Clean input
    search_term = str(search_term)
    candidates = [str(c) for c in candidates]

    str1 = to_lower(search_term)
    str2 = to_lower(candidates)

    scoring = (method or scorer).lower()
    get_rank = True if any([top_n is not None, get_rank]) else False
    if exact_match_first:
        for idx, cand in enumerate(str2):
            if str1 == cand:
                if verbose:
                    print(f"Exact match found: {candidates[idx]}")
                if return_scores:
                    return candidates[idx], idx
                return candidates[idx]
    # 中文别名映射 + 方法匹配逻辑
    methods_map = {
        "auto": fuzz.WRatio,
        "balance": fuzz.WRatio,
        "wr": fuzz.WRatio,
        "subset_match": fuzz.token_set_ratio,
        "order_free_match": fuzz.token_sort_ratio,
        "free": fuzz.token_sort_ratio,
        "partial_phrase_match": fuzz.partial_token_sort_ratio,
        "partial_match": fuzz.partial_ratio,
        "strict_match": fuzz.ratio,
        "strict": fuzz.ratio,
    } 

    # 使用 fuzzywuzzy 的 process.extractOne 找出最接近的 method key
    available_keys = list(methods_map.keys())
    scoring_key, score = process.extractOne(scoring, available_keys)
    print(f"use {scoring_key} method") if verbose else None
    # 设置得分阈值（避免误匹配）
    if score < 60:
        raise ValueError(f"无法识别的匹配方法: {scoring}. 可用方法为: {available_keys}") 
    score_func = methods_map.get(scoring, fuzz.WRatio)  # 默认 WRatio

    similarity_scores = [score_func(str1, c) for c in str2]

    if get_rank:
        sorted_indices = sorted(range(len(similarity_scores)), key=lambda i: similarity_scores[i], reverse=True)
        if top_n:
            sorted_indices = sorted_indices[:top_n]
        results = []
        for i in sorted_indices:
            if return_scores:
                results.append((candidates[i], i, similarity_scores[i]))
            else:
                results.append((candidates[i], i))
        if verbose:
            print("Top matches:")
            for r in results:
                print(r)
        return results[0] if top_n==1 else results

    else:
        best_idx = similarity_scores.index(max(similarity_scores))
        if verbose:
            print(f"Best match: {candidates[best_idx]} (Score: {similarity_scores[best_idx]})")
            suggestions = process.extract(search_term, candidates, limit=5)
            print("Suggestions:", suggestions)
        
        if return_scores:
            return candidates[best_idx], best_idx,similarity_scores[best_idx]
        return candidates[best_idx],best_idx

def imgcmp(
    img: list,
    method: str = "knn",
    thr: float = 0.75,
    detector: str = "sift",
    plot_: bool = True,
    figsize=[12, 6],
    grid_size=10,  # only for grid detector
    **kwargs,
):
    """
    Compare two images using SSIM, Feature Matching (SIFT), or KNN Matching.

    Parameters:
    - img (list): List containing two image file paths [img1, img2] or two numpy arrays.
    - method (str): Comparison method ('ssim', 'match', or 'knn').
    - detector (str): Feature detector ('sift', 'grid', 'pixel').
    - thr (float): Threshold for filtering matches.
    - plot_ (bool): Whether to display the results visually.
    - figsize (list): Size of the figure for plots.

    Returns:
    - For 'ssim': (diff, score): SSIM difference map and similarity score.
    - For 'match' or 'knn': (good_matches, len(good_matches), similarity_score): Matches and similarity score.
    """
    import cv2
    import matplotlib.pyplot as plt
    from skimage.metrics import structural_similarity as ssim

    # Load images
    if isinstance(img, list) and isinstance(img[0], str):
        image1 = cv2.imread(img[0])
        image2 = cv2.imread(img[1])
        bool_cvt = True
    else:
        image1, image2 = np.array(img[0]), np.array(img[1])
        bool_cvt = False

    if image1 is None or image2 is None:
        raise ValueError("Could not load one or both images. Check file paths.")
    methods = ["ssim", "match", "knn"]
    method = strcmp(method, methods)[0]
    if method == "ssim":
        # Convert images to grayscale
        gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)

        # Compute SSIM
        score, diff = ssim(gray1, gray2, full=True)
        print(f"SSIM Score: {score:.4f}")

        # Convert diff to 8-bit for visualization
        diff = (diff * 255).astype("uint8")

        # Plot if needed
        if plot_:
            fig, ax = plt.subplots(1, 3, figsize=figsize)
            ax[0].imshow(gray1, cmap="gray")
            ax[0].set_title("Image 1")
            ax[1].imshow(gray2, cmap="gray")
            ax[1].set_title("Image 2")
            ax[2].imshow(diff, cmap="gray")
            ax[2].set_title("Difference (SSIM)")
            plt.tight_layout()
            plt.show()

        return diff, score

    elif method in ["match", "knn"]:
        # Convert images to grayscale
        gray1 = cv2.cvtColor(image1, cv2.COLOR_BGR2GRAY)
        gray2 = cv2.cvtColor(image2, cv2.COLOR_BGR2GRAY)

        if detector == "sift":
            # SIFT detector
            sift = cv2.SIFT_create()
            keypoints1, descriptors1 = sift.detectAndCompute(gray1, None)
            keypoints2, descriptors2 = sift.detectAndCompute(gray2, None)

        elif detector == "grid":
            # Grid-based detection
            keypoints1, descriptors1 = [], []
            keypoints2, descriptors2 = [], []

            for i in range(0, gray1.shape[0], grid_size):
                for j in range(0, gray1.shape[1], grid_size):
                    patch1 = gray1[i : i + grid_size, j : j + grid_size]
                    patch2 = gray2[i : i + grid_size, j : j + grid_size]
                    if patch1.size > 0 and patch2.size > 0:
                        keypoints1.append(
                            cv2.KeyPoint(
                                j + grid_size // 2, i + grid_size // 2, grid_size
                            )
                        )
                        keypoints2.append(
                            cv2.KeyPoint(
                                j + grid_size // 2, i + grid_size // 2, grid_size
                            )
                        )
                        descriptors1.append(np.mean(patch1))
                        descriptors2.append(np.mean(patch2))

            descriptors1 = np.array(descriptors1).reshape(-1, 1)
            descriptors2 = np.array(descriptors2).reshape(-1, 1)

        elif detector == "pixel":
            # Pixel-based direct comparison
            descriptors1 = gray1.flatten()
            descriptors2 = gray2.flatten()
            keypoints1 = [
                cv2.KeyPoint(x, y, 1)
                for y in range(gray1.shape[0])
                for x in range(gray1.shape[1])
            ]
            keypoints2 = [
                cv2.KeyPoint(x, y, 1)
                for y in range(gray2.shape[0])
                for x in range(gray2.shape[1])
            ]

        else:
            raise ValueError("Invalid detector. Use 'sift', 'grid', or 'pixel'.")

        # Handle missing descriptors
        if descriptors1 is None or descriptors2 is None:
            raise ValueError("Failed to compute descriptors for one or both images.")
        # Ensure descriptors are in the correct data type
        if descriptors1.dtype != np.float32:
            descriptors1 = descriptors1.astype(np.float32)
        if descriptors2.dtype != np.float32:
            descriptors2 = descriptors2.astype(np.float32)

        # BFMatcher initialization
        bf = cv2.BFMatcher()
        if method == "match":  # Cross-check matching
            bf = cv2.BFMatcher(cv2.NORM_L2, crossCheck=True)
            matches = bf.match(descriptors1, descriptors2)
            matches = sorted(matches, key=lambda x: x.distance)

            # Filter good matches
            good_matches = [
                m for m in matches if m.distance < thr * matches[-1].distance
            ]

        elif method == "knn":  # KNN matching with ratio test
            bf = cv2.BFMatcher()
            matches = bf.knnMatch(descriptors1, descriptors2, k=2)
            # Apply Lowe's ratio test
            good_matches = [m for m, n in matches if m.distance < thr * n.distance]

        # Calculate similarity score
        similarity_score = len(good_matches) / min(len(keypoints1), len(keypoints2))
        print(f"Number of good matches: {len(good_matches)}")
        print(f"Similarity Score: {similarity_score:.4f}")
        # Handle case where no good matches are found
        if len(good_matches) == 0:
            print("No good matches found.")
            return good_matches, 0.0, None

        # Identify matched keypoints
        src_pts = np.float32([keypoints1[m.queryIdx].pt for m in good_matches]).reshape(
            -1, 1, 2
        )
        dst_pts = np.float32([keypoints2[m.trainIdx].pt for m in good_matches]).reshape(
            -1, 1, 2
        )
        # Apply the homography to image2
        try:
            # Calculate Homography using RANSAC
            homography_matrix, mask = cv2.findHomography(
                src_pts, dst_pts, cv2.RANSAC, 5.0
            )
            h, w = image1.shape[:2]
            warped_image2 = cv2.warpPerspective(image2, homography_matrix, (w, h))

            # Plot result if needed
            if plot_:
                fig, ax = plt.subplots(1, 2, figsize=figsize)
                (
                    ax[0].imshow(cv2.cvtColor(image1, cv2.COLOR_BGR2RGB))
                    if bool_cvt
                    else ax[0].imshow(image1)
                )
                ax[0].set_title("Image 1")
                (
                    ax[1].imshow(cv2.cvtColor(warped_image2, cv2.COLOR_BGR2RGB))
                    if bool_cvt
                    else ax[1].imshow(warped_image2)
                )
                ax[1].set_title("Warped Image 2")
                plt.tight_layout()
                plt.show()
        except Exception as e:
            print(e)

        # Plot matches if needed
        if plot_:
            result = cv2.drawMatches(
                image1, keypoints1, image2, keypoints2, good_matches, None, flags=2
            )
            plt.figure(figsize=figsize)
            (
                plt.imshow(cv2.cvtColor(result, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else plt.imshow(result)
            )
            plt.title(
                f"Feature Matches ({len(good_matches)} matches, Score: {similarity_score:.4f})"
            )
            plt.axis("off")
            plt.show()
        # Identify unmatched keypoints
        matched_idx1 = [m.queryIdx for m in good_matches]
        matched_idx2 = [m.trainIdx for m in good_matches]
        matched_kp1 = [kp for i, kp in enumerate(keypoints1) if i in matched_idx1]
        matched_kp2 = [kp for i, kp in enumerate(keypoints2) if i in matched_idx2]
        unmatched_kp1 = [kp for i, kp in enumerate(keypoints1) if i not in matched_idx1]
        unmatched_kp2 = [kp for i, kp in enumerate(keypoints2) if i not in matched_idx2]

        # Mark keypoints on the images
        img1_match = cv2.drawKeypoints(
            image1,
            matched_kp1,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img2_match = cv2.drawKeypoints(
            image2,
            matched_kp2,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img1_unmatch = cv2.drawKeypoints(
            image1,
            unmatched_kp1,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )
        img2_unmatch = cv2.drawKeypoints(
            image2,
            unmatched_kp2,
            None,
            color=(0, 0, 255),
            flags=cv2.DRAW_MATCHES_FLAGS_DRAW_RICH_KEYPOINTS,
        )

        if plot_:
            fig, ax = plt.subplots(1, 2, figsize=figsize)
            (
                ax[0].imshow(cv2.cvtColor(img1_unmatch, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[0].imshow(img1_unmatch)
            )
            ax[0].set_title("Unmatched Keypoints (Image 1)")
            (
                ax[1].imshow(cv2.cvtColor(img2_unmatch, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[1].imshow(img2_unmatch)
            )
            ax[1].set_title("Unmatched Keypoints (Image 2)")
            ax[0].axis("off")
            ax[1].axis("off")
            plt.tight_layout()
            plt.show()
        if plot_:
            fig, ax = plt.subplots(1, 2, figsize=figsize)
            (
                ax[0].imshow(cv2.cvtColor(img1_match, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[0].imshow(img1_match)
            )
            ax[0].set_title("Matched Keypoints (Image 1)")
            (
                ax[1].imshow(cv2.cvtColor(img2_match, cv2.COLOR_BGR2RGB))
                if bool_cvt
                else ax[1].imshow(img2_match)
            )
            ax[1].set_title("Matched Keypoints (Image 2)")
            ax[0].axis("off")
            ax[1].axis("off")
            plt.tight_layout()
            plt.show()
        return good_matches, similarity_score  # , homography_matrix

    else:
        raise ValueError("Invalid method. Use 'ssim', 'match', or 'knn'.")

def fcmp(file1, file2, kind= None, verbose=True, **kwargs):
    import pandas as pd
    import os
    from concurrent.futures import ThreadPoolExecutor
    from datetime import datetime
    import json

    # --- Compare excel files ---
    def cmp_excel(
        file1,# base
        file2,  # new
        sheet_name=None,  # list or strings; default:"common" sheet
        key_columns=None,
        ignore_columns=None,
        numeric_tolerance=0,
        ignore_case=False,
        detect_reordered_rows=False,
        verbose=True,
        **kwargs,
    ):
        """
        Compare two Excel files and identify differences across specified sheets.

        Parameters:
        - file1 (Base/Reference): str, path to the first Excel file.
        - file2: str, path to the second Excel file.
        - sheet_name: list of str, specific sheets to compare (default: all common sheets).
        - key_columns: list of str, columns to use as unique identifiers (default: None, compares all columns).
        - ignore_columns: list of str, columns to exclude from comparison (default: None).
        - numeric_tolerance: float, tolerance for numeric column differences (default: 0, exact match).
        - ignore_case: bool, whether to ignore case differences (default: False).  # Changed here
        - detect_reordered_rows: bool, whether to detect reordered rows (default: False).
        - verbose: bool, whether to print progress messages (default: True).

        Returns:
        - dict, summary of differences for each sheet.
        """
        # Define output directory based on file1 basename
        file1_basename = os.path.splitext(os.path.basename(file1))[0]
        output_dir = f"CMP_{file1_basename}"
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        # Load both files into a dictionary of DataFrames
        xl1 = pd.ExcelFile(file1)
        xl2 = pd.ExcelFile(file2)

        # Get the sheets to compare
        sheets1 = set(xl1.sheet_names)
        sheets2 = set(xl2.sheet_names)
        if sheet_name is None:
            sheet_name = list(sheets1 & sheets2)  # Compare only common sheets
        else:
            sheet_name = [sheet for sheet in sheet_name if sheet in sheets1 and sheets2]

        summary = {}
        print(f"Reference file: '{os.path.basename(file1)}'")
        def compare_sheet(sheet):
            
            if verbose:
                print(f"Comparing sheet: {sheet}...")

            # Read sheets as DataFrames
            df1 = xl1.parse(sheet).fillna("NA")
            df2 = xl2.parse(sheet).fillna("NA")

            # Handle case insensitivity
            if ignore_case:
                df1.columns = [col.lower() for col in df1.columns]
                df2.columns = [col.lower() for col in df2.columns]
                df1 = df1.applymap(lambda x: x.lower() if isinstance(x, str) else x)
                df2 = df2.applymap(lambda x: x.lower() if isinstance(x, str) else x)

            # Drop ignored columns
            if ignore_columns:
                df1 = df1.drop(
                    columns=[col for col in ignore_columns if col in df1.columns],
                    errors="ignore",
                )
                df2 = df2.drop(
                    columns=[col for col in ignore_columns if col in df2.columns],
                    errors="ignore",
                )

            # Normalize column order for comparison
            common_cols = df1.columns.intersection(df2.columns)
            df1 = df1[common_cols]
            df2 = df2[common_cols]

            # Specify key columns for comparison
            if key_columns:
                df1 = df1.set_index(key_columns)
                df2 = df2.set_index(key_columns)
            # Identify added and deleted rows based on entire row comparison, not just index
            added_rows = df2[~df2.apply(tuple, 1).isin(df1.apply(tuple, 1))]
            deleted_rows = df1[~df1.apply(tuple, 1).isin(df2.apply(tuple, 1))]

            # Detect reordered rows
            reordered_rows = pd.DataFrame()
            if detect_reordered_rows:
                # Find rows that exist in both DataFrames but are in different positions
                for idx in df1.index:
                    if idx in df2.index:
                        if not df1.loc[idx].equals(df2.loc[idx]):
                            reordered_rows = reordered_rows.append(df1.loc[idx])

            # Detect modified rows (in case of exact matches between the two files)
            aligned_df1 = df1[df1.index.isin(df2.index)]
            aligned_df2 = df2[df2.index.isin(df1.index)]

            if numeric_tolerance > 0:
                modified_rows = aligned_df1.compare(
                    aligned_df2,
                    keep_shape=False,
                    keep_equal=False,
                    result_names=["left", "right"],
                ).pipe(
                    lambda df: df[
                        ~df.apply(
                            lambda row: (
                                abs(row["left"] - row["right"]) <= numeric_tolerance
                                if pd.api.types.is_numeric_dtype(row["left"])
                                else False
                            ),
                            axis=1,
                        )
                    ]
                )
            else:
                modified_rows = aligned_df1.compare(
                    aligned_df2, keep_shape=False, keep_equal=False
                )

            # Save differences to Excel files
            sheet_dir = os.path.join(output_dir, sheet)
            os.makedirs(sheet_dir, exist_ok=True)
            added_path = os.path.join(sheet_dir, f"{sheet}_added.xlsx")
            deleted_path = os.path.join(sheet_dir, f"{sheet}_deleted.xlsx")
            modified_path = os.path.join(sheet_dir, f"{sheet}_modified.xlsx")
            reordered_path = os.path.join(sheet_dir, f"{sheet}_reordered.xlsx")

            if not added_rows.empty:
                added_rows.to_excel(added_path)
            if not deleted_rows.empty:
                deleted_rows.to_excel(deleted_path)
            if not modified_rows.empty:
                modified_rows.to_excel(modified_path)
            if not reordered_rows.empty:
                reordered_rows.to_excel(reordered_path)

            # Return the summary
            return {
                "added_rows": len(added_rows),
                "deleted_rows": len(deleted_rows),
                "modified_rows": len(modified_rows),
                "reordered_rows": len(reordered_rows),
                "added_file": added_path if not added_rows.empty else None,
                "deleted_file": deleted_path if not deleted_rows.empty else None,
                "modified_file": modified_path if not modified_rows.empty else None,
                "reordered_file": reordered_path if not reordered_rows.empty else None,
            }

        # Use ThreadPoolExecutor for parallel processing
        with ThreadPoolExecutor() as executor:
            results = executor.map(compare_sheet, sheet_name)

        # Collect results
        summary = {sheet: result for sheet, result in zip(sheet_name, results)}

        # Save JSON log
        json_path = os.path.join(output_dir, "comparison_summary.json")
        if os.path.exists(json_path):
            with open(json_path, "r") as f:
                existing_data = json.load(f)
        else:
            existing_data = {}

        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        existing_data[timestamp] = summary
        # Sort the existing data by the timestamp in descending order (latest first)
        existing_data = dict(sorted(existing_data.items(), reverse=True))

        with open(json_path, "w") as f:
            json.dump(existing_data, f, indent=4)
        if verbose:
            print(f"Comparison complete. Results saved in '{output_dir}'")

        return summary

    # --- Compare CSV files ---
    def cmp_csv(
        file1,
        file2,
        ignore_case=False,
        numeric_tolerance=0,
        ignore_columns=None,
        verbose=True,
        **kwargs,
    ):
        import pandas as pd

        # Load data and fill NaNs
        df1 = pd.read_csv(file1).fillna("NA")
        df2 = pd.read_csv(file2).fillna("NA")

        # Standardize case if needed
        if ignore_case:
            df1.columns = df1.columns.str.lower()
            df2.columns = df2.columns.str.lower()
            df1 = df1.applymap(lambda x: x.lower() if isinstance(x, str) else x)
            df2 = df2.applymap(lambda x: x.lower() if isinstance(x, str) else x)

        # Drop ignored columns
        if ignore_columns:
            ignore_columns = [col.lower() if ignore_case else col for col in ignore_columns]
            df1.drop(columns=[col for col in ignore_columns if col in df1.columns], errors="ignore", inplace=True)
            df2.drop(columns=[col for col in ignore_columns if col in df2.columns], errors="ignore", inplace=True)

        # Reset index to ensure alignment
        df1.reset_index(drop=True, inplace=True)
        df2.reset_index(drop=True, inplace=True)

        # Align DataFrames by columns
        df1, df2 = df1.align(df2, join="inner", axis=1)

        # Compare rows
        added_rows = df2[~df2.apply(tuple, axis=1).isin(df1.apply(tuple, axis=1))]
        deleted_rows = df1[~df1.apply(tuple, axis=1).isin(df2.apply(tuple, axis=1))]

        # Compare modified rows
        if numeric_tolerance > 0:
            def numeric_diff(row):
                if pd.api.types.is_numeric_dtype(row["left"]):
                    return abs(row["left"] - row["right"]) > numeric_tolerance
                return row["left"] != row["right"]

            modified_rows = df1.compare(df2, keep_shape=True, keep_equal=False)
            modified_rows = modified_rows[modified_rows.apply(numeric_diff, axis=1)]
        else:
            modified_rows = df1.compare(df2, keep_shape=True, keep_equal=False)

        # Return results
        return {
            "added_rows": len(added_rows),
            "deleted_rows": len(deleted_rows),
            "modified_rows": len(modified_rows),
            "added_file": added_rows if not added_rows.empty else pd.DataFrame(),
            "deleted_file": deleted_rows if not deleted_rows.empty else pd.DataFrame(),
            "modified_file": modified_rows if not modified_rows.empty else pd.DataFrame(),
        }

    # --- Compare JSON files ---
    def cmp_json(
        file1, file2, ignore_case=False, numeric_tolerance=0, verbose=True, **kwargs
    ):
        import json

        with open(file1, "r") as f1:
            json1 = json.load(f1)
        with open(file2, "r") as f2:
            json2 = json.load(f2)

        # Normalize case and compare JSONs
        if ignore_case:

            def normalize(obj):
                if isinstance(obj, dict):
                    return {k.lower(): normalize(v) for k, v in obj.items()}
                elif isinstance(obj, list):
                    return [normalize(item) for item in obj]
                elif isinstance(obj, str):
                    return obj.lower()
                else:
                    return obj

            json1 = normalize(json1)
            json2 = normalize(json2)

        # Compare JSONs
        def compare_json(obj1, obj2):
            if isinstance(obj1, dict) and isinstance(obj2, dict):
                added_keys = {k: obj2[k] for k in obj2 if k not in obj1}
                deleted_keys = {k: obj1[k] for k in obj1 if k not in obj2}
                modified_keys = {
                    k: (obj1[k], obj2[k])
                    for k in obj1
                    if k in obj2 and obj1[k] != obj2[k]
                }
                return added_keys, deleted_keys, modified_keys

            elif isinstance(obj1, list) and isinstance(obj2, list):
                added_items = [item for item in obj2 if item not in obj1]
                deleted_items = [item for item in obj1 if item not in obj2]
                modified_items = [
                    (item1, item2) for item1, item2 in zip(obj1, obj2) if item1 != item2
                ]
                return added_items, deleted_items, modified_items

            else:
                if obj1 != obj2:
                    return obj1, obj2, None
                else:
                    return None, None, None

        added, deleted, modified = compare_json(json1, json2)

        return {"added_keys": added, "deleted_keys": deleted, "modified_keys": modified}

    # --- Compare Text files ---
    def cmp_txt(
        file1, file2, ignore_case=False, numeric_tolerance=0, verbose=True, **kwargs
    ):
        def read_lines(file):
            with open(file, "r") as f:
                return f.readlines()

        lines1 = read_lines(file1)
        lines2 = read_lines(file2)

        if ignore_case:
            lines1 = [line.lower() for line in lines1]
            lines2 = [line.lower() for line in lines2]

        added_lines = [line for line in lines2 if line not in lines1]
        deleted_lines = [line for line in lines1 if line not in lines2]

        modified_lines = []
        if numeric_tolerance > 0:
            for line1, line2 in zip(lines1, lines2):
                if abs(float(line1) - float(line2)) > numeric_tolerance:
                    modified_lines.append((line1, line2))
        else:
            for line1, line2 in zip(lines1, lines2):
                if line1 != line2:
                    modified_lines.append((line1, line2))

        return {
            "added_lines": added_lines,
            "deleted_lines": deleted_lines,
            "modified_lines": modified_lines,
        }

    if kind is None:
        kind = os.path.splitext(file1)[1].lower()[1:]
    # Compare based on the file type
    if kind == "xlsx":
        return cmp_excel(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "csv":
        return cmp_csv(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "json":
        return cmp_json(file1=file1, file2=file2, verbose=verbose, **kwargs)

    elif kind == "txt":
        return cmp_txt(file1=file1, file2=file2, verbose=verbose, **kwargs)

    else:
        raise ValueError(f"Unsupported file type: {kind}")
def cn2pinyin(
    cn_str: Union[str, list] = None,
    sep: str = " ",
    fmt: str = "normal",  # which style you want to set
):
    from pypinyin import pinyin, Style

    """
    Converts Chinese characters to Pinyin.
    usage: 
        cn2pinyin(cn_str, sep="_", fmt="tone")
    Args:
        cn_str (str): Chinese string to convert.
        sep (str): Separator for the output Pinyin string.
        fmt (Style): "normal","tone", "tone2","tone3","finals","finals_tone","finals_tone2","finals_tone3","initials","bopomofo","bopomofo_first","cyrillic","pl",
    Returns:
        cn_str: The Pinyin representation of the Chinese string.
    """
    fmts = [
        "normal",
        "tone",
        "tone2",
        "tone3",
        "finals",
        "finals_tone",
        "finals_tone2",
        "finals_tone3",
        "initials",
        "bopomofo",
        "bopomofo_first",
        "cyrillic",
        "pl",
    ]
    fmt = strcmp(fmt, fmts)[0]
    if fmt == "normal":
        style = Style.NORMAL
    elif fmt == "tone":
        style = Style.TONE
    elif fmt == "tone2":
        style = Style.TONE2
    elif fmt == "tone3":
        style = Style.TONE3
    elif fmt == "finals":
        style = Style.FINALS
    elif fmt == "finals_tone":
        style = Style.FINALS_TONE
    elif fmt == "finals_tone2":
        style = Style.FINALS_TONE2
    elif fmt == "finals_tone3":
        style = Style.FINALS_TONE3
    elif fmt == "initials":
        style = Style.INITIALS
    elif fmt == "bopomofo":
        style = Style.BOPOMOFO
    elif fmt == "bopomofo_first":
        style = Style.BOPOMOFO_FIRST
    elif fmt == "cyrillic":
        style = Style.CYRILLIC
    elif fmt == "pl":
        style = Style.PL
    else:
        style = Style.NORMAL
    if not isinstance(cn_str, list):
        cn_str = [cn_str]
    pinyin_flat = []
    for cn_str_ in cn_str:
        pinyin_string = pinyin(cn_str_, style=style)
        pinyin_flat.append(sep.join([item[0] for item in pinyin_string]))
    if len(pinyin_flat) == 1:
        return pinyin_flat[0]
    else:
        return pinyin_flat


def counter(list_, verbose=True):
    from collections import Counter

    c = Counter(list_)
    # Print the name counts
    for item, count in c.items():
        if verbose:
            print(f"{item}: {count}")
    return c


# usage:
# print(f"Return an iterator over elements repeating each as many times as its count:\n{sorted(c.elements())}")
# print(f"Return a list of the n most common elements:\n{c.most_common()}")
# print(f"Compute the sum of the counts:\n{c.total()}")

def dict2df(dict_, fill=None, axis=0):
    """
    Convert a dictionary to a DataFrame with flexible axis and padding options.

    Parameters:
    - dict_: The dictionary to convert (keys are columns or index).
    - fill: Value to fill in case of shorter lists.
    - axis: Axis for DataFrame construction (0 for columns, 1 for rows).

    Returns:
    - DataFrame created from the dictionary.
    """
    for key, value in dict_.items():
        if not isinstance(value, list):
            dict_[key] = [value]
            print(f"'{key}' is not a list. trying to convert it to 'list'")

    # Get the maximum length of values
    len_max = max(len(value) for value in dict_.values())

    # Extend lists to match the length of the longest list
    for key, value in dict_.items():
        if isinstance(value, list):
            value.extend([fill] * (len_max - len(value)))  # Fill shorter lists
        dict_[key] = value

    # If axis=0, the dictionary keys will be treated as column names
    if axis == 0:
        return pd.DataFrame(dict_)
    # If axis=1, the dictionary keys will be treated as index names (rows)
    else:
        return pd.DataFrame(dict_).transpose()

def text2audio(
    text,
    method=None,  # "pyttsx3","gTTS"
    rate=200,
    slow=False,  # "gTTS"
    volume=1.0,
    voice=None,
    lang=None,
    gender=None,
    age=None,
    dir_save=None,
):
    """
    # sample_text = "Hello! This is a test of the pyttsx3 text-to-speech system."
    # sample_text = "这个是中文, 测试"
    # sample_text = "Hallo, ich bin echo, Wie Heissen Sie"

    # text2audio(
    #     text=sample_text,
    #     rate=150,
    #     volume=0.9,
    #     # voice=None,  # Replace with a voice name or ID available on your system
    # )
    """
    if method is not None:
        methods = ["gTTS", "pyttsx3", "google"]
        method = strcmp(method, methods)[0]
    else:
        try:
            text2audio(
                text,
                method="google",
                rate=rate,
                slow=slow,
                volume=volume,
                voice=voice,
                lang=lang,
                gender=gender,
                age=age,
                dir_save=dir_save,
            )
        except Exception as e:
            print(e)
            text2audio(
                text,
                method="pyttsx3",
                rate=rate,
                slow=slow,
                volume=volume,
                voice=voice,
                lang=lang,
                gender=gender,
                age=age,
                dir_save=dir_save,
            )

    if method == "pyttsx3":
        import pyttsx3

        try:
            engine = pyttsx3.init()
            engine.setProperty("rate", rate)
            if 0.0 <= volume <= 1.0:
                engine.setProperty("volume", volume)
            else:
                raise ValueError("Volume must be between 0.0 and 1.0")

            if gender is not None:
                gender = strcmp(gender, ["male", "female"])[0]
            if age is not None:
                if isinstance(age, (float, int)):
                    if age <= 10:
                        age = "child"
                    elif 10 < age < 18:
                        age = "senior"
                    else:
                        age = "adult"
                elif isinstance(age, str):
                    age = strcmp(age, ["child", "adult", "senior"])[0]
                else:
                    raise ValueError("age: should be in ['child', 'adult', 'senior']")
            voices = engine.getProperty("voices")
            if voice is None:
                if lang is None:
                    voice = strcmp(detect_lang(text), [v.name for v in voices])[0]
                else:
                    if run_once_within():
                        print([v.name for v in voices])
                    print(f"lang:{lang}")
                    voice = strcmp(lang, [v.name for v in voices])[0]
            selected_voice = None

            for v in voices:
                # Check if the voice matches the specified gender or age
                if voice and (voice.lower() in v.name.lower() or voice in v.id):
                    selected_voice = v
                    break
                if gender and gender.lower() in v.name.lower():
                    selected_voice = v
                if age and age.lower() in v.name.lower():
                    selected_voice = v

            if selected_voice:
                engine.setProperty("voice", selected_voice.id)
            else:
                if voice or gender or age:
                    raise ValueError(
                        f"No matching voice found for specified criteria. Available voices: {[v.name for v in voices]}"
                    )
            # Generate audio
            if dir_save:
                engine.save_to_file(text, dir_save)
                print(f"Audio saved to {dir_save}")
            else:
                engine.say(text)

            engine.runAndWait()
        except Exception as e:
            print(f"An error occurred: {e}")
            #     # Explicitly terminate the pyttsx3 engine to release resources
            try:
                engine.stop()
            except RuntimeError:
                pass
                # Safely exit the script if running interactively to avoid kernel restarts
            try:
                import sys

                sys.exit()
            except SystemExit:
                pass
    elif method in ["google", "gtts"]:
        from gtts import gTTS

        try:
            if lang is None:
                from langdetect import detect

                lang = detect(text)
            # Initialize gTTS with the provided parameters
            tts = gTTS(text=text, lang=lang, slow=slow)
        except Exception as e:
            print(f"An error occurred: {e}")

        print("not realtime reading...")
        if dir_save:
            if "." not in dir_save:
                dir_save = dir_save + ".mp3"
            tts.save(dir_save)
            print(f"Audio saved to {dir_save}")
        else:
            dir_save = "temp_audio.mp3"
            if "." not in dir_save:
                dir_save = dir_save + ".mp3"
            tts.save(dir_save)
        try:
            fopen(dir_save)
        except Exception as e:
            print(f"Error opening file: {e}")
    print("done")

# from datetime import datetime

def str2time(
    time_str: str, 
    fmt: str = "24", 
    return_obj: bool = False,
    raise_errors: bool = True,
    default: Optional[Any] = None
) -> Union[str, datetime]:
    """
    Convert a time string into the specified format with enhanced parsing capabilities.
    
    Parameters:
    - time_str (str): The time string to be converted.
    - fmt (str): The format to convert the time to. Can be:
        - '12' or '12h' for 12-hour format (02:30:45 PM)
        - '24' or '24h' for 24-hour format (14:30:45)
        - Any valid strftime format string
    - return_obj (bool): If True, returns datetime.time object instead of string
    - raise_errors (bool): If False, returns default value on error instead of raising
    - default: Value to return on error when raise_errors=False
    
    Returns:
    - str/datetime.time: The converted time string or time object
    
    Supported Input Formats:
    - 14:30:45
    - 02:30:45 PM
    - 2:30 PM
    - 2PM
    - 1430 (military time)
    - 2.30.45
    - 14-30-45
    - and many more
    
    Example Usage:
    str2time("14:30:45", fmt='12')  # returns '02:30:45 PM'
    str2time("02:30:45 PM", fmt='24')  # returns '14:30:45'
    str2time("2PM", return_obj=True)  # returns datetime.time(14, 0)
    """
    from dateutil import parser
    def preprocess_time_string(ts: str) -> str:
        """Normalize various time string formats to a parseable format"""
        if not isinstance(ts, str):
            if raise_errors:
                raise ValueError(f"Expected string, got {type(ts)}")
            return default
        
        # Clean the string
        ts = ts.strip().lower()
        
        # Handle military time (e.g., "1430" -> "14:30")
        if re.fullmatch(r'^[0-2]?\d[0-5]\d$', ts):
            ts = f"{ts[:2]}:{ts[2:]}"
        
        # Handle formats without separators (e.g., "2pm" -> "2 pm")
        ts = re.sub(r'([0-9])([a-z])', r'\1 \2', ts)
        
        # Replace various separators with colons
        ts = re.sub(r'[\.\-_]', ':', ts)
        
        # Handle AM/PM without space (e.g., "2:30pm" -> "2:30 pm")
        ts = re.sub(r'([0-9])(am|pm)', r'\1 \2', ts, flags=re.IGNORECASE)
        
        # Add seconds if missing (e.g., "14:30" -> "14:30:00")
        if re.fullmatch(r'^\d{1,2}:\d{2}$', ts):
            ts = f"{ts}:00"
        
        # Capitalize AM/PM for parsing
        ts = re.sub(r'am|pm', lambda x: x.group().upper(), ts)
        
        return ts.strip()
    
    # Determine output format
    fmt = fmt.lower()
    if fmt in ('12', '12h'):
        output_fmt = "%I:%M:%S %p"
    elif fmt in ('24', '24h'):
        output_fmt = "%H:%M:%S"
    else:
        output_fmt = fmt  # allow custom format strings
    
    try:
        processed_str = preprocess_time_string(time_str)
        
        # Try multiple parsing strategies
        for pattern in ["%H:%M:%S", "%I:%M:%S %p", "%H:%M", "%I:%M %p"]:
            try:
                time_obj = datetime.strptime(processed_str, pattern).time()
                break
            except ValueError:
                continue
        else:
            # Fallback to dateutil's parser for complex cases
            time_obj = parser.parse(processed_str).time()
        
        if return_obj:
            return time_obj
        return time_obj.strftime(output_fmt)
    
    except Exception as e:
        if raise_errors:
            raise ValueError(f"Unable to parse time string: '{time_str}'. Error: {e}")
        return default

def str2date(
    date_str: Union[str, int, float], 
    fmt: Optional[str] = "%Y-%m-%d",
    original_fmt: Optional[str] = None, 
    return_obj: bool = False,
    raise_errors: bool = True,
    default: Optional[Any] = None,
    **parser_kwargs
) -> Union[str, datetime, Dict[str, Any]]:
    """
    Convert a date string to the desired format with enhanced parsing capabilities.
    
    Parameters:
    - date_str: The input date string or numeric value (will be stringified)
    - original_fmt: The original format of the date string (optional)
    - fmt: The desired output format string. None returns datetime object
    - return_obj: If True, always returns datetime object (overrides fmt)
    - raise_errors: If False, returns default value on error instead of raising
    - default: Value to return on error when raise_errors=False
    - parser_kwargs: Additional kwargs for dateutil.parser.parse
    
    Returns:
    - str/datetime: The converted date string or datetime object
    
    Supported Input Formats:
    - 2023-05-15
    - 15/05/2023
    - May 15, 2023
    - 15 May 2023
    - 20230515
    - and many more
    
    Example Usage:
    str2date("15/05/2023", fmt="%Y-%m-%d")  # returns '2023-05-15'
    str2date("May 15, 2023", return_obj=True)  # returns datetime object
    str2date("2023-05-15", fmt="%d.%m.%Y")  # returns '15.05.2023'
    """
    from dateutil import parser
    try:
        if date_str is None:
            return None
        if not isinstance(date_str, str):
            date_str = str(date_str).strip()
        
        # Clean the string
        date_str = re.sub(r'[\s_]+', ' ', date_str.strip())
        
        # Try parsing with original format first
        if original_fmt:
            try:
                date_obj = datetime.strptime(date_str, original_fmt)
                if return_obj:
                    return date_obj
                return date_obj.strftime(fmt) if fmt else date_obj
            except ValueError:

                
                common_formats = [
                    "%Y%m%d",  # 20230515
                    "%d%m%Y",  # 15052023
                    "%m%d%Y",  # 05152023
                    "%Y-%m-%d",
                    "%d-%m-%Y",
                    "%m/%d/%Y",
                    "%d/%m/%Y",
                    "%b %d, %Y",  # May 15, 2023
                    "%d %b %Y",  # 15 May 2023
                ]
                
                for original_fmt_ in common_formats:
                    try:
                        print(f"original_fmt:'original_fmt_'")
                        date_obj = datetime.strptime(date_str, original_fmt_)
                        if return_obj:
                            return date_obj
                        return date_obj.strftime(fmt) if fmt else date_obj
                    except ValueError:
                        continue

                # if raise_errors and not parser_kwargs.get('fuzzy', False):
                #     raise
        
        # Fall back to dateutil's parser
        try:
            date_obj = parser.parse(date_str, **parser_kwargs)
            if return_obj:
                return date_obj
            return date_obj.strftime(fmt) if fmt else date_obj
        except parser.ParserError as e:
            # Try some common alternative formats
            common_formats = [
                "%Y%m%d",  # 20230515
                "%d%m%Y",  # 15052023
                "%m%d%Y",  # 05152023
                "%Y-%m-%d",
                "%d-%m-%Y",
                "%m/%d/%Y",
                "%d/%m/%Y",
                "%b %d, %Y",  # May 15, 2023
                "%d %b %Y",  # 15 May 2023
            ]
            
            for date_format in common_formats:
                try:
                    date_obj = datetime.strptime(date_str, date_format)
                    if return_obj:
                        return date_obj
                    return date_obj.strftime(fmt) if fmt else date_obj
                except ValueError:
                    continue
            
            if raise_errors:
                raise ValueError(f"Unable to parse date string: '{date_str}'. Error: {e}")
            return default
    
    except Exception as e:
        if raise_errors:
            raise ValueError(f"Unable to process date string: '{date_str}'. Error: {e}")
        return default


def datetime_parser(
    datetime_str: str,
    time_fmt: Optional[str] = None,
    date_fmt: Optional[str] = None,
    return_obj: bool = False,
    raise_errors: bool = True,
    default: Optional[Any] = None,
    **parser_kwargs
) -> Union[str, datetime, Dict[str, Any]]:
    """
    Parse a combined date and time string with flexible format handling.
    
    Parameters:
    - datetime_str: String containing both date and time
    - time_fmt: Format for time portion (uses str2time rules)
    - date_fmt: Format for date portion (uses str2date rules)
    - return_obj: If True, returns datetime object
    - raise_errors: If False, returns default on error
    - default: Value to return on error when raise_errors=False
    - parser_kwargs: Passed to dateutil.parser
    
    Returns:
    - Parsed datetime string or object
    
    Example Usage:
    >>> datetime_parser("May 15, 2023 2:30 PM")
    '2023-05-15 14:30:00'
    >>> datetime_parser("15.05.2023 14:30", date_fmt="%d.%m.%Y", time_fmt="12")
    '15.05.2023 02:30:00 PM'
    """
    try:
        # First try parsing as complete datetime
        dt_obj = parser.parse(datetime_str, **parser_kwargs)
        
        # Apply formatting if needed
        if return_obj:
            return dt_obj
        
        date_part = str2date(dt_obj.date(), fmt=date_fmt or "%Y-%m-%d", return_obj=False)
        time_part = str2time(dt_obj.time(), fmt=time_fmt or "24", return_obj=False)
        
        return f"{date_part} {time_part}"
    
    except Exception as e:
        if raise_errors:
            raise ValueError(f"Unable to parse datetime string: '{datetime_str}'. Error: {e}")
        return default



def datetime2str(
    dt_input: Union[datetime, date, time, str, int, float],
    fmt: str = "auto",
    target: str = "both",
    raise_errors: bool = True,
    default: Optional[Any] = None,
    **format_kwargs,
) -> str:
    """
    Convert datetime/date/time objects or strings to formatted strings with extensive options.

    Parameters:
    - dt_input: Input to format (datetime, date, time, or parsable string/timestamp)
    - fmt: Output format specification:
        - For dates:
            - 'iso' = ISO8601 (YYYY-MM-DD)
            - 'sql' = YYYY-MM-DD
            - 'euro' = DD.MM.YYYY
            - 'us' = MM/DD/YYYY
            - 'full' = Weekday, Month Day, Year (e.g., "Monday, January 01, 2023")
            - Any strftime format string
        - For times:
            - 'iso' = ISO8601 (HH:MM:SS)
            - '12'/'12h' = 12-hour with AM/PM
            - '24'/'24h' = 24-hour format
            - 'sql' = HH:MM:SS
            - 'short' = HH:MM
            - Any strftime format string
        - 'auto' (default) = smart formatting based on input type
    - target: What to format ('date', 'time', or 'both')
    - raise_errors: If False, returns default value on error
    - default: Value to return on error when raise_errors=False
    - format_kwargs: Additional formatting options:
        - month_format: 'full' (January), 'short' (Jan), or 'number' (01)
        - weekday_format: 'full' (Monday), 'short' (Mon), or None
        - ampm: 'lower' (am/pm), 'upper' (AM/PM), or 'title' (Am/Pm)
        - precision: for times - 'seconds', 'minutes', 'hours', or 'milliseconds'

    Returns:
    - Formatted string representation of the input

    Example Usage:
    now = datetime.now()
    datetime2str(now, fmt='full')  # "Monday, January 01, 2023 02:30:45 PM"
    datetime2str(now.date(), fmt='euro')  # "01.01.2023"
    datetime2str("14:30:45", fmt='12', target='time')  # "02:30:45 PM"
    datetime2str(1672539045, fmt='us')  # "01/01/2023" (from timestamp)
    """

    def _resolve_format(fmt: str, obj_type: type) -> str:
        """Determine the appropriate format string based on shorthand options"""
        fmt = fmt.lower()

        # Date formats
        if obj_type in (datetime, date):
            if fmt in ("auto", "iso", "sql"):
                return "%Y-%m-%d"
            elif fmt == "euro":
                return "%d.%m.%Y"
            elif fmt == "us":
                return "%m/%d/%Y"
            elif fmt == "full":
                return "%A, %B %d, %Y"
            elif fmt == "short":
                return "%b %d, %Y"

        # Time formats
        elif obj_type == time:
            if fmt in ("auto", "iso", "sql"):
                return "%H:%M:%S"
            elif fmt in ("12", "12h"):
                return "%I:%M:%S %p"
            elif fmt in ("24", "24h"):
                return "%H:%M:%S"
            elif fmt == "short":
                return "%H:%M"

        return fmt  # Return as-is if not a recognized shorthand

    def _apply_format_options(format_str: str, obj_type: type, options: dict) -> str:
        """Modify format string based on additional options"""
        # Handle month formatting
        if obj_type in (datetime, date):
            if options.get("month_format") == "short":
                format_str = format_str.replace("%B", "%b")
            elif options.get("month_format") == "number":
                format_str = format_str.replace("%B", "%m").replace("%b", "%m")

            # Handle weekday formatting
            if options.get("weekday_format") == "short":
                format_str = format_str.replace("%A", "%a")
            elif options.get("weekday_format") is None:
                format_str = re.sub(r"%[Aa],?\s*", "", format_str)

        # Handle AM/PM formatting
        if obj_type in (datetime, time) and "%p" in format_str:
            ampm = options.get("ampm", "upper")
            if ampm == "lower":
                format_str = format_str.replace("%p", "%p").lower()
            elif ampm == "title":
                format_str = format_str.replace("%p", "%p")

        # Handle time precision
        if obj_type in (datetime, time):
            precision = options.get("precision")
            if precision == "minutes":
                format_str = re.sub(r":%S(\s*%p)?$", r"\1", format_str)
            elif precision == "hours":
                format_str = re.sub(r":%M:%S(\s*%p)?$", r"\1", format_str)
            elif precision == "milliseconds":
                format_str = (
                    format_str.replace("%S", "%S.%f")
                    if "%S" in format_str
                    else format_str + ".%f"
                )

        return format_str.strip()

    try:
        # Convert input to appropriate type if it's a string or timestamp
        if isinstance(dt_input, (int, float)):
            dt_input = datetime.fromtimestamp(dt_input)
        elif isinstance(dt_input, str):
            try:
                # Try parsing as datetime first
                dt_input = parser.parse(dt_input)
            except (ValueError, TypeError):
                try:
                    # Try parsing as time
                    dt_input = parser.parse(dt_input).time()
                except (ValueError, TypeError):
                    try:
                        # Try parsing as date
                        dt_input = parser.parse(dt_input).date()
                    except (ValueError, TypeError) as e:
                        if raise_errors:
                            raise ValueError(
                                f"Could not parse input string: {dt_input}"
                            ) from e
                        return default

        # Determine what we're formatting
        if target == "auto":
            if isinstance(dt_input, time):
                target = "time"
            elif isinstance(dt_input, date) and not isinstance(dt_input, datetime):
                target = "date"
            else:
                target = "both"

        # Get the appropriate format string
        if fmt == "auto":
            if target == "time":
                fmt = "24h"
            elif target == "date":
                fmt = "iso"
            else:
                fmt = "iso"  # Will be combined later

        format_str = _resolve_format(fmt, type(dt_input))
        format_str = _apply_format_options(format_str, type(dt_input), format_kwargs)

        # Handle different target types
        if target == "date" and isinstance(dt_input, (datetime, date)):
            return dt_input.strftime(format_str)
        elif target == "time" and isinstance(dt_input, (datetime, time)):
            if isinstance(dt_input, datetime):
                return dt_input.time().strftime(format_str)
            return dt_input.strftime(format_str)
        elif target == "both" and isinstance(dt_input, datetime):
            date_part = datetime2str(
                dt_input.date(), fmt=fmt, target="date", **format_kwargs
            )
            time_part = datetime2str(
                dt_input.time(), fmt=fmt, target="time", **format_kwargs
            )
            return f"{date_part} {time_part}"
        else:
            if raise_errors:
                raise ValueError(f"Cannot format {type(dt_input)} as {target}")
            return default

    except Exception as e:
        if raise_errors:
            raise ValueError(f"Failed to format {dt_input}: {str(e)}")
        return default


# Additional convenience functions
def date2str(
    date_input: Union[date, datetime, str, int, float], fmt: str = "iso", **kwargs
) -> str:
    """Convert date objects or date strings to formatted date strings"""
    return datetime2str(date_input, fmt=fmt, target="date", **kwargs)


def time2str(
    time_input: Union[time, datetime, str, int, float], fmt: str = "24h", **kwargs
) -> str:
    """Convert time objects or time strings to formatted time strings"""
    return datetime2str(time_input, fmt=fmt, target="time", **kwargs)

def str2num(
    s: str,
    *args,
    sep: Optional[Union[str, List[str]]] = None,
    round_digits: Optional[int] = None,
    return_list: bool = True,
    handle_text: bool = True
) -> Union[float, int, List[Union[float, int]], None]:
    """
# Examples
print(str2num("123"))                # Output: 123
print(str2num("123.456", 2))         # Output: 123.46
print(str2num("one hundred and twenty three"))  # Output: 123
print(str2num("seven million"))      # Output: 7000000
print(str2num('one thousand thirty one',','))  # Output: 1,031
print(str2num("12345.6789", ","))    # Output: 12,345.6789
print(str2num("12345.6789", " ", 2)) # Output: 12 345.68
print(str2num('111113.34555',3,',')) # Output: 111,113.346
print(str2num("123.55555 sec miniuets",3)) # Output: 1.3
print(str2num("every 3,300.55 hours and 5.045555 min", sep=",", round=1))
print(str2num("five hundred fourty one"), str2num(
    "this is 5.9435 euros for 10.04499 killograme", round=3
)[0])
    Convert a string containing numeric or textual data into an integer, float, or list of numbers.

    Parameters:
    - s (str): Input string containing a number or textual representation of a number.
    - *args: Additional arguments for delimiter or rounding digits.
    - sep (str or list): Delimiter(s) to remove from the string (e.g., ',' or ['.', ',']).
    - round_digits (int): Number of decimal places to round the result to.
    - return_list (bool): Whether to return a list of numbers if multiple are found.
    - handle_text (bool): Whether to process textual numbers using the numerizer library.

    Returns:
    - Union[float, int, List[Union[float, int]], None]: Converted number(s) or None if conversion fails.
    """
    import re
    from numerizer import numerize
    if isinstance(s,list):
        return [str2num(i, sep=sep,round_digits=round_digits,return_list=return_list,handle_text=handle_text) for i in s]
    elif isinstance(s, str):
        pass
    else:
        return None

    # Merge args with explicit parameters
    if sep is None:
        sep = []
    elif isinstance(sep, str):
        sep = [sep]
    for arg in args:
        if isinstance(arg, str):
            sep.append(arg)
        elif isinstance(arg, int) and round_digits is None:
            round_digits = arg

    # Remove all specified delimiters
    for delimiter in sep:
        s = s.replace(delimiter, "")

    # Attempt conversion
    def try_convert(segment: str) -> Union[float, int, None]:
        try:
            return int(segment)
        except ValueError:
            try:
                return float(segment)
            except ValueError:
                return None

    # Handle textual numbers
    if handle_text:
        try:
            s = numerize(s)
        except Exception:
            pass

    # Extract numeric segments
    number_segments = re.findall(r"[-+]?\d*\.\d+|\d+", s)
    numbers = [try_convert(seg) for seg in number_segments if seg]
    numbers = [num for num in numbers if num is not None]

    if not numbers:
        return None  # No valid numbers found

    # Single or multiple numbers
    if len(numbers) == 1 and not return_list:
        result = numbers[0]
    else:
        result = (
            numbers[0] if len(numbers) == 1 else numbers if return_list else numbers[0]
        )

    # Apply rounding if necessary
    if round_digits is not None:
        if isinstance(result, list):
            result = [round(num + 1e-10, round_digits) for num in result]
        else:
            result = round(result + 1e-10, round_digits)

        # Convert to int if rounding to 0 digits
        if round_digits == 0:
            if isinstance(result, list):
                result = [int(num) for num in result]
            else:
                result = int(result)
    return result

# def num2str(num, *args, **kwargs):
#     delimiter = kwargs.get("sep", None)
#     round_digits = kwargs.get("round", None)

#     # Parse additional arguments
#     for arg in args:
#         if isinstance(arg, str):
#             delimiter = arg
#         elif isinstance(arg, int):
#             round_digits = arg

#     # Apply rounding if specified
#     if round_digits is not None:
#         num = round(num, round_digits)

#     # Convert number to string
#     num_str = f"{num}"

#     # Apply delimiter if specified
#     if delimiter is not None:
#         num_str = num_str.replace(".", ",")  # Replace decimal point with comma
#         num_str_parts = num_str.split(",")
#         if len(num_str_parts) > 1:
#             integer_part = num_str_parts[0]
#             decimal_part = num_str_parts[1]
#             integer_part = "{:,}".format(int(integer_part))
#             num_str = integer_part + "." + decimal_part
#         else:
#             num_str = "{:,}".format(int(num_str_parts[0]))

#     return num_str 
def num2str(
    num: Union[int, float],
    *args,
    decimal_sep: str = ".",
    thousand_sep: str = ",",
    round_digits: Optional[int] = None,
    interval:int=3,
    force_decimal: bool = False,
    unit: Optional[str] = None,
    **kwargs,
) -> str:
    """
    Convert a number to a formatted string with customizable separators.

    Parameters:
    - num: Number to convert (int or float)
    - decimal_sep: Decimal point character (default ".")
    - thousand_sep: Thousands separator character (default ",")
    - round_digits: Number of decimal places to round to
    - force_decimal: Always show decimal places even if zero (e.g., "123.00")
    - Additional args/kwargs for backward compatibility:
      - sep: Alternate way to specify thousand_sep
      - round: Alternate way to specify round_digits

    Returns:
    - Formatted number string
    """
    # Backward compatibility with old sep/round parameters
    if "sep" in kwargs:
        thousand_sep = kwargs["sep"]
    if "round" in kwargs:
        round_digits = kwargs["round"]

    # Parse additional positional arguments (legacy support)
    for arg in args:
        if isinstance(arg, str):
            thousand_sep = arg
        elif isinstance(arg, int):
            round_digits = arg

    # Apply rounding if specified
    if round_digits is not None:
        num = round(num, round_digits)

    # Handle special cases
    if isinstance(num, float):
        if num == float("inf"):
            return "∞"
        if num == float("-inf"):
            return "-∞"
        if num != num:  # NaN check
            return "NaN"

    # Initialize variables
    integer_part = ""
    decimal_part = ""

    # Split into integer and decimal parts
    if isinstance(num, int):
        integer_part = str(num)
    else:
        num_str = (
            f"{num:.10f}".rstrip("0").rstrip(".")
            if not force_decimal
            else f"{num:.10f}"
        )
        if "." in num_str:
            integer_part, decimal_part = num_str.split(".")
            if round_digits is not None:
                decimal_part = decimal_part[:round_digits]
        else:
            integer_part = num_str

    # Add thousand separators to integer part
    if thousand_sep:
        sign = "-" if integer_part.startswith("-") else ""
        digits = integer_part.lstrip("-")

        if digits:  # Only process if we have digits
            reversed_digits = digits[::-1]
            chunks = [
                reversed_digits[i : i + interval] for i in range(0, len(reversed_digits), interval)
            ]
            formatted_reversed = thousand_sep.join(chunks)
            integer_part = sign + formatted_reversed[::-1]

    # Combine parts
    result = ""
    if decimal_part or (force_decimal and round_digits is not None and round_digits > 0):
        if force_decimal and round_digits is not None:
            decimal_part = decimal_part.ljust(round_digits, '0')[:round_digits]
        result = f"{integer_part}{decimal_sep}{decimal_part}"
    else:
        result = integer_part

    # Append unit if specified
    if unit:
        result = f"{result} {unit}"

    return result
# Examples
# print(num2str(123), type(num2str(123)))  # Output: "123"
# print(num2str(123.456, round=2), type(num2str(123.456, 2)))  # Output: "123.46"
# print(num2str(7000.125, round=1), type(num2str(7000.125, 2)))  # Output: "7000.13"
# print(
#     num2str(12345333.6789, sep=","), type(num2str(12345.6789, ","))
# )  # Output: "12,345.6789"
# print(num2str(7000.00, sep=","), type(num2str(7000.00, ",")))  # Output: "7,000.00"


# Helper to convert text or list of text to HTML
def str2html(text_list, strict=False):
    if not isinstance(text_list, list):
        text_list = [text_list]

    # Define a mapping for escape sequences to their HTML representations
    escape_mapping = {
        "\\": "&bsol;",  # Backslash
        "'": "&apos;",  # Single quote
        '"': "&quot;",  # Double quote
        "\n": "<br>",  # Newline
        "\r": "",  # Carriage return (not represented in HTML)
        "\t": "&nbsp;&nbsp;&nbsp;&nbsp;",  # Tab (4 spaces)
        "\b": "",  # Backspace (not typically represented)
        "\f": "",  # Form feed (not typically represented)
        "\v": "",  # Vertical tab (not typically represented)
        "\a": "",  # Bell/alert sound (not typically represented)
        "\0": "",  # Null character (not typically represented)
    }

    # Convert text by replacing escape sequences with their HTML equivalents
    html_content = ""
    for text in text_list:
        for escape_seq, html_rep in escape_mapping.items():
            text = text.replace(escape_seq, html_rep)
        html_content += text.replace("\n", "<br>")  # Add line breaks for newlines

    if strict:
        html_content = "<html><body>\n" + html_content + "\n</body></html>"

    return html_content


def cm2px(*cm, dpi=300) -> list:
    # Case 1: When the user passes a single argument that is a list or tuple, such as cm2px([8, 5]) or inch2cm((8, 5))
    if len(cm) == 1 and isinstance(cm[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i / 2.54 * dpi for i in cm[0]]
    # Case 2: When the user passes multiple arguments directly, such as cm2px(8, 5)
    else:
        return [i / 2.54 * dpi for i in cm]


def px2cm(*px, dpi=300) -> list:
    # Case 1: When the user passes a single argument that is a list or tuple, such as px2cm([8, 5]) or inch2cm((8, 5))
    if len(px) == 1 and isinstance(px[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i * 2.54 / dpi for i in px[0]]
    # Case 2: When the user passes multiple arguments directly, such as px2cm(8, 5)
    else:
        # Here, we convert each individual argument directly to cm
        return [i * 2.54 / dpi for i in px]


def px2inch(*px, dpi=300) -> list:
    """
    px2inch: converts pixel measurements to inches based on the given dpi.
    Usage:
    px2inch(300, 600, dpi=300); px2inch([300, 600], dpi=300)
    Returns:
        list: in inches
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as px2inch([300, 600]) or px2inch((300, 600))
    if len(px) == 1 and isinstance(px[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to inches
        return [i / dpi for i in px[0]]
    # Case 2: When the user passes multiple arguments directly, such as px2inch(300, 600)
    else:
        # Here, we convert each individual argument directly to inches
        return [i / dpi for i in px]


def inch2cm(*cm) -> list:
    """
    cm2inch: converts centimeter measurements to inches.
    Usage:
    cm2inch(10, 12.7); cm2inch((10, 12.7)); cm2inch([10, 12.7])
    Returns:
        list: in inches
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as cm2inch([10, 12.7]) or cm2inch((10, 12.7))
    if len(cm) == 1 and isinstance(cm[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to inches
        return [i * 2.54 for i in cm[0]]
    # Case 2: When the user passes multiple arguments directly, such as cm2inch(10, 12.7)
    else:
        # Here, we convert each individual argument directly to inches
        return [i * 2.54 for i in cm]


def inch2px(*inch, dpi=300) -> list:
    """
    inch2px: converts inch measurements to pixels based on the given dpi.

    Usage:
    inch2px(1, 2, dpi=300); inch2px([1, 2], dpi=300)

    Parameters:
    inch : float, list, or tuple
        Single or multiple measurements in inches to convert to pixels.
    dpi : int, optional (default=300)
        Dots per inch (DPI), representing the pixel density.

    Returns:
        list: Converted measurements in pixels.
    """
    # Case 1: When the user passes a single argument that is a list or tuple, e.g., inch2px([1, 2]) or inch2px((1, 2))
    if len(inch) == 1 and isinstance(inch[0], (list, tuple)):
        return [i * dpi for i in inch[0]]

    # Case 2: When the user passes multiple arguments directly, e.g., inch2px(1, 2)
    else:
        return [i * dpi for i in inch]


def cm2inch(*inch) -> list:
    """
    Usage:
    inch2cm(8,5); inch2cm((8,5)); inch2cm([8,5])
    Returns:
        list: in centimeters
    """
    # Case 1: When the user passes a single argument that is a list or tuple, such as inch2cm([8, 5]) or inch2cm((8, 5))
    if len(inch) == 1 and isinstance(inch[0], (list, tuple)):
        # If the input is a single list or tuple, we unpack its elements and convert each to cm
        return [i / 2.54 for i in inch[0]]
    # Case 2: When the user passes multiple arguments directly, such as inch2cm(8, 5)
    else:
        # Here, we convert each individual argument directly to cm
        return [i / 2.54 for i in inch]

def img2svg(fpath):
    """
    Convert an image file to SVG format using pixels2svg.

    e.g., img2svg(fpath)
    """
    from pathlib import Path
    from pixels2svg import pixels2svg
    fpath = Path(fpath)  # Ensure it's a Path object
    
    if not fpath.exists():
        raise FileNotFoundError(f"File not found: {fpath}")
    
    if not fpath.is_file():
        raise ValueError(f"Invalid file path: {fpath} is not a file")

    ftype = fpath.suffix.lstrip(".").lower()
    output_path = fpath.with_suffix(".svg")

    try:
        pixels2svg(fpath, output_path)
        print(f"saved @: {output_path}")
    except Exception as e:
        print(f"失败了!{e}")

    return output_path

def sqlite2sql(db_path, sql_path):
    """
    Export an SQLite database to an SQL file, including schema and data for all tables.

    :param db_path: Path to the SQLite .db file
    :param output_file: Path to the output .sql file

    # Usage
        db_path = "your_database.db"  # Replace with the path to your SQLite database
        sql_path = "output.sql"  # Replace with your desired output file name
        export_sqlite_to_sql(db_path, sql_path)

    """
    import sqlite3
    try:
        # Connect to the SQLite database
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()

        with open(sql_path, 'w') as f:
            # Write a header for the SQL dump
            f.write("-- SQLite Database Dump\n")
            f.write(f"-- Source: {db_path}\n\n")
            
            # Retrieve all table names
            cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%';")
            tables = [row[0] for row in cursor.fetchall()]

            for table in tables:
                # Write the schema for the table
                cursor.execute(f"SELECT sql FROM sqlite_master WHERE type='table' AND name='{table}';")
                schema = cursor.fetchone()
                if schema:
                    f.write(f"{schema[0]};\n\n")
                
                # Write data for the table
                cursor.execute(f"SELECT * FROM {table};")
                rows = cursor.fetchall()
                if rows:
                    cursor.execute(f"PRAGMA table_info({table});")
                    column_names = [info[1] for info in cursor.fetchall()]
                    column_list = ', '.join(f'"{col}"' for col in column_names)

                    for row in rows:
                        values = ', '.join(f"'{str(val).replace('\'', '\'\'')}'" if val is not None else 'NULL' for val in row)
                        f.write(f"INSERT INTO {table} ({column_list}) VALUES ({values});\n")
                
                f.write("\n")
            
            print(f"Database exported successfully to {sql_path}")
    
    except sqlite3.Error as e:
        print(f"SQLite error: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")
    finally:
        # Ensure the connection is closed
        if conn:
            conn.close()


def sreplace(*args, **kwargs):
    """
    sreplace(text, by=None, robust=True)
    Replace specified substrings in the input text with provided replacements.
    Args:
        text (str): The input text where replacements will be made.
        by (dict, optional): A dictionary containing substrings to be replaced as keys
            and their corresponding replacements as values. Defaults to {".com": "..come", "\n": " ", "\t": " ", "  ": " "}.
        robust (bool, optional): If True, additional default replacements for newline and tab characters will be applied.
                                Default is False.
    Returns:
        str: The text after replacements have been made.
    """
    text = None
    by = kwargs.get("by", None)
    robust = kwargs.get("robust", True)

    for arg in args:
        if isinstance(arg, str):
            text = arg
        elif isinstance(arg, dict):
            by = arg
        elif isinstance(arg, bool):
            robust = arg
        else:
            Error(f"{type(arg)} is not supported")

    # Default replacements for newline and tab characters
    default_replacements = {
        "\a": "",
        "\b": "",
        "\f": "",
        "\n": "",
        "\r": "",
        "\t": "",
        "\v": "",
        "\\": "",  # Corrected here
        # "\?": "",
        "�": "",
        "\\x": "",  # Corrected here
        "\\x hhhh": "",
        "\\ ooo": "",  # Corrected here
        "\xa0": "",
        "  ": " ",
    }

    # If dict_replace is None, use the default dictionary
    if by is None:
        by = {}
    # If robust is True, update the dictionary with default replacements
    if robust:
        by.update(default_replacements)

    # Iterate over each key-value pair in the dictionary and replace substrings accordingly
    for k, v in by.items():
        text = text.replace(k, v)
    return text


# usage:
# sreplace(text, by=dict(old_str='new_str'), robust=True)


def paper_size(paper_type_str="a4"):
    df = pd.DataFrame(
        {
            "a0": [841, 1189],
            "a1": [594, 841],
            "a2": [420, 594],
            "a3": [297, 420],
            "a4": [210, 297],
            "a5": [148, 210],
            "a6": [105, 148],
            "a7": [74, 105],
            "b0": [1028, 1456],
            "b1": [707, 1000],
            "b2": [514, 728],
            "b3": [364, 514],
            "b4": [257, 364],
            "b5": [182, 257],
            "b6": [128, 182],
            "letter": [215.9, 279.4],
            "legal": [215.9, 355.6],
            "business card": [85.6, 53.98],
            "photo china passport": [33, 48],
            "passport single": [125, 88],
            "visa": [105, 74],
            "sim": [25, 15],
        }
    )
    for name in df.columns:
        if paper_type_str in name.lower():
            paper_type = name
    if not paper_type:
        paper_type = "a4"  # default
    return df[paper_type].tolist()
 
def docx2pdf(dir_docx, dir_pdf=None):
    """
    Converts .docx to .pdf. Works on Windows using docx2pdf and on Linux/macOS using LibreOffice.
    
    Parameters:
    - dir_docx: path to .docx file or directory containing .docx files
    - dir_pdf: optional output directory; if None, uses same directory as input
    """

    system = platform.system()
    is_file = os.path.isfile(dir_docx)
    is_dir = os.path.isdir(dir_docx)

    if not is_file and not is_dir:
        raise FileNotFoundError(f"Input path '{dir_docx}' does not exist.")

    if system == "Windows":
        try:
            from docx2pdf import convert
        except ImportError:
            raise ImportError("docx2pdf is not installed. Run: pip install docx2pdf")

        convert(dir_docx, dir_pdf) if dir_pdf else convert(dir_docx)

    elif system in {"Linux", "Darwin"}:
        # Check if libreoffice is available
        if shutil.which("libreoffice") is None:
            raise EnvironmentError("LibreOffice is not installed or not in PATH. Install it with: sudo apt install libreoffice")

        # Determine the output directory
        output_dir = dir_pdf or os.path.dirname(dir_docx) if is_file else dir_docx

        if is_file:
            subprocess.run([
                "libreoffice", "--headless", "--convert-to", "pdf", "--outdir", output_dir, dir_docx
            ], check=True)
        elif is_dir:
            for filename in os.listdir(dir_docx):
                if filename.lower().endswith(".docx"):
                    full_path = os.path.join(dir_docx, filename)
                    subprocess.run([
                        "libreoffice", "--headless", "--convert-to", "pdf", "--outdir", dir_pdf or dir_docx, full_path
                    ], check=True)
    else:
        raise OSError(f"Unsupported OS: {system}")


def img2pdf(dir_img, kind=None, page=None, dir_save=None, page_size="a4", dpi=300):
    import img2pdf as image2pdf

    def mm_to_point(size):
        return (image2pdf.mm_to_pt(size[0]), image2pdf.mm_to_pt(size[1]))

    def set_dpi(x):
        dpix = dpiy = x
        return image2pdf.get_fixed_dpi_layout_fun((dpix, dpiy))

    if kind is None:
        _, kind = os.path.splitext(dir_img)
    if not kind.startswith("."):
        kind = "." + kind
    if dir_save is None:
        dir_save = dir_img.replace(kind, ".pdf")
    imgs = []
    if os.path.isdir(dir_img):
        if not dir_save.endswith(".pdf"):
            dir_save += "#merged_img2pdf.pdf"
        if page is None:
            select_range = listdir(dir_img, kind=kind).fpath
        else:
            if not isinstance(page, (np.ndarray, list, range)):
                page = [page]
            select_range = listdir(dir_img, kind=kind)["fpath"][page]
        for fname in select_range:
            if not fname.endswith(kind):
                continue
            path = os.path.join(dir_img, fname)
            if os.path.isdir(path):
                continue
            imgs.append(path)
    else:
        imgs = [
            # os.path.isdir(dir_img),
            dir_img
        ]
    print(imgs)
    if page_size:
        if isinstance(page_size, str):
            pdf_in_mm = mm_to_point(paper_size(page_size))
        else:
            print("default: page_size = (210,297)")
            pdf_in_mm = mm_to_point(page_size)
            print(f"page size was set to {page_size}")
        p_size = image2pdf.get_layout_fun(pdf_in_mm)
    else:
        p_size = set_dpi(dpi)
    with open(dir_save, "wb") as f:
        f.write(image2pdf.convert(imgs, layout_fun=p_size))


# usage:
# dir_img="/Users/macjianfeng/Dropbox/00-Personal/2015-History/2012-2015_兰州大学/120901-大学课件/生物统计学 陆卫/复习题/"
# img2pdf(dir_img,kind='tif', page=range(3,7,2))


def pdf2ppt(dir_pdf, dir_ppt):
    from PyPDF2 import PdfReader
    from pptx.util import Inches
    from pptx import Presentation

    prs = Presentation()

    # Open the PDF file
    with open(dir_pdf, "rb") as f:
        reader = PdfReader(f)
        num_pages = len(reader.pages)

        # Iterate through each page in the PDF
        for page_num in range(num_pages):
            page = reader.pages[page_num]
            text = page.extract_text()

            # Add a slide for each page's content
            slide_layout = prs.slide_layouts[
                5
            ]  # Use slide layout that suits your needs
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"Page {page_num + 1}"
            slide.shapes.add_textbox(
                Inches(1), Inches(1.5), Inches(8), Inches(5)
            ).text = text

    # Save the PowerPoint presentation
    prs.save(dir_ppt)
    print(f"Conversion from {dir_pdf} to {dir_ppt} complete.")


def ssplit(text, by="space", verbose: bool =False, strict: bool =False, strip_results: bool = True, **kws):
    """# Determines the splitting strategy:
    #         - "space", "whitespace", "sp": split by whitespace (default)
    #         - "word": split into words using NLTK's word_tokenize
    #         - "sentence", "sent": split into sentences using NLTK's sent_tokenize
    #         - "sent_num", "sent_n": split every N sentences (use n=10 by default)
    #         - "char", "character": split into individual characters
    #         - "camel", "camelcase": split camelCase words
    #         - "upper", "capital": split at uppercase letters
    #         - "upper_lower", "ul": split at uppercase followed by lowercase
    #         - "lower_upper", "lu": split at lowercase followed by uppercase
    #         - "digits", "numbers", "num": split at numeric sequences
    #         - "number_words", "num_str": split at spelled-out numbers
    #         - "punctuation", "punct": split at punctuation marks
    #         - "lines", "line", "li": split by lines
    #         - "regex", "re": use custom regex pattern (provide in 'pattern' kwarg)
    #         - "fixed_len", "len": split into fixed-length chunks (provide 'length' kwarg)
    #         - "non_ascii", "lang": split at non-ASCII characters
    #         - "non_alphanum": split at consecutive non-alphanumeric chars
    #         - str: split by the exact string provided
    #         - list: split by any of the strings in the list

    Example:
    # Split by whitespace (default)
    print(ssplit("Hello world! How are you?", verbose=1))
    # Split by specific string
    print(ssplit("apple,orange,banana", by=",", verbose=1))
    # Split by multiple delimiters
    print(ssplit("apple,orange;banana-grape", by=[",", ";", "-"], verbose=1))
    print(ssplit("apple,orange,banana", by=",", verbose=1))
    print(
        ssplit(
            "1. First. 2. Second. 3. Third. 4. Fourth.",
            by="sent_num",
            n=2,
            strip_results=1,
            verbose=1,
        )
    )
    # Split into words (using NLTK)
    print(ssplit("Can't stop won't stop!", by="word"))
    # Split into sentences
    print(ssplit("First sentence. Second sentence! Third one?", by="sentence"))
    # Split every N sentences
    print(ssplit("1. First. 2. Second. 3. Third. 4. Fourth.", by="sent_num", n=2))
    # Split into characters
    print(ssplit("hello", by="char"))
    # Split by fixed length chunks
    print(ssplit("abcdefghijklmnop", by="len", length=4))
    # Split at uppercase letters
    print(ssplit("CamelCaseString", by="upper"))
    # Split camelCase words
    print(ssplit("camelCaseVariableName", by="camel"))
    # Split at numbers (your requested example)
    print(ssplit("Item1Price20Quantity5", by="digits"))
    # Split at number words
    print(ssplit("Order twenty-five items", by="number_words"))
    # Split by punctuation
    print(ssplit("Hello! How are you? I'm fine.", by="punctuation"))
    # Split by lines
    print(ssplit("Line 1\nLine 2\r\nLine 3", by="lines"))
    # Split at non-ASCII characters
    print(ssplit("English日本語Русский", by="non_ascii"))
    # Split at non-alphanumeric characters
    print(ssplit("word1@word2#word3$word4", by="non_alphanum"))
    # Split with regex pattern
    print(ssplit("a1b2c3d4", by="regex", pattern=r"\d"))
    # Keep delimiters
    print(ssplit("a,b;c.d", by=[",", ";", "."], keep_delimiters=True))
    # Case-insensitive splitting
    print(ssplit("Apple ORANGE banana", by="orange", ignore_case=True))
    # Limit number of splits
    print(ssplit("one two three four", by="space", maxsplit=2))
    # Empty string
    print(ssplit("", by="space"))
    # Only delimiters
    print(ssplit(",,,", by=","))
    # Mixed content
    print(ssplit("ID:123-Name:John Doe-Age:30", by=["ID:", "-Name:", "-Age:"]))
    # Process CSV-like string
    print(ssplit("Name;John|Age;30||City;New York", by=[";", "|"]))
    # Extract API endpoint parts
    print(ssplit("getUserByIdAndName", by="camel"))
    """
    import re

    if isinstance(text, list):
        nested_list = [ssplit(i, by=by, verbose=verbose, **kws) for i in text]
        flat_list = [item for sublist in nested_list for item in sublist]
        return flat_list

    def split_by_word_length(text, length):
        return [word for word in text.split() if len(word) == length]

    def split_by_multiple_delimiters(text, delimiters):
        regex_pattern = "|".join(map(re.escape, delimiters))
        return re.split(regex_pattern, text)

    def split_by_camel_case(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]*(?=[A-Z]|$))", text)

    def split_at_upper_fl_lower(text):
        return re.findall(r"[A-Z](?:[a-z]+|[A-Z]+(?=[A-Z]|$))", text)

    def split_at_lower_fl_upper(text):
        split_text = re.split(r"(?<=[a-z])(?=[A-Z])", text)
        return split_text

    def split_at_upper(text):
        split_text = re.split(r"(?=[A-Z])", text)
        split_text = [part for part in split_text if part]
        return split_text

    def split_by_regex_lookahead(text, pattern):
        return re.split(f"(?<={pattern})", text)

    def split_by_regex_end(text, pattern):
        return re.split(f"(?={pattern})", text)
    def split_non_ascii(text, keep_delimiters=False):
        """
        Split text at non-ASCII characters.
        
        Args:
            text: Input string to split
            keep_delimiters: If True, keeps the non-ASCII delimiters in the output
        
        Returns:
            List of strings split at non-ASCII characters
        """
        if keep_delimiters:
            # Split and keep non-ASCII delimiters
            parts = re.split(r"([^\x00-\x7F]+)", text)
            # Combine adjacent parts and delimiters
            result = []
            for i in range(0, len(parts)-1, 2):
                if parts[i]:
                    result.append(parts[i])
                if i+1 < len(parts) and parts[i+1]:
                    result.append(parts[i+1])
            if len(parts) % 2 == 1 and parts[-1]:
                result.append(parts[-1])
            return result
        else:
            # Split and discard non-ASCII delimiters
            return [part for part in re.split(r"[^\x00-\x7F]+", text) if part]
    def split_by_consecutive_non_alphanumeric(text):
        return re.split(r"\W+", text)

    def split_by_fixed_length_chunks(text, length):
        return [text[i : i + length] for i in range(0, len(text), length)]

    def split_by_sent_num(text, n=10):
        from nltk.tokenize import sent_tokenize
        from itertools import pairwise

        # split text into sentences
        text_split_by_sent = sent_tokenize(text)
        cut_loc_array = np.arange(0, len(text_split_by_sent), n)
        if cut_loc_array[-1] != len(text_split_by_sent):
            cut_loc = np.append(cut_loc_array, len(text_split_by_sent))
        else:
            cut_loc = cut_loc_array
        # get text in section (e.g., every 10 sentences)
        text_section = []
        for i, j in pairwise(cut_loc):
            text_section.append(" ".join(text_split_by_sent[i:j]))
        return text_section

    def split_general(text, by, verbose=False, ignore_case=False):
        if ignore_case:
            if verbose:
                print(f"used {by} to split, ignore_case=True")
            pattern = re.compile(re.escape(by), re.IGNORECASE)
            split_text = pattern.split(text)
            return split_text
        else:
            if verbose:
                print(f"used {by} to split, ignore_case=False")
            return text.split(by)

    def reg_split(text, pattern):
        return re.split(pattern, text)

    if ("sp" in by or "white" in by) and not strict:
        if verbose:
            print(f"splited by space")
        result= text.split()
    elif ("word" in by and "len" in by) and not strict:
        if verbose:
            print(f"split_by_word_length(text, length)")
        result=  split_by_word_length(text, **kws)  # split_by_word_length(text, length)
    # elif "," in by:
    #     if verbose:
    #         print(f"splited by ','")
    #     return text.split(",")
    elif isinstance(by, list):
        if verbose:
            print(f"split_by_multiple_delimiters: ['|','&']")
        result=  split_by_multiple_delimiters(text, by)
    elif (
        all([("digi" in by or "num" in by), not "sent" in by, not "str" in by])
        and not strict
    ):
        if verbose:
            print(f"splited by digital (numbers)")
        result=  re.split(r"(\d+)", text)
    elif all([("digi" in by or "num" in by), "str" in by]) and not strict:
        if verbose:
            print(f"Splitting by (number strings)")
        pattern = re.compile(
            r"\b((?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?)(?:[-\s]?(?:one|two|three|four|five|six|seven|eight|nine|ten|eleven|twelve|thirteen|fourteen|fifteen|sixteen|seventeen|eighteen|nineteen|twenty|thirty|forty|fifty|sixty|seventy|eighty|ninety|hundred|thousand|million|billion|trillion|and|[\d,]+(?:\.\d+)?))*)\b",
            re.IGNORECASE,
        )
        result=  re.split(pattern, text)
    elif ("pun" in by) and not strict:
        if verbose:
            print(f"splited by 标点('.!?;')")
        result=  re.split(r"[.!?;]", text)
    elif ("\n" in by or "li" in by) and not strict:
        if verbose:
            print(f"splited by lines('\n')")
        result=  text.splitlines()
    elif ("cam" in by) and not strict:
        if verbose:
            print(f"splited by camel_case")
        result=  split_by_camel_case(text)
    elif ("word" in by) and not strict:
        from nltk.tokenize import word_tokenize

        if verbose:
            print(f"splited by word")
        result=  word_tokenize(text)
    elif ("sen" in by and not "num" in by) and not strict:
        from nltk.tokenize import sent_tokenize

        if verbose:
            print(f"splited by sentence")
        result=  sent_tokenize(text)
    elif ("sen" in by and "num" in by) and not strict:
        result=  split_by_sent_num(text, **kws)
    elif ("cha" in by) and not strict:
        if verbose:
            print(f"splited by chracters")
        result=  list(text)
    elif ("up" in by or "cap" in by) and ("l" not in by) and not strict:
        if verbose:
            print(f"splited by upper case")
        result=  split_at_upper(text)
    elif ("u" in by and "l" in by) and not strict:
        if by.find("u") < by.find("l"):
            if verbose:
                print(f"splited by upper followed by lower case")
            result=  split_at_upper_fl_lower(text)
        else:
            if verbose:
                print(f"splited by lower followed by upper case")
            result=  split_at_lower_fl_upper(text)
    elif ("start" in by or "head" in by) and not strict:
        if verbose:
            print(f"splited by lookahead")
        result=  split_by_regex_lookahead(text, **kws)
    elif ("end" in by or "tail" in by) and not strict:
        if verbose:
            print(f"splited by endings")
        result=  split_by_regex_end(text, **kws)
    elif ("other" in by or "non_alp" in by) and not strict:
        if verbose:
            print(f"splited by non_alphanumeric")
        result=  split_by_consecutive_non_alphanumeric(text)
    elif ("len" in by) and not strict:
        if verbose:
            print(f"splited by fixed length")
        result=  split_by_fixed_length_chunks(text, **kws)
    elif ("re" in by or "cus" in by or "cos" in by) and not strict:
        if verbose:
            print(f"splited by customed, re; => {by}")
        result=  reg_split(text, **kws)
    elif any(["lang" in by, "eng" in by,"non_ascii" in by]) and not strict:
        result=  split_non_ascii(text)
    else:
        result=  split_general(text, by, verbose=verbose, **kws)

    if strip_results:
        result = [item.strip() for item in result if item.strip()] 
    else:
        result = [item for item in result if item]

    return result


def pdf2img(dir_pdf, dir_save=None, page=None, kind="png", verbose=True, **kws):
    from pdf2image import convert_from_path, pdfinfo_from_path

    df_dir_img_single_page = pd.DataFrame()
    dir_single_page = []
    if verbose:
        from pprint import pp

        pp(pdfinfo_from_path(dir_pdf))
    if isinstance(page, tuple) and page:
        page = list(page)
    if isinstance(page, int):
        page = [page]
    if page is None:
        page = [pdfinfo_from_path(dir_pdf)["Pages"]]
    if len(page) == 1 and page != [pdfinfo_from_path(dir_pdf)["Pages"]]:
        page = [page[0], page[0]]
    else:
        page = [1, page[0]]
    print(page)
    pages = convert_from_path(dir_pdf, first_page=page[0], last_page=page[-1], **kws)
    if dir_save is None:
        dir_save = mkdir(dirname(dir_pdf), basename(dir_pdf).split(".")[0] + "_img")
    for i, page in enumerate(pages):
        if verbose:
            print(f"processing page: {i+1}")
        if i < 9:
            dir_img_each_page = dir_save + f"page_0{i+1}.png"
        else:
            dir_img_each_page = dir_save + f"page_{i+1}.png"
        dir_single_page.append(dir_img_each_page)
        page.save(dir_img_each_page, kind.upper())
    df_dir_img_single_page["fpath"] = dir_single_page
    return df_dir_img_single_page


# dir_pdf = "/Users/macjianfeng/Dropbox/github/python/240308_Python Data Science Handbook.pdf"
# df_page = pdf2img(dir_pdf, page=[1, 5],dpi=300)
def get_encoding(fpath, alternative_encodings=None, verbose=False):
    """
    Attempt to determine the encoding of a file by trying multiple encodings.

    Parameters:
    fpath (str): The path to the file.
    alternative_encodings (list): List of encodings to try. If None, uses a default list.
    verbose (bool): If True, print detailed information about each attempted encoding.

    Returns:
    str: The encoding that successfully read the file, or None if no encoding worked.
    """
    if alternative_encodings is None:
        alternative_encodings = [
            "utf-8",
            "latin1",
            "windows-1252",
            "iso-8859-1",
            "iso-8859-2",
            "iso-8859-3",
            "iso-8859-4",
            "iso-8859-5",
            "iso-8859-6",
            "iso-8859-7",
            "iso-8859-8",
            "iso-8859-9",
            "windows-1250",
            "windows-1251",
            "windows-1253",
            "windows-1254",
            "windows-1255",
            "windows-1256",
            "windows-1257",
            "windows-1258",
            "big5",
            "gb18030",
            "shift_jis",
            "euc_jp",
            "koi8_r",
            "mac_roman",
            "mac_central_europe",
            "mac_greek",
            "mac_cyrillic",
            "mac_arabic",
            "mac_hebrew",
        ]

    if not os.path.isfile(fpath):
        raise FileNotFoundError(f"The file {fpath} does not exist.")

    for enc in alternative_encodings:
        try:
            with open(fpath, mode="r", encoding=enc) as file:
                file.read()  # Try to read the file
            if verbose:
                print(f"Successfully detected encoding: {enc}")
            return enc
        except UnicodeDecodeError:
            if verbose:
                print(f"Failed to decode with encoding: {enc}")
            continue

    # If no encoding worked
    print("No suitable encoding found.")
    return None


def unzip(dir_path, output_dir=None):
    """
    Unzips or extracts various compressed file formats (.gz, .zip, .7z, .tar, .bz2, .xz, .rar).
    If the output directory already exists, it will be replaced.

    # Example usage:
    output_dir = unzip('data.tar.gz')
    output_file = unzip('file.csv.gz')
    output_dir_zip = unzip('archive.zip')
    output_dir_7z = unzip('archive.7z')

    Parameters:
    dir_path (str): Path to the compressed file.
    output_dir (str): Directory where the extracted files will be saved.
                      If None, it extracts to the same directory as the file, with the same name.

    Returns:
    str: The path to the output directory where files are extracted.
    """

    # Set default output directory to the same as the input file
    if output_dir is None:
        output_dir = os.path.splitext(dir_path)[0]

    # If the output directory already exists, remove it and replace it
    if os.path.exists(output_dir):
        if os.path.isdir(output_dir):  # check if it is a folder
            import shutil

            shutil.rmtree(output_dir)  # remove folder
        else:
            os.remove(output_dir)  # remove file

    # Handle .tar.gz files
    if dir_path.endswith(".tar.gz") or dir_path.endswith(".tgz"):
        import tarfile

        with tarfile.open(dir_path, "r:gz") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir
    # Handle .gz files
    if dir_path.endswith(".gz") or dir_path.endswith(".gzip"):
        import gzip

        output_file = os.path.splitext(dir_path)[0]  # remove the .gz extension
        try:
            import shutil

            with gzip.open(dir_path, "rb") as gz_file:
                with open(output_file, "wb") as out_file:
                    shutil.copyfileobj(gz_file, out_file)
            print(f"unzipped '{dir_path}' to '{output_file}'")
        except FileNotFoundError:
            print(f"Error: The file '{dir_path}' was not found.")
        except PermissionError:
            print(
                f"Error: Permission denied when accessing '{dir_path}' or writing to '{output_file}'."
            )
        except Exception as e:
            try:
                import tarfile

                with tarfile.open(dir_path, "r:gz") as tar:
                    tar.extractall(path=output_file)
            except Exception as final_e:
                print(f"An final unexpected error occurred: {final_e}")
        return output_file

    # Handle .zip files
    elif dir_path.endswith(".zip"):
        import zipfile

        with zipfile.ZipFile(dir_path, "r") as zip_ref:
            zip_ref.extractall(output_dir)
        return output_dir

    # Handle .7z files (requires py7zr)
    elif dir_path.endswith(".7z"):
        import py7zr

        with py7zr.SevenZipFile(dir_path, mode="r") as z:
            z.extractall(path=output_dir)
        return output_dir

    # Handle .tar files
    elif dir_path.endswith(".tar"):
        import tarfile

        with tarfile.open(dir_path, "r") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir

    # Handle .tar.bz2 files
    elif dir_path.endswith(".tar.bz2"):
        import tarfile

        with tarfile.open(dir_path, "r:bz2") as tar_ref:
            tar_ref.extractall(output_dir)
        return output_dir

    # Handle .bz2 files
    elif dir_path.endswith(".bz2"):
        import bz2

        output_file = os.path.splitext(dir_path)[0]  # remove the .bz2 extension
        with bz2.open(dir_path, "rb") as bz_file:
            with open(output_file, "wb") as out_file:
                shutil.copyfileobj(bz_file, out_file)
        return output_file

    # Handle .xz files
    elif dir_path.endswith(".xz"):
        import lzma

        output_file = os.path.splitext(dir_path)[0]  # remove the .xz extension
        with lzma.open(dir_path, "rb") as xz_file:
            with open(output_file, "wb") as out_file:
                shutil.copyfileobj(xz_file, out_file)
        return output_file

    # Handle .rar files (requires rarfile)
    elif dir_path.endswith(".rar"):
        import rarfile

        with rarfile.RarFile(dir_path) as rar_ref:
            rar_ref.extractall(output_dir)
        return output_dir

    else:
        raise ValueError(f"Unsupported file format: {os.path.splitext(dir_path)[1]}")

def is_df_abnormal(df: pd.DataFrame, verbose=False) -> bool:
    """
    Usage
    is_abnormal = is_df_abnormal(df, verbose=1)
    True: abnormal
    False: normal
    """
    if not isinstance(df, pd.DataFrame):
        if verbose:
            print("not pd.DataFrame")
        return False
    df.columns = df.columns.astype(str)  # 把它变成str, 这样就可以进行counts运算了
    # Initialize a list to hold messages about abnormalities
    messages = []
    is_abnormal = False
    # Check the shape of the DataFrame
    actual_shape = df.shape
    messages.append(f"Shape of DataFrame: {actual_shape}")

    # Check column names
    column_names = df.columns.tolist()

    # Count of delimiters and their occurrences
    delimiter_counts = {"\t": 0, ",": 0, "\n": 0, "": 0}  # Count of empty strings

    for name in column_names:
        # Count occurrences of each delimiter
        delimiter_counts["\t"] += name.count("\t")
        delimiter_counts[","] += name.count(",")
        delimiter_counts["\n"] += name.count("\n")
        if name.strip() == "":
            delimiter_counts[""] += 1

    # Check for abnormalities based on delimiter counts
    if len(column_names) == 1 and delimiter_counts["\t"] > 1:
        messages.append("Abnormal: Column names are not split correctly.")
        is_abnormal = True
        if verbose:
            print(f'len(column_names) == 1 and delimiter_counts["\t"] > 1')
    if verbose:
        print("1", is_abnormal)
    if any(delimiter_counts[d] > 3 for d in delimiter_counts if d != ""):
        messages.append("Abnormal: Too many delimiters in column names.")
        is_abnormal = True
        if verbose:
            print(f'any(delimiter_counts[d] > 3 for d in delimiter_counts if d != "")')
    if verbose:
        print("2", is_abnormal)
    if delimiter_counts[""] > 3:
        messages.append("Abnormal: There are empty column names.")
        is_abnormal = True
        if verbose:
            print(f'delimiter_counts[""] > 3')
    if verbose:
        print("3", is_abnormal)
    if any(delimiter_counts[d] > 3 for d in ["\t", ",", "\n"]):
        messages.append("Abnormal: Some column names contain unexpected characters.")
        is_abnormal = True
        if verbose:
            print(f'any(delimiter_counts[d] > 3 for d in ["\t", ",", "\n"])')
    if verbose:
        print("4", is_abnormal)
    # # Check for missing values
    # missing_values = df.isnull().sum()
    # if missing_values.any():
    #     messages.append("Missing values in columns:")
    #     messages.append(missing_values[missing_values > 0].to_string())
    #     is_abnormal = True
    #     print(f'missing_values.any()')

    # Check data types
    data_types = df.dtypes
    # messages.append(f"Data types of columns:\n{data_types}")

    # Check for an unreasonable number of rows or columns
    if actual_shape[0] < 2 or actual_shape[1] < 2:
        messages.append(
            "Abnormal: DataFrame is too small (less than 2 rows or columns)."
        )
        is_abnormal = True
        if verbose:
            print(f"actual_shape[0] < 2 or actual_shape[1] < 2")
    if verbose:
        print("6", is_abnormal)
    # Compile results
    if verbose:
        print("\n".join(messages))
    return is_abnormal  # Data is abnormal
def decrypt_excel(fpath, password):
    # * needs a password?
    import msoffcrypto  # pip install msoffcrypto-tool
    from io import BytesIO

    # Open the encrypted Excel file
    with open(fpath, "rb") as f:
        try:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.load_key(password=password)  # Provide the password
            decrypted = BytesIO()
            office_file.decrypt(decrypted)
        except:
            office_file = msoffcrypto.OfficeFile(f)
            office_file.load_key(password=depass(password))  # Provide the password
            decrypted = BytesIO()
            office_file.decrypt(decrypted)
    decrypted.seek(0) # reset pointer to start
    return decrypted

# ! for Excel formating   
def _backup_validations(sheet, verbose=False):
    """
    Comprehensive validation backup with multiple verification layers
    """
    from openpyxl.utils import range_to_tuple, get_column_letter
    from openpyxl.worksheet.datavalidation import DataValidation
    
    backup = {
        "validations": [],
        "conditional_formatting": [],
        "merged_cells": [str(mr) for mr in sheet.merged_cells.ranges],
        "_metadata": {
            "validated_cells": set(),
            "validated_columns": set(),
            "validation_types": set()
        }
    }

    # METHOD 1: Official data_validations collection (primary source)
    for dv in sheet.data_validations.dataValidation:
        validation_data = {
            "type": dv.type,
            "formula1": dv.formula1,
            "formula2": dv.formula2,
            "allow_blank": dv.allow_blank,
            "showDropDown": dv.showDropDown,
            "showInputMessage": getattr(dv, 'showInputMessage', True),
            "showErrorMessage": getattr(dv, 'showErrorMessage', False),
            "errorTitle": dv.errorTitle,
            "error": dv.error,
            "promptTitle": dv.promptTitle,
            "prompt": dv.prompt,
            "ranges": [],
            "_source": "data_validations"
        }
        
        # Process ranges
        for rng in dv.cells.ranges:
            range_str = str(rng)
            print(range_str)
            try:
                coord = range_to_tuple(range_str)
                min_col, min_row, max_col, max_row = coord
                is_single = (min_col == max_col and min_row == max_row)
                
                # Track coverage
                for col in range(min_col, max_col + 1):
                    col_letter = get_column_letter(col)
                    backup["_metadata"]["validated_columns"].add(col_letter)
                    for row in range(min_row, max_row + 1):
                        backup["_metadata"]["validated_cells"].add(f"{col_letter}{row}")
                
                validation_data["ranges"].append({
                    "str": range_str,
                    "coord": coord,
                    "is_single": is_single
                })
            except Exception as e:
                validation_data["ranges"].append({"str": range_str, "error": str(e)})
        
        backup["validations"].append(validation_data)
        backup["_metadata"]["validation_types"].add(dv.type)

    # METHOD 2: Cell-by-cell verification (fallback)
    missing_validations = []
    for row in sheet.iter_rows():
        for cell in row:
            # Version-agnostic cell validation check
            cell_has_validation = False
            try:
                # OpenPyXL 3.0+ style
                if hasattr(cell, 'data_validation') and cell.data_validation:
                    cell_has_validation = True
                # Older versions style
                elif hasattr(cell, 'has_data_validation') and cell.has_data_validation:
                    cell_has_validation = True
            except Exception:
                continue
                
            if cell_has_validation:
                cell_ref = f"{get_column_letter(cell.column)}{cell.row}"
                if cell_ref not in backup["_metadata"]["validated_cells"]:
                    missing_validations.append(cell_ref)

    # METHOD 3: Handle any missing validations
    if missing_validations:
        if verbose:
            print(f"Found {len(missing_validations)} validations not in data_validations collection")
        
        # Group by column for more efficient processing
        from collections import defaultdict
        column_groups = defaultdict(list)
        for ref in missing_validations:
            col = ''.join(filter(str.isalpha, ref))
            column_groups[col].append(ref)
        
        # Create supplemental validations
        for col, cell_refs in column_groups.items():
            try:
                sample_cell = sheet[cell_refs[0]]
                dv = sample_cell.data_validation if hasattr(sample_cell, 'data_validation') else None
                
                if dv:
                    validation_data = {
                        "type": dv.type,
                        "formula1": getattr(dv, 'formula1', None),
                        "formula2": getattr(dv, 'formula2', None),
                        "allow_blank": getattr(dv, 'allow_blank', False),
                        "showDropDown": getattr(dv, 'showDropDown', False),
                        "showInputMessage": getattr(dv, 'showInputMessage', True),
                        "showErrorMessage": getattr(dv, 'showErrorMessage', False),
                        "ranges": cell_refs,
                        "_source": "cell_validation_fallback",
                        "_recovered": True
                    }
                    backup["validations"].append(validation_data)
                    backup["_metadata"]["validation_types"].add(dv.type)
                    
                    # Update coverage tracking
                    backup["_metadata"]["validated_columns"].add(col)
                    backup["_metadata"]["validated_cells"].update(cell_refs)
            except Exception as e:
                if verbose:
                    print(f"Failed to backup validation for {col}: {str(e)}")
 
    # METHOD 4: Comprehensive cross-sheet dropdown detection
    for row in sheet.iter_rows():
        for cell in row:
            try:
                # Skip if no data validation or not a list type
                if not hasattr(cell, 'data_validation') or not cell.data_validation:
                    continue
                    
                dv = cell.data_validation
                if dv.type != 'list':
                    continue
                    
                formula = dv.formula1
                if not formula:
                    continue

                # Standard cleaning and pattern detection
                clean_formula = formula.strip('"\'').lstrip('=')
                cell_ref = f"{get_column_letter(cell.column)}{cell.row}"
                
                # Detection patterns (ordered by priority)
                patterns = [
                    (r'^[\w\s]+!\$?[A-Za-z]+\$?\d+(?::\$?[A-Za-z]+\$?\d+)?$', "direct sheet reference"),
                    (r'INDIRECT\(["\'][\w\s]+![A-Za-z]+\d+(?::[A-Za-z]+\d+)?["\']\)', "INDIRECT sheet reference"),
                    (r'^[\w\s]+$', "potential named range"),
                ]

                detected_type = None
                for pattern, description in patterns:
                    if re.match(pattern, clean_formula, re.IGNORECASE):
                        detected_type = description
                        break

                if detected_type:
                    # Special handling for direct sheet references
                    if detected_type == "direct sheet reference":
                        # Extract sheet name for verification
                        sheet_name = clean_formula.split('!')[0]
                        # Verify the sheet exists in the workbook
                        if sheet_name not in sheet.parent.sheetnames:
                            detected_type = f"broken reference (sheet '{sheet_name}' not found)"
                    
                    validation_data = {
                        "type": "list",
                        "formula1": formula,
                        "formula2": getattr(dv, 'formula2', None),
                        "allow_blank": getattr(dv, 'allow_blank', True),
                        "showDropDown": not getattr(dv, 'showDropDown', False),  # Correct dropdown display
                        "showInputMessage": getattr(dv, 'showInputMessage', True),
                        "showErrorMessage": getattr(dv, 'showErrorMessage', False),
                        "errorTitle": getattr(dv, 'errorTitle', ""),
                        "error": getattr(dv, 'error', ""),
                        "promptTitle": getattr(dv, 'promptTitle', ""),
                        "prompt": getattr(dv, 'prompt', ""),
                        "ranges": [cell_ref],
                        "_source": "cross_sheet_detection",
                        "_detection_method": detected_type,
                        "_is_cross_sheet": True,
                        "_formula_clean": clean_formula,
                        "_sheet_name": clean_formula.split('!')[0] if '!' in clean_formula else None
                    }

                    # Check for duplicates before adding
                    is_duplicate = any(
                        v.get("_source") == "cross_sheet_detection" and
                        cell_ref in v.get("ranges", []) and
                        v.get("formula1") == formula
                        for v in backup["validations"]
                    )
                    
                    if not is_duplicate:
                        backup["validations"].append(validation_data)
                        backup["_metadata"]["validated_cells"].add(cell_ref)
                        backup["_metadata"]["validated_columns"].add(get_column_letter(cell.column))
                        backup["_metadata"]["validation_types"].add(dv.type)
                        
            except Exception as e:
                if verbose:
                    print(f"Error processing cell {cell.coordinate}: {str(e)}")

    return backup

def _restore_validations(sheet, backup,verbose=False):
    """
    恢复数据验证和条件格式规则到工作表
    
    Args:
        sheet: openpyxl的工作表对象
        backup: 从_backup_validations()获取的备份字典
    """
    from openpyxl.worksheet.datavalidation import DataValidation, DataValidationList
    from openpyxl.formatting.rule import Rule, ColorScaleRule, DataBarRule, IconSetRule
    from openpyxl.utils import get_column_letter

    # 1. 清除现有验证和格式（更可靠的方法）
    sheet.data_validations = DataValidationList()
    
    # 清除条件格式的更好方法
    if hasattr(sheet.conditional_formatting, '_rules'):
        sheet.conditional_formatting._rules.clear()
    elif hasattr(sheet.conditional_formatting, 'cf_rules'):
        sheet.conditional_formatting.cf_rules.clear()
    else:
        # 最彻底的清除方法 
        cf_ranges = list(sheet.conditional_formatting)
        for cf_range in cf_ranges:
            try:
                del sheet.conditional_formatting[cf_range]
            except TypeError:
                # Skip problematic ranges that can't be deleted
                continue

    # 2. 先恢复合并单元格
    for mr in backup.get("merged_cells", []):
        try:
            if mr and mr not in sheet.merged_cells:
                sheet.merge_cells(mr)
        except Exception as e:
            if "already merged" not in str(e):
                print(f"[合并单元格] 警告: {mr} - {str(e)}")

    # 3. 恢复数据验证规则（修复了关键参数）
    for i, val in enumerate(backup.get("validations", [])):
        try:
            dv = DataValidation(
                type=val["type"],
                formula1=val["formula1"],
                formula2=val.get("formula2"),
                allow_blank=val["allow_blank"],
                showDropDown=False,  # 强制开启, # Dropdown control (False = show dropdown)
                showInputMessage=val.get("showInputMessage", True),  
                showErrorMessage=val.get("showErrorMessage", False), 
                errorTitle=val.get("errorTitle"),
                error=val.get("error"),
                promptTitle=val.get("promptTitle"),
                prompt=val.get("prompt")
            )
            for rng in val["ranges"]:
                try:
                    # Handle both string ranges and coordinate dicts from backup
                    if isinstance(rng, dict):
                        # Restore from coordinate backup if available
                        if "coord" in rng:
                            min_col, min_row, max_col, max_row = rng["coord"]
                            if min_col == max_col and min_row == max_row:
                                dv.add(f"{get_column_letter(min_col)}{min_row}")
                            else:
                                dv.add(f"{get_column_letter(min_col)}{min_row}:"
                                      f"{get_column_letter(max_col)}{max_row}")
                        else:
                            dv.add(rng["str"])
                    else:
                        # Direct string range
                        dv.add(rng)
                except Exception as e:
                    print(f"Warning: Could not add range {rng} to validation: {str(e)}")
            sheet.add_data_validation(dv)
        except Exception as e:
            if verbose:
                print(f"Warning: Could not restore validation: {str(e)}")
 
    # 4. 恢复条件格式规则（完全重写的逻辑）
    for i, rule_data in enumerate(backup.get("conditional_formatting", [])):
        try:
            rule = None
            rule_type = rule_data.get("type", "")
            
            # 调试信息
            debug_info = f"[条件格式#{i}] 类型: {rule_type}"
            if "colorScale" in rule_data:
                debug_info += " (颜色比例)"
            elif "dataBar" in rule_data:
                debug_info += " (数据条)"
            elif "iconSet" in rule_data:
                debug_info += " (图标集)"
            print(debug_info)

            # 创建对应类型的规则
            if rule_data.get("colorScale"):
                cs = rule_data["colorScale"]
                rule = ColorScaleRule(
                    start_type=cs.get("start_type", "min"),
                    start_value=cs.get("start_value"),
                    start_color=cs.get("start_color"),
                    mid_type=cs.get("mid_type", "percentile"),
                    mid_value=cs.get("mid_value", 50),
                    mid_color=cs.get("mid_color"),
                    end_type=cs.get("end_type", "max"),
                    end_value=cs.get("end_value"),
                    end_color=cs.get("end_color")
                )
            elif rule_data.get("dataBar"):
                db = rule_data["dataBar"]
                rule = DataBarRule(
                    start_type=db.get("start_type", "min"),
                    start_value=db.get("start_value"),
                    end_type=db.get("end_type", "max"),
                    end_value=db.get("end_value"),
                    color=db.get("color", "FF638EC6"),
                    showValue=db.get("showValue", True)
                )
            elif rule_data.get("iconSet"):
                icon = rule_data["iconSet"]
                rule = IconSetRule(
                    iconSet=icon.get("iconSet", "3TrafficLights"),
                    showValue=icon.get("showValue", True),
                    values=icon.get("values", None),
                    type=icon.get("type", None),
                    reverse=icon.get("reverse", False)
                )
            else:
                # 标准规则
                rule = Rule(
                    type=rule_type,
                    dxf=rule_data.get("dxf"),
                    stopIfTrue=rule_data.get("stopIfTrue", False)
                )
                # 动态设置属性
                for attr in ['formula', 'formula1', 'formula2', 'operator', 'text']:
                    if attr in rule_data and rule_data[attr] is not None:
                        try:
                            setattr(rule, attr, rule_data[attr])
                        except AttributeError:
                            if verbose:
                                print(f"[条件格式#{i}] 警告: 无法设置属性 {attr}")

            if rule:
                # 应用规则到所有范围
                range_count = 0
                for rng in rule_data.get("ranges", []):
                    try:
                        if not isinstance(rng, str):
                            rng = str(rng)
                        if rng:  # 非空字符串检查
                            sheet.conditional_formatting.add(rng, rule)
                            range_count += 1
                    except Exception as e:
                        if verbose:
                            print(f"[条件格式#{i}] 范围错误: {rng} - {str(e)}")
                
                if range_count == 0:
                    if verbose:
                        print(f"[条件格式#{i}] 警告: 没有有效范围，规则未应用")
            else:
                if verbose:     
                    print(f"[条件格式#{i}] 错误: 无法创建规则对象")
                
        except Exception as e:
            if verbose:
                print(f"[条件格式#{i}] 严重错误: {str(e)}")

# ! fload master
def fload(fpath, kind=None, **kwargs):
    """
    Load content from a file with specified file type.
    Parameters:
        fpath (str): The file path from which content will be loaded.
        kind (str): The file type to load. Supported options: 'docx', 'txt', 'md', 'html', 'json', 'yaml', 'xml', 'csv', 'xlsx', 'pdf'.
        **kwargs: Additional parameters for 'csv' and 'xlsx' file types.
    Returns:
        content: The content loaded from the file.
    """

    def read_mplstyle(style_file):
        import matplotlib.pyplot as plt

        # Load the style file
        plt.style.use(style_file)

        # Get the current style properties
        style_dict = plt.rcParams

        # Convert to dictionary
        style_dict = dict(style_dict)
        # Print the style dictionary
        for i, j in style_dict.items():
            print(f"\n{i}::::{j}")
        return style_dict

    # #example usage:
    # style_file = "/ std-colors.mplstyle"
    # style_dict = read_mplstyle(style_file)

    def load_txt_md(fpath):
        with open(fpath, "r") as file:
            content = file.read()
        return content
    def load_html(fpath, **kwargs):
        return pd.read_html(fpath, **kwargs)

    def load_json(fpath, **kwargs):
        output = kwargs.pop("output", "json")
        if output == "json":
            import json

            with open(fpath, "r") as file:
                content = json.load(file)
            return content
        else:
            return pd.read_json(fpath, **kwargs)

    def load_yaml(fpath):
        import yaml

        with open(fpath, "r") as file:
            content = yaml.safe_load(file)
        return content

    def load_xml(fpath, fsize_thr: int = 100):
        from lxml import etree

        def load_small_xml(fpath):
            tree = etree.parse(fpath)
            root = tree.getroot()
            return etree.tostring(root, pretty_print=True).decode()

        def load_large_xml(fpath):
            xml_parts = []
            context = etree.iterparse(
                fpath, events=("start", "end"), recover=True, huge_tree=True
            )

            for event, elem in context:
                if event == "end":
                    xml_parts.append(etree.tostring(elem, pretty_print=True).decode())
                    elem.clear()
                    while elem.getprevious() is not None:
                        del elem.getparent()[0]
            del context
            return "".join(xml_parts)

        file_size = os.path.getsize(fpath) / 1024 / 1024  # in MB

        if file_size > fsize_thr:
            print(f"reading a small file:{file_size} Mb")
            return load_large_xml(fpath)
        else:
            print(f"reading a big file:{file_size} Mb")
            return load_small_xml(fpath)

    def get_comment(fpath, comment=None, encoding="utf-8", lines_to_check=5):
        """
        Detect comment characters in a file.

        Parameters:
        - fpath: str, the file path of the CSV file.
        - encoding: str, the encoding of the file (default is 'utf-8').
        - lines_to_check: int, number of lines to check for comment characters (default is 5).

        Returns:
        - str or None: the detected comment character, or None if no comment character is found.
        """
        comment_chars = [
            "#",
            "!",
            "//",
            ";",
        ]  # can use any character or string as a comment
        try:
            with open(fpath, "r", encoding=encoding) as f:
                lines = [next(f) for _ in range(lines_to_check)]
        except (UnicodeDecodeError, ValueError):
            with open(fpath, "r", encoding=get_encoding(fpath)) as f:
                lines = [next(f) for _ in range(lines_to_check)]
        for line in lines:
            for char in comment_chars:
                if line.startswith(char):
                    return char
        return None

    def _get_chunks(df_fake):
        """
        helper func for 'load_csv'
        """
        chunks = []
        for chunk in df_fake:
            chunks.append(chunk)
        return pd.concat(chunks, ignore_index=True)

    def load_csv(fpath, **kwargs):
        from pandas.errors import EmptyDataError

        engine = kwargs.pop("engine", "pyarrow")  # default: None
        sep = kwargs.pop("sep", None)  # default: ','
        index_col = kwargs.pop("index_col", None)  # default: None
        memory_map = kwargs.pop("memory_map", False)  # default: False
        skipinitialspace = kwargs.pop("skipinitialspace", False)  # default: False
        encoding = kwargs.pop("encoding", "utf-8")  # default: "utf-8"
        on_bad_lines = kwargs.pop("on_bad_lines", "skip")  # default: 'error'
        comment = kwargs.pop("comment", None)  # default: None
        fmt = kwargs.pop("fmt", False)  # default:
        chunksize = kwargs.pop("chunksize", None)  # default: None

        # check filesize
        f_size = round(os.path.getsize(fpath) / 1024 / 1024, 3)
        if f_size >= 50:  # 50 MB
            if chunksize is None:
                chunksize = 5000
                print(
                    f"file size is {f_size}MB, then set the chunksize with {chunksize}"
                )
        engine = "c" if chunksize else engine  # when chunksize, recommend 'c'
        low_memory = kwargs.pop("low_memory", True)  # default: True
        low_memory = (
            False if chunksize else True
        )  # when chunksize, recommend low_memory=False # default:
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True) and verbose:
            use_pd("read_csv", verbose=verbose)

        if comment is None:  # default: None
            comment = get_comment(
                fpath, comment=None, encoding="utf-8", lines_to_check=5
            )
        try:
            df = pd.read_csv(
                fpath,
                engine=engine,
                index_col=index_col,
                memory_map=memory_map,
                encoding=encoding,
                comment=comment,
                skipinitialspace=skipinitialspace,
                sep=sep,
                on_bad_lines=on_bad_lines,
                chunksize=chunksize,
                low_memory=low_memory,
                **kwargs,
            )
            if chunksize:
                df = _get_chunks(df)
                print(df.shape)
            if is_df_abnormal(df, verbose=0):  # raise error
                raise ValueError("the df is abnormal")
        except:
            try:
                try:
                    if engine == "pyarrow" and not chunksize:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            encoding=encoding,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    else:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            memory_map=memory_map,
                            encoding=encoding,
                            sep=sep,
                            skipinitialspace=skipinitialspace,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    if chunksize:
                        df = _get_chunks(df)
                        print(df.shape)
                    if is_df_abnormal(df, verbose=0):
                        raise ValueError("the df is abnormal")
                except (UnicodeDecodeError, ValueError):
                    encoding = get_encoding(fpath)
                    # print(f"utf-8 failed. Retrying with detected encoding: {encoding}")
                    if engine == "pyarrow" and not chunksize:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            encoding=encoding,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    else:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            index_col=index_col,
                            memory_map=memory_map,
                            encoding=encoding,
                            sep=sep,
                            skipinitialspace=skipinitialspace,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                    if chunksize:
                        df = _get_chunks(df)
                        print(df.shape)
                    if is_df_abnormal(df, verbose=0):
                        raise ValueError("the df is abnormal")
            except Exception as e:
                separators = [",", "\t", ";", "|", " "]
                for sep in separators:
                    sep2show = sep if sep != "\t" else "\\t"
                    if verbose:
                        print(f'trying with: engine=pyarrow, sep="{sep2show}"')
                    try:
                        df = pd.read_csv(
                            fpath,
                            engine=engine,
                            skipinitialspace=skipinitialspace,
                            sep=sep,
                            on_bad_lines=on_bad_lines,
                            comment=comment,
                            chunksize=chunksize,
                            low_memory=low_memory,
                            **kwargs,
                        )
                        if chunksize:
                            df = _get_chunks(df)
                            print(df.shape)
                        if not is_df_abnormal(df, verbose=0) and verbose:  # normal
                            display(df.head(2))
                            if verbose:
                                print(f"shape: {df.shape}")
                            return df
                    except:
                        pass
                else:
                    if not chunksize:
                        engines = [None, "c", "python"]
                        for engine in engines:
                            separators = [",", "\t", ";", "|", " "]
                            for sep in separators:
                                try:
                                    sep2show = sep if sep != "\t" else "\\t"
                                    if verbose:
                                        print(
                                            f"trying with: engine={engine}, sep='{sep2show}'"
                                        )
                                    # print(".")
                                    df = pd.read_csv(
                                        fpath,
                                        engine=engine,
                                        sep=sep,
                                        on_bad_lines=on_bad_lines,
                                        comment=comment,
                                        chunksize=chunksize,
                                        low_memory=low_memory,
                                        **kwargs,
                                    )
                                    # display(df.head(2))
                                    # print(f"is_df_abnormal:{is_df_abnormal(df, verbose=0)}")
                                    if chunksize:
                                        df = _get_chunks(df)
                                        print(df.shape)
                                    if not is_df_abnormal(df, verbose=0):
                                        if verbose:
                                            (
                                                display(df.head(2))
                                                if isinstance(df, pd.DataFrame)
                                                else desisplay("it is not a DataFrame")
                                            )
                                            (
                                                print(f"shape: {df.shape}")
                                                if isinstance(df, pd.DataFrame) and verbose
                                                else display("it is not a DataFrame")
                                            )
                                        return df
                                except EmptyDataError as e:
                                    continue
                            else:
                                pass
        # print(kwargs)
        # if is_df_abnormal(df,verbose=verbose):
        #     df=pd.read_csv(fpath,**kwargs)
        if verbose:
            display(df.head(2))
            print(f"shape: {df.shape}")
        return df
 
    def load_excel(fpath, **kwargs):
        engine = kwargs.get("engine", "openpyxl")
        verbose = kwargs.pop("verbose", False)
        password = kwargs.pop("password", None)
        output = kwargs.pop("output", "DataFrame").lower()
        sheet_name = kwargs.pop("sheet_name", None)
        def print_sheet_info(fpath):
            try:
                meta = pd.ExcelFile(fpath)
                print(f"n_sheet={len(meta.sheet_names)},\t'sheetname = 0 (default)':")
                [print(f"{i}:\t{name}") for i, name in enumerate(meta.sheet_names)]
            except Exception as e:
                if verbose:
                    print(f"Error retrieving sheet info: {e}")

    

        if output in ["dataframe", "df"]:
            if verbose:
                print("loading data as a DataFrame")
            if not bool(password):
                if verbose:
                    print("Reading Excel without password protection...")
                df = pd.read_excel(fpath, engine=engine, sheet_name=sheet_name, **kwargs)
                if verbose:
                    print_sheet_info(fpath)
                return df
            # Handle password-protected DataFrame case
            else:
                if verbose:
                    print("Decrypting and loading DataFrame...")
                try:
                    decrypted = decrypt_excel(fpath, password=password)
                    df = pd.read_excel(decrypted, engine=engine,sheet_name=sheet_name, **kwargs)
                except:
                    df = pd.read_excel(fpath, engine=engine,sheet_name=sheet_name, **kwargs)
                if verbose:
                    print_sheet_info(fpath)
                return df
        # Handle cases for non-dataframe output
        else:
            if verbose:
                print("loading data as a formatted workbook")
            from openpyxl import load_workbook
            try:
                if verbose:
                    print("Returning worksheet (non-DataFrame output)...")
                if password:
                    try:
                        decrypted = decrypt_excel(fpath, password=password)
                        workbook = load_workbook(decrypted, data_only=False)
                    except:
                        workbook = load_workbook(fpath, data_only=False)
                    
                else:
                    workbook = load_workbook(fpath, data_only=False) 

                # Handle sheet selection
                if sheet_name:
                    if sheet_name not in workbook.sheetnames:
                        raise ValueError(f"Sheet '{sheet_name}' not found. Available sheets: {workbook.sheetnames}")
                        
                    # Remove other sheets while preserving case sensitivity
                    for sheet in [s for s in workbook.sheetnames if s != sheet_name]:
                        del workbook[sheet]
                return workbook

            except Exception as e:
                raise Exception(f"Error loading Workbook: {str(e)}")

    def load_parquet(fpath, **kwargs):
        """
        Load a Parquet file into a Pandas DataFrame with advanced options.

        Parameters:
        - fpath (str): The file path to the Parquet file.
        - engine (str): The engine to use for reading the Parquet file (default is 'pyarrow').
        - columns (list): List of columns to load. If None, loads all columns.
        - verbose (bool): If True, prints additional information about the loading process.
        - filters (list): List of filter conditions for predicate pushdown.
        - **kwargs: Additional keyword arguments for `pd.read_parquet`.

        Returns:
        - df (DataFrame): The loaded DataFrame.
        """

        engine = kwargs.get("engine", "pyarrow")
        verbose = kwargs.pop("verbose", False)

        if run_once_within(reverse=True) and verbose:
            use_pd("read_parquet", verbose=verbose)
        try:
            df = pd.read_parquet(fpath, engine=engine, **kwargs)
            if verbose:
                if "columns" in kwargs:
                    print(f"Loaded columns: {kwargs['columns']}")
                else:
                    print("Loaded all columns.")
            if verbose:
                print(f"shape: {df.shape}")
        except Exception as e:
            print(f"An error occurred while loading the Parquet file: {e}")
            df = None

        return df

    def load_ipynb(fpath, **kwargs):
        import nbformat
        from nbconvert import MarkdownExporter

        as_version = kwargs.get("as_version", 4)
        with open(fpath, "r") as file:
            nb = nbformat.read(file, as_version=as_version)
            md_exporter = MarkdownExporter()
            md_body, _ = md_exporter.from_notebook_node(nb)
        return md_body

    def load_pdf(fpath, page="all", verbose=False, **kwargs):
        """
        Parameters:
        fpath: The path to the PDF file to be loaded.
        page (optional):
            Specifies which page or pages to extract text from. By default, it's set to "all", which means text from all
            pages will be returned. It can also be an integer to specify a single page number or a list of integers to
            specify multiple pages.
        verbose (optional):
            If True, prints the total number of pages processed.
        Functionality:
        It initializes an empty dictionary text_dict to store page numbers as keys and their corresponding text as values.
        It iterates through each page of the PDF file using a for loop.
        For each page, it extracts the text using PyPDF2's extract_text() method and stores it in text_dict with the page number incremented by 1 as the key.
        If the page parameter is an integer, it converts it into a list containing that single page number to ensure consistency in handling.
        If the page parameter is a NumPy array, it converts it to a list using the tolist() method to ensure compatibility with list operations.
        If verbose is True, it prints the total number of pages processed.
        If page is a list, it combines the text of the specified pages into a single string combined_text and returns it.
        If page is set to "all", it returns the entire text_dict containing text of all pages.
        If page is an integer, it returns the text of the specified page number.
        If the specified page is not found, it returns the string "Page is not found".
        """
        from PyPDF2 import PdfReader

        text_dict = {}
        with open(fpath, "rb") as file:
            pdf_reader = PdfReader(file)
            num_pages = len(pdf_reader.pages)
            for page_num in range(num_pages):
                if verbose:
                    print(f"processing page {page_num}")
                page_ = pdf_reader.pages[page_num]
                text_dict[page_num + 1] = page_.extract_text()
        if isinstance(page, int):
            page = [page]
        elif isinstance(page, np.ndarray):
            page = page.tolist()
        if verbose:
            print(f"total pages: {page_num}")
        if isinstance(page, list):
            combined_text = ""
            for page_num in page:
                combined_text += text_dict.get(page_num, "")
            return combined_text
        elif "all" in page.lower():
            combined_text = ""
            for i in text_dict.values():
                combined_text += i
            return combined_text
        else:
            return text_dict.get(int(page), "Page is not found")

    def load_docx(fpath):
        from docx import Document

        doc = Document(fpath)
        content = [para.text for para in doc.paragraphs]
        return content

    def load_rtf(file_path):
        from striprtf.striprtf import rtf_to_text

        try:
            with open(file_path, "r") as file:
                rtf_content = file.read()
                text = rtf_to_text(rtf_content)
                return text
        except Exception as e:
            print(f"Error loading RTF file: {e}")

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()
    kind = kind.lstrip(".").lower()
    img_types = [
        "bmp",
        "eps",
        "gif",
        "png",
        "jpg",
        "jpeg",
        "jpeg2000",
        "tiff",
        "tif",
        "icns",
        "ico",
        "im",
        "msp",
        "pcx",
        "ppm",
        "sgi",
        "spider",
        "tga",
        "webp",
    ]
    doc_types = [
        "docx",
        "pdf",
        "txt",
        "csv",
        "xlsx",
        "tsv",
        "parquet",
        "snappy",
        "md",
        "html",
        "json",
        "yaml",
        "xml",
        "ipynb",
        "mtx",
        "rtf",
    ]
    zip_types = [
        "gz",
        "zip",
        "7z",
        "rar",
        "tgz",
        "tar",
        "tar.gz",
        "tar.bz2",
        "bz2",
        "xz",
        "gzip",
    ]
    other_types = ["fcs"]
    supported_types = [*doc_types, *img_types, *zip_types, *other_types]
    if kind not in supported_types:
        print(
            f'Warning:\n"{kind}" is not in the supported list '
        )  # {supported_types}')

    if kind == "docx":
        return load_docx(fpath)
    elif kind == "txt" or kind == "md":
        return load_txt_md(fpath)
    elif kind == "html":
        return load_html(fpath, **kwargs)
    elif kind == "json":
        return load_json(fpath, **kwargs)
    elif kind == "yaml":
        return load_yaml(fpath)
    elif kind == "xml":
        return load_xml(fpath)
    elif kind in ["csv", "tsv"]:
        # verbose = kwargs.pop("verbose", False)
        # if run_once_within(reverse=True) and verbose:
        #     use_pd("read_csv")
        content = load_csv(fpath, **kwargs)
        return content
    elif kind == "pkl":
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True) and verbose:
            use_pd("read_pickle")
        try:
            res_ = pd.read_pickle(fpath, **kwargs)
        except Exception as e:
            import pickle

            with open("sgd_classifier.pkl", "rb") as f:
                res_ = pickle.load(f)
        return res_
    elif kind in ["ods", "ods", "odt"]:
        engine = kwargs.get("engine", "odf")
        kwargs.pop("engine", None)
        return load_excel(fpath, engine=engine, **kwargs)
    elif kind == "xls":
        verbose = kwargs.pop("verbose", False)
        engine = kwargs.get("engine", "xlrd")
        kwargs.pop("engine", None)
        content = load_excel(fpath, engine=engine, **kwargs)
        (
            print(f"shape: {content.shape}")
            if isinstance(content, pd.DataFrame) and verbose
            else None
        )
        display(content.head(3)) if isinstance(content, pd.DataFrame) else None
        return content
    elif kind == "xlsx":
        verbose = kwargs.pop("verbose", False)
        content = load_excel(fpath, verbose=verbose,**kwargs)
        # (
        #     display(content.head(3))
        #     if isinstance(content, pd.DataFrame) and verbose
        #     else None
        # )
        print(f"shape: {content.shape}") if isinstance(content, pd.DataFrame) and verbose else None
        return content
    elif kind == "mtx":
        from scipy.io import mmread

        verbose = kwargs.pop("verbose", False)
        dat_mtx = mmread(fpath)
        content = pd.DataFrame.sparse.from_spmatrix(dat_mtx, **kwargs)
        (
            display(content.head(3))
            if isinstance(content, pd.DataFrame) and verbose
            else None
        )
        print(f"shape: {content.shape}")
        return content
    elif kind == "ipynb":
        return load_ipynb(fpath, **kwargs)
    elif kind in ["parquet", "snappy"]:
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("read_parquet")
        return load_parquet(fpath, **kwargs)
    elif kind == "feather":
        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("read_feather")
        content = pd.read_feather(fpath, **kwargs)
        return content
    elif kind == "h5":
        content = pd.read_hdf(fpath, **kwargs)
        return content
    elif kind == "pkl":
        content = pd.read_pickle(fpath, **kwargs)
        return content
    elif kind == "pdf":
        # print('usage:load_pdf(fpath, page="all", verbose=False)')
        return load_pdf(fpath, **kwargs)
    elif kind.lower() in img_types:
        print(f'Image ".{kind}" is loaded.')
        return load_img(fpath)
    elif kind == "gz" and fpath.endswith(".soft.gz"):
        import GEOparse

        return GEOparse.get_GEO(filepath=fpath)
    elif kind.lower() in zip_types:
        from pprint import pp

        keep = kwargs.get("keep", False)
        fpath_unzip = unzip(fpath)
        if os.path.isdir(fpath_unzip):
            print(f"{fpath_unzip} is a folder. fload stoped.")
            fpath_list = os.listdir("./datasets/GSE10927_family.xml")
            print(f"{len(fpath_list)} files within the folder")
            if len(fpath_list) > 5:
                pp(fpath_list[:5])
                print("there are more ...")
            else:
                pp(fpath_list)
            return fpath_list
        elif os.path.isfile(fpath_unzip):
            print(f"{fpath_unzip} is a file.")
            content_unzip = fload(fpath_unzip, **kwargs)
            if not keep:
                os.remove(fpath_unzip)
            return content_unzip
        else:
            print(f"{fpath_unzip} does not exist or is a different type.")

    elif kind.lower() == "gmt":
        import gseapy as gp

        gene_sets = gp.read_gmt(fpath)
        return gene_sets

    elif kind.lower() == "fcs":
        import fcsparser

        # https://github.com/eyurtsev/fcsparser
        meta, data = fcsparser.parse(fpath, reformat_meta=True)
        print("meta, data = fload(*.fcs)")
        return meta, data

    elif kind == "mplstyle":
        return read_mplstyle(fpath)
    elif kind == "rtf":
        return load_rtf(fpath)

    else:
        print("direct reading...")
        try:
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    content = f.readlines()
            except UnicodeDecodeError:
                print("Failed to read as utf-8, trying different encoding...")
                with open(
                    fpath, "r", encoding=get_encoding(fpath)
                ) as f:  # Trying with a different encoding
                    content = f.readlines()
        except:
            try:
                with open(fpath, "r", encoding="utf-8") as f:
                    content = f.read()
            except UnicodeDecodeError:
                print("Failed to read as utf-8, trying different encoding...")
                with open(
                    fpath, "r", encoding=get_encoding(fpath)
                ) as f:  # Trying with a different encoding
                    content = f.read()
        return content


# Example usage
# txt_content = fload('sample.txt')
# md_content = fload('sample.md')
# html_content = fload('sample.html')
# json_content = fload('sample.json')
# yaml_content = fload('sample.yaml')
# xml_content = fload('sample.xml')
# csv_content = fload('sample.csv')
# xlsx_content = fload('sample.xlsx')
# docx_content = fload('sample.docx')


def fopen(fpath):
    import os
    import platform
    import sys

    try:
        # Check if the file exists
        if not os.path.isfile(fpath):
            print(f"Error: The file does not exist - {fpath}")
            return

        # Get the system platform
        system = platform.system()

        # Platform-specific file opening commands
        if system == "Darwin":  # macOS
            os.system(f'open "{fpath}"')
        elif system == "Windows":  # Windows
            # Ensure the path is handled correctly in Windows, escape spaces
            os.system(f'start "" "{fpath}"')
        elif system == "Linux":  # Linux
            os.system(f'xdg-open "{fpath}"')
        elif system == "Java":  # Java (or other unhandled systems)
            print(f"Opening {fpath} on unsupported system.")
        else:
            print(f"Unsupported OS: {system}")

        print(f"Successfully opened {fpath} with the default application.")
    except Exception as e:
        print(f"Error opening file {fpath}: {e}")


def fupdate(fpath, content=None, how="head"):
    """
    Update a file by adding new content at the top and moving the old content to the bottom.
    If the file is a JSON file, merge the new content with the old content.

    Parameters
    ----------
    fpath : str
        The file path where the content should be updated.
    content : str or dict, optional
        The new content to add at the top of the file (for text) or merge (for JSON).
        If not provided, the function will not add any new content.

    Notes
    -----
    - If the file at `fpath` does not exist, it will be created.
    - For text files, the new content will be added at the top, followed by the old content.
    - For JSON files, the new content will be merged with the existing JSON content.
    """
    content = content or ""
    file_ext = os.path.splitext(fpath)[1]
    how_s = ["head", "tail", "start", "end", "beginning", "stop", "last", "before"]
    how = strcmp(how, how_s)[0]
    print(how)
    add_where = "head" if how in ["head", "start", "beginning", "before"] else "tail"
    if "json" in file_ext.lower():
        old_content = fload(fpath, kind="json") if os.path.exists(fpath) else {}
        updated_content = (
            {**content, **old_content}
            if add_where == "head"
            else (
                {**old_content, **content} if isinstance(content, dict) else old_content
            )
        )
        fsave(fpath, updated_content)
    else:
        # Handle text file
        if os.path.exists(fpath):
            with open(fpath, "r") as file:
                old_content = file.read()
        else:
            old_content = ""

        # Write new content at the top followed by old content
        with open(fpath, "w") as file:
            if add_where == "head":
                file.write(content + "\n")
                file.write(old_content)
            else:
                file.write(old_content)
                file.write(content + "\n")


def fappend(fpath, content=None):
    """
    append new content at the end.
    """
    content = content or ""
    if os.path.exists(fpath):
        with open(fpath, "r") as file:
            old_content = file.read()
    else:
        old_content = ""

    with open(fpath, "w") as file:
        file.write(old_content)
        file.write(content)


def filter_kwargs(kws, valid_kwargs):
    if isinstance(valid_kwargs, dict):
        kwargs_filtered = {
            key: value for key, value in kws.items() if key in valid_kwargs.keys()
        }
    elif isinstance(valid_kwargs, list):
        kwargs_filtered = {
            key: value for key, value in kws.items() if key in valid_kwargs
        }
    return kwargs_filtered


str_space_speed = 'sapce cmp:parquet(0.56GB)<feather(1.14GB)<csv(6.55GB)<pkl=h5("26.09GB")\nsaving time: pkl=feather("13s")<parquet("35s")<h5("2m31s")<csv("58m")\nloading time: pkl("6.9s")<parquet("16.1s")=feather("15s")<h5("2m 53s")<csv(">>>30m")'


def fsave(
    fpath,
    content,
    how="overwrite",
    kind=None,
    font_name="Times",
    font_size=10,
    spacing=6,
    **kwargs,
):
    """
    Save content into a file with specified file type and formatting.
    Parameters:
        fpath (str): The file path where content will be saved.
        content (list of str or dict): The content to be saved, where each string represents a paragraph or a dictionary for tabular data.
        kind (str): The file type to save. Supported options: 'docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml'.
        font_name (str): The font name for text formatting (only applicable for 'docx', 'html', and 'pdf').
        font_size (int): The font size for text formatting (only applicable for 'docx', 'html', and 'pdf').
        spacing (int): The space after each paragraph (only applicable for 'docx').
        **kwargs: Additional parameters for 'csv', 'xlsx', 'json', 'yaml' file types.
    Returns:
        None
    """

    def save_content(fpath, content, mode="w", how="overwrite"):
        if "wri" in how.lower():
            with open(fpath, mode, encoding="utf-8") as file:
                file.write(content)
        elif "upd" in how.lower():
            fupdate(fpath, content=content)
        elif "app" in how.lower():
            fappend(fpath, content=content)

    def save_docx(fpath, content, font_name, font_size, spacing):
        import docx

        if isinstance(content, str):
            content = content.split(". ")
        doc = docx.Document()
        for i, paragraph_text in enumerate(content):
            paragraph = doc.add_paragraph()
            run = paragraph.add_run(paragraph_text)
            font = run.font
            font.name = font_name
            font.size = docx.shared.Pt(font_size)
            if i != len(content) - 1:  # Add spacing for all but the last paragraph
                paragraph.space_after = docx.shared.Pt(spacing)
        doc.save(fpath)

    def save_txt_md(fpath, content, sep="\n", mode="w"):
        # Ensure content is a single string
        if isinstance(content, list):
            content = sep.join(content)
        save_content(fpath, sep.join(content), mode)

    def save_html(fpath, content, font_name, font_size, mode="w"):
        html_content = "<html><body>"
        for paragraph_text in content:
            html_content += f'<p style="font-family:{font_name}; font-size:{font_size}px;">{paragraph_text}</p>'
        html_content += "</body></html>"
        save_content(fpath, html_content, mode)

    def save_pdf(fpath, content, font_name, font_size):
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        # pdf.add_font('Arial','',r'/System/Library/Fonts/Supplemental/Arial.ttf',uni=True)
        pdf.set_font(font_name, "", font_size)
        for paragraph_text in content:
            pdf.multi_cell(0, 10, paragraph_text)
            pdf.ln(h="")
        pdf.output(fpath, "F")

    def save_csv(fpath, data, **kwargs):
        # https://pandas.pydata.org/docs/reference/api/pandas.DataFrame.to_csv.html

        verbose = kwargs.pop("verbose", False)
        if run_once_within(reverse=True):
            use_pd("to_csv", verbose=verbose)
        kwargs_csv = dict(
            path_or_buf=None,
            sep=",",
            na_rep="",
            float_format=None,
            columns=None,
            header=True,
            index=True,
            index_label=None,
            mode="w",
            encoding="UTF-8",
            compression="infer",
            quoting=None,
            quotechar='"',
            lineterminator=None,
            chunksize=None,
            date_format=None,
            doublequote=True,
            escapechar=None,
            decimal=".",
            errors="strict",
            storage_options=None,
        )
        kwargs_valid = filter_kwargs(kwargs, kwargs_csv)
        df = pd.DataFrame(data)
        df.to_csv(fpath, **kwargs_valid)

    def save_xlsx(fpath, data, password=None,apply_format=None, **kwargs):
        import msoffcrypto
        from io import BytesIO
        import openpyxl
        from openpyxl import Workbook
        from openpyxl.worksheet.worksheet import Worksheet
        import pandas.io.formats.style

        verbose = kwargs.pop("verbose", False)
        sheet_name = kwargs.pop("sheet_name", "Sheet1")
        engine = kwargs.pop("engine", "xlsxwriter")
        mode = kwargs.pop("mode","a")
        if_sheet_exists = strcmp(kwargs.get("if_sheet_exists","overwrite"),['error', 'new', 'replace', 'overlay','overwrite'])[0]
        if_sheet_exists="overlay" if if_sheet_exists=="overwrite" else if_sheet_exists
        kwargs.pop("if_sheet_exists",None)
        if run_once_within(reverse=True):
            use_pd("to_excel", verbose=verbose)
 
        if apply_format is None:
            kwargs_format=list(extract_kwargs(format_excel).keys())[4:]
            apply_format=True if any([i in kwargs_format for i in kwargs]) else False

        if apply_format or any([
                isinstance(data, Worksheet), 
                isinstance(data, Workbook),
                isinstance(data, pd.io.formats.style.Styler)
             ]): 
            if isinstance(data, pd.io.formats.style.Styler):
                try:
                    with pd.ExcelWriter(fpath, mode=mode, engine="openpyxl") as writer:
                        # First save raw data
                        data.data.to_excel(writer, sheet_name=sheet_name, index=False)
                        # Then save style on top
                        data.to_excel(writer, sheet_name=sheet_name, index=False)
                except Exception as e:
                    print(f"Cannot save the styles, only saving raw data! Because: {e}")
                    with pd.ExcelWriter(fpath, mode=mode, engine="openpyxl", if_sheet_exists='overlay') as writer:
                        data.data.to_excel(writer, sheet_name=sheet_name, index=False)
                        data.to_excel(writer, sheet_name=sheet_name, index=False)

            format_excel(df=data, 
                        filename=fpath,
                        sheet_name=sheet_name,
                        #  password=password,
                        if_sheet_exists=if_sheet_exists,
                        mode=mode, 
                        engine=engine,
                        verbose=verbose,
                        **kwargs)
        else:
            # Remove non-relevant kwargs
            irrelevant_keys=list(extract_kwargs(format_excel).keys())[4:]
            [kwargs.pop(key, None) for key in irrelevant_keys]
            df = pd.DataFrame(data)
            # Write to Excel without password first
            temp_file = BytesIO()
            
            df.to_excel(
                temp_file,
                sheet_name=sheet_name,
                index=False,
                engine="xlsxwriter",
                **kwargs,
            )
            # If a password is provided, encrypt the file
            if password:
                temp_file.seek(0)
                office_file = msoffcrypto.OfficeFile(temp_file) 
                with open(fpath, "wb") as encrypted_file:
                    office_file.encrypt(outfile=encrypted_file,password=password)
            else: # Save the file without encryption if no password is provided
                try:
                    # Use ExcelWriter with append mode if the file exists
                    engine="openpyxl" if mode=="a" else "xlsxwriter"
                    if mode=="a":
                        with pd.ExcelWriter(fpath, engine=engine, mode=mode,if_sheet_exists=if_sheet_exists) as writer:
                            df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
                    else:
                        with pd.ExcelWriter(fpath, engine=engine, mode=mode) as writer:
                            df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
                except FileNotFoundError:
                    # If file doesn't exist, create a new one
                    df.to_excel(fpath, sheet_name=sheet_name, index=False, **kwargs)

    def save_ipynb(fpath, data, **kwargs):
        # Split the content by code fences to distinguish between code and markdown
        import nbformat

        parts = data.split("```")
        cells = []

        for i, part in enumerate(parts):
            if i % 2 == 0:
                # Even index: markdown content
                cells.append(nbformat.v4.new_markdown_cell(part.strip()))
            else:
                # Odd index: code content
                cells.append(nbformat.v4.new_code_cell(part.strip()))
        # Create a new notebook
        nb = nbformat.v4.new_notebook()
        nb["cells"] = cells
        # Write the notebook to a file
        with open(fpath, "w", encoding="utf-8") as ipynb_file:
            nbformat.write(nb, ipynb_file) 
    def save_json(fpath_fname, var_dict_or_df):
        import json
        def _convert_js(data):
            if isinstance(data, pd.DataFrame):
                return data.to_dict(orient="list")
            elif isinstance(data, np.ndarray):
                return data.tolist()
            elif isinstance(data, dict):
                return {key: _convert_js(value) for key, value in data.items()}
            return data

        serializable_data = _convert_js(var_dict_or_df)
        # Save the serializable data to the JSON file
        with open(fpath_fname, "w") as f_json:
            json.dump(serializable_data, f_json, indent=4) 

    def save_yaml(fpath, data, **kwargs):
        import yaml

        with open(fpath, "w") as file:
            yaml.dump(data, file, **kwargs)

    def save_xml(fpath, data):
        from lxml import etree

        root = etree.Element("root")
        if isinstance(data, dict):
            for key, val in data.items():
                child = etree.SubElement(root, key)
                child.text = str(val)
        else:
            raise ValueError("XML saving only supports dictionary data")
        tree = etree.ElementTree(root)
        tree.write(fpath, pretty_print=True, xml_declaration=True, encoding="UTF-8")

    def save_parquet(fpath: str, data: pd.DataFrame, **kwargs):
        engine = kwargs.pop(
            "engine", "auto"
        )  # auto先试pyarrow, 不行就转为fastparquet, {‘auto’, ‘pyarrow’, ‘fastparquet’}
        compression = kwargs.pop(
            "compression", None
        )  # Use None for no compression. Supported options: ‘snappy’, ‘gzip’, ‘brotli’, ‘lz4’, ‘zstd’
        try:
            # Attempt to save with "pyarrow" if engine is set to "auto"
            data.to_parquet(fpath, engine=engine, compression=compression, **kwargs)
            print(
                f"DataFrame successfully saved to {fpath} with engine '{engine}' and {compression} compression."
            )
        except Exception as e:
            print(
                f"Error using with engine '{engine}' and {compression} compression: {e}"
            )
            if "Sparse" in str(e):
                try:
                    # Handle sparse data by converting columns to dense
                    print("Attempting to convert sparse columns to dense format...")
                    data = data.apply(
                        lambda x: (
                            x.sparse.to_dense() if pd.api.types.is_sparse(x) else x
                        )
                    )
                    save_parquet(fpath, data=data, **kwargs)
                except Exception as last_e:
                    print(
                        f"After converted sparse columns to dense format, Error using with engine '{engine}' and {compression} compression: {last_e}"
                    )

    if kind is None:
        _, kind = os.path.splitext(fpath)
        kind = kind.lower()

    kind = kind.lstrip(".").lower()

    if kind not in [
        "docx",
        "txt",
        "md",
        "html",
        "pdf",
        "csv",
        "xlsx",
        "json",
        "xml",
        "yaml",
        "ipynb",
    ]:
        print(
            f"Warning:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
        )
    mode = kwargs.get("mode", "w")
    if kind == "docx" or kind == "doc":
        save_docx(fpath, content, font_name, font_size, spacing)
    elif kind == "txt":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "md":
        save_txt_md(fpath, content, sep="", mode=mode)
    elif kind == "html":
        save_html(fpath, content, font_name, font_size)
    elif kind == "pdf":
        save_pdf(fpath, content, font_name, font_size)
    elif kind == "csv":
        save_csv(fpath, content, **kwargs)
    elif kind == "xlsx":
        save_xlsx(fpath, content, **kwargs)
    elif kind == "json":
        save_json(fpath, content)
    elif kind == "xml":
        save_xml(fpath, content)
    elif kind == "yaml":
        save_yaml(fpath, content, **kwargs)
    elif kind == "ipynb":
        save_ipynb(fpath, content, **kwargs)
    elif kind.lower() in ["parquet", "pq", "big", "par"]:
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_parquet")
            return None
        compression = kwargs.pop(
            "compression", None
        )  # Use None for no compression. Supported options: ‘snappy’, ‘gzip’, ‘brotli’, ‘lz4’, ‘zstd’
        # fix the fpath ends
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "parquet")
        if compression is not None:
            if not fpath.endswith(compression):
                fpath = fpath + f".{compression}"
        save_parquet(fpath=fpath, data=content, compression=compression, **kwargs)
    elif kind.lower() in ["pkl", "pk", "pickle", "pick"]:
        # Pickle: Although not as efficient in terms of I/O speed and storage as Parquet or Feather,
        # Pickle is convenient if you want to preserve exact Python object types.
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_pickle")
            return None
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "pkl")
        compression = kwargs.pop("compression", None)
        if compression is not None:
            if not fpath.endswith(compression["method"]):
                fpath = fpath + f".{compression['method']}"
        if isinstance(content, pd.DataFrame):
            content.to_pickle(fpath, **kwargs)
        else:
            try:
                content = pd.DataFrame(content)
                content.to_pickle(fpath, **kwargs)
            except Exception as e:
                try:
                    import pickle

                    with open(fpath, "wb") as f:
                        pickle.dump(content, f)
                    print("done!", fpath)
                except Exception as e:
                    raise ValueError(
                        f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                    )
    elif kind.lower() in ["fea", "feather", "ft", "fe", "feat", "fether"]:
        # Feather: The Feather format, based on Apache Arrow, is designed for fast I/O operations. It's
        # optimized for data analytics tasks and is especially fast when working with Pandas.

        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_feather")
            return None
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "feather")
        if isinstance(content, pd.DataFrame):
            content.to_feather(fpath, **kwargs)
        else:
            try:
                print("trying to convert it as a DataFrame...")
                content = pd.DataFrame(content)
                content.to_feather(fpath, **kwargs)
            except Exception as e:
                raise ValueError(
                    f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                )
    elif kind.lower() in ["hd", "hdf", "h", "h5"]:
        # particularly useful for large datasets and can handle complex data structures
        verbose = kwargs.pop("verbose", False)
        if verbose:
            print(str_space_speed)
            use_pd("to_hdf")
        _fpath, _ext = os.path.splitext(fpath)
        fpath = _fpath + _ext.replace(kind, "h5")
        compression = kwargs.pop("compression", None)
        if compression is not None:
            if not fpath.endswith(compression):
                fpath = fpath + f".{compression}"
        if isinstance(content, pd.DataFrame):
            content.to_hdf(fpath, key="content", **kwargs)
        else:
            try:
                print("trying to convert it as a DataFrame...")
                content = pd.DataFrame(content)
                content.to_hdf(fpath, **kwargs)
            except Exception as e:
                raise ValueError(
                    f"content is not a DataFrame, cannot be saved as a 'pkl' format: {e}"
                )
    else:
        from . import netfinder

        try:
            netfinder.downloader(url=content, dir_save=dirname(fpath), kind=kind)
        except:
            print(
                f"Error:\n{kind} is not in the supported list ['docx', 'txt', 'md', 'html', 'pdf', 'csv', 'xlsx', 'json', 'xml', 'yaml']"
            )
def addpath(fpath):
    sys.path.insert(0, dir)


def dirname(fpath):
    """
    dirname: Extracting Directory Name from a File Path
    Args:
        fpath (str): the file or directory path
    Returns:
        str: directory, without filename
    """
    dirname_ = os.path.dirname(fpath)
    if not dirname_.endswith("/"):
        dirname_ = dirname_ + "/"
    return dirname_


def dir_name(fpath):  # same as "dirname"
    return dirname(fpath)


def basename(fpath):
    """
    basename: # Output: file.txt
    Args:
        fpath (str): the file or directory path
    Returns:
        str: # Output: file.txt
    """
    return os.path.basename(fpath)


def flist(fpath, contains="all"):
    all_files = [
        os.path.join(fpath, f)
        for f in os.listdir(fpath)
        if os.path.isfile(os.path.join(fpath, f))
    ]
    if isinstance(contains, list):
        filt_files = []
        for filter_ in contains:
            filt_files.extend(flist(fpath, filter_))
        return filt_files
    else:
        if "all" in contains.lower():
            return all_files
        else:
            filt_files = [f for f in all_files if isa(f, contains)]
            return filt_files


def sort_kind(df, by="name", ascending=True):
    if df[by].dtype == "object":  # Check if the column contains string values
        if ascending:
            sorted_index = df[by].str.lower().argsort()
        else:
            sorted_index = df[by].str.lower().argsort()[::-1]
    else:
        if ascending:
            sorted_index = df[by].argsort()
        else:
            sorted_index = df[by].argsort()[::-1]
    sorted_df = df.iloc[sorted_index].reset_index(drop=True)
    return sorted_df


def isa(content, kind):
    """
    content, kind='img'
    kinds file paths based on the specified kind.
    Args:
        content (str): Path to the file.
        kind (str): kind of file to kind. Default is 'img' for images. Other options include 'doc' for documents,
                    'zip' for ZIP archives, and 'other' for other types of files.
    Returns:
        bool: True if the file matches the kind, False otherwise.
    """
    if "img" in kind.lower() or "image" in kind.lower():
        return is_image(content)
    elif "vid" in kind.lower():
        return is_video(content)
    elif "aud" in kind.lower():
        return is_audio(content)
    elif "doc" in kind.lower():
        return is_document(content)
    elif "zip" in kind.lower():
        return is_zip(content)
    elif "dir" in kind.lower() or ("f" in kind.lower() and "d" in kind.lower()):
        return os.path.isdir(content)
    elif "code" in kind.lower():  # file
        return is_code(content)
    elif "fi" in kind.lower():  # file
        return os.path.isfile(content)
    elif "num" in kind.lower():  # file
        return isnum(content)
    elif "text" in kind.lower() or "txt" in kind.lower():  # file
        return is_text(content)
    elif "color" in kind.lower():  # file
        return is_str_color(content)
    elif "html" in kind.lower():
        import re

        if content is None or not isinstance(content, str):
            return False
        # Remove leading and trailing whitespace
        content = content.strip()
        # Check for common HTML tags using regex
        # This pattern matches anything that looks like an HTML tag
        tag_pattern = r"<[a-zA-Z][^>]*>(.*?)</[a-zA-Z][^>]*>"
        # Check for any opening and closing tags
        if re.search(tag_pattern, content):
            return True
        # Additional checks for self-closing tags
        self_closing_tags = ["img", "br", "hr", "input", "meta", "link"]
        for tag in self_closing_tags:
            if f"<{tag}" in content:
                return True
        return False
    else:
        print(f"{kind} was not set up correctly")
        return False


def get_os(full=False, verbose=False):
    """Collects comprehensive system information.
    full(bool): True, get more detailed info
    verbose(bool): True, print it
    usage:
        info = get_os(full=True, verbose=False)
    """
    import sys
    import platform
    import psutil
    # import GPUtil
    import socket
    import uuid
    import cpuinfo
    import os
    import subprocess
    from datetime import datetime, timedelta

    def get_os_type():
        os_name = sys.platform
        if "dar" in os_name:
            return "macOS"
        else:
            if "win" in os_name:
                return "Windows"
            elif "linux" in os_name:
                return "Linux"
            else:
                print(f"{os_name}, returned 'None'")
                return None

    if not full:
        return get_os_type()

    def get_os_info():
        """Get the detailed OS name, version, and other platform-specific details."""

        def get_mac_os_info():
            """Get detailed macOS version and product name."""
            try:
                sw_vers = subprocess.check_output(["sw_vers"]).decode("utf-8")
                product_name = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductName")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                product_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                build_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("BuildVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )

                # Return the formatted macOS name, version, and build
                return f"{product_name} {product_version} (Build {build_version})"
            except Exception as e:
                return f"Error retrieving macOS name: {str(e)}"

        def get_windows_info():
            """Get detailed Windows version and edition."""
            try:
                # Get basic Windows version using platform
                windows_version = platform.version()
                release = platform.release()
                version = platform.win32_ver()[0]

                # Additional information using Windows-specific system commands
                edition_command = "wmic os get caption"
                edition = (
                    subprocess.check_output(edition_command, shell=True)
                    .decode("utf-8")
                    .strip()
                    .split("\n")[1]
                )

                # Return Windows information
                return f"Windows {version} {release} ({edition})"
            except Exception as e:
                return f"Error retrieving Windows information: {str(e)}"

        def get_linux_info():
            """Get detailed Linux version and distribution info."""
            try:
                # Check /etc/os-release for modern Linux distros
                with open("/etc/os-release") as f:
                    os_info = f.readlines()

                os_name = (
                    next(line for line in os_info if line.startswith("NAME"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                os_version = (
                    next(line for line in os_info if line.startswith("VERSION"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )

                # For additional info, check for the package manager (e.g., apt, dnf)
                package_manager = "Unknown"
                if os.path.exists("/usr/bin/apt"):
                    package_manager = "APT (Debian/Ubuntu)"
                elif os.path.exists("/usr/bin/dnf"):
                    package_manager = "DNF (Fedora/RHEL)"

                # Return Linux distribution, version, and package manager
                return f"{os_name} {os_version} (Package Manager: {package_manager})"
            except Exception as e:
                return f"Error retrieving Linux information: {str(e)}"

        os_name = platform.system()

        if os_name == "Darwin":
            return get_mac_os_info()
        elif os_name == "Windows":
            return get_windows_info()
        elif os_name == "Linux":
            return get_linux_info()
        else:
            return f"Unknown OS: {os_name} {platform.release()}"

    def get_os_name_and_version():
        os_name = platform.system()
        if os_name == "Darwin":
            try:
                # Run 'sw_vers' command to get macOS details like "macOS Sequoia"
                sw_vers = subprocess.check_output(["sw_vers"]).decode("utf-8")
                product_name = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductName")
                    ][0]
                    .split(":")[1]
                    .strip()
                )
                product_version = (
                    [
                        line
                        for line in sw_vers.split("\n")
                        if line.startswith("ProductVersion")
                    ][0]
                    .split(":")[1]
                    .strip()
                )

                # Return the formatted macOS name and version
                return f"{product_name} {product_version}"

            except Exception as e:
                return f"Error retrieving macOS name: {str(e)}"

        # For Windows, we use platform to get the OS name and version
        elif os_name == "Windows":
            os_version = platform.version()
            return f"Windows {os_version}"

        # For Linux, check for distribution info using platform and os-release file
        elif os_name == "Linux":
            try:
                # Try to read Linux distribution info from '/etc/os-release'
                with open("/etc/os-release") as f:
                    os_info = f.readlines()

                # Find fields like NAME and VERSION
                os_name = (
                    next(line for line in os_info if line.startswith("NAME"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                os_version = (
                    next(line for line in os_info if line.startswith("VERSION"))
                    .split("=")[1]
                    .strip()
                    .replace('"', "")
                )
                return f"{os_name} {os_version}"

            except Exception as e:
                return f"Error retrieving Linux name: {str(e)}"

        # Default fallback (for unknown OS or edge cases)
        return f"{os_name} {platform.release()}"

    def get_system_uptime():
        """Returns system uptime as a human-readable string."""
        try:
            boot_time = datetime.fromtimestamp(psutil.boot_time())
            uptime = datetime.now() - boot_time
            return str(uptime).split(".")[0]  # Remove microseconds
        except:
            return None

    def get_active_processes(limit=10):
        try:
            processes = []
            for proc in psutil.process_iter(
                ["pid", "name", "cpu_percent", "memory_percent"]
            ):
                try:
                    processes.append(proc.info)
                except psutil.NoSuchProcess:
                    pass
            # Handle NoneType values by treating them as 0
            processes.sort(key=lambda x: x["cpu_percent"] or 0, reverse=True)
            return processes[:limit]
        except:
            return None

    def get_virtual_environment_info():
        """Checks if the script is running in a virtual environment and returns details."""
        try:
            # Check if running in a virtual environment
            if hasattr(sys, "real_prefix") or (
                hasattr(sys, "base_prefix") and sys.base_prefix != sys.prefix
            ):
                return {
                    "Virtual Environment": sys.prefix,
                    "Site-Packages Path": os.path.join(
                        sys.prefix,
                        "lib",
                        "python{}/site-packages".format(sys.version_info.major),
                    ),
                }
            else:
                return {"Virtual Environment": "Not in a virtual environment"}
        except Exception as e:
            return {"Error": str(e)}

    def get_temperatures():
        """Returns temperature sensor readings."""
        try:
            return psutil.sensors_temperatures(fahrenheit=False)
        except AttributeError:
            return {"Error": "Temperature sensors not available"}

    def get_battery_status():
        """Returns battery status."""
        try:
            battery = psutil.sensors_battery()
            if battery:
                time_left = (
                    str(timedelta(seconds=battery.secsleft))
                    if battery.secsleft != psutil.POWER_TIME_UNLIMITED
                    else "Charging/Unlimited"
                )
                return {
                    "Percentage": battery.percent,
                    "Plugged In": battery.power_plugged,
                    "Time Left": time_left,
                }
            return {"Status": "No battery detected"}
        except:
            return {"Status": "No battery detected"}

    def get_disk_io():
        """Returns disk I/O statistics."""
        disk_io = psutil.disk_io_counters()
        return {
            "Read (GB)": disk_io.read_bytes / (1024**3),
            "Write (GB)": disk_io.write_bytes / (1024**3),
            "Read Count": disk_io.read_count,
            "Write Count": disk_io.write_count,
        }

    def get_network_io():
        """Returns network I/O statistics."""
        net_io = psutil.net_io_counters()
        return {
            "Bytes Sent (GB)": net_io.bytes_sent / (1024**3),
            "Bytes Received (GB)": net_io.bytes_recv / (1024**3),
            "Packets Sent": net_io.packets_sent,
            "Packets Received": net_io.packets_recv,
        }

    def run_shell_command(command):
        """Runs a shell command and returns its output."""
        try:
            result = subprocess.run(
                command,
                shell=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            return (
                result.stdout.strip()
                if result.returncode == 0
                else result.stderr.strip()
            )
        except Exception as e:
            return f"Error running command: {e}"

    system_info = {
        "timestamp": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "os": get_os_type(),
        "system": {
            "os": get_os_info(),
            "platform": f"{platform.system()} {platform.release()}",
            "version": platform.version(),
            "machine": platform.machine(),
            "processor": platform.processor(),
            "architecture": platform.architecture()[0],
            "hostname": socket.gethostname(),
            "ip address": socket.gethostbyname(socket.gethostname()),
            "mac address": ":".join(
                ["{:02x}".format((uuid.getnode() >> i) & 0xFF) for i in range(0, 48, 8)]
            ),
            "cpu brand": cpuinfo.get_cpu_info().get("brand_raw", "Unknown"),
            "python version": platform.python_version(),
            "uptime": get_system_uptime(),
        },
        "cpu": {
            "physical cores": psutil.cpu_count(logical=False),
            "logical cores": psutil.cpu_count(logical=True),
            "max frequency (MHz)": psutil.cpu_freq().max,
            "min frequency (MHz)": psutil.cpu_freq().min,
            "current frequency (MHz)": psutil.cpu_freq().current,
            "usage per core (%)": psutil.cpu_percent(percpu=True),
            "total cpu Usage (%)": psutil.cpu_percent(),
            "load average (1m, 5m, 15m)": (
                os.getloadavg() if hasattr(os, "getloadavg") else "N/A"
            ),
        },
        "memory": {
            "total memory (GB)": psutil.virtual_memory().total / (1024**3),
            "available memory (GB)": psutil.virtual_memory().available / (1024**3),
            "used memory (GB)": psutil.virtual_memory().used / (1024**3),
            "memory usage (%)": psutil.virtual_memory().percent,
            "swap total (GB)": psutil.swap_memory().total / (1024**3),
            "swap free (GB)": psutil.swap_memory().free / (1024**3),
            "swap used (GB)": psutil.swap_memory().used / (1024**3),
            "swap usage (%)": psutil.swap_memory().percent,
        },
        "disk": {},
        "disk io": get_disk_io(),
        "network": {},
        "network io": get_network_io(),
        "gpu": [],
        # "temperatures": get_temperatures(),
        # "battery": get_battery_status(),
        "active processes": get_active_processes(),
        "environment": {
            "user": os.getenv("USER", "Unknown"),
            "environment variables": dict(os.environ),
            "virtual environment info": get_virtual_environment_info(),  # Virtual env details
            "docker running": os.path.exists("/.dockerenv"),  # Check for Docker
            "shell": os.environ.get("SHELL", "Unknown"),
            "default terminal": run_shell_command("echo $TERM"),
            "kernel version": platform.uname().release,
            "virtualization type": run_shell_command("systemd-detect-virt"),
        },
        "additional info": {
            "Shell": os.environ.get("SHELL", "Unknown"),
            "default terminal": run_shell_command("echo $TERM"),
            "kernel version": platform.uname().release,
            "virtualization type": run_shell_command("systemd-detect-virt"),
            "running in docker": os.path.exists("/.dockerenv"),
        },
    }

    # Disk Information
    for partition in psutil.disk_partitions():
        try:
            usage = psutil.disk_usage(partition.mountpoint)
            system_info["disk"][partition.device] = {
                "mountpoint": partition.mountpoint,
                "file system type": partition.fstype,
                "total size (GB)": usage.total / (1024**3),
                "used (GB)": usage.used / (1024**3),
                "free (GB)": usage.free / (1024**3),
                "usage (%)": usage.percent,
            }
        except PermissionError:
            system_info["disk"][partition.device] = "Permission Denied"

    # Network Information
    if_addrs = psutil.net_if_addrs()
    for interface_name, interface_addresses in if_addrs.items():
        system_info["network"][interface_name] = []
        for address in interface_addresses:
            if str(address.family) == "AddressFamily.AF_INET":
                system_info["network"][interface_name].append(
                    {
                        "ip address": address.address,
                        "netmask": address.netmask,
                        "broadcast ip": address.broadcast,
                    }
                )
            elif str(address.family) == "AddressFamily.AF_PACKET":
                system_info["network"][interface_name].append(
                    {
                        "mac address": address.address,
                        "netmask": address.netmask,
                        "broadcast mac": address.broadcast,
                    }
                )

    # # GPU Information
    # gpus = GPUtil.getGPUs()
    # for gpu in gpus:
    #     gpu_info = {
    #         "name": gpu.name,
    #         "load (%)": gpu.load * 100,
    #         "free memory (MB)": gpu.memoryFree,
    #         "used memory (MB)": gpu.memoryUsed,
    #         "total memory (MB)": gpu.memoryTotal,
    #         "driver version": gpu.driver,
    #         "temperature (°C)": gpu.temperature,
    #     }
    #     if hasattr(gpu, "powerDraw"):
    #         gpu_info["Power Draw (W)"] = gpu.powerDraw
    #     if hasattr(gpu, "powerLimit"):
    #         gpu_info["Power Limit (W)"] = gpu.powerLimit
    #     system_info["gpu"].append(gpu_info)

    res = system_info if full else get_os_type()
    if verbose:
        try:
            preview(res)
        except Exception as e:
            pnrint(e)
    return res


def listdir(
    rootdir,
    kind=None,
    sort_by="name",
    ascending=True,
    contains=None,  # filter filenames using re
    booster=False,  # walk in subfolders
    depth=0,  # 0: no subfolders; None: all subfolders; [int 1,2,3]: levels of subfolders
    hidden=False,  # Include hidden files/folders
    orient="list",
    output="df",  # "df", 'list','dict','records','index','series'
    verbose=False,
):
    def is_hidden(filepath):
        """Check if a file or folder is hidden."""
        system = platform.system()
        if system == "Windows":
            import ctypes

            attribute = ctypes.windll.kernel32.GetFileAttributesW(filepath)
            if attribute == -1:
                raise FileNotFoundError(f"File {filepath} not found.")
            return bool(attribute & 2)  # FILE_ATTRIBUTE_HIDDEN
        else:  # macOS/Linux: Hidden if the name starts with a dot
            return os.path.basename(filepath).startswith(".")

    def get_user():
        """Retrieve the username of the current user."""
        system = platform.system()
        if system == "Windows":
            return os.environ.get("USERNAME", "Unknown")
        else:
            import pwd

            return pwd.getpwuid(os.getuid()).pw_name

    if isinstance(kind, list):
        f_ = []
        for kind_ in kind:
            f_tmp = listdir(
                rootdir=rootdir,
                kind=kind_,
                sort_by=sort_by,
                ascending=ascending,
                contains=contains,
                depth=depth,  # walk in subfolders
                hidden=hidden,
                orient=orient,
                output=output,
                verbose=verbose,
            )
            f_.append(f_tmp)
        if f_:
            return pd.concat(f_, ignore_index=True)
    if kind is not None:
        if not kind.startswith("."):
            kind = "." + kind
    fd = [".fd", ".fld", ".fol", ".fd", ".folder"]
    i = 0
    f = {
        "name": [],
        "kind": [],
        "length": [],
        "basename": [],
        "path": [],
        "created_time": [],
        "modified_time": [],
        "last_open_time": [],
        "size": [],
        "permission": [],
        "owner": [],
        "rootdir": [],
        "fname": [],
        "fpath": [],
        "num": [],
        "os": [],
    }
    root_depth = rootdir.rstrip(os.sep).count(os.sep)
    for dirpath, dirnames, ls in os.walk(rootdir):
        current_depth = dirpath.rstrip(os.sep).count(os.sep) - root_depth
        # Check depth limit
        if depth is not None and current_depth > depth:
            dirnames[:] = []  # Prevent further traversal into subfolders
            continue

        if not hidden:
            dirnames[:] = [
                d for d in dirnames if not is_hidden(os.path.join(dirpath, d))
            ]
            ls = [i for i in ls if not is_hidden(os.path.join(dirpath, i))]

        for dirname in dirnames:
            if kind is not None and kind not in fd:  # do not check folders
                continue
            if contains and not re.search(contains, dirname):
                continue
            dirname_path = os.path.join(dirpath, dirname)
            fpath = os.path.join(os.path.dirname(dirname_path), dirname)
            try:
                stats_file = os.stat(fpath)
            except Exception as e:
                print(e)
                continue
            filename, file_extension = os.path.splitext(dirname)
            file_extension = file_extension if file_extension != "" else None
            f["name"].append(filename)
            f["kind"].append(file_extension)
            f["length"].append(len(filename))
            f["size"].append(round(os.path.getsize(fpath) / 1024 / 1024, 3))
            f["basename"].append(os.path.basename(dirname_path))
            f["path"].append(os.path.join(os.path.dirname(dirname_path), dirname))
            f["created_time"].append(
                pd.to_datetime(int(os.path.getctime(dirname_path)), unit="s")
            )
            f["modified_time"].append(
                pd.to_datetime(int(os.path.getmtime(dirname_path)), unit="s")
            )
            f["last_open_time"].append(
                pd.to_datetime(int(os.path.getatime(dirname_path)), unit="s")
            )
            f["permission"].append(stat.filemode(stats_file.st_mode)),
            f["owner"].append(get_user()),
            f["rootdir"].append(dirpath)
            f["fname"].append(filename)  # will be removed
            f["fpath"].append(fpath)  # will be removed
            i += 1
        for item in ls:
            if kind in fd:  # only check folders
                continue
            if contains and not re.search(contains, item):
                continue
            item_path = os.path.join(dirpath, item)
            fpath = os.path.join(os.path.dirname(item_path), item)
            try:
                stats_file = os.stat(fpath)
            except Exception as e:
                print(e)
                continue
            filename, file_extension = os.path.splitext(item)
            if kind is not None:
                is_folder = kind.lower() in fd and os.path.isdir(item_path)
                is_file = kind.lower() in file_extension.lower() and (
                    os.path.isfile(item_path)
                )
                if kind in [
                    ".doc",
                    ".img",
                    ".zip",
                    ".code",
                    ".file",
                    ".image",
                    ".video",
                    ".audio",
                ]:  # 选择大的类别
                    if kind != ".folder" and not isa(item_path, kind):
                        continue
                elif kind in [".all"]:
                    return flist(fpath, contains=contains)
                else:  # 精确到文件的后缀
                    if not is_folder and not is_file:
                        continue
            file_extension = file_extension if file_extension != "" else None
            f["name"].append(filename)
            f["kind"].append(file_extension)
            f["length"].append(len(filename))
            f["size"].append(round(os.path.getsize(fpath) / 1024 / 1024, 3))
            f["basename"].append(os.path.basename(item_path))
            f["path"].append(os.path.join(os.path.dirname(item_path), item))
            f["created_time"].append(
                pd.to_datetime(int(os.path.getctime(item_path)), unit="s")
            )
            f["modified_time"].append(
                pd.to_datetime(int(os.path.getmtime(item_path)), unit="s")
            )
            f["last_open_time"].append(
                pd.to_datetime(int(os.path.getatime(item_path)), unit="s")
            )
            f["permission"].append(stat.filemode(stats_file.st_mode)),
            f["owner"].append(
                os.getlogin() if platform.system() != "Windows" else "N/A"
            ),
            f["fname"].append(filename)  # will be removed
            f["fpath"].append(fpath)  # will be removed
            f["rootdir"].append(dirpath)
            i += 1

        f["num"] = i
        f["os"] = get_os()  # os.uname().machine
        # if not booster: # go deeper subfolders
        #     break
    # * convert to pd.DataFrame
    f = pd.DataFrame(f)
    f = f[
        [
            "basename",
            "name",
            "kind",
            "length",
            "size",
            "num",
            "path",
            "created_time",
            "modified_time",
            "last_open_time",
            "rootdir",
            "permission",
            "owner",
            "os",
            "fname",
            "fpath",
        ]
    ]
    if "nam" in sort_by.lower():
        f = sort_kind(f, by="name", ascending=ascending)
    elif "crea" in sort_by.lower():
        f = sort_kind(f, by="created_time", ascending=ascending)
    elif "modi" in sort_by.lower():
        f = sort_kind(f, by="modified_time", ascending=ascending)
    elif "s" in sort_by.lower() and "z" in sort_by.lower():
        f = sort_kind(f, by="size", ascending=ascending)

    if "df" in output:
        if verbose:
            display(f.head())
            print(f"shape: {f.shape}")
        return f
    else:
        from box import Box

        if "l" in orient.lower():  # list # default
            res_output = Box(f.to_dict(orient="list"))
            return res_output
        if "d" in orient.lower():  # dict
            return Box(f.to_dict(orient="dict"))
        if "r" in orient.lower():  # records
            return Box(f.to_dict(orient="records"))
        if "in" in orient.lower():  # records
            return Box(f.to_dict(orient="index"))
        if "se" in orient.lower():  # records
            return Box(f.to_dict(orient="series"))

 
def listpkg(where="env", verbose=False):
    """list all pacakages"""

    def listfunc_(lib_name, opt="call"):
        """list functions in specific lib"""
        if opt == "call":
            funcs = [
                func
                for func in dir(lib_name)
                if callable(getattr(lib_name, func))
                if not func.startswith("__")
            ]
        else:
            funcs = dir(lib_name)
        return funcs

    if any([i in where.lower() for i in ["all", "env"]]):
        import pkg_resources

        lst_pkg = [pkg.key for pkg in pkg_resources.working_set]
    else:
        lst_pkg = listfunc_(where)
    print(lst_pkg) if verbose else None
    return lst_pkg



def local_path(fpath,station=r"Q:\\IM\\AGLengerke\\Jeff\\# testing\\temp\\"):
    """copy file to a specific folder first, to aviod file conflict"""
    try:
        new_path=fbackup(fpath, station, interval=12*3600,verbose=0)
    except Exception as e:
        print(f"Path did not update because: Error:{e}")
        new_path=fpath
    return new_path
 
 


def fmove(
    src: Union[str, Path],
    dst: Union[str, Path],
    overwrite: bool = False,
    verbose: bool = True,
    filter: Optional[Union[str, List[str]]] = None,
    booster: bool = False
) -> None:
    """
    Enhanced move function with filter and booster support.
    
    Args:
        src: Source path (file/folder) or directory when using filter
        dst: Destination path
        overwrite: Whether to overwrite existing files
        verbose: Show operation details
        filter: Pattern(s) for selective moving:
            - None: Original behavior (move exact path)
            - str: File extension (e.g., '.txt') or name pattern (e.g., 'temp*')
            - List[str]: Multiple patterns
        booster: Search subdirectories when using filter
    """
    def _move_single(src: Path, dst: Path, overwrite: bool, verbose: bool) -> bool:
        """Original move logic for single file/folder"""
        try:
            dir_par_dst = os.path.dirname(str(dst))
            if not os.path.isdir(dir_par_dst):
                os.makedirs(dir_par_dst, exist_ok=True)
                
            if dst.is_dir():
                dst = dst / src.name
                
            if dst.exists():
                if overwrite:
                    if dst.is_file():
                        dst.unlink()
                    else:
                        shutil.rmtree(dst)
                else:
                    dst = dst.with_name(
                        f"{dst.stem}_{datetime.now().strftime('_%H%M%S')}{dst.suffix}"
                    )
                    
            shutil.move(str(src), str(dst))
            if verbose:
                print(f"\nDone! Moved to {dst}\n")
            return True
            
        except Exception as e:
            logging.error(f"Failed to move {src} to {dst}: {str(e)}")
            return False

    try:
        src_path = Path(src).resolve()
        dst_path = Path(dst).resolve()
        # Prevent infinite loop check
        if booster:
            try:
                if dst_path.is_relative_to(src_path):
                    logging.error("Destination cannot be inside source directory when using booster mode")
                    return None
            except ValueError:
                pass  # Different drive case
        if filter is None:
            return _move_single(src_path, dst_path, overwrite, verbose)
            
        if not src_path.exists():
            print(f"Source path '{src_path}' does not exist")
            return
            
        if not src_path.is_dir():
            print("Filter mode requires source to be a directory")
            return
            
        filters = [filter] if isinstance(filter, str) else filter
        moved_count = 0
        
        for pattern in filters:
            # Handle extension patterns (starting with .)
            if pattern.startswith('.'):
                search_pattern = f"*{pattern}"
            else:
                search_pattern = pattern
                
            matches = src_path.rglob(search_pattern) if booster else src_path.glob(search_pattern)
            
            for item in matches:
                if not item.exists():  # Skip broken symlinks
                    continue
                    
                # Calculate relative path and create destination path
                rel_path = item.relative_to(src_path)
                target_path = dst_path / rel_path
                
                # Ensure parent directory exists
                target_path.parent.mkdir(parents=True, exist_ok=True)
                
                if _move_single(item, target_path, overwrite, verbose):
                    moved_count += 1
        
        if verbose:
            print(f"\nDone! Moved {moved_count} items matching filter(s) {filters}\n")
            
    except Exception as e:
        logging.error(f"Move error: {str(e)}")


def fcopy(
    src: Union[str, Path],
    dst: Union[str, Path],
    overwrite: bool = False,
    verbose: bool = True,
    filter: Optional[Union[str, List[str]]] = None,
    booster: bool = False
) -> Optional[Path]:
    """
    Enhanced copy function with filter and booster support.
    
    Args:
        src: Source path (file/folder) or directory when using filter
        dst: Destination path
        overwrite: Whether to overwrite existing files
        verbose: Show operation details
        filter: Pattern(s) for selective copying:
            - None: Original behavior (copy exact path)
            - str: File extension (e.g., '.txt') or name pattern (e.g., 'temp*')
            - List[str]: Multiple patterns
        booster: Search subdirectories when using filter
    
    Returns:
        Path to the last copied item or None if failed
    """
    def _copy_single(src: Path, dst: Path, overwrite: bool, verbose: bool) -> Optional[Path]:
        """Original copy logic for single file/folder"""
        try:
            dir_par_dst = os.path.dirname(str(dst))
            if not os.path.isdir(dir_par_dst):
                os.makedirs(dir_par_dst, exist_ok=True)
                
            if not src.is_dir():
                if dst.is_dir():
                    dst = dst / src.name
                    
                if dst.exists():
                    if overwrite:
                        dst.unlink()
                    else:
                        dst = dst.with_name(
                            f"{dst.stem}_{datetime.now().strftime('%y%m%d_%H%M%S')}{dst.suffix}"
                        )
                        
                shutil.copy(str(src), str(dst))
                if verbose:
                    print(f"\nDone! Copied to {dst}\n")
                return dst
                
            else:
                dst = dst / src.name
                if dst.exists():
                    if overwrite:
                        shutil.rmtree(str(dst))
                    else:
                        dst = dst.with_name(
                            f"{dst.stem}_{datetime.now().strftime('%y%m%d%H%M%S')}"
                        )
                        
                shutil.copytree(str(src), str(dst))
                if verbose:
                    print(f"\nDone! Copied to {dst}\n")
                return dst
                
        except Exception as e:
            logging.error(f"Failed to copy {src} to {dst}: {str(e)}")
            return None

    try:
        src_path = Path(src).resolve()
        dst_path = Path(dst).resolve()
        # Prevent infinite loop check
        if booster:
            try:
                if dst_path.is_relative_to(src_path):
                    logging.error("Destination cannot be inside source directory when using booster mode")
                    return None
            except ValueError:
                pass  # Different drive case
        if filter is None:
            return _copy_single(src_path, dst_path, overwrite, verbose)
            
        if not src_path.exists():
            print(f"Source path '{src_path}' does not exist")
            return None
            
        if not src_path.is_dir():
            print("Filter mode requires source to be a directory")
            return None
            
        filters = [filter] if isinstance(filter, str) else filter
        copied_count = 0
        last_copied = None
        
        for pattern in filters:
            # Handle extension patterns (starting with .)
            if pattern.startswith('.'):
                search_pattern = f"*{pattern}"
            else:
                search_pattern = pattern
                
            matches = src_path.rglob(search_pattern) if booster else src_path.glob(search_pattern)
            
            for item in matches:
                if not item.exists():  # Skip broken symlinks
                    continue
                    
                # Calculate relative path and create destination path
                rel_path = item.relative_to(src_path)
                target_path = dst_path / rel_path
                
                # Ensure parent directory exists
                target_path.parent.mkdir(parents=True, exist_ok=True)
                
                result = _copy_single(item, target_path, overwrite, verbose)
                if result:
                    copied_count += 1
                    last_copied = result
        
        if verbose:
            print(f"\nDone! Copied {copied_count} items matching filter(s) {filters}\n")
        return last_copied
        
    except Exception as e:
        logging.error(f"Copy error: {str(e)}")
        return None


def delete(
    fpath: Union[str, Path],
    filter: Optional[Union[str, List[str]]] = None,
    booster: bool = False,
) -> None:
    """
    Powerful deletion function with booster mode for deep searching.

    Args:
        fpath: Path to file/folder OR parent directory when using filter
        filter: Optional filter pattern(s) for selective deletion:
            - None: Delete exact path (original behavior)
            - str:
                - If contains '.' treat as extension (e.g., '.tmp')
                - Else treat as exact filename (e.g., 'tempfile')
            - List[str]: Multiple filters to apply
        booster: Enable deep searching through all subdirectories (supercharged mode)

    Behavior:
    1. When filter is None: Original behavior (delete exact path)
    2. With filter: Delete matching items in fpath's parent directory
    3. With booster=True: Search deeply through all subdirectories
    """


    def _delete_single_path(path: Path) -> bool:
        """Delete a single file or folder, returns success status"""
        try:
            if not path.exists():
                print(f"Path '{path}' does not exist")
                return False

            if path.is_file():
                path.unlink()
                print(f"Deleted file: {path}")
                return True
            elif path.is_dir():
                shutil.rmtree(path)
                print(f"Deleted folder: {path}")
                return True
            return False
        except Exception as e:
            logging.error(f"Failed to delete {path}: {str(e)}")
            return False
    try:
        fpath = Path(fpath).resolve()

        # Original behavior when no filter
        if filter is None:
            _delete_single_path(fpath)
            return

        # Convert single filter to list for uniform processing
        filters = [filter] if isinstance(filter, str) else filter

        # Determine search root directory
        parent = fpath.parent if fpath.is_file() else fpath
        if not parent.exists():
            print(f"Path '{parent}' does not exist")
            return

        deleted_count = 0

        # Process each filter
        for pattern in filters:
            # Determine search method based on booster mode
            if booster:
                matches = parent.rglob(f"*{pattern}" if "." in pattern else pattern)
            else:
                matches = parent.glob(f"*{pattern}" if "." in pattern else pattern)

            for item in matches:
                if _delete_single_path(item):
                    deleted_count += 1

        print(f"\nDone! Deleted {deleted_count} items matching filter(s) {filters}\n")

    except Exception as e:
        logging.error(f"Deletion error: {str(e)}")
 

def rename(fpath, dst, smart=True):
    """Rename a file or folder."""
    try:
        src_kind, dst_kind = None, None
        if smart:
            dir_name_src = os.path.dirname(fpath)
            dir_name_dst = os.path.dirname(dst)
            src_kind = os.path.splitext(fpath)[1]
            dst_kind = os.path.splitext(dst)[1]
            if dir_name_dst != dir_name_src:
                dst = os.path.join(dir_name_src, dst)
            if dst_kind is not None and src_kind is not None:
                if dst_kind != src_kind:
                    dst = dst + src_kind
        if os.path.exists(fpath):
            os.rename(fpath, dst)
            print(f"Done! rename to {dst}")
        else:
            print(f"Failed: {fpath} does not exist.")
    except Exception as e:
        logging.error(f"Failed to rename {fpath} to {dst}: {e}")


def mkdir_nest(fpath: str) -> str:
    """
    Create nested directories based on the provided file path.

    Parameters:
    - fpath (str): The full file path for which the directories should be created.

    Returns:
    - str: The path of the created directory.
    """
    # Split the full path into directories
    f_slash = "/" if "mac" in get_os().lower() else "\\"
    if os.path.isdir(fpath):
        fpath = fpath + f_slash if not fpath.endswith(f_slash) else fpath
        return fpath
    dir_parts = fpath.split(f_slash)  # Split the path by the OS-specific separator

    # Start creating directories from the root to the desired path
    root_dir = os.path.splitdrive(fpath)[
        0
    ]  # Get the root drive on Windows (e.g., 'C:')
    current_path = (
        root_dir if root_dir else f_slash
    )  # Start from the root directory or POSIX '/'

    for part in dir_parts:
        if part:
            current_path = os.path.join(current_path, part)
            if not os.path.isdir(current_path):
                os.makedirs(current_path)
    if not current_path.endswith(f_slash):
        current_path += f_slash
    return current_path


def mkdir(pardir: str = None, chdir: str | list = None, overwrite=False):
    """
    Create a directory.

    Parameters:
    - pardir (str): Parent directory where the new directory will be created. If None, uses the current working directory.
    - chdir (str | list): Name of the new directory or a list of directories to create.
                          If None, a default name 'new_directory' will be used.
    - overwrite (bool): If True, overwrite the directory if it already exists. Defaults to False.

    Returns:
    - str: The path of the created directory or an error message.
    """
    rootdir = []
    pardir = mkdir_nest(pardir)
    if chdir is None:
        return pardir
    else:
        pass
    print(pardir)
    if isinstance(chdir, str):
        chdir = [chdir]
    chdir = list(set(chdir))
    if isinstance(pardir, str):  # Dir_parents should be 'str' type
        pardir = os.path.normpath(pardir)
    if "mac" in get_os().lower() or "lin" in get_os().lower():
        stype = "/"
    elif "win" in get_os().lower():
        stype = "\\"
    else:
        stype = "/"

    if os.path.isdir(pardir):
        os.chdir(pardir)  # Set current path
        # Check if subdirectories are not empty
        if chdir:
            chdir.sort()
            for folder in chdir:
                child_tmp = os.path.join(pardir, folder)
                if not os.path.isdir(child_tmp):
                    os.mkdir("./" + folder)
                    print(f"\n {folder} was created successfully!\n")
                else:
                    if overwrite:
                        shutil.rmtree(child_tmp)
                        os.mkdir("./" + folder)
                        print(f"\n {folder} overwrite! \n")
                    else:
                        print(f"\n {folder} already exists! \n")
                rootdir.append(child_tmp + stype)  # Note down
        else:
            print("\nWarning: Dir_child doesn't exist\n")
    else:
        print("\nWarning: Dir_parent is not a directory path\n")
    # Dir is the main output, if only one dir, then str type is inconvenient
    if len(rootdir) == 1:
        rootdir = rootdir[0]
        rootdir = rootdir + stype if not rootdir.endswith(stype) else rootdir

    return rootdir


def split_path(fpath):
    f_slash = "/" if "mac" in get_os().lower() else "\\"
    dir_par = f_slash.join(fpath.split(f_slash)[:-1])
    dir_ch = "".join(fpath.split(f_slash)[-1:])
    return dir_par, dir_ch
 
def figsave(*args, dpi=300, **kwargs):
    """
    Save a Matplotlib figure or image file in various formats.

    This function automatically determines whether to save a Matplotlib figure 
    or an image (PIL or NumPy array) and handles different file formats, including:
    - PDF, EPS, PNG, JPG, TIFF, ICO, EMF

    Parameters:
    -----------
    *args : str, PIL.Image, np.ndarray
        - File path (directory and/or filename) to save the figure.
        - If an image is provided (PIL or NumPy), it will be saved accordingly.
    
    dpi : int, optional (default=300)
        - Resolution (dots per inch) for saved figures.
    
    **kwargs : dict
        - Additional keyword arguments for `matplotlib.pyplot.savefig()`, including:
            - bbox_inches (str, default="tight"): Bounding box for figure.
            - pad_inches (float, default=0): Padding around figure.
            - facecolor (str, default="white"): Background color.
            - edgecolor (str, default="auto"): Edge color.

    Supported Formats:
    ------------------
    - Vector: `pdf`, `eps`, `emf`
    - Raster: `png`, `jpg`, `jpeg`, `tiff`, `tif`, `ico`
    
    Example Usage:
    --------------
    >>> figsave("output_plot.pdf")
    >>> figsave("figs/plot.png", dpi=600)
    >>> figsave("./results/figure", format="pdf")
    >>> figsave("icons/logo.ico", image)  # Save an image file as an icon

    Returns:
    --------
    None
    """
    import matplotlib.pyplot as plt
    from PIL import Image
    from pathlib import Path
    
    bbox_inches=kwargs.pop("bbox_inches","tight")
    pad_inches=kwargs.pop("pad_inches",0)
    facecolor=kwargs.pop("facecolor",'white')
    edgecolor=kwargs.pop("edgecolor",'auto')
    
    dir_save = None
    fname = None
    img = None
    for arg in args:
        if isinstance(arg, str):
            path = Path(arg)
            if path.suffix:  # Has file extension
                fname = path.name
                dir_save = path.parent
            else:
                dir_save = path
        elif isinstance(arg, (Image.Image, np.ndarray)):
            img = arg  # Store PIL image or numpy array

    dir_save = Path(dir_save) if dir_save else Path(".")
    dir_save.mkdir(parents=True, exist_ok=True)
    # Handle filename and extension
    if fname is None:
        fname = dir_save
    else:
        fname = dir_save / fname
    if fname.suffix == "":
        fname = fname.with_suffix(".pdf")  # Default format

    ftype = fname.suffix.lstrip(".").lower()

    # Save figure based on file type
    if ftype == "eps":
        plt.savefig(fname, format="eps", bbox_inches=bbox_inches)
        plt.savefig(fname.with_suffix(".pdf"), format="pdf", dpi=dpi,
                    pad_inches=pad_inches, bbox_inches=bbox_inches,
                    facecolor=facecolor, edgecolor=edgecolor) 
    elif ftype.lower() in ["jpg", "jpeg", "png", "tiff", "tif"]:
        if img is not None:  # If a PIL image is provided
            if isinstance(img, Image.Image):
                if img.mode == "RGBA":
                    img = img.convert("RGB")
                img.save(fname, format=ftype.upper(), dpi=(dpi, dpi))
            elif isinstance(img, np.ndarray):
                import cv2
                if img.ndim == 2:
                    # Grayscale image
                    Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                elif img.ndim == 3:
                    if img.shape[2] == 3:
                        # RGB image
                        img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
                        Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                    elif img.shape[2] == 4:
                        # RGBA image
                        img = cv2.cvtColor(img, cv2.COLOR_BGRA2RGBA)  # Convert BGRA to RGBA
                        Image.fromarray(img).save(fname, format=ftype.upper(), dpi=(dpi, dpi))
                    else:
                        raise ValueError("Unexpected number of channels in the image array.")
                else:
                    raise ValueError("Image array has an unexpected number of dimensions.")
        else:
            plt.savefig(fname,format=ftype.lower(),dpi=dpi,pad_inches=pad_inches,bbox_inches=bbox_inches,facecolor=facecolor,edgecolor=edgecolor)
    elif ftype.lower() in ["emf","pdf","fig"]:
        plt.savefig(fname,format=ftype.lower(),dpi=dpi,pad_inches=pad_inches,bbox_inches=bbox_inches,facecolor=facecolor,edgecolor=edgecolor)
    elif ftype.lower() == "ico":
        # Ensure the image is in a format that can be saved as an icon (e.g., 32x32, 64x64, etc.)
        if img is None:  # If no image is provided, use the matplotlib figure
            img = plt.figure() 
            print(fname)
            img.savefig(fname, 
                        format="png",
                        dpi=dpi,
                        pad_inches=pad_inches,
                        bbox_inches=bbox_inches,
                        facecolor=facecolor,
                        edgecolor=edgecolor )
            img = Image.open(fname)  # Load the saved figure image

        # Resize the image to typical icon sizes and save it as .ico
        icon_sizes = [(32, 32), (64, 64), (128, 128), (256, 256)]
        img = img.convert("RGBA") 
        img.save(fname, format="ICO", sizes=icon_sizes)
        print(f"Icon saved @: {fname} with sizes: {icon_sizes}")
    print(f"\nSaved @: dpi={dpi}\n{fname}")


def is_str_color(s):
    # Regular expression pattern for hexadecimal color codes
    if isinstance(s, str):
        import re

        color_code_pattern = r"^#([A-Fa-f0-9]{6}|[A-Fa-f0-9]{8})$"
        return re.match(color_code_pattern, s) is not None
    else:
        return True


def is_num(s):
    """
    Check if a string can be converted to a number (int or float).
    Parameters:
    - s (str): The string to check.
    Returns:
    - bool: True if the string can be converted to a number, False otherwise.
    """
    try:
        float(s)  # Try converting the string to a float
        return True
    except ValueError:
        return False


def isnum(s):
    return is_num(s)


def is_image(fpath):
    """
    Determine if a given file is an image based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized image, False otherwise.
    """
    from PIL import Image

    if isinstance(fpath, str):
        import mimetypes

        # Known image MIME types
        image_mime_types = {
            "image/jpeg",
            "image/png",
            "image/gif",
            "image/bmp",
            "image/webp",
            "image/tiff",
            "image/x-icon",
            "image/svg+xml",
            "image/heic",
            "image/heif",
        }

        # Known image file extensions
        image_extensions = {
            ".jpg",
            ".jpeg",
            ".png",
            ".gif",
            ".bmp",
            ".webp",
            ".tif",
            ".tiff",
            ".ico",
            ".svg",
            ".heic",
            ".heif",
            ".fig",
            ".jpg",
        }

        # Get MIME type using mimetypes
        mime_type, _ = mimetypes.guess_type(fpath)

        # Check MIME type
        if mime_type in image_mime_types:
            return True

        # Fallback: Check file extension
        ext = os.path.splitext(fpath)[
            -1
        ].lower()  # Get the file extension and ensure lowercase
        if ext in image_extensions:
            return True

        return False

    elif isinstance(fpath, Image.Image):
        # If the input is a PIL Image object
        return True

    return False


def is_video(fpath):
    """
    Determine if a given file is a video based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized video, False otherwise.
    """
    import mimetypes

    # Known video MIME types
    video_mime_types = {
        "video/mp4",
        "video/quicktime",
        "video/x-msvideo",
        "video/x-matroska",
        "video/x-flv",
        "video/webm",
        "video/ogg",
        "video/x-ms-wmv",
        "video/x-mpeg",
        "video/3gpp",
        "video/avi",
        "video/mpeg",
        "video/x-mpeg2",
        "video/x-ms-asf",
    }

    # Known video file extensions
    video_extensions = {
        ".mp4",
        ".mov",
        ".avi",
        ".mkv",
        ".flv",
        ".webm",
        ".ogv",
        ".wmv",
        ".mpg",
        ".mpeg",
        ".3gp",
        ".mpeg2",
        ".asf",
        ".ts",
        ".m4v",
        ".divx",
    }

    # Get MIME type using mimetypes
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type in video_mime_types:
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the file extension and ensure lowercase
    if ext in video_extensions:
        return True

    return False


def is_document(fpath):
    """
    Determine if a given file is a document based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized document, False otherwise.
    """
    import mimetypes

    # Define known MIME types for documents
    document_mime_types = {
        "text/",
        "application/pdf",
        "application/msword",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.ms-excel",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-powerpoint",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/rtf",
        "application/x-latex",
        "application/vnd.oasis.opendocument.text",
        "application/vnd.oasis.opendocument.spreadsheet",
        "application/vnd.oasis.opendocument.presentation",
    }

    # Define extensions for fallback
    document_extensions = {
        ".txt",
        ".log",
        ".csv",
        ".json",
        ".xml",
        ".pdf",
        ".doc",
        ".docx",
        ".xls",
        ".xlsx",
        ".ppt",
        ".pptx",
        ".odt",
        ".ods",
        ".odp",
        ".rtf",
        ".tex",
    }

    # Get MIME type
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type and any(
        mime_type.startswith(doc_type) for doc_type in document_mime_types
    ):
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the extension, ensure it's lowercase
    if ext in document_extensions:
        return True

    return False


def is_audio(fpath):
    """
    Determine if a given file is an audio file based on MIME type and file extension.

    Args:
        fpath (str): Path to the file.

    Returns:
        bool: True if the file is a recognized audio file, False otherwise.
    """
    import mimetypes

    # Known audio MIME types
    audio_mime_types = {
        "audio/mpeg",
        "audio/wav",
        "audio/ogg",
        "audio/aac",
        "audio/flac",
        "audio/midi",
        "audio/x-midi",
        "audio/x-wav",
        "audio/x-flac",
        "audio/pcm",
        "audio/x-aiff",
        "audio/x-m4a",
    }

    # Known audio file extensions
    audio_extensions = {
        ".mp3",
        ".wav",
        ".ogg",
        ".aac",
        ".flac",
        ".midi",
        ".m4a",
        ".aiff",
        ".pcm",
        ".wma",
        ".ape",
        ".alac",
        ".opus",
    }

    # Get MIME type using mimetypes
    mime_type, _ = mimetypes.guess_type(fpath)

    # Check MIME type
    if mime_type in audio_mime_types:
        return True

    # Fallback: Check file extension
    ext = os.path.splitext(fpath)[
        -1
    ].lower()  # Get the file extension and ensure lowercase
    if ext in audio_extensions:
        return True

    return False


def is_code(fpath):
    """
    Determine if a given file is a code file based on file extension and optionally MIME type.

    Args:
        fpath (str): Path to the file.
        check_mime (bool): Whether to perform a MIME type check in addition to file extension check.

    Returns:
        bool: True if the file is a recognized code file, False otherwise.
    """
    # Known programming and scripting file extensions
    code_extensions = {
        ".m",
        ".py",
        ".ipynb",
        ".js",
        ".html",
        ".css",
        ".java",
        ".cpp",
        ".h",
        ".cs",
        ".go",
        ".rs",
        ".sh",
        ".rb",
        ".swift",
        ".ts",
        ".json",
        ".xml",
        ".yaml",
        ".toml",
        ".bash",
        ".r",
    }

    # Check file extension
    ext = os.path.splitext(fpath)[-1].lower()
    if ext in code_extensions:
        return True
    return False


def is_zip(fpath):
    import mimetypes

    mime_type, _ = mimetypes.guess_type(fpath)
    if mime_type == "application/zip":
        return True
    else:
        return False


def adjust_spines(ax=None, spines=["left", "bottom"], distance=2):
    import matplotlib.pyplot as plt

    if ax is None:
        ax = plt.gca()
    for loc, spine in ax.spines.items():
        if loc in spines:
            spine.set_position(("outward", distance))  # outward by 2 points
            # spine.set_smart_bounds(True)
        else:
            spine.set_color("none")  # don't draw spine
    # turn off ticks where there is no spine
    if "left" in spines:
        ax.yaxis.set_ticks_position("left")
    else:
        ax.yaxis.set_ticks([])
    if "bottom" in spines:
        ax.xaxis.set_ticks_position("bottom")
    else:
        # no xaxis ticks
        ax.xaxis.set_ticks([])


# And then plot the data:


def add_colorbar(im, width=None, pad=None, **kwargs):
    # usage: add_colorbar(im, width=0.01, pad=0.005, label="PSD (dB)", shrink=0.8)
    l, b, w, h = im.axes.get_position().bounds  # get boundaries
    width = width or 0.1 * w  # get width of the colorbar
    pad = pad or width  # get pad between im and cbar
    fig = im.axes.figure  # get figure of image
    cax = fig.add_axes([l + w + pad, b, width, h])  # define cbar Axes
    return fig.colorbar(im, cax=cax, **kwargs)  # draw cbar


# =============================================================================
# # for plot figures: setting rcParams
# usage: set_pub()
# or by setting sns.set_theme...see below:
# sns.set_theme(style="ticks", rc=params)      # 白色无刻度线，有坐标轴标度
# # sns.set_theme(style="whitegrid", rc=params)# 白色＋刻度线，无坐标轴标度
# # sns.set_theme(style="white", rc=params)    # 白色无刻度线，无坐标轴标度
# # sns.set_theme(style="dark", rc=params)     # 深色无刻度线，无坐标轴标度
# =============================================================================

def list2slice(indices: Union[List[int], np.ndarray]) -> List[slice]:
    """
    Convert a list of indices to a list of slice objects for Excel formatting.

    This function takes a list of indices (which may or may not be consecutive)
    and converts them into slice objects that represent consecutive ranges.

    Args:
        indices: A list or numpy array of integers representing indices.
                Must be non-empty and contain only integers.

    Returns:
        A list of slice objects representing consecutive ranges in the input indices.

    Raises:
        ValueError: If input is empty or contains non-integer values.
        TypeError: If input is not a list or numpy array.

    Examples:
        >>> list_to_slice([1, 2, 3, 5, 6, 8])
        [slice(1, 4, None), slice(5, 7, None), slice(8, 9, None)]

        >>> list_to_slice([])
        ValueError: Input indices cannot be empty
    """
    # Input validation
    if not isinstance(indices, (list, np.ndarray)):
        raise TypeError("Input must be a list or numpy array")

    if len(indices) == 0:
        raise ValueError("Input indices cannot be empty")

    # Convert to numpy array and validate all elements are integers
    try:
        indices_arr = np.array(indices, dtype=int)
    except (ValueError, TypeError):
        raise ValueError("All elements in indices must be integers")

    # Handle single element case
    if len(indices_arr) == 1:
        return [slice(int(indices_arr[0]), int(indices_arr[0]) + 1)]

    # Sort and remove duplicates while preserving order
    unique_sorted = np.unique(indices_arr)

    # Find where the difference between consecutive elements is greater than 1
    break_points = np.where(np.diff(unique_sorted) != 1)[0] + 1

    # Split into consecutive ranges
    ranges = np.split(unique_sorted, break_points)

    # Create slice objects
    slices = []
    for r in ranges:
        start = int(r[0])
        end = int(r[-1]) + 1  # slices are exclusive of the end value
        slices.append(slice(start, end))

    return slices

def list2str(x_str):
    s = "".join(str(x) for x in x_str)
    return s


def str2list(str_):
    l = []
    [l.append(x) for x in str_]
    return l


def str2words(
    content,
    method="combined",
    custom_dict=None,
    sym_spell_params=None,
    use_threading=True,
):
    """
    Ultimate text correction function supporting multiple methods,
    lists or strings, and domain-specific corrections.

    Parameters:
        content (str or list): Input text or list of strings to correct.
        method (str): Correction method ('textblob', 'sym', 'combined').
        custom_dict (dict): Custom dictionary for domain-specific corrections.
        sym_spell_params (dict): Parameters for initializing SymSpell.

    Returns:
        str or list: Corrected text or list of corrected strings.
    """
    from textblob import TextBlob
    from symspellpy import SymSpell, Verbosity
    from functools import lru_cache
    import pkg_resources
    from concurrent.futures import ThreadPoolExecutor

    def initialize_symspell(max_edit_distance=2, prefix_length=7):
        """Initialize SymSpell for advanced spelling correction."""
        sym_spell = SymSpell(max_edit_distance, prefix_length)
        dictionary_path = pkg_resources.resource_filename(
            "symspellpy",
            # "frequency_bigramdictionary_en_243_342.txt",
            "frequency_dictionary_en_82_765.txt",
        )

        sym_spell.load_dictionary(dictionary_path, term_index=0, count_index=1)
        return sym_spell

    def segment_words(text, sym_spell):
        """Segment concatenated words into separate words."""
        segmented = sym_spell.word_segmentation(text)
        return segmented.corrected_string

    @lru_cache(maxsize=1000)  # Cache results for repeated corrections
    def advanced_correction(word, sym_spell):
        """Correct a single word using SymSpell."""
        suggestions = sym_spell.lookup(word, Verbosity.CLOSEST, max_edit_distance=2)
        return suggestions[0].term if suggestions else word

    def apply_custom_corrections(word, custom_dict):
        """Apply domain-specific corrections using a custom dictionary."""
        return custom_dict.get(word.lower(), word)

    def preserve_case(original_word, corrected_word):
        """
        Preserve the case of the original word in the corrected word.
        """
        if original_word.isupper():
            return corrected_word.upper()
        elif original_word[0].isupper():
            return corrected_word.capitalize()
        else:
            return corrected_word.lower()

    def process_string(text, method, sym_spell=None, custom_dict=None):
        """
        Process a single string for spelling corrections.
        Handles TextBlob, SymSpell, and custom corrections.
        """
        if method in ("sym", "combined") and sym_spell:
            text = segment_words(text, sym_spell)

        if method in ("textblob", "combined"):
            text = str(TextBlob(text).correct())

        corrected_words = []
        for word in text.split():
            original_word = word
            if method in ("sym", "combined") and sym_spell:
                word = advanced_correction(word, sym_spell)

            # Step 3: Apply custom corrections
            if custom_dict:
                word = apply_custom_corrections(word, custom_dict)
            # Preserve original case
            word = preserve_case(original_word, word)
            corrected_words.append(word)

        return " ".join(corrected_words)

    # Initialize SymSpell if needed
    sym_spell = None
    if method in ("sym", "combined"):
        if not sym_spell_params:
            sym_spell_params = {"max_edit_distance": 2, "prefix_length": 7}
        sym_spell = initialize_symspell(**sym_spell_params)

    # Process lists or strings
    if isinstance(content, list):
        if use_threading:
            with ThreadPoolExecutor() as executor:
                corrected_content = list(
                    executor.map(
                        lambda x: process_string(x, method, sym_spell, custom_dict),
                        content,
                    )
                )
            return corrected_content
        else:
            return [
                process_string(item, method, sym_spell, custom_dict) for item in content
            ]
    else:
        return process_string(content, method, sym_spell, custom_dict)


def load_img(fpath):
    """
    Load an image from the specified file path.
    Args:
        fpath (str): The file path to the image.
    Returns:
        PIL.Image: The loaded image.
    Raises:
        FileNotFoundError: If the specified file is not found.
        OSError: If the specified file cannot be opened or is not a valid image file.
    """
    from PIL import Image

    try:
        img = Image.open(fpath)
        return img
    except FileNotFoundError:
        raise FileNotFoundError(f"The file '{fpath}' was not found.")
    except OSError:
        raise OSError(f"Unable to open file '{fpath}' or it is not a valid image file.")


def apply_filter(img, *args, verbose=True):
    # def apply_filter(img, filter_name, filter_value=None):
    """
    Apply the specified filter to the image.
    Args:
        img (PIL.Image): The input image.
        filter_name (str): The name of the filter to apply.
        **kwargs: Additional parameters specific to the filter.
    Returns:
        PIL.Image: The filtered image.
    """
    from PIL import ImageFilter

    def correct_filter_name(filter_name):
        if all(
            [
                "b" in filter_name.lower(),
                "ur" in filter_name.lower(),
                "box" not in filter_name.lower(),
            ]
        ):
            return "BLUR"
        elif "cont" in filter_name.lower():
            return "Contour"
        elif "det" in filter_name.lower():
            return "Detail"
        elif (
            "edg" in filter_name.lower()
            and "mo" not in filter_name.lower()
            and "f" not in filter_name.lower()
        ):
            return "EDGE_ENHANCE"
        elif "edg" in filter_name.lower() and "mo" in filter_name.lower():
            return "EDGE_ENHANCE_MORE"
        elif "emb" in filter_name.lower():
            return "EMBOSS"
        elif "edg" in filter_name.lower() and "f" in filter_name.lower():
            return "FIND_EDGES"
        elif "sh" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SHARPEN"
        elif "sm" in filter_name.lower() and "mo" not in filter_name.lower():
            return "SMOOTH"
        elif "sm" in filter_name.lower() and "mo" in filter_name.lower():
            return "SMOOTH_MORE"
        elif "min" in filter_name.lower():
            return "MIN_FILTER"
        elif "max" in filter_name.lower():
            return "MAX_FILTER"
        elif "mod" in filter_name.lower():
            return "MODE_FILTER"
        elif "mul" in filter_name.lower():
            return "MULTIBAND_FILTER"
        elif "gau" in filter_name.lower():
            return "GAUSSIAN_BLUR"
        elif "box" in filter_name.lower():
            return "BOX_BLUR"
        elif "med" in filter_name.lower():
            return "MEDIAN_FILTER"
        else:
            supported_filters = [
                "BLUR",
                "CONTOUR",
                "DETAIL",
                "EDGE_ENHANCE",
                "EDGE_ENHANCE_MORE",
                "EMBOSS",
                "FIND_EDGES",
                "SHARPEN",
                "SMOOTH",
                "SMOOTH_MORE",
                "MIN_FILTER",
                "MAX_FILTER",
                "MODE_FILTER",
                "MULTIBAND_FILTER",
                "GAUSSIAN_BLUR",
                "BOX_BLUR",
                "MEDIAN_FILTER",
            ]
            raise ValueError(
                f"Unsupported filter: {filter_name}, should be one of: {supported_filters}"
            )

    for arg in args:
        if isinstance(arg, str):
            filter_name = correct_filter_name(arg)
        else:
            filter_value = arg
    if verbose:
        print(f"processing {filter_name}")
    filter_name = filter_name.upper()  # Ensure filter name is uppercase

    # Supported filters
    supported_filters = {
        "BLUR": ImageFilter.BLUR,
        "CONTOUR": ImageFilter.CONTOUR,
        "DETAIL": ImageFilter.DETAIL,
        "EDGE_ENHANCE": ImageFilter.EDGE_ENHANCE,
        "EDGE_ENHANCE_MORE": ImageFilter.EDGE_ENHANCE_MORE,
        "EMBOSS": ImageFilter.EMBOSS,
        "FIND_EDGES": ImageFilter.FIND_EDGES,
        "SHARPEN": ImageFilter.SHARPEN,
        "SMOOTH": ImageFilter.SMOOTH,
        "SMOOTH_MORE": ImageFilter.SMOOTH_MORE,
        "MIN_FILTER": ImageFilter.MinFilter,
        "MAX_FILTER": ImageFilter.MaxFilter,
        "MODE_FILTER": ImageFilter.ModeFilter,
        "MULTIBAND_FILTER": ImageFilter.MultibandFilter,
        "GAUSSIAN_BLUR": ImageFilter.GaussianBlur,
        "BOX_BLUR": ImageFilter.BoxBlur,
        "MEDIAN_FILTER": ImageFilter.MedianFilter,
    }
    # Check if the filter name is supported
    if filter_name not in supported_filters:
        raise ValueError(
            f"Unsupported filter: {filter_name}, should be one of: {[i.lower() for i in supported_filters.keys()]}"
        )

    # Apply the filter
    if filter_name.upper() in [
        "BOX_BLUR",
        "GAUSSIAN_BLUR",
        "MEDIAN_FILTER",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
    ]:
        radius = filter_value if filter_value is not None else 2
        return img.filter(supported_filters[filter_name](radius))
    elif filter_name in ["MULTIBAND_FILTER"]:
        bands = filter_value if filter_value is not None else None
        return img.filter(supported_filters[filter_name](bands))
    else:
        if filter_value is not None and verbose:
            print(
                f"{filter_name} doesn't require a value for {filter_value}, but it remains unaffected"
            )
        return img.filter(supported_filters[filter_name])


def detect_angle(image, by="median", template=None):
    """Detect the angle of rotation using various methods."""
    from sklearn.decomposition import PCA
    from skimage import transform, feature, filters, measure
    from skimage.color import rgb2gray
    from scipy.fftpack import fftshift, fft2
    import numpy as np
    import cv2

    # Convert to grayscale
    if np.array(image).shape[-1] > 3:
        image = np.array(image)[:, :, :3]
    gray_image = rgb2gray(image)

    # Detect edges using Canny edge detector
    edges = feature.canny(gray_image, sigma=2)

    # Use Hough transform to detect lines
    lines = transform.probabilistic_hough_line(edges)
    if isinstance(by, bool):
        by="mean" if by else 0
    if not lines and any(["me" in by, "pca" in by]):
        print("No lines detected. Adjust the edge detection parameters.")
        return 0
    methods = [
        "mean",
        "median",
        "pca",
        "gradient orientation",
        "template matching",
        "moments",
        "fft",
    ]
    by = strcmp(by, methods)[0]
    # Hough Transform-based angle detection (Median/Mean)
    if "me" in by.lower():
        angles = []
        for line in lines:
            (x0, y0), (x1, y1) = line
            angle = np.arctan2(y1 - y0, x1 - x0) * 180 / np.pi
            if 80 < abs(angle) < 100:
                angles.append(angle)
        if not angles:
            return 0
        if "di" in by:
            median_angle = np.median(angles)
            rotation_angle = (
                90 - median_angle if median_angle > 0 else -90 - median_angle
            )

            return rotation_angle
        else:
            mean_angle = np.mean(angles)
            rotation_angle = 90 - mean_angle if mean_angle > 0 else -90 - mean_angle

            return rotation_angle

    # PCA-based angle detection
    elif "pca" in by.lower():
        y, x = np.nonzero(edges)
        if len(x) == 0:
            return 0
        pca = PCA(n_components=2)
        pca.fit(np.vstack((x, y)).T)
        angle = np.arctan2(pca.components_[0, 1], pca.components_[0, 0]) * 180 / np.pi
        return angle

    # Gradient Orientation-based angle detection
    elif "gra" in by.lower():
        gx, gy = np.gradient(gray_image)
        angles = np.arctan2(gy, gx) * 180 / np.pi
        hist, bin_edges = np.histogram(angles, bins=360, range=(-180, 180))
        return bin_edges[np.argmax(hist)]

    # Template Matching-based angle detection
    elif "temp" in by.lower():
        if template is None:
            # Automatically extract a template from the center of the image
            height, width = gray_image.shape
            center_x, center_y = width // 2, height // 2
            size = (
                min(height, width) // 4
            )  # Size of the template as a fraction of image size
            template = gray_image[
                center_y - size : center_y + size, center_x - size : center_x + size
            ]
        best_angle = None
        best_corr = -1
        for angle in range(0, 180, 1):  # Checking every degree
            rotated_template = transform.rotate(template, angle)
            res = cv2.matchTemplate(gray_image, rotated_template, cv2.TM_CCOEFF)
            _, max_val, _, _ = cv2.minMaxLoc(res)
            if max_val > best_corr:
                best_corr = max_val
                best_angle = angle
        return best_angle

    # Image Moments-based angle detection
    elif "mo" in by.lower():
        moments = measure.moments_central(gray_image)
        angle = (
            0.5
            * np.arctan2(2 * moments[1, 1], moments[0, 2] - moments[2, 0])
            * 180
            / np.pi
        )
        return angle

    # Fourier Transform-based angle detection
    elif "fft" in by.lower():
        f = fft2(gray_image)
        fshift = fftshift(f)
        magnitude_spectrum = np.log(np.abs(fshift) + 1)
        rows, cols = magnitude_spectrum.shape
        r, c = np.unravel_index(np.argmax(magnitude_spectrum), (rows, cols))
        angle = np.arctan2(r - rows // 2, c - cols // 2) * 180 / np.pi
        return angle

    else:
        print(f"Unknown method {by}: supported methods: {methods}")
        return 0


def imgsets(
    img,
    auto: bool = True,
    size=None,
    figsize=None,
    dpi: int = 200,
    show_axis: bool = False,
    plot_: bool = True,
    verbose: bool = False,
    model: str = "isnet-general-use",
    **kwargs,
):
    """
    Apply various enhancements and filters to an image using PIL's ImageEnhance and ImageFilter modules.

    Args:
        img (PIL.Image): The input image.
        sets (dict): A dictionary specifying the enhancements, filters, and their parameters.
        show (bool): Whether to display the enhanced image.
        show_axis (bool): Whether to display axes on the image plot.
        size (tuple): The size of the thumbnail, cover, contain, or fit operation.
        dpi (int): Dots per inch for the displayed image.
        figsize (tuple): The size of the figure for displaying the image.
        auto (bool): Whether to automatically enhance the image based on its characteristics.

    Returns:
        PIL.Image: The enhanced image.

    Supported enhancements and filters:
        - "sharpness": Adjusts the sharpness of the image. Values > 1 increase sharpness, while values < 1 decrease sharpness.
        - "contrast": Adjusts the contrast of the image. Values > 1 increase contrast, while values < 1 decrease contrast.
        - "brightness": Adjusts the brightness of the image. Values > 1 increase brightness, while values < 1 decrease brightness.
        - "color": Adjusts the color saturation of the image. Values > 1 increase saturation, while values < 1 decrease saturation.
        - "rotate": Rotates the image by the specified angle.
        - "crop" or "cut": Crops the image. The value should be a tuple specifying the crop box as (left, upper, right, lower).
        - "size": Resizes the image to the specified dimensions.
        - "thumbnail": Resizes the image to fit within the given size while preserving aspect ratio.
        - "cover": Resizes and crops the image to fill the specified size.
        - "contain": Resizes the image to fit within the specified size, adding borders if necessary.
        - "fit": Resizes and pads the image to fit within the specified size.
        - "filter": Applies various filters to the image (e.g., BLUR, CONTOUR, EDGE_ENHANCE).

    Note:
        The "color" and "enhance" enhancements are not implemented in this function.
    Usage: 
    imgsets(dir_img, auto=1, color=1.5, plot_=0)
    imgsets(dir_img, color=2)
    imgsets(dir_img, pad=(300, 300), bgcolor=(73, 162, 127), plot_=0)
    imgsets(dir_img, contrast=0, color=1.2, plot_=0)
    imgsets(get_clip(), flip="tb")# flip top and bottom
    imgsets(get_clip(), contrast=1, rm=[100, 5, 2]) #'foreground_threshold', 'background_threshold' and 'erode_structure_size'
    imgsets(dir_img, rm="birefnet-portrait") # with using custom model
    """

    import matplotlib.pyplot as plt
    from PIL import ImageEnhance, ImageOps, Image

    supported_filters = [
        "BLUR",
        "CONTOUR",
        "DETAIL",
        "EDGE_ENHANCE",
        "EDGE_ENHANCE_MORE",
        "EMBOSS",
        "FIND_EDGES",
        "SHARPEN",
        "SMOOTH",
        "SMOOTH_MORE",
        "MIN_FILTER",
        "MAX_FILTER",
        "MODE_FILTER",
        "MULTIBAND_FILTER",
        "GAUSSIAN_BLUR",
        "BOX_BLUR",
        "MEDIAN_FILTER",
    ]
    # *Rembg is a tool to remove images background.
    # https://github.com/danielgatis/rembg
    rem_models = {
        "u2net": "general use cases.",
        "u2netp": "A lightweight version of u2net model.",
        "u2net_human_seg": "human segmentation.",
        "u2net_cloth_seg": "Cloths Parsing from human portrait. Here clothes are parsed into 3 category: Upper body, Lower body and Full body.",
        "silueta": "Same as u2net but the size is reduced to 43Mb.",
        "isnet-general-use": "A new pre-trained model for general use cases.",
        "isnet-anime": "A high-accuracy segmentation for anime character.",
        "sam": "any use cases.",
        "birefnet-general": "general use cases.",
        "birefnet-general-lite": "A light pre-trained model for general use cases.",
        "birefnet-portrait": "human portraits.",
        "birefnet-dis": "dichotomous image segmentation (DIS).",
        "birefnet-hrsod": "high-resolution salient object detection (HRSOD).",
        "birefnet-cod": "concealed object detection (COD).",
        "birefnet-massive": "A pre-trained model with massive dataset.",
    }
    models_support_rem = list(rem_models.keys())
    str_usage = """
    imgsets(dir_img, auto=1, color=1.5, plot_=0)
    imgsets(dir_img, color=2)
    imgsets(dir_img, pad=(300, 300), bgcolor=(73, 162, 127), plot_=0)
    imgsets(dir_img, contrast=0, color=1.2, plot_=0)
    imgsets(get_clip(), flip="tb")# flip top and bottom
    imgsets(get_clip(), contrast=1, rm=[100, 5, 2]) #'foreground_threshold', 'background_threshold' and 'erode_structure_size'
    imgsets(dir_img, rm="birefnet-portrait") # with using custom model
    """
    if run_once_within():
        print(str_usage)

    def gamma_correction(image, gamma=1.0, v_max=255):
        # adjust gama value
        inv_gamma = 1.0 / gamma
        lut = [
            int((i / float(v_max)) ** inv_gamma * int(v_max)) for i in range(int(v_max))
        ]
        return lut  # image.point(lut)

    def auto_enhance(img):
        """
        Automatically enhances the image based on its characteristics, including brightness,
        contrast, color range, sharpness, and gamma correction.

        Args:
            img (PIL.Image): The input image.

        Returns:
            dict: A dictionary containing the optimal enhancement values applied.
            PIL.Image: The enhanced image.
        """
        from PIL import Image, ImageEnhance, ImageOps, ImageFilter
        import numpy as np

        # Determine the bit depth based on the image mode
        try:
            if img.mode in ["1", "L", "P", "RGB", "YCbCr", "LAB", "HSV"]:
                bit_depth = 8
            elif img.mode in ["RGBA", "CMYK"]:
                bit_depth = 8
            elif img.mode in ["I", "F"]:
                bit_depth = 16
            else:
                raise ValueError("Unsupported image mode")
        except:
            bit_depth = 8

        # Initialize enhancement factors
        enhancements = {
            "brightness": 1.0,
            "contrast": 0,  # autocontrasted
            "color": 1.35,
            "sharpness": 1.0,
            "gamma": 1.0,
        }

        # Calculate brightness and contrast for each channel
        num_channels = len(img.getbands())
        brightness_factors = []
        contrast_factors = []
        for channel in range(num_channels):
            channel_histogram = img.split()[channel].histogram()
            total_pixels = sum(channel_histogram)
            brightness = (
                sum(i * w for i, w in enumerate(channel_histogram)) / total_pixels
            )
            channel_min, channel_max = img.split()[channel].getextrema()
            contrast = channel_max - channel_min
            # Adjust calculations based on bit depth
            normalization_factor = 2**bit_depth - 1
            brightness_factor = (
                1.0 + (brightness - normalization_factor / 2) / normalization_factor
            )
            contrast_factor = (
                1.0 + (contrast - normalization_factor / 2) / normalization_factor
            )
            brightness_factors.append(brightness_factor)
            contrast_factors.append(contrast_factor)

        # Calculate average brightness and contrast factors across channels
        enhancements["brightness"] = sum(brightness_factors) / num_channels
        # Adjust brightness and contrast
        img = ImageEnhance.Brightness(img).enhance(enhancements["brightness"])

        # # Automatic color enhancement (saturation)
        # if img.mode == "RGB":
        #     color_enhancer = ImageEnhance.Color(img)
        #     color_histogram = np.array(img.histogram()).reshape(3, -1)
        #     avg_saturation = np.mean([np.std(channel) for channel in color_histogram]) / normalization_factor
        #     print(avg_saturation)
        #     enhancements["color"] = min(0, max(0.5, 1.0 + avg_saturation))  # Clamp to a reasonable range
        #     # img = color_enhancer.enhance(enhancements["color"])

        # Adjust sharpness
        sharpness_enhancer = ImageEnhance.Sharpness(img)
        # Use edge detection to estimate sharpness need
        edges = img.filter(ImageFilter.FIND_EDGES).convert("L")
        avg_edge_intensity = np.mean(np.array(edges))
        enhancements["sharpness"] = min(
            2.0, max(0.5, 1.0 + avg_edge_intensity / normalization_factor)
        )
        # img = sharpness_enhancer.enhance(enhancements["sharpness"])

        # # Apply gamma correction
        # def gamma_correction(image, gamma):
        #     inv_gamma = 1.0 / gamma
        #     lut = [min(255, max(0, int((i / 255.0) ** inv_gamma * 255))) for i in range(256)]
        #     return image.point(lut)

        # avg_brightness = np.mean(np.array(img.convert("L"))) / 255
        # enhancements["gamma"] = min(2.0, max(0.5, 1.0 if avg_brightness > 0.5 else 1.2 - avg_brightness))
        # img = gamma_correction(img, enhancements["gamma"])

        # Return the enhancements and the enhanced image
        return enhancements

    # Load image if input is a file path
    if isinstance(img, str):
        img = load_img(img)
    img_update = img.copy()

    if auto:
        kwargs = {**auto_enhance(img_update), **kwargs}
    params = [
        "sharp",
        "color",
        "contrast",
        "bright",
        "crop",
        "rotate",
        "size",
        "resize",
        "thumbnail",
        "cover",
        "contain",
        "filter",
        "fit",
        "pad",
        "rem",
        "rm",
        "back",
        "bg_color",
        "cut",
        "gamma",
        "flip",
        "booster",
    ]
    for k, value in kwargs.items():
        k = strcmp(k, params)[0]  # correct the param name
        if "shar" in k.lower():
            enhancer = ImageEnhance.Sharpness(img_update)
            img_update = enhancer.enhance(value)
        elif all(
            ["col" in k.lower(), "bg" not in k.lower(), "background" not in k.lower()]
        ):
            # *color
            enhancer = ImageEnhance.Color(img_update)
            img_update = enhancer.enhance(value)
        elif "contr" in k.lower():
            if value and isinstance(value, (float, int)):
                enhancer = ImageEnhance.Contrast(img_update)
                img_update = enhancer.enhance(value)
            else:
                try:
                    img_update = ImageOps.autocontrast(img_update)
                    print("autocontrasted")
                except Exception as e:
                    print(f"Failed 'auto-contrasted':{e}")
        elif "bri" in k.lower():
            enhancer = ImageEnhance.Brightness(img_update)
            img_update = enhancer.enhance(value)
        elif "cro" in k.lower() or "cut" in k.lower():
            img_update = img_update.crop(value)
        elif "rota" in k.lower():
            if isinstance(value, (str,bool)):
                value = detect_angle(img_update, by=value)
                print(f"rotated by {value}°")
            img_update = img_update.rotate(value)
        elif "flip" in k.lower():
            if "l" in value and "r" in value:
                # left/right
                img_update = img_update.transpose(Image.FLIP_LEFT_RIGHT)
            elif any(["u" in value and "d" in value, "t" in value and "b" in value]):
                # up/down or top/bottom
                img_update = img_update.transpose(Image.FLIP_TOP_BOTTOM)
        elif "si" in k.lower():
            if isinstance(value, tuple):
                value = list(value)
            value = [int(i) for i in value]
            img_update = img_update.resize(value)
        elif "thum" in k.lower():
            img_update.thumbnail(value)
        elif "cover" in k.lower():
            img_update = ImageOps.cover(img_update, size=value)
        elif "contain" in k.lower():
            img_update = ImageOps.contain(img_update, size=value)
        elif "fi" in k.lower() and "t" in k.lower():  # filter
            if isinstance(value, dict):
                if verbose:
                    print(f"supported filter: {supported_filters}")
                for filter_name, filter_value in value.items():
                    img_update = apply_filter(
                        img_update, filter_name, filter_value, verbose=verbose
                    )
            else:
                img_update = ImageOps.fit(img_update, size=value)
        elif "pad" in k.lower():
            # *ImageOps.pad ensures that the resized image has the exact size specified by the size parameter while maintaining the aspect ratio.
            # size: A tuple specifying the target size (width, height).
            img_update = ImageOps.pad(img_update, size=value)
        elif "rem" in k.lower() or "rm" in k.lower() or "back" in k.lower():
            from rembg import remove, new_session

            if verbose:
                preview(rem_models)

            print(f"supported modles: {models_support_rem}")
            model = strcmp(model, models_support_rem)[0]
            session = new_session(model)
            if isinstance(value, bool):
                print(f"using model:{model}")
                img_update = remove(img_update, session=session)
            elif value and isinstance(value, (int, float, list)):
                if verbose:
                    print("https://github.com/danielgatis/rembg/blob/main/USAGE.md")
                    print(
                        f"rm=True # using default setting;\nrm=(240,10,10)\n'foreground_threshold'(240) and 'background_threshold' (10) values used to determine foreground and background pixels. \nThe 'erode_structure_size'(10) parameter specifies the size of the erosion structure to be applied to the mask."
                    )
                if isinstance(value, int):
                    value = [value]
                if len(value) < 2:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value,
                        session=session,
                    )
                elif 2 <= len(value) < 3:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        session=session,
                    )
                elif 3 <= len(value) < 4:
                    img_update = remove(
                        img_update,
                        alpha_matting=True,
                        alpha_matting_background_threshold=value[0],
                        alpha_matting_foreground_threshold=value[1],
                        alpha_matting_erode_size=value[2],
                        session=session,
                    )
            elif isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value, session=session)
            elif isinstance(value, str):
                # use custom model
                print(f"using model:{strcmp(value, models_support_rem)[0]}")
                img_update = remove(
                    img_update,
                    session=new_session(strcmp(value, models_support_rem)[0]),
                )
        elif "bg" in k.lower() and "color" in k.lower():
            from rembg import remove

            if isinstance(value, list):
                value = tuple(value)
            if isinstance(value, tuple):  # replace the background color
                if len(value) == 3:
                    value += (255,)
                img_update = remove(img_update, bgcolor=value)
        elif "boost" in k.lower():
            import torch
            from realesrgan import RealESRGANer

            if verbose:
                print("Applying Real-ESRGAN for image reconstruction...")
            if isinstance(value, bool):
                scale = 4
            elif isinstance(value, (float, int)):
                scale = value
            else:
                scale = 4

            # try:
            device = "cuda" if torch.cuda.is_available() else "cpu"
            dir_curr_script = os.path.dirname(os.path.abspath(__file__))
            model_path = dir_curr_script + "/data/RealESRGAN_x4plus.pth"
            model_RealESRGAN = RealESRGANer(
                device=device,
                scale=scale,
                model_path=model_path,
                model="RealESRGAN_x4plus",
            )
            # https://github.com/xinntao/Real-ESRGAN?tab=readme-ov-file#python-script

            img_update = model_RealESRGAN.enhance(np.array(img_update))[0]
            # except Exception as e:
            #     print(f"Failed to apply Real-ESRGAN: {e}")

        # elif "ga" in k.lower() and "m" in k.lower():
        #     img_update = gamma_correction(img_update, gamma=value)
    # Display the image if requested
    if plot_:
        if figsize is None:
            plt.figure(dpi=dpi)
        else:
            plt.figure(figsize=figsize, dpi=dpi)
        plt.imshow(img_update)
        if show_axis:
            plt.axis("on")  # Turn on axis
            plt.minorticks_on()
            plt.grid(
                which="both", linestyle="--", linewidth=0.5, color="gray", alpha=0.7
            )

        else:
            plt.axis("off")  # Turn off axis
    return img_update


def thumbnail(dir_img_list, figsize=(10, 10), dpi=100, dir_save=None, kind=".png"):
    """
    Display a thumbnail figure of all images in the specified directory.
    Args:
        dir_img_list (list): List of the Directory containing the images.
    """
    import matplotlib.pyplot as plt
    from PIL import Image

    num_images = len(dir_img_list)
    if not kind.startswith("."):
        kind = "." + kind

    if num_images == 0:
        print("No images found to display.")
        return
    grid_size = int(num_images**0.5) + 1  # Determine grid size
    fig, axs = plt.subplots(grid_size, grid_size, figsize=figsize, dpi=dpi)
    for ax, image_file in zip(axs.flatten(), dir_img_list):
        try:
            img = Image.open(image_file)
            ax.imshow(img)
            ax.axis("off")
        except:
            continue
    # for ax in axs.flatten():
    #     ax.axis('off')
    [ax.axis("off") for ax in axs.flatten()]
    plt.tight_layout()
    if dir_save is None:
        plt.show()
    else:
        if basename(dir_save):
            fname = basename(dir_save) + kind
        else:
            fname = "_thumbnail_" + basename(dirname(dir_save)[:-1]) + ".png"
        if dirname(dir_img_list[0]) == dirname(dir_save):
            figsave(dirname(dir_save[:-1]), fname)
        else:
            figsave(dirname(dir_save), fname)

# search and fine the director of the libary, which installed at local
def dir_lib(lib_oi):
    """
    # example usage:
    # dir_lib("seaborn")
    """
    import site

    # Get the site-packages directory
    f = listdir(site.getsitepackages()[0], "folder")

    # Find Seaborn directory within site-packages
    dir_list = []
    for directory in f.fpath:
        if lib_oi in directory.lower():
            dir_list.append(directory)

    if dir_list != []:
        print(f"{lib_oi} directory:", dir_list)
    else:
        print(f"Cannot find the {lib_oi} in site-packages directory.")
    return dir_list

 
class FileInfo:
    def __init__(self, size, creation_time, ctime, mod_time, mtime, parent_dir, fname, kind, owner, extra_info=None):
        self.size = size
        self.creation_time = creation_time
        self.ctime = ctime
        self.mod_time = mod_time
        self.mtime = mtime
        self.parent_dir = parent_dir
        self.fname = fname
        self.kind = kind
        self.owner = owner
        if extra_info:
            for key, value in extra_info.items():
                setattr(self, key, value)
        print("To show the result: 'finfo(fpath).show()'")

    def __repr__(self):
        return (
            f"FileInfo(size={self.size} MB,  "
            f"ctime='{self.ctime}',  mtime='{self.mtime}', "
            f"parent_dir='{self.parent_dir}', fname='{self.fname}', kind='{self.kind}', owner='{self.owner}')"
        )

    def __str__(self):
        return (
            f"FileInfo:\n"
            f"  Size: {self.size} MB\n" 
            f"  CTime: {self.ctime}\n" 
            f"  MTime: {self.mtime}\n"
            f"  Parent Directory: {self.parent_dir}\n"
            f"  File Name: {self.fname}\n"
            f"  Kind: {self.kind}\n"
            f"  Owner: {self.owner}"
        )

    def show(self):
        return {
            "size": self.size, 
            "ctime": self.ctime, 
            "mtime": self.mtime,
            "parent_dir": self.parent_dir,
            "fname": self.fname,
            "kind": self.kind,
            "owner": self.owner,
            **{
                key: getattr(self, key)
                for key in vars(self)
                if key not in ["size",  "ctime", "mtime", "parent_dir", "fname", "kind", "owner"]
            },
        }

def _format_timestamp(timestamp,fmt='%Y-%m-%d %H:%M:%S'):
    """Convert timestamp to a human-readable format."""
    import datetime
    return datetime.datetime.fromtimestamp(timestamp).strftime(fmt)

def get_file_owner(file_path):
    """Retrieve file owner information cross-platform."""
    import platform
    if platform.system() == "Windows":
        import win32security
        import win32con
    else:
        import pwd
        import grp
    if platform.system() == "Windows":
        try:
            sd = win32security.GetFileSecurity(file_path, win32security.OWNER_SECURITY_INFORMATION)
            owner_sid = sd.GetSecurityDescriptorOwner()
            owner_name, domain, _ = win32security.LookupAccountSid(None, owner_sid)
            return f"{domain}\\{owner_name}"
        except Exception as e:
            return f"Could not retrieve owner: {e}"
    else:
        try:
            file_stat = os.stat(file_path)
            user_id = file_stat.st_uid
            return pwd.getpwuid(user_id).pw_name
        except Exception as e:
            return f"Could not retrieve owner: {e}"

def finfo(fpath, output='json', verbose=False):
    """Retrieve detailed metadata about a file."""
    import os
    if not os.path.exists(fpath):
        raise FileNotFoundError(f"File '{fpath}' not found.")

    fname, fmt = os.path.splitext(fpath)
    dir_par = os.path.dirname(fpath) + "/"

    data = {
        "size": round(os.path.getsize(fpath) / 1024 / 1024, 3), 
        "ctime": _format_timestamp(os.path.getctime(fpath)), 
        "mtime": _format_timestamp(os.path.getmtime(fpath)),
        "parent_dir": dir_par,
        "fname": fname.replace(dir_par, ""),
        "kind": fmt,
        "owner": get_file_owner(fpath),
    }

    # Extract PDF metadata if applicable
    extra_info = {}
    if data["kind"].lower() == ".pdf":
        try:
            from pdf2image import pdfinfo_from_path
            extra_info = pdfinfo_from_path(fpath)
        except ImportError:
            extra_info = {"Warning": "pdf2image module not installed, cannot extract PDF metadata."}
    if output=="json":
        if verbose:
            print(data)
        return data
    return FileInfo(
        size=data["size"], 
        ctime=data["ctime"], 
        mtime=data["mtime"],
        parent_dir=data["parent_dir"],
        fname=data["fname"],
        kind=data["kind"],
        owner=data["owner"],
        extra_info=extra_info,
    )

def color2rgb(
    color_input: str | tuple | list | None, 
    alpha: float | None = None
) -> tuple | None:
    """
    Ultimate color conversion utility with support for multiple formats and transparency.
    
    Parameters:
    -----------
    color_input : str | tuple | list | None
        Supported formats:
        - Hex strings ("#RRGGBB", "#RGB")
        - Named colors ("red", "blue")
        - RGB tuples ((0.2, 0.4, 0.6))
        - RGBA tuples ((0.2, 0.4, 0.6, 0.8))
        - HTML/CSS colors ("cornflowerblue")
        - CSS formats:
          - rgb(100,200,50)
          - rgba(100,200,50,0.8)
          - hsl(120,60%,70%)
          - hsla(120,60%,70%,0.8)
    alpha : float | None, optional
        Opacity value (0.0-1.0). If provided, adds/overrides alpha channel.
    
    Returns:
    --------
    tuple | None
        (R, G, B) or (R, G, B, A) tuple in 0-1 range, or None if invalid
    """
    from matplotlib import colors as mcolors
    import re
    
    if color_input is None:
        return None
    
    # Case 1: Already in RGB/RGBA tuple format
    if isinstance(color_input, (tuple, list)):
        if 3 <= len(color_input) <= 4:
            if all(0 <= x <= 1 for x in color_input):
                if alpha is not None and len(color_input) == 3:
                    return (*color_input, alpha)
                return tuple(color_input)
    
    # Case 2: String input
    if isinstance(color_input, str):
        # Remove whitespace and make lowercase
        color_str = color_input.strip().lower()
        
        # Handle CSS rgb/rgba format
        if color_str.startswith(('rgb(', 'rgba(')):
            try:
                nums = list(map(float, re.findall(r"[\d.]+", color_str)))
                if 3 <= len(nums) <= 4:
                    rgb = tuple(x/255 if i < 3 else x for i, x in enumerate(nums))
                    if alpha is not None:
                        return (*rgb[:3], alpha)
                    return rgb[:4] if len(rgb) == 4 else rgb[:3]
            except:
                pass
        
        # Handle CSS hsl/hsla format
        elif color_str.startswith(('hsl(', 'hsla(')):
            try:
                nums = list(map(float, re.findall(r"[\d.]+", color_str)))
                if 3 <= len(nums) <= 4:
                    h, s, l = nums[0]/360, nums[1]/100, nums[2]/100
                    rgb = mcolors.hsv_to_rgb((h, s, l))
                    if len(nums) == 4:
                        rgb += (nums[3],)
                    if alpha is not None:
                        return (*rgb[:3], alpha)
                    return rgb[:4] if len(rgb) == 4 else rgb[:3]
            except:
                pass
        
        # Standard hex/named color processing
        try:
            rgb = mcolors.to_rgba(color_str)
            if alpha is not None:
                return (*rgb[:3], alpha)
            return rgb if len(rgb) == 4 and rgb[3] != 1 else rgb[:3]
        except ValueError:
            pass
    
    # Fallback for invalid colors
    print(f"Warning: Invalid color format '{color_input}'")
    return None
 
def color2hex(
    color_input: str | tuple | list | dict | int | None,
    keep_alpha: bool = False,
    force_long: bool = False,
    uppercase: bool = False,
    prefix: str = "#",
    allow_short: bool = True
) -> str | None:
    """
    Ultimate color to hex converter with comprehensive format support.
    
    Parameters:
    -----------
    color_input : str | tuple | list | dict | int | None
        Input color in any of these formats:
        - Hex strings ("#RRGGBB", "#RGB", "RRGGBB", "RGB")
        - Named colors ("red", "blue", "transparent")
        - RGB/RGBA tuples ((0.2, 0.4, 0.6), (255, 0, 0), (100, 100, 100, 0.5))
        - CSS formats:
          - rgb(100,200,50)
          - rgba(100,200,50,0.8)
          - hsl(120,60%,70%)
          - hsla(120,60%,70%,0.8)
        - Integer RGB (0xFF0000 for red)
        - Dictionary {"r": 255, "g": 0, "b": 0} or {"h": 0, "s": 100, "l": 50}
    keep_alpha : bool, optional
        Whether to include alpha channel in hex format (#RRGGBBAA)
    force_long : bool, optional
        Force 6/8-digit hex even when 3/4-digit would be possible
    uppercase : bool, optional
        Use uppercase hex characters (False for lowercase)
    prefix : str, optional
        Prefix for hex string ("#" for CSS, "0x" for programming, "" for raw)
    allow_short : bool, optional
        Allow shortened 3/4-digit hex when possible
    
    Returns:
    --------
    str | None
        Hex color string or None if invalid
        
    Examples:
    ---------
    >>> color2hex((0.5, 0.2, 0.8)) → "#7f33cc"
    >>> color2hex("rgb(127, 51, 204)") → "#7f33cc"
    >>> color2hex((0.2, 0.4, 0.6, 0.8), True) → "#336699cc"
    >>> color2hex(0xFF0000, uppercase=True) → "#FF0000"
    >>> color2hex({"r": 255, "g": 165, "b": 0}, prefix="") → "ffa500"
    >>> color2hex("hsl(120, 100%, 50%)") → "#00ff00"
    """
    from matplotlib import colors as mcolors
    import re
    
    def to_rgba(color) -> tuple | None:
        """Internal conversion to RGBA tuple"""
        # Handle None
        if color is None:
            return None
            
        # Handle integer RGB
        if isinstance(color, int):
            if color < 0:
                return None
            return (
                (color >> 16) & 0xFF,
                (color >> 8) & 0xFF,
                color & 0xFF,
                255
            )
        
        # Handle dictionary formats
        if isinstance(color, dict):
            keys = set(color.keys())
            if {'r','g','b'}.issubset(keys):
                return (
                    color['r'] / 255 if color['r'] > 1 else color['r'],
                    color['g'] / 255 if color['g'] > 1 else color['g'],
                    color['b'] / 255 if color['b'] > 1 else color['b'],
                    color.get('a', 1.0)
                )
            elif {'h','s','l'}.issubset(keys):
                return mcolors.hsv_to_rgb((
                    color['h'] / 360,
                    color['s'] / 100,
                    color['l'] / 100
                )) + (color.get('a', 1.0),)
            return None
        
        # Handle string formats
        if isinstance(color, str):
            color = color.strip().lower()
            
            # Handle hex without prefix
            if re.match(r'^[0-9a-f]{3,8}$', color):
                return mcolors.to_rgba(f"#{color}")
                
            # Handle CSS functions
            if color.startswith(('rgb(', 'rgba(', 'hsl(', 'hsla(')):
                try:
                    return mcolors.to_rgba(color)
                except ValueError:
                    return None
                    
            # Handle named colors (including 'transparent')
            try:
                return mcolors.to_rgba(color)
            except ValueError:
                return None
        
        # Handle tuple/list formats
        if isinstance(color, (tuple, list)):
            if len(color) in (3, 4):
                # Normalize values
                normalized = []
                for i, v in enumerate(color):
                    if i < 3:  # RGB channels
                        if isinstance(v, int):
                            normalized.append(v / 255 if v > 1 else v)
                        else:
                            normalized.append(float(v))
                    else:  # Alpha channel
                        normalized.append(float(v))
                return tuple(normalized)
        
        return None
    
    # Convert input to RGBA
    rgba = to_rgba(color_input)
    if rgba is None:
        return None
    
    # Extract components
    components = []
    for i, c in enumerate(rgba):
        if i == 3 and not keep_alpha:
            break
        components.append(round(c * 255 if c <= 1 else c))
    
    # Determine if we can use short format
    use_short = (allow_short and 
                not force_long and 
                len(components) in (3, 4) and
                all((x % 17 == 0) for x in components[:3]))
    
    # Format the hex string
    if use_short:
        short_components = [x//17 for x in components[:3]] + components[3:]
        hex_str = "".join(f"{x:1x}" for x in short_components)
    else:
        hex_str = "".join(f"{x:02x}" for x in components)
    
    # Apply case and prefix
    if uppercase:
        hex_str = hex_str.upper()
    
    return f"{prefix}{hex_str}"
# ! format excel file
 
def hex2argb(color):
    """
    Convert a color name or hex code to aARGB format required by openpyxl.

    :param color: A color in the format: 'blue', '#RRGGBB', 'RRGGBB', 'aARRGGBB'
    :return: A hex color code in the format aARRGGBB.

    Example:
        print(hex2argb("blue"))      # Output: FF0000FF
        print(hex2argb("FFFF00"))    # Output: FFFFFF00
        print(hex2argb("#DF4245"))   # Output: FFDf4245
        print(hex2argb("FF00FF00"))  # Output: FF00FF00 (already in aARGB format)
    """
    import matplotlib.colors as mcolors
    import re
    color = color.lower().replace(" ", "") # 'light blue'
    # Convert color name (e.g., "blue") to hex
    if color.lower() in mcolors.CSS4_COLORS:
        color = mcolors.CSS4_COLORS[color.lower()].lstrip("#")
    color = color.lstrip("#").upper()# Remove '#' if present
    
    # Validate hex format
    if not re.fullmatch(r"[A-F0-9]{6,8}", color):
        raise ValueError(f"格式错误: {color}, 应该使用 RRGGBB, #RRGGBB, or aARRGGBB format.")

    # If already in aARRGGBB format (8 chars), return as is
    if len(color) == 8:
        return color
    
    # If in RRGGBB format, add FF (full opacity) as alpha
    return f"FF{color}"

def extract_kwargs(func):
    import inspect

    # Get the signature of the function
    signature = inspect.signature(func)
    # Extract parameters that are kwargs (parameters with default values or **kwargs)
    kwargs = {
        param.name: param.default
        for param in signature.parameters.values()
        if param.default is not inspect.Parameter.empty
    }

    return kwargs

def df2workbook(dataframes: dict, filename: str, mode: str = 'w',**kwargs):
    """
    将单个DataFrame化成workbook

    Parameters:
    - dataframes: dict of {sheet_name: DataFrame}
    - filename: str, path to the Excel file
    - mode: 'w' = write new file, 'a' = append to existing file
    """
    import pandas as pd
    import os
    from openpyxl import load_workbook
    if mode not in ['w', 'a']:
        raise ValueError("Mode must be 'w' (write) or 'a' (append).")

    if mode == 'a' and os.path.exists(filename):
        # Append mode using openpyxl
        with pd.ExcelWriter(filename, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
            for sheet_name, df in dataframes.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)
    else:
        # Write mode or file does not exist
        with pd.ExcelWriter(filename, engine='xlsxwriter') as writer:
            for sheet_name, df in dataframes.items():
                df.to_excel(writer, sheet_name=sheet_name, index=False)

# ! excel formatting
def copy_format(
    fpath: str, target_path: str, sheet_name, copy_values: bool = True, save:bool=True
):
    """
    ULTIMATE worksheet formatting copier with complete validation and conditional formatting support

    Features:
    - Full style copying (fonts, borders, fills, alignment)
    - Number formats and comments
    - Merged cells and freeze panes
    - Auto-filters
    - Data validations (including cross-sheet)
    - Conditional formatting
    - Optimized performance

    Parameters:
    - fpath: Path to source workbook
    - target_path: Path to target workbook
    - sheet_name: Single string (same name) or tuple/list (source_name, target_name)
    - copy_values: Whether to copy cell values (default: True)

    Returns: Modified target worksheet

    例子:
    1. copy_format("source.xlsx", "target.xlsx", sheet_name="same sheet_name")
    2. copy_format("source.xlsx", "target.xlsx", sheet_name=["sheet_name1", "sheet_name2"])
    3. copy_format("source.xlsx", "target.xlsx", sheet_name=("sheet_name1", "sheet_name2"))
    """
    import openpyxl
    from openpyxl import load_workbook,Workbook
    from openpyxl.worksheet.worksheet import Worksheet

    from openpyxl.styles import Font, Border, PatternFill, Alignment
    from copy import copy
    import warnings

    # --- Helper Functions ---
    def safe_copy_style(obj):
        """Deep copy style objects without recursion issues"""
        if obj is None:
            return None

        if isinstance(obj, (Font, Border, PatternFill, Alignment)):
            new_obj = obj.__class__()
            for attr in dir(obj):
                if not attr.startswith("_") and not callable(getattr(obj, attr)):
                    try:
                        setattr(new_obj, attr, getattr(obj, attr))
                    except (AttributeError, TypeError):
                        pass
            return new_obj
        return copy(obj)

    # --- Main Execution ---
    try:
        # Load workbooks with optimized settings
        wb_source = load_workbook(
            fpath, read_only=False, keep_vba=False, data_only=True
        )
        if isinstance(target_path, Worksheet) or isinstance(target_path, Workbook):
            wb_target=target_path
        else:
            wb_target = load_workbook(target_path, read_only=False, keep_vba=False)

        # Handle sheet names
        if isinstance(sheet_name, (list, tuple)) and len(sheet_name) == 2:
            src_name, tgt_name = sheet_name
        else:
            src_name = tgt_name = str(sheet_name)

        # Verify sheet existence
        if src_name not in wb_source.sheetnames:
            raise ValueError(f"Source sheet '{src_name}' not found")
        if tgt_name not in wb_target.sheetnames:
            raise ValueError(f"Target sheet '{tgt_name}' not found")

        src_ws = wb_source[src_name]
        tgt_ws = wb_target[tgt_name]

        # Backup validations and conditional formatting FIRST
        backup = _backup_validations(src_ws)

        # --- Copy Core Formatting ---
        # 1. Column dimensions
        for col, dim in src_ws.column_dimensions.items():
            tgt_dim = tgt_ws.column_dimensions[col]
            tgt_dim.width = dim.width
            tgt_dim.hidden = dim.hidden
            if dim.has_style:
                tgt_dim.font = safe_copy_style(dim.font)
                tgt_dim.fill = safe_copy_style(dim.fill)
                tgt_dim.border = safe_copy_style(dim.border)
                tgt_dim.alignment = safe_copy_style(dim.alignment)

        # 2. Row dimensions
        for row, dim in src_ws.row_dimensions.items():
            tgt_dim = tgt_ws.row_dimensions[row]
            tgt_dim.height = dim.height
            tgt_dim.hidden = dim.hidden
            if dim.has_style:
                tgt_dim.font = safe_copy_style(dim.font)
                tgt_dim.fill = safe_copy_style(dim.fill)
                tgt_dim.border = safe_copy_style(dim.border)
                tgt_dim.alignment = safe_copy_style(dim.alignment)

        # 3. Cell styles and values
        for row in src_ws.iter_rows():
            for cell in row:
                tgt_cell = tgt_ws.cell(
                    row=cell.row,
                    column=cell.column,
                    value=cell.value if copy_values else None,
                )
                if cell.has_style:
                    tgt_cell.font = safe_copy_style(cell.font)
                    tgt_cell.fill = safe_copy_style(cell.fill)
                    tgt_cell.border = safe_copy_style(cell.border)
                    tgt_cell.alignment = safe_copy_style(cell.alignment)
                    tgt_cell.number_format = cell.number_format
                    if cell.comment:
                        tgt_cell.comment = copy(cell.comment)

        # --- Copy Structural Elements ---
        # 1. Merged cells
        tgt_ws.merged_cells.ranges = []
        for merged_range in src_ws.merged_cells.ranges:
            tgt_ws.merge_cells(str(merged_range))

        # 2. Freeze panes and auto-filter
        tgt_ws.freeze_panes = src_ws.freeze_panes
        if src_ws.auto_filter.ref:
            tgt_ws.auto_filter.ref = src_ws.auto_filter.ref

        # --- Restore Advanced Features ---
        _restore_validations(tgt_ws, backup)

        # Save with optimization
        if save:
            wb_target.save(target_path)
            return tgt_ws

    except Exception as e:
        warnings.warn(f"Format copying failed: {str(e)}")
        raise
    finally:
        # Ensure workbooks are closed
        if "wb_source" in locals():
            wb_source.close()
        if "wb_target" in locals():
            wb_target.close()
# ! =========(below) interact with worrkbook and DataFrame=========== 
import pandas as pd
from openpyxl import load_workbook 
from openpyxl.workbook.workbook import Workbook 
from openpyxl.utils import get_column_letter 

class DataFrameAlignExcel:
    """
    A powerful tool for updating Excel files with data from DataFrames with various matching strategies.

    Features:
    - Accepts either file path or open Workbook object
    - Multiple matching strategies (exact, contains, starts_with, ends_with, regex)
    - Multiple value update strategies (overwrite, add, subtract, multiply, divide, append)
    - Support for multiple worksheets
    - Automatic column creation
    - Value normalization options
    - Detailed logging and dry-run mode
    - Progress reporting
    - Data validation
    - make_backup functionality
    """

    def __init__(self, fpath: Union[str, Workbook], df: pd.DataFrame = None):
        """
        Initialize the DataFrameAlignExcel.

        Args:
            fpath: Path to the Excel file (str) or open Workbook object
            df: Optional DataFrame to use for updates
        """
        self.fpath_or_wb = fpath
        self.df = df
        self.wb = None
        self.backup_path = None
        self.log = []
        self.owns_workbook = (
            False  # Track whether we created the workbook or it was passed in
        )

    def load_workbook(self) -> None:
        """Load the Excel workbook if a path was provided."""
        if isinstance(self.fpath_or_wb, str):
            if not os.path.exists(self.fpath_or_wb):
                raise FileNotFoundError(f"Excel file not found: {self.fpath_or_wb}")
            self.wb = load_workbook(self.fpath_or_wb)
            self.owns_workbook = True
        elif isinstance(self.fpath_or_wb, Workbook):
            self.wb = self.fpath_or_wb
            self.owns_workbook = False
        else:
            raise TypeError(
                "fpath must be either a string path or an openpyxl Workbook object"
            )

    def create_make_backup(self) -> None:
        """Create a make_backup of the original Excel file (only if we loaded from a file)."""
        if not isinstance(self.fpath_or_wb, str):
            self.log.append(
                "Skipping make_backup - working with Workbook object directly"
            )
            return

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        self.backup_path = os.path.join(
            os.path.dirname(self.fpath_or_wb),
            f"backup_{timestamp}_{os.path.basename(self.fpath_or_wb)}",
        )
        self.wb.save(self.backup_path)
        self.log.append(f"Created make_backup at: {self.backup_path}")

    def save_workbook(self, dir_save: str = None) -> None:
        """
        Save the workbook to a file.

        Args:
            dir_save: Optional path to save to. If None and we loaded from a file,
                       saves to the original path.
        """
        if self.wb is None:
            raise ValueError("No workbook loaded")

        if dir_save is None:
            if isinstance(self.fpath_or_wb, str):
                dir_save = self.fpath_or_wb
            else:
                dir_save = datetime.now().strftime("%Y%m%d_%H%M%S") + ".xlsx"
                print(
                    f"No save path provided and original input was a Workbook object, so used : {dir_save}"
                )
            self.wb.save(dir_save)
            self.log.append(f"Saved workbook to: {dir_save}")

    def normalize_value(self, value, clean_keys: str = "strip_split_first") -> str:
        """
        Normalize a value based on the specified method.

        Args:
            value: Value to normalize
            clean_keys: One of:
                - 'strip': just strip whitespace
                - 'strip_lower': strip and lowercase
                - 'strip_split_first': strip and take first part before comma
                - 'strip_split_last': strip and take last part after comma
                - None: no normalization

        Returns:
            Normalized value
        """
        if value is None:
            return None

        value = str(value)

        if clean_keys is None:
            return value

        if clean_keys == "strip":
            return value.strip()
        elif clean_keys == "strip_lower":
            return value.strip().lower()
        elif clean_keys == "strip_split_first":
            return value.strip().split(",")[0].strip()
        elif clean_keys == "strip_split_last":
            parts = value.strip().split(",")
            return parts[-1].strip() if len(parts) > 1 else value.strip()
        else:
            warnings.warn(f"Unknown clean_keys: {clean_keys}. Using 'strip'.")
            return value.strip()

    def find_column_index(self, ws, header_row: int, column_name: str, max_search_columns: int = 100) -> int:
        """
        Efficiently find the column index (1-based) for a given column name,
        considering only non-empty cells and limiting search range.

        Args:
            ws: Worksheet object
            header_row: Row number containing headers (1-based)
            column_name: Column name to find
            max_search_columns: Max number of columns to search (to prevent infinite loops)

        Returns:
            Column index (1-based), or -1 if not found
        """
        row_iter = ws.iter_rows(min_row=header_row, max_row=header_row, max_col=max_search_columns, values_only=False)
        for row in row_iter:
            for cell in row:
                if cell.value and str(cell.value).strip().lower() == column_name.lower():
                    return cell.column
            break  # Only process the header row
        return -1
    # def find_column_index(self, ws, header_row: int, column_name: str, max_search_columns: int = 100) -> int:
    #     """
    #     Find the column index (1-based) for a given column name.
    #     If not found, return the last non-empty header column index.

    #     Args:
    #         ws: Worksheet object
    #         header_row: Row number containing headers (1-based)
    #         column_name: Column name to find
    #         max_search_columns: Max number of columns to search

    #     Returns:
    #         Column index (1-based)
    #     """
    #     row_iter = ws.iter_rows(min_row=header_row, max_row=header_row, max_col=max_search_columns, values_only=False)
    #     last_non_empty_col = -1

    #     for row in row_iter:
    #         for cell in row:
    #             if cell.value and str(cell.value).strip():
    #                 last_non_empty_col = cell.column
    #                 if str(cell.value).strip().lower() == column_name.lower():
    #                     return cell.column
    #         break  # Only one row being read

    #     return last_non_empty_col

    def update_values(
        self,
        df: pd.DataFrame = None,
        sheet_name: Union[str, int, List[Union[str, int]]] = 0,
        header_row: int = 1,
        column_match: Union[Dict[str, str], List[Tuple[str, str]]] = None,
        column_mapping: Union[Dict[str, str], List[Tuple[str, str]]] = None,
        clean_keys: str = "strip_split_first",
        match_method: str = "exact",
        update_strategy: str = "overwrite",
        create_missing_columns: bool = True,
        preview_only: bool = False,
        show_progress: bool = True,
        skip_no_match: bool = True,
        make_backup: bool = True,
        dir_save: str = None,
        row_max=500
    ) -> Dict[str, int]:
        """
        Update Excel with values from DataFrame.

        Args:
            df: DataFrame containing update data (if None, uses self.df)
            sheet_name: Sheet name(s) to update (str, int, or list of these)
            header_row: Row number containing headers (1-based)
            column_match: Dict or list of tuples mapping DataFrame columns to Excel columns for matching
                          e.g., {'SampleID': 'ID'} or [('SampleID', 'ID'), ('Batch', 'Lot')]
            column_mapping: Dict or list of tuples mapping DataFrame columns to Excel columns to update
                          e.g., {'Vials': 'Qty'} or [('Vials', 'Qty'), ('Status', 'State')]
            clean_keys: How to normalize matching values (see normalize_value())
            match_method: How to match values ('exact', 'contains', 'starts_with', 'ends_with', 'regex')
            update_strategy: How to update values ('overwrite', 'add', 'subtract', 'multiply', 'divide', 'append')
            create_missing_columns: Whether to create columns that don't exist
            preview_only: If True, don't actually update the Excel file
            show_progress: If True, print progress updates
            skip_no_match: If True, skip rows where match columns don't match
            make_backup: If True, create a make_backup before updating (only if working with file path)
            dir_save: Optional path to save to. If None and we loaded from a file,
                      saves to the original path. Ignored if preview_only=True.

        Returns:
            Dictionary with update statistics
        """
        # Initialize
        start_time = datetime.now()
        if df is None:
            df = self.df
        if df is None:
            raise ValueError("No DataFrame provided")

        if not isinstance(column_match, (dict, list)) or not column_match:
            raise ValueError(
                "column_match must be a non-empty dict or list of tuples"
            )

        if not isinstance(column_mapping, (dict, list)) or not column_mapping:
            raise ValueError("column_mapping must be a non-empty dict or list of tuples")

        # Convert match/update columns to consistent format
        if isinstance(column_match, dict):
            column_match = list(column_match.items())
        if isinstance(column_mapping, dict):
            column_mapping = list(column_mapping.items())

        # Load workbook if not already loaded
        if self.wb is None:
            self.load_workbook()

        # Create make_backup (only if we're working with a file path)
        if not preview_only:
            self.create_make_backup()

        # Prepare statistics
        stats = {
            "processed_sheet_names":[],
            "processed_sheets": 0,
            "total_updates": 0,
            "skipped_rows": 0,
            "created_columns": 0,
        }

        # Normalize sheet names
        if not isinstance(sheet_name, list):
            sheet_name = [sheet_name]

        # Process each sheet
        for sheet in sheet_name:
            try:
                if isinstance(sheet, str):
                    ws = self.wb[sheet]
                elif isinstance(sheet, int):
                    ws = self.wb.worksheets[sheet]
                else:
                    ws = self.wb.active

                sheet_name = ws.title
                self.log.append(f"\nProcessing sheet: {sheet_name}")

                # Prepare matching data
                match_dict = {}
                for df_col, excel_col in column_match:
                    if clean_keys:
                        match_dict[excel_col] = dict(
                            zip(
                                df[df_col].apply(
                                    lambda x: self.normalize_value(x, clean_keys)
                                ),
                                df.index,
                            )
                        )
                    else:
                        match_dict[excel_col] = dict(zip(df[df_col], df.index))

                # Find or create update columns
                update_col_indices = {}
                for df_col, excel_col in column_mapping: 
                    col_idx = self.find_column_index(ws, header_row, excel_col)
                    if col_idx == -1:
                        if create_missing_columns:
                            # Find last column
                            last_col = max(
                                [cell.column for cell in ws[header_row] if cell.value is not None], default=0
                            )
                            col_idx = last_col + 1
                            ws.cell(row=header_row, column=col_idx, value=excel_col)
                            update_col_indices[excel_col] = col_idx
                            stats["created_columns"] += 1
                            self.log.append(
                                f"Created new column '{excel_col}' at position {col_idx}"
                            )
                        else:
                            raise ValueError(
                                f"Column '{excel_col}' not found and create_missing_columns=False"
                            )
                    else:
                        update_col_indices[excel_col] = col_idx

                # Process rows 
                for row in ws.iter_rows(min_row=header_row + 1): 
                    match_values = {}
                    match_failed = False

                    for excel_col in match_dict.keys():
                        col_idx = self.find_column_index(ws, header_row, excel_col)
                        if col_idx == -1:
                            if skip_no_match:
                                match_failed = True
                                break
                            else:
                                raise ValueError(
                                    f"Match column '{excel_col}' not found in sheet"
                                )

                        cell_value = row[
                            col_idx - 1
                        ].value  # -1 because iter_rows returns 0-based list
                        if clean_keys:
                            cell_value = self.normalize_value(cell_value, clean_keys)

                        match_values[excel_col] = cell_value

                    if match_failed:
                        stats["skipped_rows"] += 1
                        continue

                    # Find matching DataFrame row
                    df_index = None
                    for excel_col, value in match_values.items():
                        if value in match_dict[excel_col]:
                            if df_index is None:
                                df_index = match_dict[excel_col][value]
                            elif df_index != match_dict[excel_col][value]:
                                # Multiple match columns point to different rows - skip
                                df_index = None
                                break

                    if df_index is None:
                        stats["skipped_rows"] += 1
                        continue

                    # Update cells
                    for df_col, excel_col in column_mapping:
                        col_idx = update_col_indices[excel_col]
                        cell = row[
                            col_idx - 1
                        ]  # -1 because iter_rows returns 0-based list
                        new_value = df.at[df_index, df_col]

                        # Apply update strategy
                        if update_strategy == "overwrite":
                            cell.value = new_value
                        elif update_strategy in (
                            "add",
                            "subtract",
                            "multiply",
                            "divide",
                        ):
                            try:
                                old_value = (
                                    float(cell.value) if cell.value is not None else 0
                                )
                                new_value = (
                                    float(new_value) if new_value is not None else 0
                                )
                                if update_strategy == "add":
                                    cell.value = old_value + new_value
                                elif update_strategy == "subtract":
                                    cell.value = old_value - new_value
                                elif update_strategy == "multiply":
                                    cell.value = old_value * new_value
                                elif update_strategy == "divide":
                                    cell.value = (
                                        old_value / new_value
                                        if new_value != 0
                                        else old_value
                                    )
                            except (ValueError, TypeError):
                                if skip_no_match:
                                    continue
                                raise ValueError(
                                    f"Could not perform {update_strategy} operation on non-numeric values"
                                )
                        elif update_strategy == "append":
                            separator = ", " if cell.value else ""
                            cell.value = (
                                f"{cell.value}{separator}{new_value}"
                                if cell.value
                                else new_value
                            )
                        else:
                            raise ValueError(
                                f"Unknown update_strategy: {update_strategy}"
                            )

                        stats["total_updates"] += 1

                stats["processed_sheets"] += 1
                stats["processed_sheet_names"].append(sheet_name)
            except Exception as e:
                self.log.append(f"Error processing sheet {sheet}: {str(e)}")
                if (
                    not preview_only
                    and self.backup_path
                    and isinstance(self.fpath_or_wb, str)
                ):
                    self.log.append("Restoring from make_backup due to error")
                    self.wb = load_workbook(self.backup_path)
                raise

        # Save changes if not dry run
        if not preview_only:
            self.save_workbook(dir_save)
            if not make_backup:
                if os.path.exists(self.backup_path):
                    os.remove(self.backup_path)
        else:
            self.log.append("\nDry run complete - no changes saved") 

        # Print summary
        summary = (
            f"\nUpdate Summary:\n"
            f"\tProcessed {stats["processed_sheets"]} sheetnames: {stats['processed_sheet_names']}\n"
            f"\tTotal updates: {stats['total_updates']}\n"
            f"\tSkipped rows: {stats['skipped_rows']}\n" 
        )
        self.log.append(summary)

        if show_progress:
            print(summary)

        return stats

    def get_log(self) -> str:
        """Get the operation log as a string."""
        return "\n".join(self.log)

    def close(self) -> None:
        """Close the workbook if we own it."""
        if self.wb is not None and self.owns_workbook:
            self.wb.close()
            self.wb = None


DFToExcelMapping = Union[Dict[str, str], List[Tuple[str, str]]]
def df_align(
    fpath: Union[str, Workbook],
    df: pd.DataFrame,
    sheet_name: Union[str, int, List[Union[str, int]]] = 0,
    header_row: int = 1,
    column_match: DFToExcelMapping = None,
    column_mapping: DFToExcelMapping = None,
    clean_keys: str = "strip_split_first",
    match_method: str = "exact",
    update_strategy: str = "overwrite",
    create_missing_columns: bool = True,
    preview_only: bool = False,
    show_progress: bool = True,
    skip_no_match: bool = True,
    make_backup: bool = True,
    dir_save: str = None,
) -> Dict[str, int]:
    """
    wb = fload(
        dir_aml,
        password="XBuzwVk4xsC2361cHzyi9JFgfJHaTSerjBOQ0JAJU24=",
        sheet_name=0,
        header=1,
        output="bit",
    )
    ws = wb[wb.sheetnames[0]]
    df_align(
        fpath=wb,
        df=df_,
        sheet_name=None,
        header_row=2,
        column_match={"SampleID": "SampleID"},# key是 df中的列名, value是 excel中,
        column_mapping={"Vials": "Vials", "Vials_": "Total Vials"}, # key是 df中的列名, value是 excel中,
    )
    """
    updater = DataFrameAlignExcel(fpath, df)
    try:
        result = updater.update_values(
            sheet_name=sheet_name,
            header_row=header_row,
            column_match=column_match,
            column_mapping=column_mapping,
            clean_keys=clean_keys,
            match_method=match_method,
            update_strategy=update_strategy,
            create_missing_columns=create_missing_columns,
            preview_only=preview_only,
            show_progress=show_progress,
            skip_no_match=skip_no_match,
            make_backup=make_backup,
            dir_save=dir_save,
        )
        return result
    finally:
        updater.close()


# ! =========(Above) interact with worrkbook and DataFrame===========
def set_sheet_visible(
    fpath: str,
    sheet_name: Union[int, str, None,list] = 1,
    show: Union[bool, str] = True,
    exclude: Union[List[str], None,list,int] = None,
    verbose: bool = False
) -> None:
    """
    Modify sheet visibility in an Excel workbook.
    set_sheet_visible(fpath=dir_data_collection,sheet_name=None,show=1,verbose=1)
    Args:
        fpath (str): Path to the Excel workbook.
        sheet_name (int | str | None): Index or name of the sheet to apply visibility to.
                                       If None, all sheets are considered.
        show (bool | str): Visibility mode. Can be:
                           - True -> visible
                           - False -> veryHidden
                           - 'visible', 'hidden', 'veryHidden' as str
        exclude (list[str] | None): List of sheet names to exclude from changes.
        verbose (bool): If True, logs actions.
    """

    try: 
        wb = fload(fpath, output="bit", get_validations=1)
    except Exception as e:
        raise FileNotFoundError(f"Unable to load workbook: {e}")

    sheet_names = wb.sheetnames
    if verbose:
        print("Workbook loaded with sheets:")
        for i, name in enumerate(sheet_names):
            print(f"  [{i}] {name}")

    excludes=[]
    if exclude is None:
        exclude=[]
    if ~isinstance(exclude, list):
        exclude = [exclude]
    for exclude_ in exclude: 
        if isinstance(exclude_, str):
            excludes.append(strcmp(exclude_, sheet_names)[0])
        elif isinstance(exclude_, int):
            if 0 <= exclude_ < len(sheet_names):
                excludes.append(sheet_names[exclude_])
            else:
                raise IndexError(f"sheet_name index {exclude_} is out of range:0~{len(sheet_names)-1}.")

    # Resolve the sheet_name target
    target_indices = []
    if not isinstance(sheet_name,list):
        sheet_name=[sheet_name]
    for sheet_name_ in sheet_name:
        if sheet_name_ is None:
            target_indices = list(range(len(sheet_names)))
            break
        elif isinstance(sheet_name_, int):
            if 0 <= sheet_name_ < len(sheet_names):
                target_indices.append(sheet_name_)
            else:
                raise IndexError(f"sheet_name index {sheet_name_} is out of range :0~{len(sheet_names)-1}.")
        elif isinstance(sheet_name_, str):
            idx = strcmp(sheet_name_, sheet_names)[1]
            if idx == -1:
                raise ValueError(f"Sheet '{sheet_name_}' not found.")
            target_indices.append(idx)

    # Map show argument to valid state
    valid_states = ["veryHidden", "visible", "hidden"]
    if isinstance(show, str):
        if show not in valid_states:
            raise ValueError(f"Invalid show value '{show}'. Must be one of {valid_states}")
        state = show
    else:
        state = "visible" if show else "veryHidden" 
    # Modify sheet visibility
    for idx in target_indices:
        ws= wb[sheet_names[idx]]
        if ws.title in excludes:
            if verbose:
                print(f"Skipping excluded sheet: '{ws.title}'")
            continue
        ws.sheet_state = state 
    # Ensure at least one sheet is visible
    visible_sheets = [s for s in wb.worksheets if s.sheet_state == "visible"]
    not_visible_sheets = [s for s in wb.worksheets if s.sheet_state != "visible"]
    if not visible_sheets:
        fallback_sheet = wb.worksheets[0]
        fallback_sheet.sheet_state = "visible"
        if verbose:
            print(f"No visible sheets found. Setting '{fallback_sheet.title}' to visible.")
    if verbose:
        print(f"visible sheets:{[s.title for s in visible_sheets]}")

    try:
        wb.save(fpath) 
    except Exception as e:
        raise IOError(f"Error saving workbook: {e}") 


def format_excel(
    df: pd.DataFrame=None,
    filename:str=None,
    sheet_name:Union[str, int]=0,
    insert_img:dict=None,# {"A1":img_path}
    usage:bool=False,
    text_color:Union[dict,bool]=False, # dict: set the text color
    bg_color:Union[dict,bool]=False, # dict: set the back_ground color
    cell:Union[dict, list]=None,  # dict: or list for multiple locs setting:
    width:Union[bool, dict]=None,  # dict
    width_factor:int=1,# width * width_factor
    width_padding:int=2,# width + width_padding
    width_max=None,
    height:Union[bool, dict]=None,  # dict e.g., {2: 50, 3: 25},  keys are columns
    height_factor:int=1,
    height_padding:int=2,
    height_max=None,
    merge:tuple=None,  # tuple e.g.,  (slice(0, 1), slice(1, 3)),
    shade:Union[dict, list]=None,  # dict
    comment:Union[dict, list]=None,  # dict e.g., {(2, 4): "This is a comment"},
    comment_always_visible:bool=True,# always display comment
    link:Union[dict, list]=None,  # dict e.g., {(2, 2): "https://example.com"},
    protect:dict=None,  # dict, protect sheet
    protect_file:dict={},
    number_format:dict=None,  # dict: e.g., {1:"0.00", 2:"#,##0",3:"0%",4:"$#,##0.00"}
    data_validation=None,  # dict
    template:dict={},# e.g., template=dict(path="xx.xlsx",sheet_name=['sheet_name1',"sheet_name2"])
    apply_filter:bool=False, # add filter 
    freeze :str= False,#"A2",
    conditional_format:dict=None,  # dict
    verbose:bool=False,
    **kwargs,
):
    """
    Parameters:
        df : pandas.DataFrame, optional
            DataFrame to be written to the Excel file.
        filename : str, optional
            Path to the output Excel file.
        sheet_name : str or int, default 0
            Name or index of the sheet where data will be written.
        insert_img : dict, optional
            Dictionary specifying image insert locations, e.g., {"A1": "path/to/image.png"}.
        usage : bool, default False
            If True, display usage examples.
        cell : dict or list, optional
            Specifies cell formatting options.
        width : dict, optional
            Dictionary specifying column widths, e.g., {1: 20, 2: 30}.
        width_factor : int, default 2
            Additional factor to adjust column width dynamically.
        height : dict, optional
            Dictionary specifying row heights, e.g., {2: 50, 3: 25}.
        height_max : int, default 25
            Maximum row height allowed.
        merge : tuple, optional
            Specifies cell merging, e.g., (slice(0, 1), slice(1, 3)).
        shade : dict, optional
            Dictionary defining cell shading/styling.
        comment : dict, optional
            Dictionary adding comments, e.g., {(2, 4): "This is a comment"}.
        comment_always_visible : bool, default True
            Whether comments should always be visible.
        link : dict, optional
            Dictionary specifying hyperlinks, e.g., {(2, 2): "https://example.com"}.
        protect : dict, optional
            Dictionary defining cell protection settings.
        number_format : dict, optional
            Dictionary specifying number formats, e.g., {1: "0.00", 2: "#,##0"}.
        data_validation : dict, optional
            Dictionary setting data validation rules.
        apply_filter : bool, default True
            Whether to apply filters to the header row.
        freeze : str, optional
            Cell reference (e.g., "A2") to freeze rows/columns.
        conditional_format : dict, optional
            Dictionary defining conditional formatting rules.
        verbose : bool, default False
            Whether to print detailed execution logs.
        **kwargs : dict
            Additional parameters for advanced customization. 
    """

    usage_str="""
        Formats an Excel file with various styling options.
        Usage:
        fsave(
                dir_save,
                fload(dir_save, output="bit", sheet_name=sheet_name),
                sheet_name=sheet_name,
                if_sheet_exists="overlay",
                mode="a",
                width_factor=0,
                height={1: 50},
                cell=[
                    {
                        (slice(0, 1), slice(0, df_exists.shape[1])): {
                            "fill": {
                                "start_color": "61AFEF",  # Starting color
                                "end_color": "61AFEF",  # Ending color (useful for gradients)
                                "fill_type": "solid",  # Fill type (solid, gradient, etc.)
                            },
                            "font": {
                                "name": "Arial",  # Font name
                                "size": 11,  # Font size
                                "bold": True,  # Bold text
                                "italic": False,  # Italic text
                                # "underline": "single",  # Underline (single, double)
                                "color": "#000000",  # Font color
                            },
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                        }
                    },
                    {
                        (
                            slice(0, df_exists.shape[0]),
                            slice(0, df_exists.shape[1]),
                        ): {
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                        }
                    },
                    {
                        (slice(0, df_exists.shape[0]), slice(2, 3)): {
                            "alignment": {
                                "horizontal": "left",  # Horizontal alignment (left, center, right)
                            },
                        }
                    },
                    {
                        (slice(0, df_exists.shape[0]), slice(7, 8)): {
                            "alignment": {
                                "horizontal": "left",  # Horizontal alignment (left, center, right)
                            },
                        }
                    },
                ],
                password=False,  # depass("ogB3B7y3xR9iuH4QIQbyy6VXG14I0A8DlsTxyiGqg1U="),
            )
            """
    if verbose:
        print(usage_str)
    import pandas as pd
    from datetime import datetime
    import openpyxl
    from openpyxl import load_workbook,Workbook
    from openpyxl.worksheet.worksheet import Worksheet

    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.worksheet.datavalidation import DataValidation
    from openpyxl.comments import Comment
    from openpyxl.formatting.rule import ColorScaleRule, DataBarRule, IconSetRule,IconSet
    from openpyxl.workbook.protection import WorkbookProtection

    def _clean_excel_str(s):
        if isinstance(s, str):
            return re.sub(r'[\x00-\x1F\x7F]', '', s)  # remove control chars
        return s
    def _escape_excel_formula(cell):
        if isinstance(cell, str) and cell.startswith(('=', '+', '-')):
            return "'" + cell  # add single quote to escape
        return cell 

    def convert_indices_to_range(row_slice, col_slice):
        """Convert numerical row and column slices to Excel-style range strings."""
        start_row = row_slice.start + 1
        end_row = row_slice.stop if row_slice.stop is not None else None
        start_col = col_slice.start + 1
        end_col = col_slice.stop if col_slice.stop is not None else None

        start_col_letter = get_column_letter(start_col)
        end_col_letter = get_column_letter(end_col) if end_col else None
        return (
            f"{start_col_letter}{start_row}:{end_col_letter}{end_row}"
            if end_col_letter
            else f"{start_col_letter}{start_row}"
        )


    def is_merged_cell(ws, cell):
        """Check if a cell is part of any merged range."""
        for merged_range in ws.merged_cells.ranges:
            if cell.coordinate in merged_range:
                return True
        return False

    def apply_auto_width(ws, width_factor=1.2, width_padding=2, width_max=50):
        """
        Automatically adjust column widths based on content length,
        with complete protection against merged cell errors.
        
        Args:
            ws: Worksheet object
            width_factor: Multiplier for content length (default 1.2)
            width_padding: Additional padding (default 2)
            width_max: Maximum column width (default 50)
        """
        # First build a set of all merged cell coordinates
        merged_coords = set()
        for merged_range in ws.merged_cells.ranges:
            for row in ws.iter_rows(min_row=merged_range.min_row,
                                max_row=merged_range.max_row,
                                min_col=merged_range.min_col,
                                max_col=merged_range.max_col):
                for cell in row:
                    merged_coords.add(cell.coordinate)
        
        for col in ws.columns:
            if not col:
                continue
                
            col_letter = get_column_letter(col[0].column)
            max_length = 0
            
            for cell in col:
                # Skip merged cells entirely
                if cell.coordinate in merged_coords:
                    continue
                    
                try:
                    if cell.value is not None:
                        # Handle both single-line and multi-line content
                        cell_value = str(cell.value)
                        lines = cell_value.split('\n')
                        current_max = max(len(line) for line in lines)
                        max_length = max(max_length, current_max)
                except Exception as e:
                    print(f"Skipping cell {cell.coordinate} due to error: {e}")
                    continue
            
            # Calculate width with constraints
            adjusted_width = min(
                max(1, (max_length * width_factor) + width_padding),
                width_max if width_max is not None else float('inf')
            )
            
            ws.column_dimensions[col_letter].width = adjusted_width

    def apply_color_to_worksheet(ws=None, sheet_name=None, conditions=None, cell_idx=None,where="text"):
        """
        Apply text color formatting to a specific cell range in an openpyxl workbook based on conditions.

        Parameters:
        ws : worrksheet
            The openpyxl workbook object to style.
        sheet_name : str
            The name of the sheet to style.
        conditions : dict
            Dictionary defining conditions for text or background coloring.
                Example:
                {
                    ">10": "#FF0000",           # Red if value is greater than 10
                    "contains:Error": "#FFFF00", # Yellow if text contains 'Error'
                    "startswith:Warn": "#FFA500" # Orange if text starts with 'Warn'
                }
        cell_idx : tuple, optional
            A tuple of slices defining the selected row and column range (only for DataFrame).
        where : str, default="text"
            "text" -> Apply color to text, "bg" -> Apply color to background.

        Returns:
        Workbook
            The workbook with applied formatting.
        """ 
        def evaluate_condition(value, condition):
            """Evaluate the condition dynamically."""
            if not isinstance(conditions, dict):
                raise ValueError(f"condition必须是dict格式:e.g., {'x>=20':'#DD0531', 'startswith:Available':'#DD0531'}")
            try:
                if "x" in condition and re.search(r"[<>=!]=*", condition):
                    expr = condition.replace("x", str(value))
                    return eval(expr)
                elif condition.startswith("startswith:") or condition.startswith("startwith:"):
                    return value.startswith(condition.split(":", 1)[1])
                elif condition.startswith("endswith:") or condition.startswith("endwith:"):
                    return value.endswith(condition.split(":", 1)[1])
                elif condition.startswith("contains:") or condition.startswith("contain:") or condition.startswith("include:"):
                    return condition.split(":", 1)[1] in value
                elif condition.startswith("matches:") or condition.startswith("match:"):
                    return re.search(condition.split(":", 1)[1], value) is not None 
                else:
                    expr = condition
                return False
            except Exception as e: 
                return False

        def apply_condition_to_cell_text_color(cell, value):
            """Apply color to a cell if it matches any condition."""
            for condition, color in conditions.items():
                if evaluate_condition(value, condition):
                    # Apply color to font
                    cell.font = openpyxl.styles.Font(
                        color=openpyxl.styles.Color(rgb=hex2argb(color))
                    )
                    return
        def apply_condition_to_cell_bg_color(cell, value):
            """Apply background color to a cell if it matches any condition."""
            for condition, color in conditions.items():
                if evaluate_condition(value, condition):
                    if not isinstance(color,list):
                        color=[color]
                    if len(color)==1:
                        cell.fill = PatternFill(
                            start_color=hex2argb(color[0]),
                            end_color=hex2argb(color[0]),
                            fill_type="solid"
                        )
                    elif len(color)==2:
                        cell.fill = PatternFill(
                            start_color=hex2argb(color[0]),
                            end_color=hex2argb(color[1]),
                            fill_type="solid"
                        )
                    return
        if isinstance(cell_idx, tuple):
            # If cell_idx is provided, select a range based on the slice
            row_slice, col_slice = cell_idx
            rows = list(
                ws.iter_rows(
                    min_row=row_slice.start + 1,
                    max_row=row_slice.stop,
                    min_col=col_slice.start + 1,
                    max_col=col_slice.stop,
                )
            )
            for row in rows:
                for cell in row:
                    if where=="text":
                        apply_condition_to_cell_text_color(cell, cell.value)
                    elif where=="bg":
                        apply_condition_to_cell_bg_color(cell, cell.value)
        else:
            # If no cell_idx is provided, apply to all cells
            for row in ws.iter_rows():
                for cell in row:
                    if where=="text":
                        apply_condition_to_cell_text_color(cell, cell.value)
                    elif where=="bg":
                        apply_condition_to_cell_bg_color(cell,cell.value)
        return ws
    
    def apply_format(ws, cell, cell_range):
        """Apply cell formatting to a specified range."""
        # Get all merged cell coordinates first
        merged_cells = set()
        for merged_range in ws.merged_cells.ranges:
            for coord in merged_range.cells:
                merged_cells.add(coord)
        cell_font, cell_fill, cell_alignment, border = None, None, None, None
        kws_cell = ["font", "fill", "alignment", "border"]
        for K, _ in cell.items():
            if strcmp(K, kws_cell)[0] == "font":
                #! font
                font_color = "000000"
                font_name = "Arial"
                font_underline = "none"
                font_size = 11
                font_bold = False
                font_strike = False
                font_italic = False
                kws_font = [
                    "name",
                    "size",
                    "bold",
                    "underline",
                    "color",
                    "strike",
                    "italic",
                ]
                for k_, v_ in cell.get(K, {}).items():
                    if strcmp(k_, kws_font)[0] == "name":
                        font_name = v_
                    elif strcmp(k_, kws_font)[0] == "size":
                        font_size = v_
                    elif strcmp(k_, kws_font)[0] == "bold":
                        font_bold = v_
                    elif strcmp(k_, kws_font)[0] == "underline":
                        font_underline = strcmp(v_, ["none", "single", "double"])[0]
                    elif strcmp(k_, kws_font)[0] == "color":
                        font_color = hex2argb(v_)
                    elif strcmp(k_, kws_font)[0] == "strike":
                        font_strike = v_
                    elif strcmp(k_, kws_font)[0] == "italic":
                        font_italic = v_

                cell_font = Font(
                    name=font_name,
                    size=font_size,
                    bold=font_bold,
                    italic=font_italic,
                    underline=font_underline,
                    strike=font_strike,
                    color=font_color,
                )

            if strcmp(K, kws_cell)[0] == "fill":
                #! fill
                kws_fill = ["start_color", "end_color", "fill_type", "color"]
                kws_fill_type = [
                    "darkVertical",
                    "lightDown",
                    "lightGrid",
                    "solid",
                    "darkDown",
                    "lightGray",
                    "lightUp",
                    "gray0625",
                    "lightVertical",
                    "lightHorizontal",
                    "darkHorizontal",
                    "gray125",
                    "darkUp",
                    "mediumGray",
                    "darkTrellis",
                    "darkGray",
                    "lightTrellis",
                    "darkGrid",
                ]
                start_color, end_color, fill_type = (
                    "FFFFFF",
                    "FFFFFF",
                    "solid",
                )  # default
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_fill)[0] == "color":
                        start_color, end_color = hex2argb(v), hex2argb(v)
                        break
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_fill)[0] == "start_color":
                        start_color = hex2argb(v)
                    elif strcmp(k, kws_fill)[0] == "end_color":
                        end_color = hex2argb(v)
                    elif strcmp(k, kws_fill)[0] == "fill_type":
                        fill_type = strcmp(v, kws_fill_type)[0]
                cell_fill = PatternFill(
                    start_color=start_color,
                    end_color=end_color,
                    fill_type=fill_type,
                )

            if strcmp(K, kws_cell)[0] == "alignment":
                #! alignment
                # default
                align_horizontal = "general"
                align_vertical = "center"
                align_rot = 0
                align_wrap = False
                align_shrink = False
                align_indent = 0
                kws_align = [
                    "horizontal",
                    "ha",
                    "vertical",
                    "va",
                    "text_rotation",
                    "rotat",
                    "rot",
                    "wrap_text",
                    "wrap",
                    "shrink_to_fit",
                    "shrink",
                    "indent",
                ]
                for k, v in cell.get(K, {}).items():
                    if strcmp(k, kws_align)[0] in ["horizontal", "ha"]:
                        align_horizontal = strcmp(
                            v, ["general", "left", "right", "center"]
                        )[0]
                    elif strcmp(k, kws_align)[0] in ["vertical", "va"]:
                        align_vertical = strcmp(v, ["top", "center", "bottom"])[0]
                    elif strcmp(k, kws_align)[0] in ["text_rotation", "rotat", "rot"]:
                        align_rot = v
                    elif strcmp(k, kws_align)[0] in ["wrap_text", "wrap"]:
                        align_wrap = v
                    elif strcmp(k, kws_align)[0] in [
                        "shrink_to_fit",
                        "shrink",
                        "wrap_text",
                        "wrap",
                    ]:
                        align_shrink = v
                    elif strcmp(k, kws_align)[0] in ["indent"]:
                        align_indent = v
                cell_alignment = Alignment(
                    horizontal=align_horizontal,
                    vertical=align_vertical,
                    text_rotation=align_rot,
                    wrap_text=align_wrap,
                    shrink_to_fit=align_shrink,
                    indent=align_indent,
                )

            if strcmp(K, kws_cell)[0] == "border":
                #! border
                kws_border = [
                    "color_left",
                    "color_l",
                    "color_right",
                    "color_r",
                    "color_top",
                    "color_t",
                    "color_bottom",
                    "color_b",
                    "color_diagonal",
                    "color_d",
                    "color_outline",
                    "color_o",
                    "color_vertical",
                    "color_v",
                    "color_horizontal",
                    "color_h",
                    "color",
                    "style_left",
                    "style_l",
                    "style_right",
                    "style_r",
                    "style_top",
                    "style_t",
                    "style_bottom",
                    "style_b",
                    "style_diagonal",
                    "style_d",
                    "style_outline",
                    "style_o",
                    "style_vertical",
                    "style_v",
                    "style_horizontal",
                    "style_h",
                    "style",
                ]
                # * border color
                border_color_l, border_color_r, border_color_t, border_color_b = (
                    "FF000000",
                    "FF000000",
                    "FF000000",
                    "FF000000",
                )
                border_color_d, border_color_o, border_color_v, border_color_h = (
                    "FF000000",
                    "FF000000",
                    "FF000000",
                    "FF000000",
                )
                # get colors config
                for k, v in cell.get(K, {}).items():
                    print(k, v,strcmp(k, kws_border)[0])
                    if strcmp(k, kws_border)[0] in ["color"]:
                        border_color_all = hex2argb(v)
                        # 如果设置了color,表示其它的所有的都设置成为一样的
                        # 然后再才开始自己定义其它的color
                        (
                            border_color_l,
                            border_color_r,
                            border_color_t,
                            border_color_b,
                        ) = (
                            border_color_all,
                            border_color_all,
                            border_color_all,
                            border_color_all,
                        )
                        (
                            border_color_d,
                            border_color_o,
                            border_color_v,
                            border_color_h,
                        ) = (
                            border_color_all,
                            border_color_all,
                            border_color_all,
                            border_color_all,
                        )
                    elif strcmp(k, kws_border)[0] in ["color_left", "color_l"]:
                        border_color_l = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_right", "color_r"]:
                        border_color_r = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_top", "color_t"]:
                        border_color_t = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_bottom", "color_b"]:
                        border_color_b = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_diagonal", "color_d"]:
                        border_color_d = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_outline", "color_o"]:
                        border_color_o = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_vertical", "color_v"]:
                        border_color_v = hex2argb(v)
                    elif strcmp(k, kws_border)[0] in ["color_horizontal", "color_h"]:
                        border_color_h = hex2argb(v)
                # *border style
                border_styles = [
                    "thin",
                    "medium",
                    "thick",
                    "dotted",
                    "dashed",
                    "hair",
                    "mediumDashed",
                    "dashDot",
                    "dashDotDot",
                    "slantDashDot",
                    "none",
                ]
                border_style_l, border_style_r, border_style_t, border_style_b = (
                    None,
                    None,
                    None,
                    None,
                )
                border_style_d, border_style_o, border_style_v, border_style_h = (
                    None,
                    None,
                    None,
                    None,
                )
                # get styles config
                for k, v in cell.get(K, {}).items():
                    # if not "style" in k:
                    #     break
                    if strcmp(k, kws_border)[0] in ["style"]:
                        border_style_all = strcmp(v, border_styles)[0]
                        # 如果设置了style,表示其它的所有的都设置成为一样的
                        # 然后再才开始自己定义其它的style
                        (
                            border_style_l,
                            border_style_r,
                            border_style_t,
                            border_style_b,
                        ) = (
                            border_style_all,
                            border_style_all,
                            border_style_all,
                            border_style_all,
                        )
                        (
                            border_style_d,
                            border_style_o,
                            border_style_v,
                            border_style_h,
                        ) = (
                            border_style_all,
                            border_style_all,
                            border_style_all,
                            border_style_all,
                        )
                    elif strcmp(k, kws_border)[0] in ["style_left", "style_l"]:
                        border_style_l = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_right", "style_r"]:
                        border_style_r = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_top", "style_t"]:
                        border_style_t = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_bottom", "style_b"]:
                        border_style_b = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_diagonal", "style_d"]:
                        border_style_d = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_outline", "style_o"]:
                        border_style_o = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_vertical", "style_v"]:
                        border_style_v = strcmp(v, border_styles)[0]
                    elif strcmp(k, kws_border)[0] in ["style_horizontal", "style_h"]:
                        border_style_h = strcmp(v, border_styles)[0]
                # * apply border config
                border = Border(
                    left=Side(border_style=border_style_l, color=border_color_l),
                    right=Side(border_style=border_style_r, color=border_color_r),
                    top=Side(border_style=border_style_t, color=border_color_t),
                    bottom=Side(border_style=border_style_b, color=border_color_b),
                    diagonal=Side(border_style=border_style_d, color=border_color_d),
                    diagonal_direction=0,
                    outline=Side(border_style=border_style_o, color=border_color_o),
                    vertical=Side(border_style=border_style_v, color=border_color_v),
                    horizontal=Side(border_style=border_style_h, color=border_color_h),
                )

        #! final apply configs
        for row in ws[cell_range]:
            for cell_ in row:
                if cell_.coordinate in merged_cells:
                    continue  # Skip merged cells
                if cell_font:
                    cell_.font = cell_font
                if cell_fill:
                    cell_.fill = cell_fill
                if cell_alignment:
                    cell_.alignment = cell_alignment
                if border:
                    cell_.border = border
    def generate_unique_sheet_name(wb, sheet_name):
        """Generate a unique sheet name if the given name already exists in the workbook."""
        if sheet_name not in wb.sheetnames:
            return sheet_name
        counter = 1
        unique_name = f"{sheet_name}_{counter}"
        while unique_name in wb.sheetnames:
            counter += 1
            unique_name = f"{sheet_name}_{counter}"
        return unique_name


    # if it is already worksheet format
    if isinstance(df, pd.DataFrame):
        pass
    elif isinstance(df, Worksheet) or isinstance(df, Workbook):
        pass
    elif df is None:
        if any(filename):
            df = fload(filename, output="bit")
    else:
        try:
            print(f"is loading file {os.path.basename(df)}")
            df = fload(df)
        except Exception as e:
            print(e) 
    if isinstance(df, tuple):
        df, original_validations=df
    if filename is None:
        filename = str(datetime.now().strftime("%y%m%d_%H.xlsx"))
    
    kwargs.pop("format", None)  # 更好地跟fsave结合使用
    kwargs.pop("sheet_name", 0)  # 更好地跟df.to_excel结合使用
    # 只有openpyxl才支持 append
    mode = strcmp(kwargs.get("mode", "a"), ["a", "w","auto"])[0]
    # print(f'mode="{mode}"')
    kwargs.pop("mode", None)
    engine = strcmp(kwargs.get("engine", "openpyxl"), ["xlsxwriter", "openpyxl"])[0]
    # corr engine
    engine="openpyxl" if mode=="a" else "xlsxwriter" 
    # print(f'engine="{engine}"')
    if_sheet_exists=kwargs.get("if_sheet_exists","replace")
    # 通常是不需要保存index的
    index = kwargs.get("index", False)
    # header https://pandas.pydata.org/docs/reference/api/pandas.DataFrame.to_excel.html
    header=kwargs.pop("header",True) # bool or list of str, default True: Write out the column names. If a list of string is given it is assumed to be aliases for the column names.
    password = kwargs.pop("password", None)  # Use kwargs if provided
    
    kwargs.pop("password", None)
    kwargs.pop("header", None)
    kwargs.pop("index", None)
    kwargs.pop("if_sheet_exists", None)
 

    if isinstance(df, Workbook):
        """打开Sheet_name指定的表格，如果该表不存在，则创建一个新的或从现有文件中加载数据"""
        wb = df
        
        # Check if sheet exists in input workbook
        try:
            if isinstance(sheet_name, int):
                ws = wb.worksheets[sheet_name]
            else:
                ws = wb[sheet_name]
        except Exception as e:
            # print(f'mode="{mode}"')
            if not os.path.exists(filename) or mode=="w":
                # ws=wb.active
                # ws.title = sheet_name
                ws = wb.create_sheet(title=sheet_name) 
            else:# file exists
                wb = load_workbook(filename) 
                # with pd.ExcelWriter(filename, mode="a", engine=engine, if_sheet_exists=if_sheet_exists) as writer:
                #     for ws in wb.worksheets:  # Iterate through worksheets in the input workbook
                #         ws_df = pd.DataFrame(ws.values)
                #         ws_df.to_excel(writer,sheet_name=sheet_name,index=index,header=header,**kwargs) 
                #         print(3)

                with pd.ExcelWriter(filename, mode="a", engine=engine, if_sheet_exists=if_sheet_exists) as writer:
                    for ws in wb.worksheets:
                        # Convert to DataFrame
                        data = list(ws.values)
                        df = pd.DataFrame(data)
                        try:
                            # Clean up NaNs in header
                            df.columns = df.iloc[0]
                            df = df[1:]
                            df.columns = df.columns.fillna("Unnamed")

                            # Optional: Sanitize cell values
                            df = df.applymap(lambda x: '' if pd.isnull(x) else x)

                            # Write to new sheet or overwrite
                            df.to_excel(writer, sheet_name=ws.title, index=False)
                            print(3)
                        except Exception as e:
                            print(e)

                # 重新打开刚更新过的数据
                wb = load_workbook(filename) 
                # print(4)
                # if sheet_name in wb.sheetnames:
                #     ws = wb[sheet_name]
                #     if not sheet_name==sheet_name:
                #         wb.remove(wb[sheet_name]) 
                # else:
                #     raise KeyError(f"Worksheet {sheet_name} does not exist.") 
                    
        # Handle case where we need to write data to file
        if os.path.exists(filename) and mode != "w":
            try:
                # Backup all sheets
                sheet_data = {}
                for sheet in wb.worksheets:
                    data = list(sheet.values)
                    df = pd.DataFrame(data[1:], columns=data[0])
                    sheet_data[sheet.title] = df
                    
                # Write all sheets back to file
                with pd.ExcelWriter(filename, engine=engine, mode=mode,if_sheet_exists=if_sheet_exists) as writer:
                    for name, df in sheet_data.items():
                        df.to_excel(writer, sheet_name=name, index=index, header=header, **kwargs)
                        
                # 重新打开刚更新过的数据
                wb = load_workbook(filename) 
                if sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    if not sheet_name==sheet_name:
                        wb.remove(wb[sheet_name]) 
                else:
                    raise KeyError(f"Worksheet {sheet_name} does not exist.") 
            except Exception as e:
                print(f"Error saving workbook: {e}")

    else:
        if not os.path.exists(filename) or mode=="w": # or overwrite
            try:
                df = df.applymap(_clean_excel_str)
            except Exception as e:
                print(f"trying to _clean_excel_str:{e}")
            try:
                df = df.applymap(_escape_excel_formula)
            except Exception as e:
                print(f"trying to _escape_excel_formula:{e}")
            try:
                if isinstance(df, pd.io.formats.style.Styler):
                    try:
                        df.to_excel(filename, sheet_name=sheet_name, index=False, engine="openpyxl")
                        print(f"in format_excel: styled file: {os.path.basename(filename)}")
                    except Exception as e: 
                        print(f"in format_excel: cannot save the styles, only save the raw data!,because: {e}")
                        with pd.ExcelWriter(filename, mode="w", engine=engine) as writer:
                            df.data.to_excel(writer, sheet_name=sheet_name, index=False)
                else:
                    with pd.ExcelWriter(filename, mode="w", engine=engine) as writer:
                        if isinstance(df, dict):
                            for name_, df_ in df.items():
                                df_.to_excel(writer, sheet_name=name_, index=index, header=header, **kwargs)
                        else:
                            df.to_excel(writer, sheet_name=sheet_name, index=index, header=header,**kwargs)
            except Exception as e:
                print(f"Failed to save Excel file: {e}")
            wb = load_workbook(filename)
            if sheet_name in wb.sheetnames:
                wb[sheet_name].sheet_state = 'visible'
            else:
                wb.active.sheet_state = 'visible'  # Fallback to the active sheet
            try:
                if isinstance(sheet_name, str):
                    ws = wb[sheet_name]
                elif isinstance(sheet_name, int):
                    ws = wb.worksheets[sheet_name]
                else:
                    ws = wb.worksheets[sheet_name]  # the index of worksheets
            except Exception as e: 
                ws = wb.create_sheet(title=sheet_name)
        else:# file exists
            wb = load_workbook(filename)
            with pd.ExcelWriter(filename, mode="a", engine=engine, if_sheet_exists=if_sheet_exists) as writer:
                if isinstance(df, dict):
                    for name_, df_ in df.items():
                        df_.to_excel(writer, sheet_name=name_, index=index, header=header, **kwargs)
                else:
                    df.to_excel(writer, sheet_name=sheet_name, index=index, header=header,**kwargs)
            wb = load_workbook(filename) 
            if sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                if not sheet_name==sheet_name:
                    wb.remove(wb[sheet_name]) 
            else:
                print(e)
                ws = wb.create_sheet(title=sheet_name) 
 
    # *  Copy Formatting when available
    if template:
        if isinstance(template, dict):
            path_source=template.get("path")
            sheet_name_source=template.get("sheet_name",sheet_name)
        else:
            path_source=template
            sheet_name_source=sheet_name
        if isinstance(sheet_name_source, (list,tuple)):
            pass
        else:
            if sheet_name_source!=sheet_name:
                sheet_name_source=list(sheet_name_source).append(sheet_name)
        try: 
            copy_format(path_source, wb, sheet_name=sheet_name_source, save=False,copy_values=False)
            print(f"Copied formatting from sheet_name '{sheet_name}'in file '{os.path.basename(path_source)}'")
        except Exception as e:
            print(e)
            copy_format(path_source, wb, sheet_name=[pd.ExcelFile(path_source).sheet_names[0],sheet_name], save=False,copy_values=False)
            print(f"Alternatively, Copied formatting from the 1st Sheet '{sheet_name}'in file '{os.path.basename(path_source)}'")
    # ! Apply Text color
    if text_color:
        if verbose:
            text_color_str="""
            text_color=[
                {
                    (slice(1, 2), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=8": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
                {
                    (slice(3, df.shape[0] + 1), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=10": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
            ],
            """
            print(text_color_str)
        if not isinstance(text_color, list):
            text_color=[text_color]
        for text_color_ in text_color:
            for indices, dict_text_conditions in text_color_.items():
                ws = apply_color_to_worksheet(ws, sheet_name=sheet_name, conditions=dict_text_conditions, cell_idx=indices,where="text")
    # ! Apply Text color
    if bg_color:
        if verbose:
            bg_color_str="""
            bg_color=[
                {
                    (slice(1, 2), slice(0, 3)): {
                        "x>20": ["#DD0531","#35B20C"],  # Numbers > 20 → red
                        "x<=8": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
                {
                    (slice(3, df.shape[0] + 1), slice(0, 3)): {
                        "x>20": "#DD0531",  # Numbers > 20 → red
                        "x<=10": "#35B20C",  # Numbers ≤ 10 → blue
                        "'x'!='available'": "#0510DD",  # 'available' → green
                        "10<x<=30": "#EAB107",  # 10 < value ≤ 30 → orange
                        "10<=x<30": "#C615BE",  # 10 ≤ value < 30 → purple
                    }
                },
            ],
            """
            print(bg_color_str)
        if not isinstance(bg_color, list):
            bg_color=[bg_color]
        for bg_color_ in bg_color:
            for indices, dict_text_conditions in bg_color_.items():
                ws = apply_color_to_worksheet(ws, sheet_name=sheet_name, conditions=dict_text_conditions, cell_idx=indices,where="bg")
    # !Apply cell formatting
    if cell:
        if not isinstance(cell, list):
            cell = [cell]
        for cell_ in cell:
            for indices, format_options in cell_.items():
                cell_range = convert_indices_to_range(*indices)
                apply_format(ws, format_options, cell_range)

        if verbose:
            cell_tmp="""cell=[
                    {
                        (slice(0, 1), slice(0, len(df.columns))): {
                            "font": {
                                "name": "Calibri",  # Font name
                                "size": 14,  # Font size
                                "bold": True,  # Bold text
                                "italic": False,  # Italic text
                                "underline": "single",  # Underline (single, double)
                                "color": "#FFFFFF",  # Font color
                            },
                            "fill": {
                                "start_color": "a1cd1e",  # Starting color
                                "end_color": "4F81BD",  # Ending color (useful for gradients)
                                "fill_type": "solid",  # Fill type (solid, gradient, etc.)
                            },
                            "alignment": {
                                "horizontal": "center",  # Horizontal alignment (left, center, right)
                                "vertical": "center",  # Vertical alignment (top, center, bottom)
                                "wrap_text": True,  # Wrap text in the cell
                                "shrink_to_fit": True,  # Shrink text to fit within cell
                                "text_rotation": 0,  # Text rotation angle
                            },
                            "border": {
                                "left": "thin",
                                "right": "thin",
                                "top": "thin",
                                "bottom": "thin",
                                "color": "000000",  # Border color
                            },
                        }
                    },{}]"""
            print(cell_tmp)
    # !Apply cell shading
    if shade:
        if not isinstance(shade, list):
            shade = [shade]
        for shade_ in shade:
            for indices, shading in shade_.items():
                cell_range = convert_indices_to_range(*indices)
                fill = PatternFill(
                    start_color=hex2argb(shading.get("bg_color", "FFFFFF")),
                    end_color=hex2argb(shading.get("end_color", "FFFFFF")),
                    fill_type=shading.get("fill_type", "solid"),
                    patternType=shading.get("pattern_type", "solid"),
                    fgColor=hex2argb(shading.get("fg_color", "0000FF")),
                )
                for row in ws[cell_range]:
                    for cell in row:
                        cell.fill = fill
        if verbose:
            shade_temp="""shade={        
                        (slice(1, 4), slice(1, 3)): {
                        "bg_color": "#63C187",  # Background color
                        "pattern_type": "solid",  # Fill pattern (e.g., solid, darkGrid, lightGrid)
                        "fg_color": "#0000FF",  # Foreground color, used in patterns
                        "end_color": "0000FF",  # End color, useful for gradients
                        "fill_type": "solid",  # Type of fill (solid, gradient, etc.)
                    }}"""
            print(shade_temp)
    # !number formatting
    if number_format:
        if not isinstance(number_format, list):
            number_format = [number_format]
        for number_format_ in number_format:
            for col_idx, fmt in number_format_.items():
                col_letter = get_column_letter(col_idx)
                for cell in ws[col_letter][1:]:  # Skip the header
                    cell.number_format = fmt
        if verbose:
            number_format_temp="""number_format={
                                1: "0.00",  # Two decimal places for column index 1
                                2: "#,##0",  # Thousands separator
                                3: "0%",  # Percentage format
                                4: "$#,##0.00",  # Currency format
                            }""" 
            print(number_format_temp)
    
    if freeze:
        if isinstance(freeze,bool):
            freeze='A2'
        ws.freeze_panes = freeze  # Freeze everything above and to the left of A2
    if apply_filter:
        if isinstance(apply_filter, bool):
            # Default: Apply filter to the entire first row (header)
            filter_range = f"A1:{get_column_letter(ws.max_column)}1"
            ws.auto_filter.ref = filter_range
            if not freeze:
                ws.freeze_panes = "A2"  # Freeze everything above and to the left of A2
        elif isinstance(apply_filter, tuple):
            row_slice, col_slice = apply_filter
            # Extract the start and end indices for rows and columns
            start_row, end_row = row_slice.start, row_slice.stop
            start_col_idx, end_col_idx = col_slice.start, col_slice.stop

            # Ensure valid row and column indices
            if start_row < 1: start_row = 1  # Row should not be less than 1
            if end_row > ws.max_row: end_row = ws.max_row  # Ensure within sheet's row limits
            if start_col_idx < 1: start_col_idx = 1  # Column should not be less than 1
            if end_col_idx > ws.max_column: end_col_idx = ws.max_column  # Ensure within sheet's column limits
            
            # Get column letters based on indices
            start_col = get_column_letter(start_col_idx)
            end_col = get_column_letter(end_col_idx)

            # Define the filter range based on specific rows and columns
            filter_range = f"{start_col}{start_row}:{end_col}{end_row}"

            # Apply the filter
            ws.auto_filter.ref = filter_range
            if freeze:
                ws.freeze_panes = freeze  # Freeze everything above and to the left of A2
    # !widths 
    if isinstance(width, bool):
        width=None if width else False
    if isinstance(height,bool):
        height=None if height else False

    merged_cells = set()
    for merged_range in ws.merged_cells.ranges:
        for row in ws.iter_rows(min_row=merged_range.min_row,
                            max_row=merged_range.max_row,
                            min_col=merged_range.min_col,
                            max_col=merged_range.max_col):
            for cell in row:
                merged_cells.add(cell.coordinate)
    if width is None or width == {}:  # automatic adjust width
        print("auto-width")
        for col in ws.columns:
            if not col:
                continue
            try:
                col_letter = get_column_letter(col[0].column)
                
                # Skip entire column if any cell is merged
                if any(cell.coordinate in merged_cells for cell in col):
                    continue
                    
                max_length = 0
                for cell in col:
                    try:
                        if cell.value:
                            cell_value = str(cell.value)
                            if '\n' in cell_value:
                                max_line_length = max(len(line) for line in cell_value.split('\n'))
                                max_length = max(max_length, max_line_length)
                            else:
                                max_length = max(max_length, len(cell_value))
                    except:
                        pass
                        
                adjusted_width = (max_length * width_factor) + width_padding
                if width_max is not None:
                    adjusted_width = min(adjusted_width, width_max)
                ws.column_dimensions[col_letter].width = max(5, adjusted_width)
                
            except Exception as e:
                print(f"Error adjusting width for column: {e}")
                continue
    elif isinstance(width, (int, float)): # set all columns to this value
        print("set to fixed width {}".format(width))
        for col in ws.columns:
            column = get_column_letter(col[0].column)
            ws.column_dimensions[column].width = width * width_factor + width_padding
    elif isinstance(width, dict):  # custom widths per column
        for col_idx, width_ in width.items():
            col_letter = get_column_letter(col_idx)
            ws.column_dimensions[col_letter].width = width_
 
    # !heights
    if height is None or height=={}:  # automatic adust height
        for row in ws.iter_rows(min_row=1, max_row=ws.max_row):
            max_height = 0
            for cell in row:
                if cell.value:
                    lines = str(cell.value).split("\n")
                    max_line_length = max(len(line) for line in lines)
                    estimated_height = 15 * len(lines)
                    if max_line_length > 20:
                        estimated_height += 5 * (max_line_length // 20)
                    max_height = max(max_height, estimated_height)
            height_adj=max_height*height_factor+height_padding
            if height_max is not None:
                height_adj=min(height_adj, height_max)
                    
            ws.row_dimensions[row[0].row].height = height_adj
    elif isinstance(height,bool) and not height:
        pass 
    elif isinstance(height, (int, float)):
        # set all rows to the same fixed height
        for i in range(1, ws.max_row + 1):
            ws.row_dimensions[i].height = height*height_factor+height_padding

    else:
        for row, height_ in height.items():
            ws.row_dimensions[row].height = height_

    # !Merge cells using slice indices
    if merge:
        if isinstance(merge, tuple):
            merge = [merge]
        for indices in merge:
            # Ensure indices are slice objects
            if len(indices) == 2:
                row_slice, col_slice = indices
                merge_range = convert_indices_to_range(row_slice, col_slice)
                ws.merge_cells(merge_range)
            elif len(indices) == 4:
                start_row, start_col, end_row, end_col = indices
                # Convert column numbers to letters (e.g., 1 -> 'A')
                start_cell = f"{get_column_letter(start_col)}{start_row}"
                end_cell = f"{get_column_letter(end_col)}{end_row}"
                merge_range = f"{start_cell}:{end_cell}"
                ws.merge_cells(merge_range)
            else:
                raise ValueError(
                    f"两种方式: 1. format: (start_row, start_col, end_row, end_col), 2. format: (slice(0, 3), slice(1, 2))"
                )

    # !Add comment
    if comment:
        if not isinstance(comment, list):
            comment = [comment]
        for comment_ in comment:
            if not isinstance(comment_, dict):
                raise TypeError("Each item in the `comments` list must be a dictionary.")
        
            for (row, col), comment_str in comment_.items():
                if not isinstance(row, int) or not isinstance(col, int):
                    raise ValueError("Row and column indices must be integers.")
                if not isinstance(comment_str, str):
                    raise ValueError("Comment text must be a string.")

                comment_curr = Comment(comment_str, "Author")
                comment_curr.visible = comment_always_visible
                # if comment_always_visible:
                #     comment_curr.width = 200  # Adjust width
                #     comment_curr.height = 100  # Adjust height
                ws.cell(row=row + 1, column=col + 1).comment = comment_curr
        if verbose:
            comment_tmp="""comment=[
                        {(0, 0): "This is a comment for A1"},
                        {(1, 1): "This is a comment for B2"},
                        {(2, 2): "This is a comment for C3"},
                    ]"""
            print(comment_tmp)
    # !Add link
    if link:
        if not isinstance(link, list):
            link = [link]
        for link_ in link:
            for (row, col), link_str in link_.items():
                ws.cell(row=row + 1, column=col + 1).hyperlink = link_str
        if verbose:
            print('link={(2, 2): "https://example.com"}')
    # !Apply data validation
    if data_validation:
        # default_dropdown_settings = {
        #     "type":"list",
        #     "allow_blank": True,
        #     "showDropDown": False,
        #     "showErrorMessage": False,
        #     "errorTitle": "",
        #     "error": "",
        # }
        # for indices, validation in data_validation.items():
        #     cell_range = convert_indices_to_range(*indices)
        #     # Preprocess formula1 if it's given as a list
        #     if "formula1" in validation and isinstance(validation["formula1"], list):
        #         validation["formula1"] = '"' + ",".join(str(i) for i in validation["formula1"]) + '"'
        #     for key, value in default_dropdown_settings.items():
        #         validation.setdefault(key, value)
        #     dv = DataValidation(**validation)
        #     ws.add_data_validation(dv)
        #     dv.add(cell_range)
        if verbose:
            print("""data_validation={
                        (slice(1, 2), slice(2, 10)): {
                            "type": "list",
                            "formula1": '"Option1,Option2,Option3"',  # List of options
                            "allow_blank": True,
                            "showDropDown": False,
                            "showErrorMessage": True,
                            "errorTitle": "Invalid input",
                            "error": "Please select a valid option.",
                        }
                    }"""
                  )
    if data_validation:
        # Default settings based on type
        default_settings_by_type = {
            "list": {
                "allow_blank": True,
                "showDropDown": False,
                "showErrorMessage": False,
                "errorTitle": "",
                "error": "",
            },
            "whole": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid number",
                "error": "Please enter a valid whole number.",
            },
            "decimal": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid decimal",
                "error": "Please enter a valid decimal number.",
            },
            "date": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid date",
                "error": "Please enter a valid date.",
            },
            "time": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid time",
                "error": "Please enter a valid time.",
            },
            "textLength": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid text length",
                "error": "Please enter text with correct length.",
            },
            "custom": {
                "allow_blank": True,
                "showErrorMessage": True,
                "errorTitle": "Invalid input",
                "error": "Your input does not match the required format.",
            },
        }

        for indices, validation in data_validation.items():
            cell_range = convert_indices_to_range(*indices)

            # Get default settings based on type
            validation_type = validation.get("type", "list")
            # Preprocess formula1 if it's given as a list (for list type)
            validation["type"]=validation_type
            print(validation)
            if validation.get("type") == "list" and "formula1" in validation and isinstance(validation["formula1"], list):
                print('yes')
                validation["formula1"] = '"' + ",".join(str(i) for i in validation["formula1"]) + '"'

            default_settings = default_settings_by_type.get(validation_type, {})

            # Apply defaults (only if not already set)
            for key, value in default_settings.items():
                validation.setdefault(key, value)
            print(validation)
            dv = DataValidation(**validation)
            ws.add_data_validation(dv)
            dv.add(cell_range)

    # !Protect sheet with a password
    # Fetch the password
    # if all([password is not None, any([protect, isinstance(password, (str, list, tuple)) and any(password)])]):  # Check if protection options are provided
    #     from openpyxl.worksheet.protection import SheetProtection

    #     if protect is None:
    #         protect = {}  # Ensure protect is always a dictionary
    #     print(protect)
    #     # Default to 'protect' password if not explicitly set
    #     password = password or protect.get("password")

    #     # Create a SheetProtection object with the necessary settings
    #     protection = SheetProtection(
    #         password=password,
    #         sheet=True,  # Enable sheet protection
    #         objects=protect.get("objects", True),
    #         scenarios=protect.get("scenarios", False),
    #         formatCells=protect.get("formatCells", False),
    #         formatColumns=protect.get("formatColumns", False),
    #         formatRows=protect.get("formatRows", False),
    #         insertColumns=protect.get("insertColumns", True),
    #         insertRows=protect.get("insertRows", True),
    #         deleteColumns=protect.get("deleteColumns", True),
    #         deleteRows=protect.get("deleteRows", True),
    #         selectLockedCells=protect.get("selectLockedCells", False),
    #         selectUnlockedCells=protect.get("selectUnlockedCells", False)
    #     )

    #     # Apply the protection to the worksheet
    #     ws.protection = protection

    if bool(protect):  # Check if protection options are provided
        from openpyxl.worksheet.protection import SheetProtection

        if protect is None:
            protect = {}  # Ensure protect is always a dictionary

        # Default to 'protect' password if not explicitly set
        password = password or protect.get("password")

        # Apply sheet protection with password if a password is provided
        if password:
            ws.protection.sheet = True
            ws.protection.set_password(password)  # Directly set the password
        else:
            ws.protection.sheet = protect.get("sheet", False)  # Use 'protect' settings if no password 
        # 创建 SheetProtection 对象，并传入各项保护配置
        protection = SheetProtection(
            password=password,  # 设置密码，用户在 Excel 中取消保护时需要输入
            sheet=True,  # 启用工作表保护，必须为 True 其余选项才生效

            # 以下为保护选项（True 表示允许操作，False 表示禁止操作）：
            objects=protect.get("objects", False),  # 是否允许编辑嵌入对象（如图表、按钮等）
            scenarios=protect.get("scenarios", False),  # 是否允许访问“方案管理器”（Excel 的假设分析功能）

            # 单元格格式设置权限：
            formatCells=protect.get("formatCells", False),  # 是否允许更改单元格格式（字体、颜色、边框等）
            formatColumns=protect.get("formatColumns", False),  # 是否允许调整列宽
            formatRows=protect.get("formatRows", False),  # 是否允许调整行高

            # 插入与删除权限：
            insertColumns=protect.get("insertColumns", False),  # 是否允许插入新列
            insertRows=protect.get("insertRows", False),  # 是否允许插入新行
            deleteColumns=protect.get("deleteColumns", False),  # 是否允许删除列
            deleteRows=protect.get("deleteRows", False),  # 是否允许删除行

            # 选择单元格权限：
            selectLockedCells=protect.get("selectLockedCells", False),  # 是否允许选择被锁定的单元格
            selectUnlockedCells=protect.get("selectUnlockedCells", True)  # 是否允许选择未被锁定的单元格
        )
        ws.protection = protection
    # !conditional formatting
    if conditional_format:
        if not isinstance(conditional_format, list):
            conditional_format = [conditional_format]
        print(f"conditional_format dict setting: 'color_scale', 'data_bar' and 'icon_set[not-ready]'")
        for conditional_format_ in conditional_format:
            for indices, rules in conditional_format_.items():
                cell_range = convert_indices_to_range(*indices)
                if not isinstance(rules, list):
                    rules=[rules]
                for rule in rules:
                    # Handle color scale
                    if "color_scale" in rule:
                        if verbose:
                            color_scale_tmp="""
                                    conditional_format={
                                        (slice(1, df.shape[0] + 1), slice(1, 2)):
                                            {
                                                "color_scale": {
                                                    "start_type": "min",
                                                    "start_value": 0,
                                                    "start_color": "#74ADE9",
                                                    "mid_type": "percentile",
                                                    "mid_value": 50,
                                                    "mid_color": "74ADE9",
                                                    "end_type": "max",
                                                    "end_value": 100,
                                                    "end_color": "#B62833",
                                                }
                                            }}
                            """
                            print(color_scale_tmp)
                        color_scale = rule["color_scale"] 

                        color_scale_rule = ColorScaleRule(
                            start_type=color_scale.get("start_type", "min"),
                            start_value=color_scale.get("start_value",None),
                            start_color=hex2argb(color_scale.get("start_color", "#74ADE9")),
                            mid_type=color_scale.get("mid_type","percentile"),
                            mid_value=color_scale.get("mid_value",None),
                            mid_color=hex2argb(color_scale.get("mid_color", "FFFFFF")),
                            end_type=color_scale.get("end_type", "max"),
                            end_value=color_scale.get("end_value",None),
                            end_color=hex2argb(color_scale.get("end_color", "#B62833")),
                        )
                        ws.conditional_formatting.add(cell_range, color_scale_rule)
                    # Handle data bar
                    if "data_bar" in rule:
                        if verbose:
                            data_bar_tmp="""
                                    conditional_format={
                                        (slice(1, df.shape[0] + 1), slice(1, 2)):
                                            {
                                                "data_bar": {
                                                    "start_type": "min",
                                                    "start_value": None,
                                                    "end_type": "max",
                                                    "end_value": None,
                                                    "color": "F6C9CE",
                                                    "show_value": True,
                                                }
                                            }}
                            """
                            print(data_bar_tmp)
                        data_bar = rule["data_bar"]
                        bar_color = hex2argb(data_bar.get("color", "638EC6"))

                        data_bar_rule = DataBarRule(
                            start_type=data_bar.get("start_type", "min"),
                            start_value=data_bar.get("start_value",None),
                            end_type=data_bar.get("end_type", "max"),
                            end_value=data_bar.get("end_value",None),
                            color=bar_color,
                            showValue=data_bar.get("show_value", True),
                        )
                        ws.conditional_formatting.add(cell_range, data_bar_rule)

                    # Handle icon setse
                    if "icon_set" in rule:
                        icon_set = rule["icon_set"]
                        icon_set_rule = IconSet(
                            iconSet=icon_set.get("iconSet", "3TrafficLights1"),  # Corrected
                            showValue=icon_set.get("show_value", True),        # Corrected
                            reverse=icon_set.get("reverse", False)            # Corrected
                        )
                        ws.conditional_formatting.add(cell_range, icon_set_rule)
                    # Handle text-based conditions
                    if "text_color" in rule: # not work
                        from openpyxl.styles.differential import DifferentialStyle
                        from openpyxl.formatting.rule import Rule
                        from openpyxl.styles import PatternFill

                        # Extract the fill properties from the rule
                        fill = rule.get("fill", {})
                        start_color = fill.get("start_color", "FFFFFF")  # Default to white if not specified
                        end_color = fill.get("end_color", "FFFFFF")  # Default to white if not specified
                        fill_type = fill.get("fill_type", "solid")  # Default to solid fill if not specified

                        # Extract the text condition or default to a space if 'text' is not provided
                        text = rule.get("text", " ")  

                        # Create the DifferentialStyle using the extracted fill settings
                        dxf = DifferentialStyle(
                            fill=PatternFill(start_color=start_color, end_color=end_color, fill_type=fill_type)
                        )
                        
                        # Create the text rule based on the text condition
                        text_rule = Rule(
                            type="containsText",  # The type of condition
                            operator=rule.get("operator", "equal"),  # Default operator is "equal"
                            text=text,
                            dxf=dxf,  # Apply the fill color from DifferentialStyle
                        )
                        ws.conditional_formatting.add(cell_range, text_rule)
        if verbose:
            conditional_format_temp="""
                    conditional_format={
                            (slice(1, 3), slice(1, 4)): [
                                {
                                    "data_bar": {
                                        "start_type": "min",
                                        "start_value": 100,
                                        "end_type": "max",
                                        "end_value": None,
                                        "color": "F6C9CE",
                                        "show_value": True,
                                    }
                                },
                                {
                                    "color_scale": {
                                        "start_type": "min",
                                        "start_value": 0,
                                        "start_color": "#74ADE9",
                                        "mid_type": "percentile",
                                        "mid_value": 50,
                                        "mid_color": "74ADE9",
                                        "end_type": "max",
                                        "end_value": 100,
                                        "end_color": "#B62833",
                                    }
                                },
                            ]
                        }
            """
            print(conditional_format_temp)
    if insert_img:
        if not isinstance(insert_img, dict):
            raise ValueError(f'insert_img 需要dict格式: e.g., insert_img={"A1":"example.png"}')
        try:
            from openpyxl import drawing
            from PIL import Image
            import PIL
            for img_cell, img_data in insert_img.items():
                img_width = img_height = None
                pil_img=img_path = None
                if isinstance(img_data, dict):
                    if "path" in img_data:
                        img_ = drawing.image.Image(img_data["path"])# File path
                    elif "image" in img_data:
                        img_ = drawing.image.Image(img_data["image"])# PIL Image object
                    elif "array" in img_data:
                        img_ = drawing.image.Image(Image.fromarray(img_data["array"]))# Convert NumPy array to PIL Image

                    img_width = img_data.get("width", None)
                    img_height = img_data.get("height", None)
                elif isinstance(img_data, str):
                    img_ = drawing.image.Image(img_data)# Direct file path
                elif isinstance(img_data, (PIL.Image.Image,PIL.PngImagePlugin.PngImageFile)):
                    img_ = drawing.image.Image(img_data)# Direct PIL Image object
                elif isinstance(img_data, np.ndarray):
                    img_ = drawing.image.Image(Image.fromarray(img_data))# Convert NumPy array to PIL Image
                elif pil_img:
                    img_ = drawing.image.Image(pil_img)

                # Set width and height if provided
                if img_width is not None:
                    img_.width = img_width
                if img_height is not None:
                    img_.height = img_height
                ws.add_image(img_, img_cell)  
                print(f"✅ at column '{img_cell}' inserted image ==>File: {os.path.basename(filename)}")

        except Exception as e:
            print(e)
    if protect_file and isinstance(protect_file,dict):
        pass_=protect_file.get("password", None) 
        wb.security=WorkbookProtection(
                workbookPassword= enpass(pass_) if len(pass_)==len(enpass('a')) else pass_,
                lockStructure=protect_file.get("lock_structure", True),
                lockWindows=protect_file.get("lock_window", True),
                revisionsPassword=protect_file.get("warning", ""),# Add this line to suppress warnings
            )
        print(f"file is protected with password: {enpass(protect_file.get("password", None))[:6]}")

    # ungroup sheets
    for sheet in wb.worksheets:
        sheet.sheet_view.tabSelected = False
    # !Save the workbook
    try:
        wb.save(filename)
    except Exception as e:
         print(f"Error saving workbook: {str(e)}")
    # Replace your final save operation with this:
     # try:
         # # Create a temporary file for safer saving
         # temp_filename = filename + '.tmp'
         # wb.save(temp_filename)
        
         # # If save succeeds, replace original file
         # if os.path.exists(filename):
             # os.remove(filename)
         # os.rename(temp_filename, filename)
        
     # except Exception as e:
         # print(f"Error saving workbook: {str(e)}")
         # if os.path.exists(temp_filename):
             # os.remove(temp_filename)
         # raise


def preview(var):
    """Master function to preview formatted variables in Jupyter."""
    from bs4 import BeautifulSoup
    from IPython.display import display, HTML, Markdown

    if isinstance(var, str):
        if isa(var, "html"):
            display(HTML(var))  # Render as HTML
        # Check if it's a valid markdown
        elif var.startswith("#"):
            display(Markdown(var))
        else:
            # Otherwise, display as plain text
            print(var)
    elif isinstance(var, BeautifulSoup):
        preview(str(var))
    elif isinstance(var, pd.DataFrame):
        # Display pandas DataFrame
        display(var)

    elif isinstance(var, list) or isinstance(var, dict):
        import json

        # Display JSON
        json_str = json.dumps(var, indent=4)
        display(Markdown(f"```json\n{json_str}\n```"))

    elif isinstance(var, bytes):
        # Display image if it's in bytes format
        display(Image(data=var))

    elif isinstance(var, str) and (var.endswith(".png") or var.endswith(".jpg")):
        # Display image from file path
        display(Image(filename=var))

    elif isinstance(var, dict):
        import json

        # Handle dictionary formatting
        json_str = json.dumps(var, indent=4)
        display(Markdown(f"```json\n{json_str}\n```"))

    else:
        # If the format is not recognized, print a message
        print("Format not recognized or unsupported.")


# # Example usages:
# preview("This is a plain text message.")
# preview("# This is a Markdown header")
# preview(pd.DataFrame({"Name": ["Alice", "Bob"], "Age": [25, 30]}))
# preview({"key": "value", "numbers": [1, 2, 3]})

def df2db(df: pd.DataFrame, db_url: str, table_name: str, if_exists: str = "replace"):
    """
    Save a pandas DataFrame to a SQL database (SQLite, MySQL, PostgreSQL).

    Usage:
    case 1. SQLite:
        df2db(df, 'sqlite:///sample_manager.db', 'users', 'replace')
    case 2. MySQL:
        df2db(df, 'mysql+mysqlconnector://user:password@localhost/mydatabase', 'users', 'replace')
    case 3. PostgreSQL:
        df2db(df, 'postgresql://user:password@localhost/mydatabase', 'users', 'replace')

    Parameters:
    - df (pd.DataFrame): The DataFrame to save.
    - db_url (str): The SQLAlchemy connection string to the database.
    - table_name (str): The name of the table to save the DataFrame to.
    - if_exists (str): What to do if the table already exists. Options:
        - 'replace': Drop the table before inserting new values.
        - 'append': Append new values to the table.
        - 'fail': Do nothing if the table exists. Default is 'replace'.
    """
    from sqlalchemy import create_engine
    from sqlalchemy.exc import SQLAlchemyError

    try:
        # Create the SQLAlchemy engine based on the provided db_url
        engine = create_engine(db_url)

        # Save the DataFrame to the SQL database (table)
        df.to_sql(table_name, con=engine, index=False, if_exists=if_exists)

        print(f"Data successfully saved to {db_url} in the {table_name} table.")

    except SQLAlchemyError as e:
        # Handle SQLAlchemy-related errors (e.g., connection issues, query issues)
        print(f"Error saving DataFrame to database: {str(e)}")

    except Exception as e:
        # Handle other unexpected errors
        print(f"An unexpected error occurred: {str(e)}")

def _df_outlier(
    data,
    columns=None,
    method=["zscore", "iqr", "percentile", "iforest"],
    min_outlier_method=3,  # 至少两种方法检查出outlier
    zscore_threshold=3,
    iqr_threshold=1.5,
    lower_percentile=5,
    upper_percentile=95,
):
    from scipy.stats import zscore
    from sklearn.ensemble import IsolationForest
    from sklearn.preprocessing import StandardScaler

    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0
    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()
    # Separate numeric and non-numeric columns
    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    # if columns is not None:
    #     numeric_data = numeric_data[columns]
    if numeric_data.empty:
        raise ValueError("Input data must contain numeric columns.")

    outliers_df = pd.DataFrame(index=numeric_data.index)
    if isinstance(method, str):
        method = [method]

    # Z-score method
    if "zscore" in method:
        z_scores = np.abs(zscore(numeric_data))
        outliers_df["zscore"] = np.any(z_scores > zscore_threshold, axis=1)

    # IQR method
    if "iqr" in method:
        Q1 = numeric_data.quantile(0.25)
        Q3 = numeric_data.quantile(0.75)
        IQR = Q3 - Q1
        lower_bound = Q1 - iqr_threshold * IQR
        upper_bound = Q3 + iqr_threshold * IQR
        outliers_df["iqr"] = (
            (numeric_data < lower_bound) | (numeric_data > upper_bound)
        ).any(axis=1)

    # Percentile method
    if "percentile" in method:
        lower_bound = numeric_data.quantile(lower_percentile / 100)
        upper_bound = numeric_data.quantile(upper_percentile / 100)
        outliers_df["percentile"] = (
            (numeric_data < lower_bound) | (numeric_data > upper_bound)
        ).any(axis=1)

    # Isolation Forest method
    if "iforest" in method:
        # iforest method cannot handle NaNs, then fillna with mean
        numeric_data_ = numeric_data.fillna(numeric_data.mean())
        scaler = StandardScaler()
        scaled_data = scaler.fit_transform(numeric_data_)
        iso_forest = IsolationForest(contamination=0.05)
        outliers_df["iforest"] = iso_forest.fit_predict(scaled_data) == -1

    # Combine all outlier detections
    if len(method) == 4:  # all method are used:
        outliers_df["outlier"] = outliers_df.sum(axis=1) >= min_outlier_method
    else:
        outliers_df["outlier"] = outliers_df.any(axis=1)

    # Handling Outliers: Remove or Winsorize or Replace with NaN
    processed_data = numeric_data.copy()

    processed_data.loc[outliers_df["outlier"]] = np.nan

    return processed_data

def df_group(
    data: pd.DataFrame,
    columns: Union[str, list, None] = None,
    by: str = None,
    param: Dict[str, Any] = None,
    sep: Union[str, list] = [", ",","],
    dropna: bool = True,
    unique: bool = False,
    astype: type = str,
    merge: List[str] = None,
    merge_symbo_column:str=' & ',
    merge_symbo_cell:str='[]',# ["{}","()","[]"]
) -> pd.DataFrame:
    """
    Groups a dataframe based on a specified column and applies aggregation functions dynamically.
    
    Parameters:
    data (pd.DataFrame): The dataframe to be grouped.
    columns (Union[str, list, None]): Columns to select; if None, all columns are selected.
    by (str): The column name to group by.
    param (dict): A dictionary specifying aggregation rules.
    sep (Union[str, list]): Separator for concatenated values. when sep is a list, then sep[0] used for general, sep[1] used in the merging
    dropna (bool): Whether to drop NaN values before aggregation.
    unique (bool): Whether to apply uniqueness before concatenation.
    astype (type): Data type to cast values before aggregation.
    merge (List[str]): List of columns to merge into a single paired column.
    merge_symbo_column:str: indicate in the columns, default ("&")
    merge_symbo_cell:str=default: '{}' or can select from ["{}","()","[]"]

    Usage:
    data = pd.DataFrame({
                        "Cage Tag": [1, 1, 2, 2, 3],
                        "Physical Tag": ["A1", "A2", "B1", "B2", "C1"],
                        "Sex": ["M", "F", "M", "M", "F"],
                        "Date of Birth": ["2021-06-01", "2021-06-02", "2021-07-01", "2021-07-02", "2021-08-01"],
                        "Age": [34, 35, 30, 31, 29],
                        "State": ["Mating", "Resting", "Mating", "Resting", "Mating"],
                        "Owner": ["Dr. Smith", "Dr. Smith", "Dr. Brown", "Dr. Brown", "Dr. Lee"],
                        })
    display(data)
    result = df_group(data, 
                    #   columns=["Sex", "Date of Birth", "Age", "State"], 
                    by="Cage Tag", 
                    merge=["Age", "State"], 
                    merge_symbo_column="|",
                    #   astype=str
                    # sep=[',',    '_'],
                    merge_symbo_cell=None
                    )  
    result

    
    """
    if param is None:
        param = {}
    if columns is None:
        columns = data.columns.tolist()
    elif isinstance(columns, str):
        columns = [columns]
    if not isinstance(sep, list):
        sep = [sep]
    sep.extend(sep) if len(sep)==1 else None

    # Merge specified columns into a single column
    if merge:
        merge_col_name = merge_symbo_column.join(merge)
        # data[merge_col_name] = data[merge].apply(lambda row: tuple(map(astype, row.dropna() if dropna else row)), axis=1)
        if merge_symbo_cell is None: 
            data[merge_col_name] = data[merge].apply(lambda row: f"{sep[1].join(map(astype, row.dropna()))}" if dropna else f"{sep[1].join(map(astype, row))}", axis=1)
        elif len(merge_symbo_cell)==2:
            data[merge_col_name] = data[merge].apply(lambda row: f"{merge_symbo_cell[0]}{sep[1].join(map(astype, row.dropna()))}{merge_symbo_cell[1]}" if dropna else f"{merge_symbo_cell[0]}{sep[1].join(map(astype, row))}{merge_symbo_cell[1]}", axis=1)
        else:
            data[merge_col_name] = data[merge].apply(lambda row: f"[{sep[1].join(map(astype, row.dropna()))}]" if dropna else f"[{sep[1].join(map(astype, row))}]", axis=1)
        columns.append(merge_col_name)
        
    default_aggregations = {
        col: (lambda x: sep[0].join(map(astype, x.dropna().unique() if unique else x.dropna())) if dropna else sep[0].join(map(astype, x.unique() if unique else x)))
        for col in columns if col != by and (merge is None or col not in merge)
    } 
    aggregation_rules = {**default_aggregations, **param}
    
    grouped_df = data.groupby(by).agg(aggregation_rules).reset_index() 
    return grouped_df

def df_outlier(
    data,
    columns=None,
    method=["zscore", "iqr", "percentile", "iforest"],
    min_outlier_method=2,  # 至少两种方法检查出outlier
    zscore_threshold=3,
    iqr_threshold=1.5,
    lower_percentile=5,
    upper_percentile=95,
):
    """
    Usage:
    data_out = df_outlier(
        data,
        columns=["income"],
        method="iforest",
        min_outlier_method=1)

    Advanced outlier detection and handling function.

    Parameters:
    - data: DataFrame, the input data (numerical).
    - method: List, the outlier detection method to use. Options: 'zscore', 'iqr', 'percentile', 'iforest'.
    - zscore_threshold: float, threshold for Z-score outlier detection (default 3).
    - iqr_threshold: float, threshold for IQR method (default 1.5).
    - lower_percentile: float, lower percentile for percentile-based outliers (default 5).
    - upper_percentile: float, upper percentile for percentile-based outliers (default 95).
    - keep_nan: bool, whether to replace outliers with NaN (default True).
    - plot: bool, whether to visualize the outliers (default False).
    - min_outlier_method: int, minimum number of method that need to flag a row as an outlier (default 2).
    - inplace: bool, whether to modify the original `data` DataFrame (default False).

    Returns:
    - processed_data: DataFrame with outliers handled based on method (if winsorize/remove is True).
    """
    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()

    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    _outlier_df_tmp = pd.DataFrame()
    for col in numeric_data.columns:
        _outlier_df_tmp = pd.concat(
            [
                _outlier_df_tmp,
                _df_outlier(
                    data=data,
                    columns=[col],
                    method=method,
                    min_outlier_method=min_outlier_method,  # 至少两种方法检查出outlier
                    zscore_threshold=zscore_threshold,
                    iqr_threshold=iqr_threshold,
                    lower_percentile=lower_percentile,
                    upper_percentile=upper_percentile,
                ),
            ],
            axis=1,
            # join="inner",
        )
    processed_data = pd.concat([_outlier_df_tmp, non_numeric_data], axis=1)
    processed_data = processed_data[col_names_org]
    return processed_data


def df_extend(data: pd.DataFrame, column, axis=0, sep=None, prefix="col"):
    """
    Extend a DataFrame by the list elecments in the column.

    Parameters:
    ----------
    data : pd.DataFrame
        The input DataFrame to be extended.

    column : str
        The name of the column to be split.

    axis : int, optional
        The axis along which to expand the DataFrame.
        - 0 (default): Expand the specified column into multiple rows.
        - 1: Expand the specified column into multiple columns.

    sep : str, optional
        The separator used to split the values in the specified column.
        Must be provided for the function to work correctly.
    """

    data = data.copy()
    mask = data[column].str.contains(sep, na=False)
    data = data.copy()
    if mask.any():
        data[column] = data[column].apply(
            lambda x: x.split(sep) if isinstance(x, str) else x
        )  # Only split if x is a string

        # Strip spaces from each item in the lists
        data[column] = data[column].apply(
            lambda x: [item.strip() for item in x] if isinstance(x, list) else x
        )

    data = data.explode(column, ignore_index=True)
    return data


def df_cycle(data: pd.DataFrame, columns=None, max_val=None, inplace=False):
    """
    Purpose: transforms a datetime feature (like month or day) into a cyclic encoding for use in machine learning models, particularly neural networks.
    Usage:
        data = pd.DataFrame({'month': [1, 4, 7, 10, 12]})  # Just months as an example
        # df_cycle month cyclically
        data = df_cycle(data, 'month', 12)
    """
    if columns is None:
        columns = list(
            data.select_dtypes(include=np.number).columns
        )  # If no columns specified, use all columns
    if max_val is None:
        max_val = np.max(
            data[columns]
        )  # If no max_val specified, use the maximum value across all columns
    if isinstance(columns, str):
        columns = [
            columns
        ]  # If a single column name is provided as a string, convert it to a list

    # Check if inplace is True, so we modify the original dataframe
    if inplace:
        # Modify the data in place, no return statement needed
        for col in columns:
            data[col + "_sin"] = np.sin(2 * np.pi * data[col] / max_val)
            data[col + "_cos"] = np.cos(2 * np.pi * data[col] / max_val)
    else:
        # If inplace is False, return the modified dataframe
        new_data = data.copy()
        for col in columns:
            new_data[col + "_sin"] = np.sin(2 * np.pi * new_data[col] / max_val)
            new_data[col + "_cos"] = np.cos(2 * np.pi * new_data[col] / max_val)
        return new_data


# ! DataFrame
def df_astype(
    data: pd.DataFrame,
    columns: Optional[Union[str, List[str]]] = None,
    astype: str = None,  # "datetime",
    skip_row: Union[str, list] = None,
    original_fmt:str=None,
    fmt: Optional[str] = None,
    inplace: bool = False,
    errors: str = "coerce",  # Can be "ignore", "raise", or "coerce"
    verbose:bool=False,
    **kwargs,
) -> Optional[pd.DataFrame]:
    """
    Convert specified columns of a DataFrame to a specified type (e.g., datetime, float, int, numeric, timedelta).
    If columns is None, all columns in the DataFrame will be converted.

    Parameters:
    - df: DataFrame containing the columns to convert.
    - columns: Either a single column name, a list of column names, or None to convert all columns.
    - astype: The target type to convert the columns to ('datetime', 'float', 'int', 'numeric', 'timedelta', etc.).
    - fmt: Optional; format to specify the datetime format (only relevant for 'datetime' conversion).
    - inplace: Whether to modify the DataFrame in place or return a new one. Defaults to False.
    - errors: Can be "ignore", "raise", or "coerce"
    - **kwargs: Additional keyword arguments to pass to the conversion function (e.g., errors='ignore' for pd.to_datetime or pd.to_numeric).

    Returns:
    - If inplace=False: DataFrame with the specified columns (or all columns if columns=None) converted to the specified type.
    - If inplace=True: The original DataFrame is modified in place, and nothing is returned.
    """
    astypes = [
        "datetime",
        "timedelta",
        "numeric",
        "int",
        "int8",
        "int16",
        "int32",
        "int64",
        "uint8",
        "uint16",
        "uint32",
        "uint64",
        "float",
        "float16",
        "float32",
        "float64",
        "complex",
        "complex64",
        "complex128",
        "str",
        "string",
        "bool",
        "datetime64",
        "datetime64[ns]",
        "timedelta64",
        "timedelta64[ns]",
        "category",
        "object",
        "Sparse",
        "hour",
        "minute",
        "second",
        "time",
        "week",
        "date",
        "day",
        "month",
        "year",
        "circular",
    ]
    # If inplace is False, make a copy of the DataFrame
    if not inplace:
        data = data.copy()
    if skip_row is not None:
        data = data.drop(index=skip_row, errors="ignore")
    # If columns is None, apply to all columns
    if columns is None:
        columns = data.columns.tolist()
    # correct the astype input
    if isinstance(astype, str):
        astype = strcmp(astype, astypes)[0]
        if verbose:
            print(f"converting as type: {astype}") 
    elif isinstance(astype, dict):
        for col, dtype in astype.items():
            dtype = "date" if dtype == "day" else dtype
            target_dtype = strcmp(dtype, astypes)[0]
            try:
                if target_dtype == "datetime":
                    data[col] = pd.to_datetime(data[col], format=original_fmt, errors=errors)
                elif target_dtype == "timedelta":
                    data[col] = pd.to_timedelta(data[col], errors=errors)
                else:
                    data[col] = data[col].astype(target_dtype)
            except Exception as e:
                if verbose:
                    print(f"Error converting column '{col}' to {target_dtype}: {e}")
        return data if not inplace else None
    # Ensure columns is a list
    if isinstance(columns, str):
        columns = [columns]

    # Convert specified columns
    for column in columns:
        try:
            if astype in [
                "datetime",
                "hour",
                "minute",
                "second",
                "time",
                "week",
                "date",
                "month",
                "year",
            ]:
                kwargs.pop("errors", None)
                # convert it as type: datetime
                if isinstance(column, int):
                    data.iloc[:, column] = pd.to_datetime(data.iloc[:, column], format=original_fmt, errors=errors, **kwargs) if original_fmt is not None else pd.to_datetime(data[column], errors=errors, **kwargs)

                    try:
                        if fmt is not None:
                            # data[column] = data[column].apply(lambda x: f"{x:{fmt}}")
                            data[column] = data[column].apply(
                                lambda x: x.strftime(fmt) if pd.notnull(x) else None
                            )
                    except Exception as e:
                        print(f"设置格式的时候有误: {e}")

                    # further convert:
                    if astype == "time":
                        data.iloc[:, column] = data.iloc[:, column].dt.time
                    elif astype == "month":
                        data.iloc[:, column] = data.iloc[:, column].dt.month
                    elif astype == "year":
                        data.iloc[:, column] = data.iloc[:, column].dt.year
                    elif astype == "date" or astype == "day":
                        data.iloc[:, column] = data.iloc[:, column].dt.date
                    elif astype == "hour":
                        data.iloc[:, column] = data.iloc[:, column].dt.hour
                    elif astype == "minute":
                        data.iloc[:, column] = data.iloc[:, column].dt.minute
                    elif astype == "second":
                        data.iloc[:, column] = data.iloc[:, column].dt.second
                    elif astype == "week":
                        data.iloc[:, column] = data.iloc[:, column].dt.day_name()
                else:
                    data[column] = (
                        pd.to_datetime(
                            data[column], format=original_fmt, errors=errors, **kwargs
                        )
                        if original_fmt is not None
                        else pd.to_datetime(data[column], errors=errors, **kwargs)
                    )

                    try:
                        if fmt is not None:
                            # data[column] = data[column].apply(lambda x: f"{x:{fmt}}")
                            data[column] = data[column].apply(
                                lambda x: x.strftime(fmt) if pd.notnull(x) else None
                            )
                    except Exception as e:
                        print(f"设置格式的时候有误: {e}")
                    # further convert:
                    if astype == "time":
                        data[column] = data[column].dt.time
                    elif astype == "month":
                        data[column] = data[column].dt.month
                    elif astype == "year":
                        data[column] = data[column].dt.year
                    elif astype == "date":
                        data[column] = data[column].dt.date
                    elif astype == "hour":
                        data[column] = data[column].dt.hour
                    elif astype == "minute":
                        data[column] = data[column].dt.minute
                    elif astype == "second":
                        data[column] = data[column].dt.second
                    elif astype == "week":
                        data[column] = data[column].dt.day_name()

            elif astype == "numeric":
                kwargs.pop("errors", None)
                data[column] = pd.to_numeric(data[column], errors=errors, **kwargs)
                # print(f"Successfully converted '{column}' to numeric.")
            elif astype == "timedelta":
                kwargs.pop("errors", None)
                data[column] = pd.to_timedelta(data[column], errors=errors, **kwargs)
                # print(f"Successfully converted '{column}' to timedelta.")
            elif astype == "circular":
                max_val = kwargs.get("max_val", None)
                data[column] = df_cycle(data=data, columns=column, max_val=max_val)
            else:
                # Convert to other types (e.g., float, int)
                if astype == "int":
                    data[column] = data[column].astype("float").astype("int")
                else:
                    data[column] = data[column].astype(astype)
                # print(f"Successfully converted '{column}' to {astype}.")

        except Exception as e:
            if verbose:
                print(f"Error converting '{column}' to {astype}: {e}")
    try:
        if verbose:
            display(data.info()[:10])
    except:
        pass
    return data


def calculate_age(birthdate,fmt="%d.%m.%Y", exact=True, default=None, return_years_months=False):
    """
    计算年龄（精确到年或年月）

    参数:
    -----------
    birthdate : datetime-like, str
        出生日期，可以是datetime对象、pandas.Timestamp或可解析的日期字符串
    exact : bool, 可选
        是否计算精确年龄（考虑月份和日），False则只计算年份差 (默认: True)
    default : any, 可选
        当输入无效时返回的默认值 (默认: None)
    return_years_months : bool, 可选
        是否返回包含年和月的字典 (默认: False)

    返回:
    --------
    int or dict or default
        返回年龄（整数），或{'years': int, 'months': int}字典，或默认值
    例子: 
    data = {
        "birthdate": [
            "04/20/19",
            "04/20/33",
            "04/20/29",
            "04/20/30",
            "04/20/68",
            "04/20/69",
            "04/20/99",
        ],
    }
    df = pd.DataFrame(data)
    # display(df)
    # 转换日期（使用默认threshold=30）
    df["converted"] = df_date(
        df[["birthdate"]], fmt="%m/%d/%y", astype="date", verbose=True, century_threshold=25
    )["birthdate"]
    # 显示结果
    # display(df)
    df["age"] = df["converted"].apply(lambda x: calculate_age(x))
    df
    """
    from datetime import datetime
    import pandas as pd

    # 处理空值
    if pd.isnull(birthdate):
        return default

    # 统一转换为datetime对象
    try:
        if not isinstance(birthdate, (datetime, pd.Timestamp)):
            birthdate = pd.to_datetime(birthdate,format=fmt, errors="raise")
    except Exception as e:
        print(f"日期解析错误: {e}")
        return default

    today = datetime.today()

    # 计算基本年龄
    try:
        if not exact:
            # 简单年份差计算
            age = today.year - birthdate.year
        else:
            # 精确年龄计算（考虑月份和日）
            age = (
                today.year
                - birthdate.year
                - ((today.month, today.day) < (birthdate.month, birthdate.day))
            )

        if return_years_months:
            # 计算完整的年和月
            months = (today.year - birthdate.year) * 12 + today.month - birthdate.month
            if today.day < birthdate.day:
                months -= 1
            years = months // 12
            months = months % 12
            return {"years": years, "months": months}
        else:
            return age

    except Exception as e:
        print(f"年龄计算错误: {e}")
        return default


def df_date(df, keywords=["date", "datum", "时间", "日期"], fmt="%d.%m.%y",original_fmt="%Y-%m-%d",
                 astype="str", century_threshold=None, inplace=False, verbose=False,**kwargs):
    """
    自动识别并转换数据框中与日期相关的列格式
    
    参数:
        df (pd.DataFrame): 输入数据框
        keywords (list): 用于识别日期列的关键词列表
        fmt (str): 日期格式字符串
        astype (str): 目标类型 ("str"或"datetime")
        century_threshold (int): 两位数年份的世纪阈值, default: current year+5, e.g., in year 2025, then the thr=30
        inplace (bool): 是否原地修改数据框
        verbose (bool): 是否显示详细转换信息
        
    返回:
        pd.DataFrame: 转换后的数据框
 
    # 测试数据
    data = {
        "birthdate": ["04/20/29", "04/20/30", "04/20/68", "04/20/69", "04/20/99"],
        "description": [
            "低于阈值(29<30)→2029年",
            "等于阈值(30=30)→1930年",
            "68→2068年(POSIX标准)",
            "69→1969年(POSIX标准)",
            "99→1999年",
        ],
    }
    df = pd.DataFrame(data)
    display(df)

    df_date(
        df,
        fmt="%m/%d/%y",
        keywords=["date", "conv"],
        astype="date",
        verbose=True,
        inplace=True,
        century_threshold=30,
    )
    # 显示结果
    display(df)
    df["age"] = df["birthdate"].apply(lambda x: calculate_age(x))
    df
    """
    import pandas as pd
    from datetime import datetime

    def parse_custom_date(date_str, fmt, century_threshold=None, verbose=False):
        """
        增强版日期解析函数，正确处理两位数年份的世纪问题
        
        参数:
            date_str (str): 要解析的日期字符串
            fmt (str): 日期格式字符串 (例如 "%d/%m/%y")
            century_threshold (int): 两位数年份的世纪阈值 (0-99)
                                    - 小于此值视为20XX年
                                    - 大于等于此值视为19XX年
            verbose (bool): 是否显示调试信息
            
        返回:
            datetime or pd.NaT: 解析后的datetime对象或NaT(解析失败时)
        """
        # 检查空值
        if pd.isnull(date_str):
            if verbose:
                print("输入为空值")
            return pd.NaT
        if century_threshold is None:
            century_threshold=int(str(datetime.today().year)[-2:])+5
        # 清理字符串
        date_str = str(date_str).strip()
        if verbose:
            print(f"\nConverting '{date_str}' (format: '{fmt}')")
        
        try:
            # 首先尝试直接解析
            dt = datetime.strptime(date_str, fmt)
            print(dt)
            
            # 处理两位数年份的情况
            if '%y' in fmt.lower() and not '%Y' in fmt.lower():
                current_year_last2 = datetime.now().year % 100
                year_full = dt.year
                
                # 计算正确的世纪
                if (dt.year % 100) < century_threshold:
                    year_full = 2000 + (dt.year % 100)
                else:
                    year_full = 1900 + (dt.year % 100)
                if verbose:
                    print(f"Handle 2-digit years: {dt.year} -> {year_full} (century_threshold={century_threshold})")
                
                # 确保年份在合理范围内
                if 1900 <= year_full <= 2100:
                    dt = dt.replace(year=year_full)
                else:
                    if verbose:
                        print(f"年份超出合理范围: {year_full}")
                    return pd.NaT
            
            if verbose:
                print(f"Done! {dt}\n")
            return dt
            
        except ValueError as e:
            if verbose:
                print(f"解析失败: {e}")
            return pd.NaT
    if century_threshold is None:
        century_threshold=int(str(datetime.today().year)[-2:])+5
    if not inplace:
        df = df.copy()
    
    # 找出名称中包含关键词的列
    date_cols = [
        col for col in df.columns
        if any(kw.lower() in col.lower() for kw in keywords)
    ]
    
    astype=strcmp(astype,["datetime","str"])[0]
    print(f"\n识别到日期相关列: {date_cols}\n")
    for col in date_cols:
        try:
            if astype == "datetime":
                # 转换为datetime对象
                if verbose:
                    print(f"将列 '{col}' 转为datetime (格式: '{fmt}')")
                
                # 应用自定义解析函数
                df[col] = df[col].apply(
                    lambda x: parse_custom_date(x, fmt, century_threshold, verbose)
                )
                
            elif astype == "str":
                # 转换为格式化字符串
                if pd.api.types.is_datetime64_any_dtype(df[col]):
                    if verbose:
                        print(f"将列 '{col}' 转为字符串 (格式: '{fmt}')")
                    df[col] = df[col].dt.strftime(fmt)
                else:
                    # 非日期类型尝试先转日期再转字符串
                    try:
                        temp = pd.to_datetime(df[col], errors='raise')
                        df[col] = temp.dt.strftime(fmt)
                        if verbose:
                            print(f"将列 '{col}' 通过中间转换转为字符串")
                    except:
                        raise ValueError(f"列 '{col}' 无法转换为日期格式")
            else:
                raise ValueError(f"无效的astype参数: {astype} (只接受'str'或'datetime')")
                
        except Exception as e:
            if verbose:
                print(f"列 '{col}' 转换失败: {str(e)}")
    try: 
        df_ = df_astype(
            df, columns=date_cols, astype="date", fmt=fmt, inplace=inplace, verbose=verbose, **kwargs
        )
        return df_
    except Exception as e:
        print(f"got final errors: {e}, 但是尝试其它办法")
    
    return df

# pd extensions 
from typing import Union, List, Optional, Literal
@pd.api.extensions.register_dataframe_accessor("column")
class Column:
    def __init__(self, pandas_obj):
        self._df = pandas_obj
    
    def __call__(
        self, 
        keywords: Optional[Union[str, List[str]]] = None, 
        output: Literal['columns', 'list', 'dataframe', 'df'] = 'columns',
        ignore_case: bool = True,
    ) -> Union[List[str], pd.DataFrame]:
        """
        Get columns matching specified keyword(s)
        
        Parameters:
        -----------
        keywords : str or list of str, optional
            Keyword(s) to search in column names. If None, returns all columns.
        ignore_case : bool, default True
            Whether to ignore case when matching keywords
        output : str, default 'columns'
            Output format:
            - 'columns' or 'list': returns list of column names
            - 'dataframe' or 'df': returns DataFrame with selected columns
            
        Returns:
        --------
        Union[List[str], pd.DataFrame]
            Either list of column names or filtered DataFrame

        df = pd.DataFrame({
            'OrderDate': ['2023-01-01', '2023-01-02'],
            'delivery_time': ['10:00', '12:00'],
            'CustomerID': [1, 2],
            'DATE_RECEIVED': ['2023-01-03', '2023-01-04']
        })

        # Get all columns
        df.column()  
        # Returns: ['OrderDate', 'delivery_time', 'CustomerID', 'DATE_RECEIVED']

        # Get date-related columns (case insensitive)
        df.column("date")  
        # Returns: ['OrderDate', 'DATE_RECEIVED']

        # Get multiple keyword matches
        df.column(["date", "time"])  
        # Returns: ['OrderDate', 'delivery_time', 'DATE_RECEIVED']

        # Exact case matching
        df.column("DATE", ignore_case=False)  
        # Returns: ['DATE_RECEIVED']
        """
        # Get matching columns
        if keywords is None:
            cols = self._df.columns.tolist()
        else:
            if isinstance(keywords, str):
                keywords = [keywords]
            
            if ignore_case:
                cols = [
                    col for col in self._df.columns
                    if any(kw.lower() in col.lower() for kw in keywords)
                ]
            else:
                cols = [
                    col for col in self._df.columns
                    if any(kw in col for kw in keywords)
                ]
        output=strcmp(output.lower(),['columns', 'list', 'dataframe', 'df','all'])[0]
        # Handle output format
        if output in ('columns', 'list'):
            return cols
        elif output in ('dataframe', 'df','all'):
            return self._df[cols]
        else:
            raise ValueError(
                f"Invalid output format '{output}'. "
                "Use 'columns'/'list' or 'dataframe'/'df'"
            )
@pd.api.extensions.register_dataframe_accessor("select")
class Select:
    def __init__(self, pandas_obj):
        self._df = pandas_obj
    
    def __call__(
        self, 
        keywords: Optional[Union[str, List[str]]] = None, 
        output: Literal['columns', 'list', 'dataframe', 'df'] = 'df',
        ignore_case: bool = True,
    ) -> Union[List[str], pd.DataFrame]:
        """
        Get columns matching specified keyword(s)
        
        Parameters:
        -----------
        keywords : str or list of str, optional
            Keyword(s) to search in column names. If None, returns all columns.
        ignore_case : bool, default True
            Whether to ignore case when matching keywords
        output : str, default 'columns'
            Output format:
            - 'columns' or 'list': returns list of column names
            - 'dataframe' or 'df': returns DataFrame with selected columns
            
        Returns:
        --------
        Union[List[str], pd.DataFrame]
            Either list of column names or filtered DataFrame

        df = pd.DataFrame({
            'OrderDate': ['2023-01-01', '2023-01-02'],
            'delivery_time': ['10:00', '12:00'],
            'CustomerID': [1, 2],
            'DATE_RECEIVED': ['2023-01-03', '2023-01-04']
        })

        # Get all columns
        df.select()  
        # Returns: ['OrderDate', 'delivery_time', 'CustomerID', 'DATE_RECEIVED']

        # Get date-related columns (case insensitive)
        df.select("date")  
        # Returns: ['OrderDate', 'DATE_RECEIVED']

        # Get multiple keyword matches
        df.select(["date", "time"])  
        # Returns: ['OrderDate', 'delivery_time', 'DATE_RECEIVED']

        # Exact case matching
        df.select("DATE", ignore_case=False)  
        # Returns: ['DATE_RECEIVED']
        """
        # Get matching columns
        if keywords is None:
            cols = self._df.columns.tolist()
        else:
            if isinstance(keywords, str):
                keywords = [keywords]
            
            if ignore_case:
                cols = [
                    col for col in self._df.columns
                    if any(kw.lower() in col.lower() for kw in keywords)
                ]
            else:
                cols = [
                    col for col in self._df.columns
                    if any(kw in col for kw in keywords)
                ]
        output=strcmp(output.lower(),['columns', 'list', 'dataframe', 'df','all'])[0]
        # Handle output format
        if output in ('columns', 'list'):
            return cols
        elif output in ('dataframe', 'df','all'):
            return self._df[cols]
        else:
            raise ValueError(
                f"Invalid output format '{output}'. "
                "Use 'columns'/'list' or 'dataframe'/'df'"
            ) 

# =================df.apply_style(rules)========================= 
@pd.api.extensions.register_dataframe_accessor("apply_style")
class ApplyStyleAccessor:
    def __init__(self, pandas_obj: pd.DataFrame):
        self._obj = pandas_obj

    def __call__(
        self,
        style_rules: List[Dict],
        hover: Optional[Dict] = None,
        caption: Optional[str] = None,
        table_styles: Optional[List[Dict]] = None,
        default_alignment: Optional[Dict] = None,
        verbose:bool = False
    ):
        """


        Parameters:
        - style_rules: List of conditional formatting rules
        - hover: Dict of hover effects (e.g., {'color': 'red', 'background-color': 'yellow'})
        - caption: Table caption text
        - table_styles: List of table-wide styles

        Style Rule Properties:
        - Text: color, font-size, font-weight, font-style, font-family, text-decoration
        - Background: background-color, background-gradient, opacity
        - Borders: border, border-top, border-bottom, border-left, border-right
        - Alignment: text-align, vertical-align
        - Padding: padding, padding-top, padding-bottom, padding-left, padding-right
        - Transformation: text-transform, letter-spacing
        - Special: cursor, display, white-space

        """
        return _apply_style_core(
            self._obj,
            style_rules,
            hover=hover,
            caption=caption,
            table_styles=table_styles,
            default_alignment=default_alignment,
            verbose=verbose
        )


def _apply_style_core(
    df: pd.DataFrame,
    style_rules: List[Dict],
    hover: Optional[Dict] = None,
    caption: Optional[str] = None,
    table_styles: Optional[List[Dict]] = None,
    default_alignment: Optional[Dict] = None,
    verbose: bool = False
):
    if verbose:
        print("""
data = {
    "Name": ["Alice", "Bob", "Charlie", "David"],
    "Age": [25, 30, 35, 40],
    "Score": [92, 88, 76, 95],
    "Status": ["Active", "Inactive", "Pending", "Active"],
    "Join Date": pd.to_datetime(
        ["2020-01-01", "2019-05-15", "2021-03-10", "2018-11-20"]
    ),
}
df = pd.DataFrame(data)

# Comprehensive styling example
styled_df = df.apply_style(
    [
        {
            "column": "Score",
            "operator": ">=",
            "value": 90,
            "color": "red",
            "background-color": "green",
            "font-weight": "bold",
            "border": "2px solid darkgreen",
            "text-align": "center",
            "padding": "5px 10px",
        },
        {
            "column": "Status",
            "operator": "==",
            "value": "Active",
            "background-color": "lightblue",
            "text-decoration": "underline",
            "font-style": "italic",
        },
        {
            "column": "Age",
            "operator": ">",
            "value": 30,
            "color": "darkred",
            "font-size": "14px",
            "border-left": "3px solid orange",
        },
        {
            "index": 0,
            "operator": "==",
            "value": "Alice",
            "background-color": "lightyellow",
            "font-family": "Arial",
        },
    ],
    default_alignment={
        "header": {
            "text-align": "left",
            "vertical-align": "middle",
            "font-weight": "bold",
        },
        "cells": {"text-align": "center", "vertical-align": "top"},
    },
    # hover={
    #     "background-color": "#F3F3BA",  # Light yellow on hover
    #     "transition": "all 3s ease",  # Smooth transition
    # },
    caption="Employee Performance Data",
    table_styles=[
        {
            "selector": "th",
            "props": [
                ("background-color", "#333"),
                ("color", "white"),
                ("font-weight", "bold"),
                ("text-align", "center"),
            ],
        },
        {
            "selector": "tr:nth-child(even)",
            "props": [("background-color", "#f9f9f9")],
        },
    ],
)

display(styled_df)
delete("styled9.xlsx")
styled_df.to_excel("styled9.xlsx", sheet_name="style9", index=False)


# example:
df = pd.DataFrame(
    {
        "Event": ["Conference", "Workshop", "Deadline"],
        "Date": [
            "2023-06-15",
            "2022-07-01",
            "2026-05-31",
        ],  # pd.to_datetime(["2023-06-15", "2023-07-01", "2023-05-31"]),
    }
)

styled_df = df.apply_style(
    [
        # {
        #     "column": "Date",
        #     "operator": "before",
        #     "value": "2023-06-01",
        #     "background-color": "#ffcccc",
        #     "text": "⚠️ Past deadline",
        # },
        # {
        #     "column": "Date",
        #     "operator": "after",
        #     "value": pd.Timestamp.now(),
        #     "background-color": "#e6ffe6",
        #     "text": "↑ Upcoming",
        # },
        {
            "column": "Date",
            "operator": "between",
            "value": ["2022-06-01", "2025-06-30"],
            "border": "2px solid darkred",
        },
    ]
)
styled_df
              """)
    def _preprocess_cfg(style_rules: List[Dict]) -> List[Dict]:
        """Normalize and validate style rules with comprehensive CSS support"""
        if not isinstance(style_rules, list):
            style_rules = [style_rules]

        corrected_rules = []
        for rule in style_rules:
            normalized_rule = {}
            for key, value in rule.items():
                # Convert common aliases to standard CSS properties 
                prop_aliases = {
                    # Background properties
                    "bg": "background",
                    "bg_color": "background-color",
                    "background_color": "background-color",
                    "bg_image": "background-image",
                    "bg_repeat": "background-repeat",
                    "bg_position": "background-position",
                    "bg_size": "background-size",
                    "bg_attachment": "background-attachment",
                    "bg_gradient": "background-gradient",
                    "bg_blend_mode": "background-blend-mode",
                    
                    # Text properties
                    "color":"color",
                    "text_color": "color",
                    "text_align": "text-align",
                    "text_decoration": "text-decoration",
                    "text_transform": "text-transform",
                    "text_indent": "text-indent",
                    "text_shadow": "text-shadow",
                    "text_overflow": "text-overflow",
                    "letter_spacing": "letter-spacing",
                    "word_spacing": "word-spacing",
                    "line_height": "line-height",
                    "white_space": "white-space",
                    "direction": "direction",
                    
                    # Font properties
                    "font": "font",
                    "font_size": "font-size",
                    "font_weight": "font-weight",
                    "font_style": "font-style",
                    "font_family": "font-family",
                    "font_variant": "font-variant",
                    "font_stretch": "font-stretch",
                    
                    # Border properties (full set)
                    "border": "border",
                    "border_top": "border-top",
                    "border_right": "border-right",
                    "border_bottom": "border-bottom",
                    "border_left": "border-left",
                    "border_color": "border-color",
                    "border_width": "border-width",
                    "border_style": "border-style",
                    "border_radius": "border-radius",
                    "border_top_color": "border-top-color",
                    "border_right_color": "border-right-color",
                    "border_bottom_color": "border-bottom-color",
                    "border_left_color": "border-left-color",
                    "border_top_width": "border-top-width",
                    "border_right_width": "border-right-width",
                    "border_bottom_width": "border-bottom-width",
                    "border_left_width": "border-left-width",
                    "border_top_style": "border-top-style",
                    "border_right_style": "border-right-style",
                    "border_bottom_style": "border-bottom-style",
                    "border_left_style": "border-left-style",
                    "border_top_left_radius": "border-top-left-radius",
                    "border_top_right_radius": "border-top-right-radius",
                    "border_bottom_right_radius": "border-bottom-right-radius",
                    "border_bottom_left_radius": "border-bottom-left-radius",
                    
                    # Padding properties
                    "padding": "padding",
                    "padding_top": "padding-top",
                    "padding_right": "padding-right",
                    "padding_bottom": "padding-bottom",
                    "padding_left": "padding-left",
                    "padding_x": "padding-left",  # Horizontal (both left/right)
                    "padding_y": "padding-top",   # Vertical (both top/bottom)
                    
                    # Margin properties
                    "margin": "margin",
                    "margin_top": "margin-top",
                    "margin_right": "margin-right",
                    "margin_bottom": "margin-bottom",
                    "margin_left": "margin-left",
                    "margin_x": "margin-left",    # Horizontal (both left/right)
                    "margin_y": "margin-top",     # Vertical (both top/bottom)
                    
                    # Display properties
                    "display": "display",
                    "visibility": "visibility",
                    "opacity": "opacity",
                    "overflow": "overflow",
                    "overflow_x": "overflow-x",
                    "overflow_y": "overflow-y",
                    "z_index": "z-index",
                    
                    # Position properties
                    "position": "position",
                    "top": "top",
                    "right": "right",
                    "bottom": "bottom",
                    "left": "left",
                    
                    # Flexbox properties
                    "flex": "flex",
                    "flex_direction": "flex-direction",
                    "flex_wrap": "flex-wrap",
                    "flex_flow": "flex-flow",
                    "flex_grow": "flex-grow",
                    "flex_shrink": "flex-shrink",
                    "flex_basis": "flex-basis",
                    "justify_content": "justify-content",
                    "align_items": "align-items",
                    "align_self": "align-self",
                    "align_content": "align-content",
                    
                    # Grid properties
                    "grid": "grid",
                    "grid_template": "grid-template",
                    "grid_template_columns": "grid-template-columns",
                    "grid_template_rows": "grid-template-rows",
                    "grid_template_areas": "grid-template-areas",
                    "grid_column": "grid-column",
                    "grid_row": "grid-row",
                    "grid_area": "grid-area",
                    "grid_gap": "grid-gap",
                    "grid_column_gap": "grid-column-gap",
                    "grid_row_gap": "grid-row-gap",
                    
                    # Animation/transition properties
                    "transition": "transition",
                    "transition_property": "transition-property",
                    "transition_duration": "transition-duration",
                    "transition_timing": "transition-timing-function",
                    "transition_delay": "transition-delay",
                    "animation": "animation",
                    
                    # Special parameters (non-CSS)
                    "apply_to": "apply_to",
                    "operator": "operator",
                    "value": "value"
                }
                if key.lower() in ["column","index","columns","indexs","operator","op","value","apply_to"]:
                    normalized_key = key.lower().replace("-", "_")
                else:
                    normalized_key = strcmp(key.lower().replace("-", "_"),list(prop_aliases.keys()))[0]
                # print(f"normalized_key={normalized_key}")
                normalized_key = prop_aliases.get(normalized_key, normalized_key)
                normalized_rule[normalized_key] = value
            # print(normalized_rule)
            # Validate required fields
            if not all(k in normalized_rule for k in ("operator", "value")):
                raise ValueError("Each rule must contain 'operator' and 'value'")
            if not any(k in normalized_rule for k in ("column", "index")):
                raise ValueError("Each rule must contain either 'column' or 'index'")
            if "apply_to" not in normalized_rule:
                normalized_rule["apply_to"] = "cell"
            corrected_rules.append(normalized_rule)

        return corrected_rules

    def _apply_column_styles(x: pd.Series) -> List[str]:
        """Apply styles to each cell in a column or entire column/row"""
        styles = [""] * len(x)
        column_name = x.name
        
        # Check for column-level rules first
        for rule in column_rules:
            if column_name != rule["column"]:
                continue
                
            apply_to = rule.get("apply_to", "cell")
            
            # For column-level application
            if apply_to == "column":
                # Check if any cell in column meets condition
                if any(_evaluate_condition(val, rule["operator"], rule["value"]) for val in x):
                    return [_build_style_string(rule)] * len(x)
                continue
                
            # For row-level application (applied in _apply_index_styles)
            if apply_to == "row":
                continue
                
            # Default cell-level application
            for i, value in enumerate(x):
                try:
                    if _evaluate_condition(value, rule["operator"], rule["value"]):
                        styles[i] = _build_style_string(rule)
                except Exception as e:
                    print(f"Error applying style for {column_name}: {e}")
        
        return styles

    def _apply_index_styles(x: pd.Series) -> List[str]:
        """Apply styles to each cell in an index or entire row"""
        styles = [""] * len(x)
        index_name = x.name
        
        # Check for index-level rules first
        for rule in index_rules:
            if index_name != rule["index"]:
                continue
                
            apply_to = rule.get("apply_to", "cell")
            
            # For row-level application
            if apply_to == "row":
                # Check if any cell in row meets condition
                if any(_evaluate_condition(val, rule["operator"], rule["value"]) for val in x):
                    return [_build_style_string(rule)] * len(x)
                continue
                
            # Default cell-level application
            for i, value in enumerate(x):
                try:
                    if _evaluate_condition(value, rule["operator"], rule["value"]):
                        styles[i] = _build_style_string(rule)
                except Exception as e:
                    print(f"Error applying style for {index_name}: {e}")
        
        return styles
    def _evaluate_condition(value, op: str, val) -> bool:
        """
        Ultimate condition evaluator with comprehensive list support for datetime comparisons
        """
        import pandas as pd
        import datetime
        op = op.lower()
        original_value = value
        
        # Helper function to safely convert to datetime
        def try_convert_to_datetime(obj, verbose=False):
            if verbose:
                print(f"Input type: {type(obj)}")
            
            # Handle list/tuple inputs
            if isinstance(obj, (list, tuple)):
                try:
                    return [pd.to_datetime(item) for item in obj]
                except:
                    if verbose:
                        print("List conversion failed")
                    return None
            
            # Handle single value inputs
            if isinstance(obj, (pd.Timestamp, datetime.datetime, datetime.date)):
                if verbose:
                    print("Already datetime type")
                return obj
            if isinstance(obj, str):
                try:
                    return pd.to_datetime(obj)
                except:
                    if verbose:
                        print("String conversion failed")
                    return None
            return None
        
        # Convert both values to datetime if possible
        value_dt = try_convert_to_datetime(value)
        val_dt = try_convert_to_datetime(val)
        
        # Special handling for 'between' operator with lists
        if op == "between":
            if not isinstance(val, (list, tuple)) or len(val) != 2:
                return False
                
            # Try numeric comparison first
            try:
                num_value = float(original_value)
                start = float(val[0])
                end = float(val[1])
                return start <= num_value <= end
            except (ValueError, TypeError):
                pass
                
            # Try datetime comparison
            value_dt = try_convert_to_datetime(original_value)
            range_dt = try_convert_to_datetime(val)
            if value_dt and all(x is not None for x in range_dt):
                return range_dt[0] <= value_dt <= range_dt[1]
                
            # Fall back to string comparison
            try:
                return str(val[0]) <= str(original_value) <= str(val[1])
            except:
                return False 
        
        # Handle datetime comparisons if both values are datetime-like
        if value_dt is not None and val_dt is not None:
            try:
                if op in ("before", "<", "earlier"):
                    return value_dt < val_dt
                elif op in ("after", ">", "later"):
                    return value_dt > val_dt
                elif op in ("on_or_before", "<=", "before_or_on"):
                    return value_dt <= val_dt
                elif op in ("on_or_after", ">=", "after_or_on"):
                    return value_dt >= val_dt
                elif op == "same_day":
                    return value_dt.date() == val_dt.date()
            except Exception as e:
                print(f"Datetime comparison error: {e}")
        
        # Fall back to regular comparison
        try:
            if op in ("==", "=", "is"):
                return original_value == val
            elif op in ("!=", "~="):
                return original_value != val
            elif op == ">":
                return original_value > val
            elif op == "<":
                return original_value < val
            elif op == ">=":
                return original_value >= val
            elif op == "<=":
                return original_value <= val
            elif op in ("in", "include", "including"):
                if isinstance(val, (list, tuple, set)):
                    return original_value in val
                return isinstance(original_value, str) and (str(val) in original_value)
            elif op == "not in":
                if isinstance(val, (list, tuple, set)):
                    return original_value not in val
                return isinstance(original_value, str) and (str(val) not in original_value)
            elif op == "contains":
                return isinstance(original_value, str) and str(val) in original_value
            elif op == "not contains":
                return isinstance(original_value, str) and str(val) not in original_value
            elif op == "startswith":
                return isinstance(original_value, str) and original_value.startswith(str(val))
            elif op == "endswith":
                return isinstance(original_value, str) and original_value.endswith(str(val))
            elif op == "isna":
                return pd.isna(original_value)
            elif op == "notna":
                return not pd.isna(original_value)
        except Exception as e:
            print(f"Comparison error: {e}")
        
        return False
    def _build_style_string(rule: Dict) -> str:
        """Build CSS string from rule dictionary"""
        style_parts = []
        for prop, val in rule.items():
            if prop in ["column", "index", "operator", "value","apply_to"]:
                continue
            style_parts.append(f"{prop}: {val}")
        return "; ".join(style_parts)
    def _build_style_props(rule: Dict) -> Dict:
        """Convert style rule to properties dictionary"""
        return {k: v for k, v in rule.items() 
                if k not in ["column", "index", "operator", "value", "apply_to"]}
    def _apply_default_alignment(styler, alignment: Dict):
        """Apply default alignment to headers and cells"""
        if not alignment:
            return styler

        # Process header alignment
        header_align = alignment.get("header", {})
        if header_align:
            if "align" in header_align:
                # Handle combined alignment shortcut
                align = header_align.pop("align")
                header_align.setdefault("text-align", align)
                header_align.setdefault("vertical-align", align)

            header_styles = [
                {
                    "selector": "th",
                    "props": [
                        (k.replace("_", "-"), v) for k, v in header_align.items()
                    ],
                }
            ]
            styler = styler.set_table_styles(header_styles, overwrite=False)

        # Process cell alignment
        cell_align = alignment.get("cells", {})
        if cell_align:
            if "align" in cell_align:
                # Handle combined alignment shortcut
                align = cell_align.pop("align")
                cell_align.setdefault("text-align", align)
                cell_align.setdefault("vertical-align", align)

            cell_styles = [
                {
                    "selector": "td",
                    "props": [(k.replace("_", "-"), v) for k, v in cell_align.items()],
                }
            ]
            styler = styler.set_table_styles(cell_styles, overwrite=False)

        return styler
    def _apply_styles(df: pd.DataFrame, styler):
        """Core function to apply all styles"""
        processed_rules = _preprocess_cfg(style_rules)
        
        # Apply column-level styles
        for rule in [r for r in processed_rules if r.get("apply_to") == "column"]:
            col = rule["column"]
            mask = df[col].apply(lambda x: _evaluate_condition(x, rule["operator"], rule["value"]))
            if mask.any():
                props = _build_style_props(rule)
                styler = styler.set_properties(subset=pd.IndexSlice[:, col], **props)
        
        # Apply row-level styles
        for rule in [r for r in processed_rules if r.get("apply_to") == "row"]:
            if "column" in rule:
                col = rule["column"]
                mask = df[col].apply(lambda x: _evaluate_condition(x, rule["operator"], rule["value"]))
            else:  # index-based
                idx = rule["index"]
                mask = df.index == idx
                
            for i in df.index[mask]:
                props = _build_style_props(rule)
                styler = styler.set_properties(subset=pd.IndexSlice[i, :], **props)
        
        # Apply cell-level styles
        for rule in [r for r in processed_rules if r.get("apply_to", "cell") == "cell"]:
            if "column" in rule:
                col = rule["column"]
                mask = df[col].apply(lambda x: _evaluate_condition(x, rule["operator"], rule["value"]))
                for i in df.index[mask]:
                    props = _build_style_props(rule)
                    styler = styler.set_properties(subset=pd.IndexSlice[i, col], **props)
            else:  # index-based
                idx = rule["index"]
                props = _build_style_props(rule)
                styler = styler.set_properties(subset=pd.IndexSlice[idx, :], **props)
        
        return styler

    # Create base styler
    styler = df.style
    styler = _apply_styles(df, styler)
    if default_alignment:
        styler = _apply_default_alignment(styler, default_alignment)
 
    # Apply additional table styling
    if hover:
        styler = styler.set_properties(**hover, subset=pd.IndexSlice[:, :])

    if caption:
        styler = styler.set_caption(caption)

    if table_styles:
        styler = styler.set_table_styles(table_styles)

    return styler

# =================df.apply_style(rules)=========================

# ! DataFrame
def df_sort_values(data, column, by=None, ascending=True, inplace=True, **kwargs):
    """
    Sort a DataFrame by a specified column based on a custom order or by count.

    Parameters:
    - data: DataFrame to be sorted.
    - column: The name of the column to sort by.
    - by: List specifying the custom order for sorting or 'count' to sort by frequency.
    - ascending: Boolean or list of booleans, default True.
                 Sort ascending vs. descending.
    - inplace: If True, perform operation in place and return None.
    - **kwargs: Additional arguments to pass to sort_values.

    Returns:
    - Sorted DataFrame if inplace is False, otherwise None.
    """
    if column not in data.columns:
        raise ValueError(f"Column '{column}' does not exist in the DataFrame.")

    if isinstance(by, str) and "count" in by.lower():
        # Count occurrences of each value in the specified column
        value_counts = data[column].value_counts()

        # Determine the order based on counts
        count_ascending = kwargs.pop("count_ascending", ascending)
        sorted_counts = value_counts.sort_values(
            ascending=count_ascending
        ).index.tolist()

        # Convert to a categorical type with the new order
        data[column] = pd.Categorical(
            data[column], categories=sorted_counts, ordered=True
        )
        # Set ascending to count_ascending for sorting
        ascending = count_ascending  # Adjust ascending for the final sort
    elif isinstance(by, list):
        # Convert the specified column to a categorical type with the custom order
        data[column] = pd.Categorical(data[column], categories=by, ordered=True)
    else:
        raise ValueError("Custom order must be a list or 'count'.")

    try:
        if inplace:  # replace the original
            data.sort_values(column, ascending=ascending, inplace=True, **kwargs)
            print(f"Successfully sorted DataFrame by '{column}'")
            return None
        else:
            sorted_df = data.sort_values(column, ascending=ascending, **kwargs)
            print(f"Successfully sorted DataFrame by '{column}' using custom order.")
            return sorted_df
    except Exception as e:
        print(f"Error sorting DataFrame by '{column}': {e}")
        return data


# # Example usage:
# # Sample DataFrame
# data = {
#     "month": ["March", "January", "February", "April", "December"],
#     "Amount": [200, 100, 150, 300, 250],
# }
# df_month = pd.DataFrame(data)

# # Define the month order
# month_order = [
#     "January",
#     "February",
#     "March",
#     "April",
#     "May",
#     "June",
#     "July",
#     "August",
#     "September",
#     "October",
#     "November",
#     "December",
# ]
# display(df_month)
# sorted_df_month = df_sort_values(df_month, "month", month_order, ascending=True)
# display(sorted_df_month)
# df_sort_values(df_month, "month", month_order, ascending=True, inplace=True)
# display(df_month)


def df_merge(
    df1: pd.DataFrame,
    df2: pd.DataFrame,
    use_index: bool = False,
    columns: list = ["col_left", "col_right"],
    how: str = "left",
    fuzz:bool= False,
    verbose:bool=True,
) -> pd.DataFrame:
    """
    Merges two DataFrames based on either the index or shared columns with matching data types.
    usage:
        #(1) if the index are the same
            df_merged = df_merge(df1, df2, use_index=True(defalut), how='outer')
        #(2) if there are shaed columns, then based on shared columns
            df_merged = df_merge(df1, df2, how='outer')
        #(3) if columns: then based on the specific columns
            df_merged = df_merge(df1, df2, columns=["col_left", "col_right"],how='outer')
    Parameters:
    - df1 (pd.DataFrame): The first DataFrame.
    - df2 (pd.DataFrame): The second DataFrame.
    - use_index (bool): If True, first try to merge by index if they are comparable; otherwise, fall back to column-based merge.
    - how (str): Type of merge to perform: 'inner', 'outer', 'left', or 'right'. Default is 'inner'.
    'inner': only the rows that have matching values in both DataFrames (intersection)
    'outer': keeps all rows from both DataFrames and fills in missing values with NaN
    'left': keeps all rows from the left DataFrame and matches rows from the right DataFrame
    'right': keeps all rows from the right DataFrame and matches rows from the left DataFrame, filling with NaN if there is no match.

    Returns:
    - pd.DataFrame: The merged DataFrame.
    """
    if fuzz:
        # try to merge them based on the columns, but not exact the same, but similar
        df1_copy=df1.copy()
        for i in df1_copy[columns[0]]:
            if verbose:
                print(f"double checking: {i}--------->{strcmp(i,df2[columns[1]].tolist())[0]}")
            df1_copy.loc[df1_copy[columns[0]]==i,columns[0]+"_corr"]=strcmp(i,df2[columns[1]].tolist() )[0]
        columns[0]=columns[0]+"_corr"
        df_=df_merge(df1_copy,df2,columns=columns)
        # drop the 'corr' column
        df_ = df_.drop(columns=[columns[0]])
        return df_
        
    # 1. Check if indices are comparable (same length and types)
    if use_index:
        print(f"Merging based on index using '{how}' join...")
        df_merged = pd.merge(df1, df2, left_index=True, right_index=True, how=how)
        return df_merged

    # 2. Find common columns with the same dtype
    common_columns = df1.columns.intersection(df2.columns)
    shared_columns = []
    for col in common_columns:
        try:
            if df1[col].dtype == df2[col].dtype:
                shared_columns.append(col)
        except Exception as e:
            print(e)
            pass
            
    if not isinstance(columns, list):
        columns = [columns]
    if len(columns) != 2:
        raise ValueError(
            "'columns':list shoule be a list: columns=['col_left','col_right']"
        )
    if all(columns):
        if verbose:
            print(f"Merging based on columns: {columns} using '{how}' join...")
        df_merged = pd.merge(df1, df2, left_on=columns[0], right_on=columns[1], how=how)
    elif shared_columns:
        if verbose:
            print(
                f"Merging based on shared columns: {shared_columns} using '{how}' join..."
            )
        df_merged = pd.merge(df1, df2, on=shared_columns, how=how)
    else:
        raise ValueError(
            "No common columns with matching data types to merge on, and indices are not comparable."
        )
    return df_merged


def df_drop_duplicates(
    data: pd.DataFrame,
    by: Union[
        str, List[str]
    ] = "index",  # Options: 'index', or column name(s) for 'rows'
    keep="first",  # Options: 'first', 'last', or False (drop all duplicates)
    ignore_index=True,
    inplace: bool = False,
    verbose=True,
):
    """
    data (pd.DataFrame): DataFrame to drop duplicates from.
    by (str): Specify by to drop duplicates:
                 - 'index': Drop duplicates based on the DataFrame index.
                 - Column name(s) for row-wise duplicate checking.
    keep (str): Which duplicates to keep:
        'first',
        'last',
        False (drop all duplicates).
    inplace (bool): Whether to modify the original DataFrame in place.
    """
    original_shape = data.shape
    if by == "index":
        # Drop duplicates in the index
        result = data[~data.index.duplicated(keep=keep)]
    else:
        # Drop duplicates row-wise based on column(s)
        result = data.drop_duplicates(subset=by, keep=keep, ignore_index=ignore_index)
    if original_shape != result.shape or verbose:
        print(f"\nshape:{original_shape} (before drop_duplicates)")
        print(f"shape:{result.shape} (after drop_duplicates)")
    if inplace:
        # Modify the original DataFrame in place
        data.drop(data.index, inplace=True)  # Drop all rows first
        data[data.columns] = result  # Refill the DataFrame
        return None
    else:
        return result


#! fillna()
def df_fillna(
    data: pd.DataFrame,
    method: str = "knn",
    axis: int = 0,  # column-wise
    constant: float = None,
    n_neighbors: int = 5,  # KNN-specific
    max_iter: int = 10,  # Iterative methods specific
    inplace: bool = False,
    random_state: int = 1,
) -> pd.DataFrame:
    """
    Fill missing values in a DataFrame using specified imputation method.

    Parameters:
    data (pd.DataFrame): The DataFrame to fill missing values.
    method (str): The imputation method to use. Options are:
        - 'mean': Replace missing values with the mean of the column.
        - 'median': Replace missing values with the median of the column.
        - 'most_frequent': Replace missing values with the most frequent value in the column.
        - 'constant': Replace missing values with a constant value provided by the `constant` parameter.
        - 'knn': Use K-Nearest Neighbors imputation; replaces missing values based on the values of the nearest neighbors
        - 'iterative': Use Iterative imputation; each feature with missing values as a function of other features and estimates them iteratively
        - 'mice' (Multivariate Imputation by Chained Equations): A special case of iterative imputation.
        # - 'missforest': A random forest-based imputation method. Uses a random forest model to predict and fill missing values
        # - 'softimpute': Matrix factorization imputation.A matrix factorization technique where missing values are imputed by
        #       reconstructing the data matrix using low-rank approximation
        # - EM (Expectation-Maximization): Often used in advanced statistics to estimate missing values in a probabilistic framework.
        # - 'svd': Use IterativeSVD (matrix factorization via Singular Value Decomposition).

    axis (int): The axis along which to impute:
        - 0: Impute column-wise (default).
        - 1: Impute row-wise.
    constant (float, optional): Constant value to use for filling NaNs if method is 'constant'.
    inplace (bool): If True, modify the original DataFrame. If False, return a new DataFrame.

    """
    if isinstance(data, pd.Series):
        data = pd.DataFrame(data)
    # handle None
    for col in data.columns:
        data[col] = data[col].apply(lambda x: np.nan if x is None else x)

    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0

    col_names_org = data.columns.tolist()
    index_names_org = data.index.tolist()
    # Separate numeric and non-numeric columns
    numeric_data = data.select_dtypes(include=[np.number])
    non_numeric_data = data.select_dtypes(exclude=[np.number])

    if data.empty:
        raise ValueError("Input DataFrame is empty.")

    # Validate method
    methods = [
        "mean",
        "median",
        "most_frequent",
        "constant",
        "knn",
        "iterative",
    ]  # ,"missforest","softimpute","svd"]
    method = strcmp(method, methods)[0]

    # If using constant method, ask for a constant value
    if constant is not None:
        method = "constant"
        try:
            constant = float(constant)
        except ValueError:
            raise ValueError("Constant value must be a number.")

    # Initialize SimpleImputer with the chosen method
    if method == "constant":
        from sklearn.impute import SimpleImputer

        imputer = SimpleImputer(strategy=method, fill_value=constant)
    elif method == "knn":
        from sklearn.impute import KNNImputer

        imputer = KNNImputer(n_neighbors=n_neighbors)
    elif method == "iterative" or method == "mice":
        from sklearn.experimental import enable_iterative_imputer
        from sklearn.impute import IterativeImputer

        imputer = IterativeImputer(max_iter=max_iter, random_state=random_state)
    else:  # mean, median, most_frequent
        from sklearn.impute import SimpleImputer

        imputer = SimpleImputer(strategy=method)

    # Fit and transform the data
    if axis == 0:
        # Impute column-wise
        imputed_data = imputer.fit_transform(numeric_data)
    elif axis == 1:
        # Impute row-wise
        imputed_data = imputer.fit_transform(numeric_data.T)
    else:
        raise ValueError("Invalid axis. Use 0 for columns or 1 for rows.")

    imputed_data = pd.DataFrame(
        imputed_data if axis == 0 else imputed_data.T,
        index=numeric_data.index if axis == 0 else numeric_data.columns,
        columns=numeric_data.columns if axis == 0 else numeric_data.index,
    )
    for col in imputed_data.select_dtypes(include=[np.number]).columns:
        imputed_data[col] = imputed_data[col].astype(numeric_data[col].dtype)

    # Handle non-numeric data imputation
    if not non_numeric_data.empty:
        from sklearn.impute import SimpleImputer

        if method == "constant":
            non_numeric_imputer = SimpleImputer(
                strategy="constant", fill_value=constant
            )
        else:
            non_numeric_imputer = SimpleImputer(strategy="most_frequent")

        # Impute non-numeric columns column-wise (axis=0)
        imputed_non_numeric = non_numeric_imputer.fit_transform(non_numeric_data)

        # Convert imputed non-numeric array back to DataFrame with original index and column names
        imputed_non_numeric_df = pd.DataFrame(
            imputed_non_numeric,
            index=non_numeric_data.index,
            columns=non_numeric_data.columns,
        )
    else:
        imputed_non_numeric_df = pd.DataFrame(index=data.index)

    imputed_data = pd.concat([imputed_data, imputed_non_numeric_df], axis=1).reindex(
        columns=data.columns
    )

    if inplace:
        # Modify the original DataFrame
        data[:] = imputed_data[col_names_org]
        return None
    else:
        # Return the modified DataFrame
        return imputed_data[col_names_org]


# # example
# data = {
#     "A": [1, 2, np.nan, 4, 5],
#     "B": [np.nan, 2, 3, 4, np.nan],
#     "C": [1, np.nan, 3, 4, 5],
#     "D": [1, 2, 3, 4, np.nan],
# }

# # Define a function to test each imputation method
# methods = [
#     "mean",
#     "median",
#     "most_frequent",
#     "constant",
#     "knn",
#     "iterative",
#     # "missforest",
#     # "softimpute",
#     # "svd",
# ]

# # Create a dictionary to hold results
# results = {}

# for method_name in methods:
#     print(method_name)
#     display(df)
#     display(df_fillna(data=df, method=method_name, inplace=False, axis=0))
def df_cut(
    df: pd.DataFrame,
    column: str,
    *,
    new_col_name: Optional[str] = None,
    bins: Optional[
        Union[int, List[float], Dict[str, Union[float, str, pd.Timestamp]]]
    ] = None,
    range_start: Optional[Union[float, str, pd.Timestamp]] = None,
    range_end: Optional[Union[float, str, pd.Timestamp]] = None,
    step: Optional[Union[float, str, pd.Timedelta]] = None,
    labels: Optional[List[str]] = None,
    label_format: Optional[Union[str, Callable[[float, float], str]]] = None,
    include_overflow: bool = True,
    include_underflow: bool = False,
    right: bool = False,
    drop_original: bool = False,
    precision: int = 2,
    show_count: bool = False,
    symbol_count: str = "n=",
    show_percentage: bool = False,
    symbol_percentage: str = "%",
    show_total_count: bool = False,
    symbol_total_count: str = "∑n=",
    sep_between: str = " | ",
    sort_labels: bool = True,
    na_action: str = "keep",
    na_fill_value: Optional[str] = None,
    dtype: Optional[Union[str, pd.CategoricalDtype]] = None,
    ordered: bool = True,
    inplace: bool = False,
    datetime_format: str = "%Y-%m-%d",
    categorical_agg: str = "count",
) -> Optional[pd.DataFrame]:
    """
            Enhanced binning function that works with numeric, datetime, and categorical columns.

            Features:
            - Automatic type detection (numeric, datetime, categorical)
            - Flexible bin specification (number of bins, explicit edges, or range+step)
            - Customizable labels with formatting
            - Count and percentage display options
            - NA value handling
            square bracket: means inclusive
            parenthesis: means exclusive
            Parameters:
            -----------
            df : pd.DataFrame
                Input DataFrame containing the column to bin
            column : str
                Name of column to bin
            new_col_name : str, optional
                Name for binned column (default: f"{column}_binned")
            bins : int, list, or dict, optional
                - int: Number of equal-width bins
                - list: Explicit bin edges
                - dict: {'start': x, 'end': y, 'step': z} for range specification
            range_start : float or datetime-like, optional
                Start value for bin range (required if bins is None or dict)
            range_end : float or datetime-like, optional
                End value for bin range (default: max of column)
            step : float or timedelta-like, optional
                Step size for bin creation (required if bins is None or dict)
            labels : list of str, optional
                Custom labels for bins (must match number of bins)
            label_format : str or callable, optional
                Format string or function for bin labels
            include_overflow : bool, default True
                Include catch-all bin for values above range_end
            include_underflow : bool, default False
                Include catch-all bin for values below range_start
            right : bool, default False
                Whether bins include the right edge
            drop_original : bool, default False
                Drop original column after binning
            precision : int, default 2
                Decimal precision for numeric bin labels
            show_count : bool, default False
                Show count of items in each bin
            show_percentage : bool, default False
                Show percentage of items in each bin
            show_total_count : bool, default False
                Show total count in labels
            na_action : str, default 'keep'
                How to handle NA values ('keep', 'drop', or 'fill')
            na_fill_value : str, optional
                Value to fill NAs with if na_action='fill'
            dtype : dtype or CategoricalDtype, optional
                Output dtype for binned column
            ordered : bool, default True
                Whether bins are ordered
            inplace : bool, default False
                Modify DataFrame in place
            datetime_format : str, default "%Y-%m-%d"
                Format string for datetime labels
            categorical_agg : str, default 'count'
                For categorical data: 'count' or 'ratio'

            Returns:
            --------
            pd.DataFrame or None
                Returns modified DataFrame unless inplace=True

            Examples:
            --------
            # Numeric binning
            df_cut(df, 'age', bins=5)
            df_cut(df, 'price', range_start=0, range_end=1000, step=100)

            # Datetime binning
            df_cut(df, 'date', bins={'start': '2023-01-01', 'end': '2023-12-31', 'step': '1M'})

            # Categorical binning
            df_cut(df, 'category', bins=5, categorical_agg='ratio')

    # Sample datetime data
    dates = pd.date_range("2020-01-01", "2023-12-31", freq="D")
    df = pd.DataFrame(
        {
            "order_date": np.random.choice(dates, 500),
            "delivery_time": np.random.randint(1, 72, 500),  # hours
        }
    )
    # Example 1: Monthly bins
    # Monthly binning with exact month boundaries
    df_cut(
        df,
        "order_date",
        bins={"start": "2019-01-01", "end": "2023-12-31", "step": "1Y"},
        datetime_format="%Y-%m-%d",
        label_format="%m-%d",
        show_count=True,
        show_percentage=True,
        show_total_count=True,
    )
    # Weekly binning
    df_cut(
        df,
        "order_date",
        bins={"start": "2019-01-01", "end": "2023-12-31", "step": "1W"},
        label_format="%Y-%m-%d",
        datetime_format="%Y-%m-%d",
        show_count=True,
        show_percentage=True,
        show_total_count=True,
    )


    # Sample numeric data
    df = pd.DataFrame(
        {"price": np.random.uniform(10, 1000, 1000), "age": np.random.randint(18, 80, 1000)}
    )

    # Example 1: Equal-width bins
    df_cut(df, "price", bins=5, show_count=True)

    # Example 2: Custom range with step
    df_cut(
        df,
        "price",
        range_start=0,
        range_end=1000,
        step=200,
        label_format="${left:.0f}-${right:.0f}",
        show_percentage=True,
    )
    df_cut(
        df,
        "price",
        bins={"start": 0, "end": 1000, "step": 200},
        # label_format="${left:.0f}-${right:.0f}",
        show_percentage=True,
    )
    """
    from pandas.api.types import is_numeric_dtype, is_datetime64_any_dtype

    def _process_time_step(step: Union[str, int, float, pd.Timedelta]) -> str:
        """Convert step to pandas frequency string."""
        if isinstance(step, pd.Timedelta):
            return step.freqstr if step.freqstr else str(step)

        if isinstance(step, (int, float)):
            return f"{step}S"  # Interpret numbers as seconds

        if isinstance(step, str):
            step = step.strip().lower()
            match = re.match(r"(\d*\.?\d+)?\s*([a-z]+)", step)
            if not match:
                raise ValueError(f"Invalid time step format: {step}")

            num_part, unit_part = match.groups()
            num = float(num_part) if num_part else 1.0

            unit_map = {
                "y": "Y",
                "yr": "Y",
                "yrs": "Y",
                "year": "Y",
                "years": "Y",
                "m": "M",
                "mo": "M",
                "mon": "M",
                "month": "M",
                "months": "M",
                "w": "W",
                "wk": "W",
                "wks": "W",
                "week": "W",
                "weeks": "W",
                "d": "D",
                "day": "D",
                "days": "D",
                "h": "H",
                "hr": "H",
                "hrs": "H",
                "hour": "H",
                "hours": "H",
                "min": "T",
                "mins": "T",
                "minute": "T",
                "minutes": "T",
                "s": "S",
                "sec": "S",
                "secs": "S",
                "second": "S",
                "seconds": "S",
            }

            if unit_part not in unit_map:
                raise ValueError(f"Unknown time unit: {unit_part}")

            freq = unit_map[unit_part]
            if num.is_integer():
                num = int(num)
            return f"{num}{freq}"

        raise TypeError(f"Unsupported step type: {type(step)}")


    def _process_datetime_column(
        col: pd.Series,
        bins: Optional[Union[int, List[pd.Timestamp]]],
        range_start: Optional[Union[str, pd.Timestamp]],
        range_end: Optional[Union[str, pd.Timestamp]],
        step: Optional[Union[str, pd.Timedelta]],
        labels: Optional[List[str]],
        label_format: Optional[Union[str, Callable]],
        datetime_format: str,
        right: bool,
        include_underflow: bool,
        include_overflow: bool,
    ) -> Tuple[pd.Categorical, List[str]]:
        """Process datetime column with accurate counting."""
        col = pd.to_datetime(col)

        # Handle bin edges
        if bins is None:
            if step is None:
                raise ValueError("Step must be provided for datetime binning")

            # Convert step to pandas frequency string
            step_freq = _process_time_step(step)

            # Set default range if needed
            range_start = (
                pd.to_datetime(range_start) if range_start is not None else col.min()
            )
            range_end = pd.to_datetime(range_end) if range_end is not None else col.max()

            # Generate bins
            try:
                bin_edges = pd.date_range(start=range_start, end=range_end, freq=step_freq)
                if len(bin_edges) == 0:
                    bin_edges = pd.date_range(start=range_start, end=range_end, periods=2)
                elif bin_edges[-1] < range_end:
                    bin_edges = bin_edges.append(pd.DatetimeIndex([range_end]))
            except ValueError as e:
                raise ValueError(f"Invalid frequency specification: {step_freq}") from e
        elif isinstance(bins, int):
            bin_edges = pd.date_range(start=col.min(), end=col.max(), periods=bins + 1)
        else:
            bin_edges = pd.to_datetime(bins)

        # Add overflow/underflow bins
        if include_underflow:
            bin_edges = bin_edges.insert(0, pd.Timestamp.min)
        if include_overflow:
            bin_edges = bin_edges.append(pd.DatetimeIndex([pd.Timestamp.max]))

        # Perform the cut - this is where we ensure proper binning
        binned = pd.cut(
            col.astype("int64"),  # Convert to nanoseconds for precise binning
            bins=bin_edges.astype("int64"),
            right=right,
            include_lowest=True,
        )

        # Generate labels if not provided
        if labels is None:
            labels = []
            for i in range(len(bin_edges) - 1):
                left = bin_edges[i]
                right_ = bin_edges[i + 1]

                # Handle special cases
                if left == pd.Timestamp.min:
                    left_str = "<"
                else:
                    left_str = left.strftime(datetime_format)

                if right_ == pd.Timestamp.max:
                    right_str = ">"
                else:
                    right_str = right_.strftime(datetime_format)

                # Apply label formatting
                if callable(label_format):
                    label = label_format(left, right_)
                elif isinstance(label_format, str):
                    try:
                        if left != pd.Timestamp.min and right_ != pd.Timestamp.max:
                            label = f"{left.strftime(label_format)}-{right_.strftime(label_format)}"
                        else:
                            label = f"{left_str}-{right_str}"
                    except (ValueError, AttributeError):
                        label = f"{left_str}-{right_str}"
                else:
                    label = f"{left_str}-{right_str}"

                labels.append(label)

        return binned, labels


    def _process_categorical_column(
        col: pd.Series,
        bins: Optional[Union[int, List[str]]],
        labels: Optional[List[str]],
        categorical_agg: str,
    ) -> Tuple[pd.Categorical, List[str]]:
        value_counts = col.value_counts(normalize=(categorical_agg == "ratio"))

        if bins is not None and isinstance(bins, int):
            top_categories = value_counts.head(bins).index
            binned = col.where(col.isin(top_categories), "Other")
        elif isinstance(bins, list):
            binned = col.where(col.isin(bins), "Other")
        else:
            binned = col

        binned = binned.astype("category")

        if labels is not None:
            binned = binned.cat.rename_categories(dict(zip(binned.cat.categories, labels)))

        return binned, list(binned.cat.categories)


    def _process_numeric_column(
        col: pd.Series,
        bins: Optional[Union[int, List[float]]],
        range_start: Optional[float],
        range_end: Optional[float],
        step: Optional[float],
        labels: Optional[List[str]],
        label_format: Optional[Union[str, Callable]],
        precision: int,
        right: bool,
        include_underflow: bool,
        include_overflow: bool,
    ) -> Tuple[pd.Categorical, List[str]]:
        if bins is None:
            if range_start is None or step is None:
                raise ValueError("If bins not provided, must set range_start and step")
            if range_end is None:
                range_end = col.max()

            bin_edges = list(np.arange(range_start, range_end + step, step))
        elif isinstance(bins, int):
            bin_edges = np.linspace(col.min(), col.max(), bins + 1).tolist()
        else:
            bin_edges = list(bins)

        # Add overflow/underflow bins if needed
        if include_underflow and not np.isinf(bin_edges[0]):
            bin_edges.insert(0, float("-inf"))
        if include_overflow and not np.isinf(bin_edges[-1]):
            bin_edges.append(float("inf"))

        # Generate labels if not provided
        if labels is None:
            labels = []
            for i in range(len(bin_edges) - 1):
                left = round(bin_edges[i], precision)
                right_ = round(bin_edges[i + 1], precision)

                if label_format:
                    label = (
                        label_format(left, right_)
                        if callable(label_format)
                        else label_format.format(left=left, right=right_)
                    )
                else:
                    if np.isinf(left) and left < 0:
                        label = f"<{right_}"
                    elif np.isinf(right_):
                        label = f">{left}"
                    else:
                        label = f"[{left}, {right_}{']' if right else ')'}"

                labels.append(label)

        binned = pd.cut(
            col, bins=bin_edges, labels=labels, right=right, include_lowest=True
        )
        return binned, labels


    def _handle_na_values(
        col: pd.Series, na_action: str, na_fill_value: Optional[str]
    ) -> pd.Series:
        if na_action == "drop":
            return col.dropna()
        elif na_action == "fill" and na_fill_value is not None:
            return col.fillna(na_fill_value)
        return col


    def _add_statistical_labels(
        binned: pd.Categorical,
        labels: List[str],
        show_count: bool,
        show_percentage: bool,
        show_total_count: bool,
        symbol_count: str,
        symbol_percentage: str,
        symbol_total_count: str,
        sep_between: str, 
    ) -> List[str]:
        """Add statistical information with accurate counts."""
        # Get counts by matching the exact bin intervals
        value_counts = binned.value_counts()
        total = len(binned.dropna())

        new_labels = []
        for i, (label, category) in enumerate(zip(labels, binned.cat.categories)):
            count = value_counts.get(category, 0)
            parts = [label]

            if show_count:
                parts.append(f"{symbol_count}{count}")
            if show_percentage:
                percentage = (count / total * 100) if total > 0 else 0
                parts.append(f"{percentage:.1f}{symbol_percentage}")
            if show_total_count:
                parts.append(f"{symbol_total_count}{total}")

            # Ensure unique labels 
            new_label = sep_between.join(parts)                
            if new_label in new_labels:
                new_label = f"{new_label}_{i}"
            new_labels.append(new_label)

        return new_labels


    def _sort_bin_labels(binned: pd.Categorical, labels: List[str]) -> pd.Categorical:
        try:
            # Attempt to sort by the underlying intervals
            sorted_categories = sorted(binned.cat.categories)
            binned = binned.cat.reorder_categories(sorted_categories, ordered=True)
        except Exception:
            # If sorting fails (e.g., string labels), fallback to given label order
            binned = binned.cat.set_categories(labels, ordered=True)
        return binned
    # Input validation
    if column not in df.columns:
        raise ValueError(f"Column '{column}' not found in DataFrame")

    if not inplace:
        df = df.copy()

    col_data = df[column]

    # Determine column type
    if is_datetime64_any_dtype(col_data):
        col_type = "datetime"
        col_data = pd.to_datetime(col_data)
    elif isinstance(col_data.dtype, pd.CategoricalDtype) or col_data.dtype == "object":
        col_type = "categorical"
    elif is_numeric_dtype(col_data):
        col_type = "numeric"
    else:
        raise TypeError(f"Unsupported column type: {col_data.dtype}")

    # Handle dictionary bin specification
    if isinstance(bins, dict):
        range_start = bins.get("start", range_start)
        range_end = bins.get("end", range_end)
        step = bins.get("step", step)
        bins = None

    # Process based on column type
    if col_type == "datetime":
        binned, bin_labels = _process_datetime_column(
            col_data,
            bins,
            range_start,
            range_end,
            step,
            labels,
            label_format,
            datetime_format,
            right,
            include_underflow,
            include_overflow,
        )
    elif col_type == "categorical":
        binned, bin_labels = _process_categorical_column(
            col_data, bins, labels, categorical_agg
        )
    else:
        binned, bin_labels = _process_numeric_column(
            col_data,
            bins,
            range_start,
            range_end,
            step,
            labels,
            label_format,
            precision,
            right,
            include_underflow,
            include_overflow,
        )

    # Handle NA values
    binned = _handle_na_values(binned, na_action, na_fill_value)

    # Add statistical information to labels if requested
    if show_count or show_percentage or show_total_count:
        bin_labels = _add_statistical_labels(
            binned,
            bin_labels,
            show_count,
            show_percentage,
            show_total_count,
            symbol_count,
            symbol_percentage,
            symbol_total_count,
            sep_between,
        )
        binned = binned.cat.rename_categories(
            dict(zip(binned.cat.categories, bin_labels))
        )

    # Sort labels if requested
    if sort_labels and not right and len(bin_labels) > 1:
        binned = _sort_bin_labels(binned, bin_labels)

    # Create final output column
    new_col = new_col_name or f"{column}_binned"
    df[new_col] = binned.astype(dtype) if dtype else binned

    if drop_original:
        df.drop(columns=[column], inplace=True)

    return None if inplace else df



def df_encoder(
    data: pd.DataFrame,
    method: str = "dummy",  #'dummy', 'onehot', 'ordinal', 'label', 'target', 'binary'
    columns=None,
    target_column=None,  # Required for 'target' encoding method
    **kwargs,
) -> pd.DataFrame:
    """
    Methods explained:
    - 'dummy': pandas' `get_dummies` to create dummy variables for categorical columns, which is another form of one-hot encoding, but with a simpler interface.

    - 'onehot': One-hot encoding is used when there is no inherent order in categories. It creates a binary column for each category and is useful for nominal categorical variables. However, it increases dimensionality significantly if there are many unique categories.

    - 'ordinal': Ordinal encoding is used when there is an inherent order in the categories. It assigns integers to categories based on their order. Use this when the categories have a ranking (e.g., 'low', 'medium', 'high').

    - 'label': Label encoding is used for converting each unique category to a numeric label. It can be useful when working with algorithms that can handle categorical data natively (e.g., decision trees). However, it might introduce unintended ordinal relationships between the categories.

    - 'target': Target encoding is used when you encode a categorical feature based on the mean of the target variable. This is useful when there is a strong correlation between the categorical feature and the target variable. It is often used in predictive modeling to capture relationships that are not directly encoded in the feature.

    - 'binary': Binary encoding is a more efficient alternative to one-hot encoding when dealing with high-cardinality categorical variables. It converts categories into binary numbers and then splits them into multiple columns, reducing dimensionality compared to one-hot encoding.
    """

    # Select categorical columns
    categorical_cols = data.select_dtypes(exclude=np.number).columns.tolist()
    methods = ["dummy", "onehot", "ordinal", "label", "target", "binary"]
    method = strcmp(method, methods)[0]

    if columns is None:
        columns = categorical_cols

    # pd.get_dummies()
    if method == "dummy":
        dtype = kwargs.pop("dtype", int)
        drop_first = kwargs.pop("drop_first", True)
        try:
            encoded_df = pd.get_dummies(
                data[columns], drop_first=drop_first, dtype=dtype, **kwargs
            )
            return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)
        except Exception as e:
            # print(f"Warning, 没有进行转换, 因为: {e}")
            return data
    # One-hot encoding
    elif method == "onehot":
        from sklearn.preprocessing import OneHotEncoder

        encoder = OneHotEncoder(drop="first", sparse_output=False, **kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        encoded_df = pd.DataFrame(
            encoded_data,
            columns=encoder.get_feature_names_out(columns),
            index=data.index,
        )
        return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)

    # Ordinal encoding
    elif method == "ordinal":
        from sklearn.preprocessing import OrdinalEncoder

        encoder = OrdinalEncoder(**kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        encoded_df = pd.DataFrame(encoded_data, columns=columns, index=data.index)
        return pd.concat([data.drop(columns, axis=1), encoded_df], axis=1)

    # Label encoding
    elif method == "label":
        from sklearn.preprocessing import LabelEncoder

        encoder = LabelEncoder()
        # Apply LabelEncoder only to non-numeric columns
        non_numeric_columns = [
            col for col in columns if not pd.api.types.is_numeric_dtype(data[col])
        ]

        if not non_numeric_columns:
            return data
        encoded_data = data[non_numeric_columns].apply(
            lambda col: encoder.fit_transform(col)
        )
        return pd.concat([data.drop(non_numeric_columns, axis=1), encoded_data], axis=1)

    # Target encoding (Mean of the target for each category)
    elif method == "target":
        if target_column is None:
            raise ValueError("target_column must be provided for target encoding.")
        from category_encoders import TargetEncoder

        encoder = TargetEncoder(cols=columns, **kwargs)
        encoded_data = encoder.fit_transform(data[columns], data[target_column])
        return pd.concat([data.drop(columns, axis=1), encoded_data], axis=1)

    # Binary encoding (for high-cardinality categorical variables)
    elif method == "binary":
        from category_encoders import BinaryEncoder

        encoder = BinaryEncoder(cols=columns, **kwargs)
        encoded_data = encoder.fit_transform(data[columns])
        return pd.concat([data.drop(columns, axis=1), encoded_data], axis=1)


def df_scaler(
    data: pd.DataFrame,  # should be numeric dtype
    scaler=None,
    method="standard",
    columns=None,  # default, select all numeric col/row
    feature_range=None,  # specific for 'minmax'
    vmin=0,
    vmax=1,
    inplace=False,
    verbose=False,  # show usage
    axis=0,  # defalut column-wise
    return_scaler: bool = False,  # True: return both: return df, scaler
    **kwargs,
):
    """
    df_scaler(data, scaler="standard", inplace=False, axis=0, verbose=True)

    Parameters:
    - data: pandas DataFrame to be scaled.
    - method: Scaler type ('standard', 'minmax', 'robust'). Default is 'standard'.
    - columns: List of columns (for axis=0) or rows (for axis=1) to scale.
               If None, all numeric columns/rows will be scaled.
    - inplace: If True, modify the DataFrame in place. Otherwise, return a new DataFrame.
    - axis: Axis along which to scale. 0 for column-wise, 1 for row-wise. Default is 0.
    - verbose: If True, prints logs of the process.
    - kwargs: Additional arguments to be passed to the scaler.
    """
    if verbose:
        print('df_scaler(data, scaler="standard", inplace=False, axis=0, verbose=True)')
    if scaler is None:
        methods = ["standard", "minmax", "robust", "maxabs"]
        method = strcmp(method, methods)[0]
        if method == "standard":
            from sklearn.preprocessing import StandardScaler

            if verbose:
                print(
                    "performs z-score normalization: This will standardize each feature to have a mean of 0 and a standard deviation of 1."
                )
                print(
                    "Use when the data is approximately normally distributed (Gaussian).\nWorks well with algorithms sensitive to feature distribution, such as SVMs, linear regression, logistic regression, and neural networks."
                )
            scaler = StandardScaler(**kwargs)
        elif method == "minmax":
            from sklearn.preprocessing import MinMaxScaler

            if feature_range is None:
                feature_range = (vmin, vmax)
            if verbose:
                print(
                    "don't forget to define the range: e.g., 'feature_range=(0, 1)'. "
                )
                print(
                    "scales the features to the range [0, 1]. Adjust feature_range if you want a different range, like [-1, 1]."
                )
                print(
                    "Use when the data does not follow a normal distribution and you need all features in a specific range (e.g., [0, 1]).\nIdeal for algorithms that do not assume a particular distribution, such as k-nearest neighbors and neural networks."
                )
            scaler = MinMaxScaler(feature_range=feature_range, **kwargs)
        elif method == "robust":
            from sklearn.preprocessing import RobustScaler

            if verbose:
                print(
                    "scales the data based on the median and interquartile range, which is robust to outliers."
                )
                print(
                    "Use when the dataset contains outliers.\nThis method is useful because it scales based on the median and the interquartile range (IQR), which are more robust to outliers than the mean and standard deviation."
                )
            scaler = RobustScaler(**kwargs)
        elif method == "maxabs":
            from sklearn.preprocessing import MaxAbsScaler

            if verbose:
                print(
                    "This scales each feature by its maximum absolute value, resulting in values within the range [-1, 1] for each feature."
                )
                print(
                    "Use for data that is already sparse or when features have positive or negative values that need scaling without shifting the data.\nOften used with sparse data (data with many zeros), where preserving zero entries is essential, such as in text data or recommendation systems."
                )
            scaler = MaxAbsScaler(**kwargs)
        if axis not in [0, 1]:
            raise ValueError("Axis must be 0 (column-wise) or 1 (row-wise).")
    if verbose:
        print(scaler)
    if axis == 0:
        # Column-wise scaling (default)
        if columns is None:
            columns = data.select_dtypes(include=np.number).columns.tolist()
        non_numeric_columns = data.columns.difference(columns)

        # scaled_data = scaler.fit_transform(data[columns])
        if scaler is None or not hasattr(scaler, "mean_"):
            scaled_data = scaler.fit_transform(data[columns])
        else:
            scaled_data = scaler.transform(data[columns])

        if inplace:
            data[columns] = scaled_data
            print("Original DataFrame modified in place (column-wise).")
        else:
            scaled_df = pd.concat(
                [
                    pd.DataFrame(scaled_data, columns=columns, index=data.index),
                    data[non_numeric_columns],
                ],
                axis=1,
            )
            scaled_df = scaled_df[data.columns]  # Maintain column order
            if return_scaler:
                return scaled_df, scaler
            else:
                return scaled_df

    elif axis == 1:
        # Row-wise scaling
        if columns is None:
            columns = data.index.tolist()
        numeric_rows = data.loc[columns].select_dtypes(include=np.number)
        if numeric_rows.empty:
            raise ValueError("No numeric rows to scale.")

        print(f"Scaling rows")

        # scaled_data = scaler.fit_transform(
        #     numeric_rows.T
        # ).T  # Transpose for scaling and then back
        scaled_data = (
            scaler.fit_transform(numeric_rows.T).T
            if scaler is None or not hasattr(scaler, "mean_")
            else scaler.transform(numeric_rows.T).T
        )

        if inplace:
            data.loc[numeric_rows.index] = scaled_data
            print("Original DataFrame modified in place (row-wise).")
        else:
            scaled_df = data.copy()
            scaled_df.loc[numeric_rows.index] = scaled_data
            if return_scaler:
                return scaled_df, scaler
            else:
                return scaled_df


def df_special_characters_cleaner(
    data: pd.DataFrame, where=["column", "content", "index"]
) -> pd.DataFrame:
    """
    to clean special characters:
    usage:
        df_special_characters_cleaner(data=df, where='column')
    """
    if not isinstance(where, list):
        where = [where]
    where_to_clean = ["column", "content", "index"]
    where_ = [strcmp(i, where_to_clean)[0] for i in where]

    # 1. Clean column names by replacing special characters with underscores
    if "column" in where_:
        try:
            data.columns = data.columns.str.replace(r"[^\w\s]", "_", regex=True)
        except Exception as e:
            print(e)

    # 2. Clean only object-type columns (text columns)
    try:
        if "content" in where_:
            for col in data.select_dtypes(include=["object"]).columns:
                data[col] = data[col].str.replace(r"[^\w\s]", "", regex=True)
        if data.index.dtype == "object" and index in where_:
            data.index = data.index.str.replace(r"[^\w\s]", "_", regex=True)
    except:
        pass
    return data


def df_cluster(
    data: pd.DataFrame,
    columns: Optional[list] = None,
    n_clusters: Optional[int] = None,
    range_n_clusters: Union[range, np.ndarray] = range(2, 11),
    scale: bool = True,
    plot: Union[str, list] = "all",
    inplace: bool = True,
    ax=None,
):
    from sklearn.preprocessing import StandardScaler
    from sklearn.cluster import KMeans
    from sklearn.metrics import silhouette_score, silhouette_samples
    import seaborn as sns
    import numpy as np
    import pandas as pd
    import matplotlib.pyplot as plt

    """
    Performs clustering analysis on the provided feature matrix using K-Means.

    Parameters:
        X (np.ndarray):
            A 2D numpy array or DataFrame containing numerical feature data,
            where each row corresponds to an observation and each column to a feature.

        range_n_clusters (range):
            A range object specifying the number of clusters to evaluate for K-Means clustering.
            Default is range(2, 11), meaning it will evaluate from 2 to 10 clusters.

        scale (bool):
            A flag indicating whether to standardize the features before clustering.
            Default is True, which scales the data to have a mean of 0 and variance of 1.

        plot (bool):
            A flag indicating whether to generate visualizations of the clustering analysis.
            Default is True, which will plot silhouette scores, inertia, and other relevant plots.
    Returns:
        tuple: 
            A tuple containing the modified DataFrame with cluster labels, 
            the optimal number of clusters, and the Axes object (if any).
    """
    X = data[columns].values if columns is not None else data.values

    silhouette_avg_scores = []
    inertia_scores = []

    # Standardize the features
    if scale:
        scaler = StandardScaler()
        X = scaler.fit_transform(X)

    for n_cluster in range_n_clusters:
        kmeans = KMeans(n_clusters=n_cluster, random_state=1)
        cluster_labels = kmeans.fit_predict(X)

        silhouette_avg = silhouette_score(X, cluster_labels)
        silhouette_avg_scores.append(silhouette_avg)
        inertia_scores.append(kmeans.inertia_)
        print(
            f"For n_clusters = {n_cluster}, the average silhouette_score is : {silhouette_avg:.4f}"
        )

    # Determine the optimal number of clusters based on the maximum silhouette score
    if n_clusters is None:
        n_clusters = range_n_clusters[np.argmax(silhouette_avg_scores)]
    print(f"n_clusters = {n_clusters}")

    # Apply K-Means Clustering with Optimal Number of Clusters
    kmeans = KMeans(n_clusters=n_clusters, random_state=1)
    cluster_labels = kmeans.fit_predict(X)

    if plot:
        # ! Interpreting the plots from your K-Means clustering analysis
        # ! 1. Silhouette Score and Inertia vs Number of Clusters
        # Description:
        # This plot has two y-axes: the left y-axis represents the Silhouette Score, and the right y-axis
        # represents Inertia.
        # The x-axis represents the number of clusters (k).

        # Interpretation:

        # Silhouette Score:
        # Ranges from -1 to 1, where a score close to 1 indicates that points are well-clustered, while a
        # score close to -1 indicates that points might be incorrectly clustered.
        # A higher silhouette score generally suggests that the data points are appropriately clustered.
        # Look for the highest value to determine the optimal number of clusters.

        # Inertia:
        # Represents the sum of squared distances from each point to its assigned cluster center.
        # Lower inertia values indicate tighter clusters.
        # As the number of clusters increases, inertia typically decreases, but the rate of decrease
        # may slow down, indicating diminishing returns for additional clusters.

        # Optimal Number of Clusters:
        # You can identify an optimal number of clusters where the silhouette score is maximized and
        # inertia starts to plateau (the "elbow" point).
        # This typically suggests that increasing the number of clusters further yields less meaningful
        # separations.
        if ax is None:
            _, ax = plt.subplots(figsize=inch2cm(10, 6))
        color = "tab:blue"
        ax.plot(
            range_n_clusters,
            silhouette_avg_scores,
            marker="o",
            color=color,
            label="Silhouette Score",
        )
        ax.set_xlabel("Number of Clusters")
        ax.set_ylabel("Silhouette Score", color=color)
        ax.tick_params(axis="y", labelcolor=color)
        # add right axis: inertia
        ax2 = ax.twinx()
        color = "tab:red"
        ax2.set_ylabel("Inertia", color=color)
        ax2.plot(
            range_n_clusters,
            inertia_scores,
            marker="x",
            color=color,
            label="Inertia",
        )
        ax2.tick_params(axis="y", labelcolor=color)

        plt.title("Silhouette Score and Inertia vs Number of Clusters")
        plt.xticks(range_n_clusters)
        plt.grid()
        plt.axvline(x=n_clusters, linestyle="--", color="r", label="Optimal n_clusters")
        # ! 2. Elbow Method Plot
        # Description:
        # This plot shows the Inertia against the number of clusters.

        # Interpretation:
        # The elbow point is where the inertia begins to decrease at a slower rate. This point suggests that
        # adding more clusters beyond this point does not significantly improve the clustering performance.
        # Look for a noticeable bend in the curve to identify the optimal number of clusters, indicated by the
        # vertical dashed line.
        # Inertia plot
        plt.figure(figsize=inch2cm(10, 6))
        plt.plot(range_n_clusters, inertia_scores, marker="o")
        plt.title("Elbow Method for Optimal k")
        plt.xlabel("Number of clusters")
        plt.ylabel("Inertia")
        plt.grid()
        plt.axvline(
            x=np.argmax(silhouette_avg_scores) + 2,
            linestyle="--",
            color="r",
            label="Optimal n_clusters",
        )
        plt.legend()
        # ! Silhouette Plots
        # 3. Silhouette Plot for Various Clusters
        # Description:
        # This horizontal bar plot shows the silhouette coefficient values for each sample, organized by cluster.

        # Interpretation:
        # Each bar represents the silhouette score of a sample within a specific cluster. Longer bars indicate
        # that the samples are well-clustered.
        # The height of the bars shows how similar points within the same cluster are to one another compared to
        # points in other clusters.
        # The vertical red dashed line indicates the average silhouette score for all samples.
        # You want the majority of silhouette values to be above the average line, indicating that most points
        # are well-clustered.

        # 以下代码不用再跑一次了
        # n_clusters = (
        #     np.argmax(silhouette_avg_scores) + 2
        # )  # Optimal clusters based on max silhouette score
        # kmeans = KMeans(n_clusters=n_clusters, random_state=1)
        # cluster_labels = kmeans.fit_predict(X)
        silhouette_vals = silhouette_samples(X, cluster_labels)

        plt.figure(figsize=inch2cm(10, 6))
        y_lower = 10
        for i in range(n_clusters):
            # Aggregate the silhouette scores for samples belonging to cluster i
            ith_cluster_silhouette_values = silhouette_vals[cluster_labels == i]

            # Sort the values
            ith_cluster_silhouette_values.sort()

            size_cluster_i = ith_cluster_silhouette_values.shape[0]
            y_upper = y_lower + size_cluster_i

            # Create a horizontal bar plot for the silhouette scores
            plt.barh(range(y_lower, y_upper), ith_cluster_silhouette_values, height=0.5)

            # Label the silhouette scores
            plt.text(-0.05, (y_lower + y_upper) / 2, str(i + 2))
            y_lower = y_upper + 10  # 10 for the 0 samples

        plt.title("Silhouette Plot for the Various Clusters")
        plt.xlabel("Silhouette Coefficient Values")
        plt.ylabel("Cluster Label")
        plt.axvline(x=np.mean(silhouette_vals), color="red", linestyle="--")

        df_clusters = pd.DataFrame(
            X, columns=[f"Feature {i+1}" for i in range(X.shape[1])]
        )
        df_clusters["Cluster"] = cluster_labels
        # ! pairplot of the clusters
        # Overview of the Pairplot
        # Axes and Grid:
        # The pairplot creates a grid of scatter plots for each pair of features in your dataset.
        # Each point in the scatter plots represents a sample from your dataset, colored according to its cluster assignment.

        # Diagonal Elements:
        # The diagonal plots usually show the distribution of each feature. In this case, since X.shape[1] <= 4,
        # there will be a maximum of four features plotted against each other. The diagonal could display histograms or
        # kernel density estimates (KDE) for each feature.

        # Interpretation of the Pairplot

        # Feature Relationships:
        # Look at each scatter plot in the off-diagonal plots. Each plot shows the relationship between two features. Points that
        # are close together in the scatter plot suggest similar values for those features.
        # Cluster Separation: You want to see clusters of different colors (representing different clusters) that are visually distinct.
        # Good separation indicates that the clustering algorithm effectively identified different groups within your data.
        # Overlapping Points: If points from different clusters overlap significantly in any scatter plot, it indicates that those clusters
        # might not be distinct in terms of the two features being compared.
        # Cluster Characteristics:
        # Shape and Distribution: Observe the shape of the clusters. Are they spherical, elongated, or irregular? This can give insights
        # into how well the K-Means (or other clustering methods) has performed:
        # Spherical Clusters: Indicates that clusters are well defined and separated.
        # Elongated Clusters: May suggest that the algorithm is capturing variations along specific axes but could benefit from adjustments
        # in clustering parameters or methods.
        # Feature Influence: Identify which features contribute most to cluster separation. For instance, if you see that one feature
        # consistently separates two clusters, it may be a key factor for clustering.
        # Diagonal Histograms/KDE:
        # The diagonal plots show the distribution of individual features across all samples. Look for:
        # Distribution Shape: Is the distribution unimodal, bimodal, skewed, or uniform?
        # Concentration: Areas with a high density of points may indicate that certain values are more common among samples.
        # Differences Among Clusters: If you see distinct peaks in the histograms for different clusters, it suggests that those clusters are
        # characterized by specific ranges of feature values.
        # Example Observations
        # Feature 1 vs. Feature 2: If there are clear, well-separated clusters in this scatter plot, it suggests that these two features
        # effectively distinguish between the clusters.
        # Feature 3 vs. Feature 4: If you observe significant overlap between clusters in this plot, it may indicate that these features do not
        # provide a strong basis for clustering.
        # Diagonal Plots: If you notice that one cluster has a higher density of points at lower values for a specific feature, while another
        # cluster is concentrated at higher values, this suggests that this feature is critical for differentiating those clusters.

        # Pairplot of the clusters
        # * 为什么要限制到4个features?
        # 2 features=1 scatter plot            # 3 features=3 scatter plots
        # 4 features=6 scatter plots           # 5 features=10 scatter plots
        # 6 features=15 scatter plots          # 10 features=45 scatter plots
        # Pairplot works well with low-dimensional data, 如果维度比较高的话, 子图也很多,失去了它的意义
        if X.shape[1] <= 6:
            plt.figure(figsize=(8, 4))
            sns.pairplot(df_clusters, hue="Cluster", palette="tab10")
            plt.suptitle("Pairplot of Clusters", y=1.02)

    # Add cluster labels to the DataFrame or modify in-place
    if inplace:  # replace the oringinal data
        data["Cluster"] = cluster_labels
        return None, n_clusters, kmeans, ax  # Return None when inplace is True
    else:
        data_copy = data.copy()
        data_copy["Cluster"] = cluster_labels
        return data_copy, n_clusters, kmeans, ax


# example:
# clustering_features = [marker + "_log" for marker in markers]
# df_cluster(data, columns=clustering_features, n_clusters=3,range_n_clusters=np.arange(3, 7))

"""
# You're on the right track, but let's clarify how PCA and clustering (like KMeans) work, especially 
# in the context of your dataset with 7 columns and 23,121 rows.

# Principal Component Analysis (PCA)
# Purpose of PCA:
# PCA is a dimensionality reduction technique. It transforms your dataset from a high-dimensional space 
# (in your case, 7 dimensions corresponding to your 7 columns) to a lower-dimensional space while 
# retaining as much variance (information) as possible.
# How PCA Works:
# PCA computes new features called "principal components" that are linear combinations of the original 
# features.
# The first principal component captures the most variance, the second captures the next most variance 
# (orthogonal to the first), and so on.
# If you set n_components=2, for example, PCA will reduce your dataset from 7 columns to 2 columns. 
# This helps in visualizing and analyzing the data with fewer dimensions.
# Result of PCA:
# After applying PCA, your original dataset with 7 columns will be transformed into a new dataset with 
# the specified number of components (e.g., 2 or 3).
# The transformed dataset will have fewer columns but should capture most of the important information 
# from the original dataset.

# Clustering (KMeans)
# Purpose of Clustering:
# Clustering is used to group data points based on their similarities. KMeans, specifically, partitions 
# your data into a specified number of clusters (groups).
# How KMeans Works:
# KMeans assigns each data point to one of the k clusters based on the feature space (original or 
# PCA-transformed).
# It aims to minimize the variance within each cluster while maximizing the variance between clusters.
# It does not classify the data in each column independently; instead, it considers the overall similarity 
# between data points based on their features.
# Result of KMeans:
# The output will be cluster labels for each data point (e.g., which cluster a particular observation 
# belongs to).
# You can visualize how many groups were formed and analyze the characteristics of each cluster.

# Summary
# PCA reduces the number of features (columns) in your dataset, transforming it into a lower-dimensional
# space.
# KMeans then classifies data points based on the features of the transformed dataset (or the original 
# if you choose) into different subgroups (clusters).
# By combining these techniques, you can simplify the complexity of your data and uncover patterns that 
# might not be visible in the original high-dimensional space. Let me know if you have further questions!
"""


def df_reducer(
    data: pd.DataFrame,
    columns: Optional[List[str]] = None,
    method: str = "umap",  # 'pca', 'umap'
    n_components: int = 2,  # Default for umap, but 50 for PCA
    umap_neighbors: int = 15,  # UMAP-specific
    umap_min_dist: float = 0.1,  # UMAP-specific
    tsne_perplexity: int = 30,  # t-SNE-specific
    hue: str = None,  # lda-specific
    scale: bool = True,
    fill_missing: bool = True,
    size=2,  # for plot marker size
    markerscale=4,  # for plot, legend marker size scale
    edgecolor="none",  # for plot,
    legend_loc="best",  # for plot,
    bbox_to_anchor=None,
    ncols=1,
    debug: bool = False,
    inplace: bool = True,  # replace the oringinal data
    plot_: bool = False,  # plot scatterplot, but no 'hue',so it is meaningless
    random_state=1,
    ax=None,
    figsize=None,
    verbose=True,
    **kwargs,
) -> pd.DataFrame:
    dict_methods = {
        #!Linear Dimensionality Reduction: For simplifying data with techniques that assume linearity.
        "pca": "pca(Principal Component Analysis): \n\tUseful for reducing dimensionality of continuous data while retaining variance. Advantage: Simplifies data, speeds up computation, reduces noise. Limitation: Assumes linear relationships, may lose interpretability in transformed dimensions.",
        "lda": "lda(Linear Discriminant Analysis):\n\tUseful for supervised dimensionality reduction when class separability is important. Advantage: Enhances separability between classes, can improve classification performance. Limitation: Assumes normal distribution and equal class covariances, linear boundaries only.",
        "factor": "factor(Factor Analysis):\n\tSuitable for datasets with observed and underlying latent variables. Advantage: Reveals hidden structure in correlated data, dimensionality reduction with interpretable factors. Limitation: Assumes factors are linear combinations, less effective for nonlinear data.",
        "svd": "svd(Singular Value Decomposition):\n\tSuitable for matrix decomposition, dimensionality reduction in tasks like topic modeling or image compression. Advantage: Efficient, preserves variance, useful in linear transformations. Limitation: Assumes linear relationships, sensitive to noise, may not capture non-linear structure.",
        #! Non-linear Dimensionality Reduction (Manifold Learning)
        "umap": "umap(Uniform Manifold Approximation and Projection):\n\tBest for high-dimensional data visualization (e.g., embeddings). Advantage: Captures complex structure while preserving both local and global data topology. Limitation: Non-deterministic results can vary, sensitive to parameter tuning.",
        "tsne": "tsne(t-Distributed Stochastic Neighbor Embedding):\n\tt-SNE excels at preserving local structure (i.e., clusters), but it often loses global. relationships, causing clusters to appear in arbitrary proximities to each other.  Ideal for clustering and visualizing high-dimensional data, especially for clear cluster separation. Advantage: Captures local relationships effectively. Limitation: Computationally intensive, does not preserve global structure well, requires parameter tuning.",
        "mds": "mds(Multidimensional Scaling):\n\tAppropriate for visualizing pairwise similarity or distance in data. Advantage: Maintains the perceived similarity or dissimilarity between points. Limitation: Computationally expensive for large datasets, less effective for complex, high-dimensional structures.",
        "lle": "lle(Locally Linear Embedding):\n\tUseful for non-linear dimensionality reduction when local relationships are important (e.g., manifold learning). Advantage: Preserves local data structure, good for manifold-type data. Limitation: Sensitive to noise and number of neighbors, not effective for global structure.",
        "kpca": "kpca(Kernel Principal Component Analysis):\n\tGood for non-linear data with complex structure, enhancing separability. Advantage: Extends PCA to capture non-linear relationships. Limitation: Computationally expensive, sensitive to kernel and parameter choice, less interpretable.",
        "ica": "ica(Independent Component Analysis):\n\tEffective for blind source separation (e.g., EEG, audio signal processing).is generally categorized under Non-linear Dimensionality Reduction, but it also serves a distinct role in Blind Source Separation. While ICA is commonly used for dimensionality reduction, particularly in contexts where data sources need to be disentangled (e.g., separating mixed signals like EEG or audio data), it focuses on finding statistically independent components rather than maximizing variance (like PCA) or preserving distances (like MDS or UMAP). Advantage: Extracts independent signals/components, useful in mixed signal scenarios. Limitation: Assumes statistical independence, sensitive to noise and algorithm choice.",
        #! Anomaly Detection: Specialized for detecting outliers or unusual patterns
        "isolation_forest": "Isolation Forest:\n\tDesigned for anomaly detection, especially in high-dimensional data. Advantage: Effective in detecting outliers, efficient for large datasets. Limitation: Sensitive to contamination ratio parameter, not ideal for highly structured or non-anomalous data.",
        #! more methods
        "truncated_svd": "Truncated Singular Value Decomposition (SVD):\n\tEfficient for large sparse datasets, useful for feature reduction in natural language processing (e.g., Latent Semantic Analysis). Advantage: Efficient in memory usage for large datasets. Limitation: Limited in non-linear transformation.",
        "spectral_embedding": "Spectral Embedding:\n\tBased on graph theory, it can be useful for clustering and visualization, especially for data with connected structures. Advantage: Preserves global structure, good for graph-type data. Limitation: Sensitive to parameter choice, not ideal for arbitrary non-connected data.",
        "autoencoder": "Autoencoder:\n\tA neural network-based approach for complex feature learning and non-linear dimensionality reduction. Advantage: Can capture very complex relationships. Limitation: Computationally expensive, requires neural network expertise for effective tuning.",
        "nmf": "Non-negative Matrix Factorization:\n\tEffective for parts-based decomposition, commonly used for sparse and non-negative data, e.g., text data or images. Advantage: Interpretability with non-negativity, efficient with sparse data. Limitation: Less effective for negative or zero-centered data.",
        "umap_hdbscan": "UMAP + HDBSCAN:\n\tCombination of UMAP for dimensionality reduction and HDBSCAN for density-based clustering, suitable for cluster discovery in high-dimensional data. Advantage: Effective in discovering clusters in embeddings. Limitation: Requires careful tuning of both UMAP and HDBSCAN parameters.",
        "manifold_learning": "Manifold Learning (Isomap, Hessian LLE, etc.):\n\tMethods designed to capture intrinsic geometrical structure. Advantage: Preserves non-linear relationships in low dimensions. Limitation: Computationally expensive and sensitive to noise.",
    }

    from sklearn.preprocessing import StandardScaler
    from sklearn.impute import SimpleImputer

    if plot_:
        import matplotlib.pyplot as plt
        import seaborn as sns
    # Check valid method input
    methods = [
        "pca",
        "umap",
        "umap_hdbscan",
        "tsne",
        "factor",
        "isolation_forest",
        "manifold_learning",
        "lda",
        "kpca",
        "ica",
        "mds",
        "lle",
        "svd",
        "truncated_svd",
        "spectral_embedding",
        # "autoencoder","nmf",
    ]
    method = strcmp(method, methods)[0]
    if run_once_within(reverse=True):
        print(f"support methods:{methods}")

    if verbose:
        print(f"\nprocessing with using {dict_methods[method]}:")
    xlabel, ylabel = None, None
    if columns is None:
        columns = data.select_dtypes(include="number").columns.tolist()
    if hue is None:
        hue = data.select_dtypes(exclude="number").columns.tolist()
        print(f"auto select the non-number as 'hue':{hue}")
    if isinstance(hue, list):
        print("Warning: hue is a list, only select the 1st one")
        hue = hue[0]
    if not any(hue):
        # Select columns if specified, else use all columns
        X = data[columns].values if columns else data.values
    else:
        # Select columns to reduce and hue for LDA
        try:
            X = data[columns].values if columns else data.drop(columns=[hue]).values
            y = data[hue].values
        except:
            pass
    print(X.shape)
    # Handle missing values
    if fill_missing:
        imputer = SimpleImputer(strategy="mean")
        X = imputer.fit_transform(X)

    # Optionally scale the data
    if scale:
        scaler = StandardScaler()
        X = scaler.fit_transform(X)

    # Apply PCA if selected
    if method == "pca":
        from sklearn.decomposition import PCA

        pca = PCA(n_components=n_components)
        X_reduced = pca.fit_transform(X)

        # Additional PCA information
        explained_variance = pca.explained_variance_ratio_
        singular_values = pca.singular_values_
        loadings = pca.components_.T * np.sqrt(pca.explained_variance_)

        if debug:
            print(f"PCA completed: Reduced to {n_components} components.")
            print(f"Explained Variance: {explained_variance}")
            print(f"Singular Values: {singular_values}")

        # Plot explained variance if debug=True
        if debug:
            # Plot explained variance
            cumulative_variance = np.cumsum(explained_variance)
            plt.figure(figsize=(8, 5))
            plt.plot(
                range(1, len(cumulative_variance) + 1), cumulative_variance, marker="o"
            )
            plt.title("Cumulative Explained Variance by Principal Components")
            plt.xlabel("Number of Principal Components")
            plt.ylabel("Cumulative Explained Variance")
            plt.axhline(y=0.95, color="r", linestyle="--", label="Threshold (95%)")
            plt.axvline(
                x=n_components,
                color="g",
                linestyle="--",
                label=f"n_components = {n_components}",
            )
            plt.legend()
            plt.grid()
            plt.show()

        # Prepare reduced DataFrame with additional PCA info
        pca_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"PC_{i+1}" for i in range(n_components)],
        )
        # pca_df["Explained Variance"] = np.tile(explained_variance[:n_components], (pca_df.shape[0], 1))
        # pca_df["Singular Values"] = np.tile(singular_values[:n_components], (pca_df.shape[0], 1))
        # Expand explained variance to multiple columns if needed
        for i in range(n_components):
            pca_df[f"Explained Variance PC_{i+1}"] = np.tile(
                format(explained_variance[i] * 100, ".3f") + "%", (pca_df.shape[0], 1)
            )
        for i in range(n_components):
            pca_df[f"Singular Values PC_{i+1}"] = np.tile(
                singular_values[i], (pca_df.shape[0], 1)
            )
        if hue:
            pca_df[hue] = y
    elif method == "lda":
        from sklearn.discriminant_analysis import LinearDiscriminantAnalysis

        if "hue" not in locals() or hue is None:
            raise ValueError(
                "LDA requires a 'hue' col parameter to specify class labels."
            )

        lda_reducer = LinearDiscriminantAnalysis(n_components=n_components)
        X_reduced = lda_reducer.fit_transform(X, y)

        # Prepare reduced DataFrame with additional LDA info
        lda_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"LDA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print(f"LDA completed: Reduced to {n_components} components.")
            print("Class separability achieved by LDA.")
        if hue:
            lda_df[hue] = y
    # Apply UMAP if selected
    elif method == "umap":
        import umap

        umap_reducer = umap.UMAP(
            n_neighbors=umap_neighbors,
            min_dist=umap_min_dist,
            n_components=n_components,
        )
        X_reduced = umap_reducer.fit_transform(X)

        # Additional UMAP information
        embedding = umap_reducer.embedding_
        trustworthiness = umap_reducer._raw_data[:, :n_components]

        if debug:
            print(f"UMAP completed: Reduced to {n_components} components.")
            print(f"Embedding Shape: {embedding.shape}")
            print(f"Trustworthiness: {trustworthiness}")

        # Prepare reduced DataFrame with additional UMAP info
        umap_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"UMAP_{i+1}" for i in range(n_components)],
        )
        umap_df["Embedding"] = embedding[:, 0]  # Example of embedding data
        umap_df["Trustworthiness"] = trustworthiness[:, 0]  # Trustworthiness metric
        if hue:
            umap_df[hue] = y
    elif method == "tsne":
        from sklearn.manifold import TSNE

        tsne = TSNE(
            n_components=n_components,
            perplexity=tsne_perplexity,
            random_state=random_state,
        )
        X_reduced = tsne.fit_transform(X)
        tsne_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"tSNE_{i+1}" for i in range(n_components)],
        )
        tsne_df["Perplexity"] = np.tile(
            f"Perplexity: {tsne_perplexity}", (tsne_df.shape[0], 1)
        )
        if hue:
            tsne_df[hue] = y
    # Apply Factor Analysis if selected
    elif method == "factor":
        from sklearn.decomposition import FactorAnalysis

        factor = FactorAnalysis(n_components=n_components, random_state=random_state)
        X_reduced = factor.fit_transform(X)
        # Factor Analysis does not directly provide explained variance, but we can approximate it
        fa_variance = factor.noise_variance_
        # Prepare reduced DataFrame with additional Factor Analysis info
        factor_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"Factor_{i+1}" for i in range(n_components)],
        )
        factor_df["Noise Variance"] = np.tile(
            format(np.mean(fa_variance) * 100, ".3f") + "%", (factor_df.shape[0], 1)
        )
        if hue:
            factor_df[hue] = y
    # Apply Isolation Forest for outlier detection if selected
    elif method == "isolation_forest":
        from sklearn.decomposition import PCA
        from sklearn.ensemble import IsolationForest

        # Step 1: Apply PCA for dimensionality reduction to 2 components
        pca = PCA(n_components=n_components)
        X_pca = pca.fit_transform(X)

        explained_variance = pca.explained_variance_ratio_
        singular_values = pca.singular_values_

        # Prepare reduced DataFrame with additional PCA info
        iso_forest_df = pd.DataFrame(
            X_pca, index=data.index, columns=[f"PC_{i+1}" for i in range(n_components)]
        )

        isolation_forest = IsolationForest(
            n_estimators=100, contamination="auto", random_state=1
        )
        isolation_forest.fit(X)
        anomaly_scores = isolation_forest.decision_function(
            X
        )  # Anomaly score: larger is less anomalous
        # Predict labels: 1 (normal), -1 (anomaly)
        anomaly_labels = isolation_forest.fit_predict(X)
        # Add anomaly scores and labels to the DataFrame
        iso_forest_df["Anomaly Score"] = anomaly_scores
        iso_forest_df["Anomaly Label"] = anomaly_labels
        # add info from pca
        for i in range(n_components):
            iso_forest_df[f"Explained Variance PC_{i+1}"] = np.tile(
                format(explained_variance[i] * 100, ".3f") + "%",
                (iso_forest_df.shape[0], 1),
            )
        for i in range(n_components):
            iso_forest_df[f"Singular Values PC_{i+1}"] = np.tile(
                singular_values[i], (iso_forest_df.shape[0], 1)
            )
        if hue:
            iso_forest_df[hue] = y
    # * Apply Kernel PCA if selected
    elif method == "kpca":
        from sklearn.decomposition import KernelPCA

        kpca = KernelPCA(
            n_components=n_components, kernel="rbf", random_state=random_state
        )
        X_reduced = kpca.fit_transform(X)

        # Prepare reduced DataFrame with KPCA info
        kpca_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"KPCA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Kernel PCA completed with RBF kernel.")
        if hue:
            kpca_df[hue] = y
    # * Apply ICA if selected
    elif method == "ica":
        from sklearn.decomposition import FastICA

        ica = FastICA(n_components=n_components, random_state=random_state)
        X_reduced = ica.fit_transform(X)

        # Prepare reduced DataFrame with ICA info
        ica_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"ICA_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Independent Component Analysis (ICA) completed.")
        if hue:
            ica_df[hue] = y
    # * Apply MDS if selected
    elif method == "mds":
        from sklearn.manifold import MDS

        mds = MDS(n_components=n_components, random_state=random_state)
        X_reduced = mds.fit_transform(X)

        # Prepare reduced DataFrame with MDS info
        mds_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"MDS_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Multidimensional Scaling (MDS) completed.")
        if hue:
            mds_df[hue] = y
    # * Apply Locally Linear Embedding (LLE) if selected
    elif method == "lle":
        from sklearn.manifold import LocallyLinearEmbedding

        lle = LocallyLinearEmbedding(
            n_components=n_components,
            n_neighbors=umap_neighbors,
            random_state=random_state,
        )
        X_reduced = lle.fit_transform(X)

        # Prepare reduced DataFrame with LLE info
        lle_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"LLE_{i+1}" for i in range(n_components)],
        )
        if debug:
            print("Locally Linear Embedding (LLE) completed.")
        if hue:
            lle_df[hue] = y
    # * Apply Singular Value Decomposition (SVD) if selected
    elif method == "svd":
        # Using NumPy's SVD for dimensionality reduction
        U, s, Vt = np.linalg.svd(X, full_matrices=False)
        X_reduced = U[:, :n_components] * s[:n_components]

        # Prepare reduced DataFrame with SVD info
        svd_df = pd.DataFrame(
            X_reduced,
            index=data.index,
            columns=[f"SVD_{i+1}" for i in range(n_components)],
        )
        colname_met = "SVD_"
        if hue:
            svd_df[hue] = y
        if debug:
            print("Singular Value Decomposition (SVD) completed.")
    elif method == "truncated_svd":
        from sklearn.decomposition import TruncatedSVD

        svd = TruncatedSVD(n_components=n_components, random_state=random_state)
        X_reduced = svd.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"SVD Component {i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "SVD Component "

        if debug:
            print("Truncated SVD completed.")
            print("Explained Variance Ratio:", svd.explained_variance_ratio_)
        if hue:
            reduced_df[hue] = y

    elif method == "spectral_embedding":
        from sklearn.manifold import SpectralEmbedding

        spectral = SpectralEmbedding(
            n_components=n_components, random_state=random_state
        )
        X_reduced = spectral.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Dimension_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Dimension_"

        if debug:
            print("Spectral Embedding completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "autoencoder":
        from tensorflow.keras.models import Model
        from tensorflow.keras.layers import Input, Dense

        input_dim = X.shape[1]
        input_layer = Input(shape=(input_dim,))
        encoded = Dense(n_components * 2, activation="relu")(input_layer)
        encoded = Dense(n_components, activation="relu")(encoded)
        autoencoder = Model(input_layer, encoded)
        autoencoder.compile(optimizer="adam", loss="mean_squared_error")
        autoencoder.fit(X, X, epochs=50, batch_size=256, shuffle=True, verbose=0)

        X_reduced = autoencoder.predict(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Score_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Score_"

        if debug:
            print("Autoencoder reduction completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "nmf":
        from sklearn.decomposition import NMF

        nmf = NMF(n_components=n_components, random_state=random_state)
        X_reduced = nmf.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"NMF_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "NMF_"

        if debug:
            print("Non-negative Matrix Factorization completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "umap_hdbscan":
        import umap
        import hdbscan

        umap_model = umap.UMAP(
            n_neighbors=umap_neighbors,
            min_dist=umap_min_dist,
            n_components=n_components,
        )
        X_umap = umap_model.fit_transform(X)

        clusterer = hdbscan.HDBSCAN()
        clusters = clusterer.fit_predict(X_umap)

        reduced_df = pd.DataFrame(
            X_umap,
            columns=[f"UMAP_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        reduced_df["Cluster"] = clusters
        colname_met = "UMAP_"
        if debug:
            print("UMAP + HDBSCAN reduction and clustering completed.")
        if hue:
            reduced_df[hue] = y

    elif method == "manifold_learning":
        from sklearn.manifold import Isomap

        isomap = Isomap(n_components=n_components)
        X_reduced = isomap.fit_transform(X)
        reduced_df = pd.DataFrame(
            X_reduced,
            columns=[f"Manifold_{i+1}" for i in range(n_components)],
            index=data.index,
        )
        colname_met = "Manifold_"

        if debug:
            print("Manifold Learning (Isomap) completed.")
        if hue:
            reduced_df[hue] = y

    #! Return reduced data and info as a new DataFrame with the same index
    if method == "pca":
        reduced_df = pca_df
        colname_met = "PC_"
        xlabel = f"PC_1 ({pca_df['Explained Variance PC_1'].tolist()[0]})"
        ylabel = f"PC_2 ({pca_df['Explained Variance PC_2'].tolist()[0]})"
    elif method == "umap":
        reduced_df = umap_df
        colname_met = "UMAP_"
    elif method == "tsne":
        reduced_df = tsne_df
        colname_met = "tSNE_"
    elif method == "factor":
        reduced_df = factor_df
        colname_met = "Factor_"
    elif method == "isolation_forest":
        reduced_df = iso_forest_df  # Already a DataFrame for outliers
        colname_met = "PC_"
        if plot_:
            ax = sns.scatterplot(
                data=iso_forest_df[iso_forest_df["Anomaly Label"] == 1],
                x="PC_1",
                y="PC_2",
                label="normal",
                c="b",
            )
            ax = sns.scatterplot(
                ax=ax,
                data=iso_forest_df[iso_forest_df["Anomaly Label"] == -1],
                x="PC_1",
                y="PC_2",
                c="r",
                label="outlier",
                marker="+",
                s=30,
            )
    elif method == "lda":
        reduced_df = lda_df
        colname_met = "LDA_"
    elif method == "kpca":
        reduced_df = kpca_df
        colname_met = "KPCA_"
    elif method == "ica":
        reduced_df = ica_df
        colname_met = "ICA_"
    elif method == "mds":
        reduced_df = mds_df
        colname_met = "MDS_"
    elif method == "lle":
        reduced_df = lle_df
        colname_met = "LLE_"
    elif method == "svd":
        reduced_df = svd_df
        colname_met = "SVD_"
    # Quick plots
    if plot_ and (not method in ["isolation_forest"]):
        from .plot import plotxy, figsets, get_color

        # if ax is None:
        #     if figsize is None:
        #         _, ax = plt.subplots(figsize=cm2inch(8, 8))
        #     else:
        #         _, ax = plt.subplots(figsize=figsize)
        # else:
        #     ax = ax.cla()
        xlabel = f"{colname_met}1" if xlabel is None else xlabel
        ylabel = f"{colname_met}2" if ylabel is None else ylabel
        palette = get_color(len(flatten(data[hue], verbose=0)))

        reduced_df = reduced_df.sort_values(by=hue)
        print(flatten(reduced_df[hue]))
        ax = plotxy(
            data=reduced_df,
            x=colname_met + "1",
            y=colname_met + "2",
            hue=hue,
            palette=palette,
            # size=size,
            edgecolor=edgecolor,
            kind_=[
                "joint",
                #    "kde",
                "ell",
            ],
            kws_kde=dict(
                hue=hue,
                levels=2,
                common_norm=False,
                fill=True,
                alpha=0.05,
            ),
            kws_joint=dict(kind="scatter", joint_kws=dict(s=size)),
            kws_ellipse=dict(alpha=0.1, lw=1, label=None),
            verbose=False,
            **kwargs,
        )
        figsets(
            legend=dict(
                loc=legend_loc,
                markerscale=markerscale,
                bbox_to_anchor=bbox_to_anchor,
                ncols=ncols,
                fontsize=8,
            ),
            xlabel=xlabel if xlabel else None,
            ylabel=ylabel if ylabel else None,
        )

    if inplace:
        # If inplace=True, add components back into the original data
        for col_idx in range(n_components):
            data.loc[:, f"{colname_met}{col_idx+1}"] = reduced_df.iloc[:, col_idx]
        # Add extra info for PCA/UMAP
        if method == "pca":
            for i in range(n_components):
                data.loc[:, f"Explained Variance PC_{i+1}"] = reduced_df.loc[
                    :, f"Explained Variance PC_{i+1}"
                ]
            for i in range(n_components):
                data.loc[:, f"Singular Values PC_{i+1}"] = reduced_df.loc[
                    :, f"Singular Values PC_{i+1}"
                ]
        elif method == "umap":
            for i in range(n_components):
                data.loc[:, f"UMAP_{i+1}"] = reduced_df.loc[:, f"UMAP_{i+1}"]
            data.loc[:, "Embedding"] = reduced_df.loc[:, "Embedding"]
            data.loc[:, "Trustworthiness"] = reduced_df.loc[:, "Trustworthiness"]

        return None  # No return when inplace=True

    return reduced_df


# example:
# df_reducer(data=data_log, columns=markers, n_components=2)
 


def get_df_format(data, threshold_unique=0.5, verbose=False, sample_size=1000):
    """
    Detect whether a DataFrame is in long or wide format with optimized performance and accuracy.
    
    Parameters:
    - data (pd.DataFrame): DataFrame to analyze
    - threshold_unique (float): Threshold for categorical column detection (0-1)
    - verbose (bool): Whether to print diagnostic messages
    - sample_size (int): Maximum number of rows/columns to sample for large datasets
    
    Returns:
    - "long" if detected as long format
    - "wide" if detected as wide format
    - "uncertain" if format is ambiguous
    """
    import pandas as pd
    import numpy as np
    from scipy.stats import entropy
    from sklearn.cluster import AgglomerativeClustering
    from sklearn.preprocessing import StandardScaler
    from sklearn.metrics import pairwise_distances
    from collections import Counter
    import re
    # ----- Initial Setup and Sampling -----
    n_rows, n_cols = data.shape
    if verbose:
        print(f"Initial shape: {n_rows} rows, {n_cols} columns")
    
    # Sample data if too large
    if n_rows > sample_size:
        data = data.sample(n=sample_size, random_state=42)
        n_rows = sample_size
    if n_cols > sample_size:
        data = data.iloc[:, :sample_size]
        n_cols = sample_size
    
    # Early exit for tiny datasets
    if n_rows < 3 or n_cols < 3:
        return "uncertain"
    
    long_score = 0
    wide_score = 0
    
    # ----- Feature Extraction -----
    # Basic statistics
    row_col_ratio = n_rows / n_cols if n_cols != 0 else float('inf')
    
    # Column types
    numeric_cols = data.select_dtypes(include=np.number).columns
    cat_cols = data.select_dtypes(include=['object', 'category']).columns
    other_cols = [col for col in data.columns if col not in numeric_cols and col not in cat_cols]
    
    # Unique value analysis
    unique_counts = data.nunique(dropna=False)
    duplicate_ratio = 1 - unique_counts / n_rows
    
    # Missing values
    missing_per_row = data.isna().sum(axis=1)
    missing_per_col = data.isna().sum()
    
    # Column name patterns
    col_names = data.columns.astype(str)
    has_suffix = sum(bool(re.search(r'(_\d+|\d+_?$)', col)) for col in col_names)
    has_time = sum(bool(re.search(r'(^time|^date|^year|^month|^day|^t\d+)', col.lower())) for col in col_names)
    
    # ----- Scoring Rules -----
    
    # 1. Row-Column Ratio (weighted)
    if row_col_ratio > 5:
        long_score += 3
        if verbose: print(f"High row/col ratio ({row_col_ratio:.1f}) → long +3")
    elif row_col_ratio < 0.2:
        wide_score += 3
        if verbose: print(f"Low row/col ratio ({row_col_ratio:.1f}) → wide +3")
    elif row_col_ratio > 2:
        long_score += 1
        if verbose: print(f"Moderate row/col ratio ({row_col_ratio:.1f}) → long +1")
    elif row_col_ratio < 0.5:
        wide_score += 1
        if verbose: print(f"Moderate row/col ratio ({row_col_ratio:.1f}) → wide +1")
    
    # 2. Duplication Patterns
    high_dupe_cols = sum(duplicate_ratio > 0.3)
    if high_dupe_cols > 0.6 * n_cols:
        wide_score += 2
        if verbose: print(f"Many columns ({high_dupe_cols}/{n_cols}) with duplicates → wide +2")
    elif high_dupe_cols < 0.2 * n_cols:
        long_score += 1
        if verbose: print(f"Few columns ({high_dupe_cols}/{n_cols}) with duplicates → long +1")
    
    # 3. Categorical Column Analysis
    if len(cat_cols) > 0:
        # Entropy analysis
        cat_entropies = []
        for col in cat_cols:
            counts = data[col].value_counts(normalize=True, dropna=False)
            cat_entropies.append(entropy(counts))
        
        avg_cat_entropy = np.mean(cat_entropies) if cat_entropies else 0
        if avg_cat_entropy < 1.2:
            long_score += 2
            if verbose: print(f"Low categorical entropy ({avg_cat_entropy:.2f}) → long +2")
        elif avg_cat_entropy > 2:
            wide_score += 1
            if verbose: print(f"High categorical entropy ({avg_cat_entropy:.2f}) → wide +1")
        
        # Entity identifier detection
        if len(cat_cols) >= 2 and n_rows > 10:
            dup_rows = data.duplicated(subset=cat_cols.tolist()[:2], keep=False).sum()
            if dup_rows > 0.3 * n_rows:
                long_score += 2
                if verbose: print(f"Duplicate rows in categorical cols ({dup_rows}/{n_rows}) → long +2")
    
    # 4. Column Name Patterns
    if has_suffix > 0.4 * n_cols:
        wide_score += 2
        if verbose: print(f"Many suffix patterns ({has_suffix}/{n_cols}) → wide +2")
    if has_time > 0.3 * n_cols:
        wide_score += 1
        if verbose: print(f"Time-like columns ({has_time}/{n_cols}) → wide +1")
    
    # 5. Numeric Column Analysis (only if enough numeric columns)
    if len(numeric_cols) > 2:
        # Correlation analysis
        corr_matrix = data[numeric_cols].corr().abs()
        avg_corr = corr_matrix.values[np.triu_indices_from(corr_matrix, k=1)].mean()
        
        if avg_corr > 0.5:
            wide_score += 2
            if verbose: print(f"High numeric correlation ({avg_corr:.2f}) → wide +2")
        elif avg_corr < 0.2:
            long_score += 1
            if verbose: print(f"Low numeric correlation ({avg_corr:.2f}) → long +1")
        
        # Entropy analysis
        try:
            numeric_data = data[numeric_cols].dropna()
            if len(numeric_data) > 10:
                numeric_entropy = numeric_data.apply(lambda x: entropy(pd.cut(x, bins=min(10, len(x.unique())).value_counts(normalize=True))))
                if numeric_entropy.mean() < 1.5:
                    wide_score += 1
                    if verbose: print(f"Low numeric entropy ({numeric_entropy.mean():.2f}) → wide +1")
        except Exception as e:
            if verbose: print(f"Numeric entropy failed: {str(e)}")
    
    # 6. Missing Value Patterns
    missing_row_std = missing_per_row.std()
    if missing_row_std < 1 and missing_per_row.mean() > 0.1 * n_cols:
        wide_score += 1
        if verbose: print(f"Uniform missing pattern (std={missing_row_std:.2f}) → wide +1")
    elif missing_per_row.mean() < 0.05 * n_cols:
        long_score += 1
        if verbose: print(f"Few missing values → long +1")
    
    # 7. Advanced Clustering (only for medium/large datasets)
    if len(numeric_cols) > 3 and n_rows > 10 and n_cols > 5:
        try:
            # Efficient clustering with sampling
            sample_data = data[numeric_cols].sample(n=min(100, n_rows), random_state=42)
            scaled_data = StandardScaler().fit_transform(sample_data.dropna())
            
            if scaled_data.shape[0] > 5:
                # Column clustering
                col_dist = pairwise_distances(scaled_data.T)
                col_clusters = AgglomerativeClustering(n_clusters=2, 
                                                     affinity='precomputed', 
                                                     linkage='complete').fit(col_dist)
                cluster_counts = Counter(col_clusters.labels_)
                if max(cluster_counts.values()) > 0.7 * len(numeric_cols):
                    wide_score += 2
                    if verbose: print(f"Column clustering shows dominant group → wide +2")
                
                # Row clustering
                row_clusters = AgglomerativeClustering(n_clusters=2).fit(scaled_data)
                row_cluster_counts = Counter(row_clusters.labels_)
                if max(row_cluster_counts.values()) > 0.8 * scaled_data.shape[0]:
                    wide_score += 1
                    if verbose: print(f"Row clustering shows homogeneity → wide +1")
        except Exception as e:
            if verbose: print(f"Clustering skipped: {str(e)}")
    
    # ----- Decision Logic -----
    score_diff = long_score - wide_score
    abs_diff = abs(score_diff)
    
    if verbose:
        print(f"\nFinal scores - Long: {long_score}, Wide: {wide_score}")
    
    if abs_diff >= 3:
        return "long" if score_diff > 0 else "wide"
    elif abs_diff >= 1:
        # Additional tie-breakers
        if score_diff == 0:
            if row_col_ratio > 1.5:
                return "long"
            elif row_col_ratio < 0.67:
                return "wide"
            elif len(cat_cols) > len(numeric_cols):
                return "long"
            else:
                return "wide"
        return "long" if score_diff > 0 else "wide"
    else:
        return "uncertain"
#! ========== workbook, worksheet, wb,ws =============

import openpyxl
from openpyxl.utils import get_column_letter
from openpyxl.styles import NamedStyle
from openpyxl.worksheet.datavalidation import DataValidation
import re
# import copy


def ws_get_header(ws, column=None, header=1, return_idx=True):
    header_list = [cell.value for cell in ws[header]]
    if column is None:  # return all values
        return (
            [(i, idx) for i, idx in zip(header_list, range(1, len(header_list) + 1))]
            if return_idx
            else header_list
        )
    if isinstance(column, list):
        return [
            ws_get_header(ws, column=col, header=header, return_idx=return_idx)
            for col in column
        ]
    return (
        strcmp(column, header_list)[0],
        header_list.index((strcmp(column, header_list)[0])) + 1,
    )


def ws_insert_column(ws, column, after=None, before=None, header=None, header_style=None):
    """
    Insert a new column while preserving formatting and data validation

    Args:
        ws: Worksheet object
        column: Name for the new column header
        after: Insert after this column name
        before: Insert before this column name
        header: Header text (defaults to column)
        header_style: Style to apply to header cell
    """
    if after is not None and before is not None:
        raise ValueError("Cannot specify both 'after' and 'before'")

    if after is not None:
        _, col_idx = ws_get_header(ws, after, header=header)
        insert_idx = col_idx + 1
    elif before is not None:
        _, col_idx = ws_get_header(ws, before, header=header)
        insert_idx = col_idx
    else:
        raise ValueError("Must specify either 'after' or 'before'")

    # Insert column
    ws.insert_cols(insert_idx)

    # Set header
    header_text = header if header is not None else column
    header_cell = ws.cell(row=header, column=insert_idx, value=header_text)

    # Apply header style if provided
    if header_style:
        if isinstance(header_style, str):
            header_style = ws.parent.named_styles[header_style]
        header_cell.style = header_style

    return insert_idx


def ws_write_column(ws, column, values, start_row=2, 
                    # preserve_formatting=True, 
                    header=1):
    """
    Write values to a column while optionally preserving formatting

    Args:
        ws: Worksheet object
        column: Name of column to write to
        values: List of values to write
        start_row: Starting row (default 2, assuming row 1 is header)
        preserve_formatting: Whether to keep existing formatting (default True)
    """
    _, col_idx = ws_get_header(ws, column, header=header)

    # Get reference column for formatting (previous column)
    ref_col_idx = col_idx - 1 if col_idx > 1 else col_idx + 1
    ref_col = ws[get_column_letter(ref_col_idx)]

    for row_idx, value in enumerate(values, start_row):
        cell = ws.cell(row=row_idx, column=col_idx, value=value)

        # if preserve_formatting and row_idx <= ws.max_row:
        #     # Copy formatting from reference cell in same row
        #     ref_cell = ws.cell(row=row_idx, column=ref_col_idx)
        #     if ref_cell.has_style:
        #         cell.font = copy.copy(ref_cell.font)
        #         cell.border = copy.copy(ref_cell.border)
        #         cell.fill = copy.copy(ref_cell.fill)
        #         cell.number_format = ref_cell.number_format
        #         cell.protection = copy.copy(ref_cell.protection)
        #         cell.alignment = copy.copy(ref_cell.alignment)


def ws_delete_column(ws, column, header=1):
    """
    Delete a column by name while preserving data validation rules

    Args:
        ws: Worksheet object
        column: Name of column to delete
    """
    _, col_idx = ws_get_header(ws, column, header=header)
    ws.delete_cols(col_idx)


# def copy_column_formatting(source_col, target_col):
#     """Copy formatting from one column to another"""
#     for source_cell, target_cell in zip(source_col, target_col):
#         if source_cell.has_style:
#             target_cell.font = copy.copy(source_cell.font)
#             target_cell.border = copy.copy(source_cell.border)
#             target_cell.fill = copy.copy(source_cell.fill)
#             target_cell.number_format = source_cell.number_format
#             target_cell.protection = copy.copy(source_cell.protection)
#             target_cell.alignment = copy.copy(source_cell.alignment)


# def copy_data_validation(source_ws, target_ws, column, header=1):
#     """
#     Copy data validation rules from one worksheet to another for a specific column

#     Args:
#         source_ws: Source worksheet
#         target_ws: Target worksheet
#         column: Column name to copy validation for
#     """
#     _, source_col_idx = ws_get_header(source_ws, column, header=header)
#     _, target_col_idx = ws_get_header(target_ws, column, header=header)

#     # Get data validation from source column
#     source_col = source_ws[get_column_letter(source_col_idx)]
#     for cell in source_col:
#         if cell.data_validation:
#             # Apply same validation to target column
#             target_cell = target_ws.cell(row=cell.row, column=target_col_idx)
#             target_cell.data_validation = copy.copy(cell.data_validation)


def apply_number_format(ws, column, number_format, header=1):
    """
    Apply number formatting to an entire column

    Args:
        ws: Worksheet object
        column: Column name to format
        number_format: Excel number format string (e.g., "0.00%")
    """
    _, col_idx = ws_get_header(ws, column, header=header)
    col_letter = get_column_letter(col_idx)

    for row in ws.iter_rows(min_col=col_idx, max_col=col_idx):
        for cell in row:
            cell.number_format = number_format


def ws_auto_width(ws):
    """Adjust column widths based on content"""
    for column in ws.columns:
        max_length = 0
        column_letter = get_column_letter(column[0].column)

        for cell in column:
            try:
                if len(str(cell.value)) > max_length:
                    max_length = len(str(cell.value))
            except:
                pass

        adjusted_width = (max_length + 2) * 1.2
        ws.column_dimensions[column_letter].width = adjusted_width


def create_style(wb, style_name, **style_attrs):
    """
    Create a named style in the workbook

    Args:
        wb: Workbook object
        style_name: Name for the new style
        **style_attrs: Style attributes (font, fill, border, etc.)
    """
    if style_name in wb.named_styles:
        return wb.named_styles[style_name]

    style = NamedStyle(name=style_name)
    for attr, value in style_attrs.items():
        setattr(style, attr, value)

    wb.add_named_style(style)
    return style


def create_data_validation(validation_type, formula1, formula2=None, **kwargs):
    """
    Create a data validation rule

    Args:
        validation_type: Excel validation type (e.g., "list", "whole", "decimal")
        formula1: First formula (e.g., list range for dropdown)
        formula2: Second formula (for between validations)
        **kwargs: Additional validation options (errorTitle, errorMessage, etc.)
    """
    dv = DataValidation(type=validation_type, formula1=formula1, formula2=formula2)
    for key, value in kwargs.items():
        setattr(dv, key, value)
    return dv

#! ========== workbook, worksheet, wb,ws =============

def plot_cluster(
    data: pd.DataFrame,
    labels: np.ndarray,
    metrics: dict = None,
    cmap="tab20",
    true_labels: Optional[np.ndarray] = None,
) -> None:
    """
    Visualize clustering results with various plots.

    Parameters:
    -----------
    data : pd.DataFrame
        The input data used for clustering.
    labels : np.ndarray
        Cluster labels assigned to each point.
    metrics : dict
        Dictionary containing evaluation metrics from evaluate_cluster function.
    true_labels : Optional[np.ndarray], default=None
        Ground truth labels, if available.
    """
    import seaborn as sns
    from sklearn.metrics import silhouette_samples
    import matplotlib.pyplot as plt

    if metrics is None:
        metrics = evaluate_cluster(data=data, labels=labels, true_labels=true_labels)

    # 1. Scatter Plot of Clusters
    plt.figure(figsize=(15, 6))
    plt.subplot(1, 3, 1)
    plt.scatter(data.iloc[:, 0], data.iloc[:, 1], c=labels, cmap=cmap, s=20)
    plt.title("Cluster Scatter Plot")
    plt.xlabel("Component 1")
    plt.ylabel("Component 2")
    plt.colorbar(label="Cluster Label")
    plt.grid()

    # 2. Silhouette Plot
    if "Silhouette Score" in metrics:
        silhouette_vals = silhouette_samples(data, labels)
        plt.subplot(1, 3, 2)
        y_lower = 10
        for i in range(len(set(labels))):
            # Aggregate the silhouette scores for samples belonging to the current cluster
            cluster_silhouette_vals = silhouette_vals[labels == i]
            cluster_silhouette_vals.sort()
            size_cluster_i = cluster_silhouette_vals.shape[0]
            y_upper = y_lower + size_cluster_i

            plt.fill_betweenx(np.arange(y_lower, y_upper), 0, cluster_silhouette_vals)
            plt.text(-0.05, y_lower + 0.5 * size_cluster_i, str(i))
            y_lower = y_upper + 10  # 10 for the 0 samples

        plt.title("Silhouette Plot")
        plt.xlabel("Silhouette Coefficient Values")
        plt.ylabel("Cluster Label")
        plt.axvline(x=metrics["Silhouette Score"], color="red", linestyle="--")
        plt.grid()

    # 3. Metrics Plot
    plt.subplot(1, 3, 3)
    metric_names = ["Davies-Bouldin Index", "Calinski-Harabasz Index"]
    metric_values = [
        metrics["Davies-Bouldin Index"],
        metrics["Calinski-Harabasz Index"],
    ]

    if true_labels is not None:
        metric_names += ["Homogeneity Score", "Completeness Score", "V-Measure"]
        metric_values += [
            metrics["Homogeneity Score"],
            metrics["Completeness Score"],
            metrics["V-Measure"],
        ]

    plt.barh(metric_names, metric_values, color="lightblue")
    plt.title("Clustering Metrics")
    plt.xlabel("Score")
    plt.axvline(x=0, color="gray", linestyle="--")
    plt.grid()
    plt.tight_layout()


def evaluate_cluster(
    data: pd.DataFrame, labels: np.ndarray, true_labels: Optional[np.ndarray] = None
) -> dict:
    """
    Evaluate clustering performance using various metrics.

    Parameters:
    -----------
    data : pd.DataFrame
        The input data used for clustering.
    labels : np.ndarray
        Cluster labels assigned to each point.
    true_labels : Optional[np.ndarray], default=None
        Ground truth labels, if available.

    Returns:
    --------
    metrics : dict
        Dictionary containing evaluation metrics.

    1. Silhouette Score:
    The silhouette score measures how similar an object is to its own cluster (cohesion) compared to
    how similar it is to other clusters (separation). The score ranges from -1 to +1:
    +1: Indicates that the data point is very far from the neighboring clusters and well clustered.
    0: Indicates that the data point is on or very close to the decision boundary between two neighboring
    clusters.
    -1: Indicates that the data point might have been assigned to the wrong cluster.

    Interpretation:
    A higher average silhouette score indicates better-defined clusters.
    If the score is consistently high (above 0.5), it suggests that the clusters are well separated.
    A score near 0 may indicate overlapping clusters, while negative scores suggest points may have
    been misclassified.

    2. Davies-Bouldin Index:
    The Davies-Bouldin Index (DBI) measures the average similarity ratio of each cluster with its
    most similar cluster. The index values range from 0 to ∞, with lower values indicating better clustering.
    It is defined as the ratio of within-cluster distances to between-cluster distances.

    Interpretation:
    A lower DBI value indicates that the clusters are compact and well-separated.
    Ideally, you want to minimize the Davies-Bouldin Index. If your DBI value is above 1, this indicates
    that your clusters might not be well-separated.

    3. Adjusted Rand Index (ARI):
    The Adjusted Rand Index (ARI) is a measure of the similarity between two data clusterings. The ARI
    score ranges from -1 to +1:
    1: Indicates perfect agreement between the two clusterings.
    0: Indicates that the clusterings are no better than random.
    Negative values: Indicate less agreement than expected by chance.

    Interpretation:
    A higher ARI score indicates better clustering, particularly if it's close to 1.
    An ARI score of 0 or lower suggests that the clustering results do not represent the true labels
    well, indicating a poor clustering performance.

    4. Calinski-Harabasz Index:
    The Calinski-Harabasz Index (also known as the Variance Ratio Criterion) evaluates the ratio of
    the sum of between-cluster dispersion to within-cluster dispersion. Higher values indicate better clustering.

    Interpretation:
    A higher Calinski-Harabasz Index suggests that clusters are dense and well-separated. It is typically
    used to validate the number of clusters, with higher values favoring more distinct clusters.

    5. Homogeneity Score:
    The homogeneity score measures how much a cluster contains only members of a single class (if true labels are provided).
    A score of 1 indicates perfect homogeneity, where all points in a cluster belong to the same class.

    Interpretation:
    A higher homogeneity score indicates that the clustering result is pure, meaning the clusters are composed
    of similar members. Lower values indicate mixed clusters, suggesting poor clustering performance.

    6. Completeness Score:
    The completeness score evaluates how well all members of a given class are assigned to the same cluster.
    A score of 1 indicates perfect completeness, meaning all points in a true class are assigned to a single cluster.

    Interpretation:
    A higher completeness score indicates that the clustering effectively groups all instances of a class together.
    Lower values suggest that some instances of a class are dispersed among multiple clusters.

    7. V-Measure:
    The V-measure is the harmonic mean of homogeneity and completeness, giving a balanced measure of clustering performance.

    Interpretation:
    A higher V-measure score indicates that the clusters are both homogenous (pure) and complete (cover all members of a class).
    Scores closer to 1 indicate better clustering quality.
    """
    from sklearn.metrics import (
        silhouette_score,
        davies_bouldin_score,
        adjusted_rand_score,
        calinski_harabasz_score,
        homogeneity_score,
        completeness_score,
        v_measure_score,
    )

    metrics = {}
    unique_labels = set(labels)
    if len(unique_labels) > 1 and len(unique_labels) < len(data):
        # Calculate Silhouette Score
        try:
            metrics["Silhouette Score"] = silhouette_score(data, labels)
        except Exception as e:
            metrics["Silhouette Score"] = np.nan
            print(f"Silhouette Score calculation failed: {e}")

        # Calculate Davies-Bouldin Index
        try:
            metrics["Davies-Bouldin Index"] = davies_bouldin_score(data, labels)
        except Exception as e:
            metrics["Davies-Bouldin Index"] = np.nan
            print(f"Davies-Bouldin Index calculation failed: {e}")

        # Calculate Calinski-Harabasz Index
        try:
            metrics["Calinski-Harabasz Index"] = calinski_harabasz_score(data, labels)
        except Exception as e:
            metrics["Calinski-Harabasz Index"] = np.nan
            print(f"Calinski-Harabasz Index calculation failed: {e}")

        # Calculate Adjusted Rand Index if true labels are provided
        if true_labels is not None:
            try:
                metrics["Adjusted Rand Index"] = adjusted_rand_score(
                    true_labels, labels
                )
            except Exception as e:
                metrics["Adjusted Rand Index"] = np.nan
                print(f"Adjusted Rand Index calculation failed: {e}")

            # Calculate Homogeneity Score
            try:
                metrics["Homogeneity Score"] = homogeneity_score(true_labels, labels)
            except Exception as e:
                metrics["Homogeneity Score"] = np.nan
                print(f"Homogeneity Score calculation failed: {e}")

            # Calculate Completeness Score
            try:
                metrics["Completeness Score"] = completeness_score(true_labels, labels)
            except Exception as e:
                metrics["Completeness Score"] = np.nan
                print(f"Completeness Score calculation failed: {e}")

            # Calculate V-Measure
            try:
                metrics["V-Measure"] = v_measure_score(true_labels, labels)
            except Exception as e:
                metrics["V-Measure"] = np.nan
                print(f"V-Measure calculation failed: {e}")
    else:
        # Metrics cannot be computed with 1 cluster or all points as noise
        metrics["Silhouette Score"] = np.nan
        metrics["Davies-Bouldin Index"] = np.nan
        metrics["Calinski-Harabasz Index"] = np.nan
        if true_labels is not None:
            metrics["Adjusted Rand Index"] = np.nan
            metrics["Homogeneity Score"] = np.nan
            metrics["Completeness Score"] = np.nan
            metrics["V-Measure"] = np.nan

    return metrics


def df_qc(
    data: pd.DataFrame,
    columns=None,
    skim=False,
    plot_=True,
    max_cols=20,  # only for plots
    hue=None,
    output=False,
    verbose=True,
    dir_save=None,
):
    """
    Usage example:
    df = pd.DataFrame(...)  # Your DataFrameres_qc = df_qc(df)
    """
    from statsmodels.stats.outliers_influence import variance_inflation_factor
    from scipy.stats import skew, kurtosis, entropy

    pd.options.display.max_seq_items = 10
    #! display(data.select_dtypes(include=[np.number]).describe())
    #!skim
    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    if skim:
        try:
            import skimpy

            skimpy.skim(data)
        except:
            numerical_data = data.select_dtypes(include=[np.number])
            skimpy.skim(numerical_data)
    # Fill completely NaN columns with a default value (e.g., 0)
    data = data.copy()
    data.loc[:, data.isna().all()] = 0
    res_qc = {}
    print(f"⤵ data.shape:{data.shape}\n⤵ data.sample(10):")
    display(data.sample(10).style.background_gradient(cmap="coolwarm", axis=1))

    # Missing values
    res_qc["missing_values"] = data.isnull().sum()
    res_qc["missing_percentage"] = round(
        (res_qc["missing_values"] / len(data)) * 100, 2
    )
    res_qc["rows_with_missing"] = data.isnull().any(axis=1).sum()

    # Data types and unique values
    res_qc["data_types"] = data.dtypes
    res_qc["unique_counts"] = (
        data.select_dtypes(exclude=np.number).nunique().sort_values()
    )
    res_qc["unique_values"] = data.select_dtypes(exclude=np.number).apply(
        lambda x: x.unique()
    )
    res_qc["constant_columns"] = [
        col for col in data.columns if data[col].nunique() <= 1
    ]

    # Duplicate rows and columns
    res_qc["duplicate_rows"] = data.duplicated().sum()
    res_qc["duplicate_columns"] = data.columns[data.columns.duplicated()].tolist()

    # Empty columns
    res_qc["empty_columns"] = [col for col in data.columns if data[col].isnull().all()]

    # outliers
    data_outliers = df_outlier(data)
    outlier_num = data_outliers.isna().sum() - data.isnull().sum()
    res_qc["outlier_num"] = outlier_num[outlier_num > 0]
    outlier_percentage = round((outlier_num / len(data_outliers)) * 100, 2)
    res_qc["outlier_percentage"] = outlier_percentage[outlier_percentage > 0]
    try:
        # Correlation and multicollinearity (VIF)
        if any(data.dtypes.apply(pd.api.types.is_numeric_dtype)):
            numeric_df = data.select_dtypes(include=[np.number]).dropna()
            corr_matrix = numeric_df.corr()
            high_corr_pairs = [
                (col1, col2)
                for col1 in corr_matrix.columns
                for col2 in corr_matrix.columns
                if col1 != col2 and abs(corr_matrix[col1][col2]) > 0.9
            ]
            res_qc["high_correlations"] = high_corr_pairs

            # VIF for multicollinearity check
            numeric_df = data.select_dtypes(include=[np.number]).dropna()
            if isinstance(numeric_df.columns, pd.MultiIndex):
                numeric_df.columns = [
                    "_".join(col).strip() if isinstance(col, tuple) else col
                    for col in numeric_df.columns
                ]

            vif_data = pd.DataFrame()
            res_qc["vif"] = vif_data
            if numeric_df.shape[1] > 1 and not numeric_df.empty:
                vif_data["feature"] = numeric_df.columns.tolist()
                vif_data["VIF"] = [
                    round(variance_inflation_factor(numeric_df.values, i), 2)
                    for i in range(numeric_df.shape[1])
                ]
                res_qc["vif"] = vif_data[
                    vif_data["VIF"] > 5
                ]  # Typically VIF > 5 indicates multicollinearity
    except Exception as e:
        print(e)
    # Skewness and Kurtosis
    skewness = data.skew(numeric_only=True)
    kurtosis_vals = data.kurt(numeric_only=True)
    res_qc["skewness"] = skewness[abs(skewness) > 1]
    res_qc["kurtosis"] = kurtosis_vals[abs(kurtosis_vals) > 3]

    # Entropy for categorical columns (higher entropy suggests more disorder)
    categorical_cols = data.select_dtypes(include=["object", "category"]).columns
    res_qc["entropy_categoricals"] = {
        col: entropy(data[col].value_counts(normalize=True), base=2)
        for col in categorical_cols
    }

    # dtypes counts
    res_qc["dtype_counts"] = data.dtypes.value_counts()

    # Distribution Analysis (mean, median, mode, std dev, IQR for numeric columns)
    
    distribution_stats = data.select_dtypes(include=[np.number]).describe().T
    iqr = data.select_dtypes(include=[np.number]).apply(
        lambda x: x.quantile(0.75) - x.quantile(0.25)
    )
    distribution_stats["IQR"] = iqr
    res_qc["distribution_analysis"] = distribution_stats

    # Variance Check: Identify low-variance columns
    variance_threshold = 0.01
    low_variance_cols = [
        col
        for col in data.select_dtypes(include=[np.number]).columns
        if data[col].var() < variance_threshold
    ]
    res_qc["low_variance_features"] = low_variance_cols

    # Categorical columns and cardinality
    categorical_cols = data.select_dtypes(include=["object", "category"]).columns
    high_cardinality = {
        col: data[col].nunique() for col in categorical_cols if data[col].nunique() > 50
    }
    res_qc["high_cardinality_categoricals"] = high_cardinality

    # Feature-type inconsistency (mixed types in columns)
    inconsistent_types = {}
    for col in data.columns:
        unique_types = set(type(val) for val in data[col].dropna())
        if len(unique_types) > 1:
            inconsistent_types[col] = unique_types
    res_qc["inconsistent_types"] = inconsistent_types

    # Text length analysis for text fields
    text_lengths = {}
    for col in categorical_cols:
        text_lengths[col] = {
            "avg_length": data[col].dropna().apply(len).mean(),
            "length_variance": data[col].dropna().apply(len).var(),
        }
    res_qc["text_length_analysis"] = text_lengths

    # Summary statistics
    res_qc["summary_statistics"] = data.describe().T.style.background_gradient(
        cmap="coolwarm", axis=0
    )

    # Automated warnings
    warnings = []
    if res_qc["duplicate_rows"] > 0:
        warnings.append("Warning: Duplicate rows detected.")
    if len(res_qc["empty_columns"]) > 0:
        warnings.append("Warning: Columns with only NaN values detected.")
    if len(res_qc["constant_columns"]) > 0:
        warnings.append("Warning: Columns with a single constant value detected.")
    if len(high_corr_pairs) > 0:
        warnings.append("Warning: Highly correlated columns detected.")
    if len(res_qc["vif"]) > 0:
        warnings.append("Warning: Multicollinearity detected in features.")
    if len(high_cardinality) > 0:
        warnings.append("Warning: High cardinality in categorical columns.")
    if len(inconsistent_types) > 0:
        warnings.append("Warning: Columns with mixed data types detected.")
    res_qc["warnings"] = warnings

    # Report generation
    if verbose:
        print("\n⤵  Summary Statistics:")
        display(res_qc["summary_statistics"])
        print("\n⤵  Data Types:")
        display(res_qc["data_types"])
        if any(res_qc["missing_values"][res_qc["missing_values"] > 0]):
            print(" ⤵  Missing Values Counts:")
            display(
                pd.DataFrame(
                    {
                        "missing_values": res_qc["missing_values"][
                            res_qc["missing_values"] > 0
                        ],
                        "missing_percent(%)": res_qc["missing_percentage"][
                            res_qc["missing_percentage"] > 0
                        ],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )
            # print(res_qc["missing_percentage"][res_qc["missing_percentage"] > 0])
            print("\n⤵  Rows with Missing Values:", res_qc["rows_with_missing"])

        (
            print("\n⤵  Constant Columns:", res_qc["constant_columns"])
            if any(res_qc["constant_columns"])
            else None
        )
        (
            print("⤵  Duplicate Rows:", res_qc["duplicate_rows"])
            if res_qc["duplicate_rows"]
            else None
        )
        (
            print("⤵  Duplicate Columns:", res_qc["duplicate_columns"])
            if any(res_qc["duplicate_columns"])
            else None
        )

        if any(res_qc["outlier_num"]):
            print("\n⤵  Outlier Report:")
            display(
                pd.DataFrame(
                    {
                        "outlier_num": res_qc["outlier_num"][res_qc["outlier_num"] > 0],
                        "outlier_percentage(%)": res_qc["outlier_percentage"][
                            res_qc["outlier_percentage"] > 0
                        ],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )

        if any(res_qc["unique_counts"]):
            print("\n⤵  Unique Values per Column:")
            display(
                pd.DataFrame(
                    {
                        "unique_counts": res_qc["unique_counts"],
                        "unique_values": res_qc["unique_values"],
                    }
                ).style.background_gradient(cmap="coolwarm", axis=0)
            )

        if res_qc["empty_columns"]:
            print("\n⤵  Empty Columns:", res_qc["empty_columns"])

        if any(res_qc["high_correlations"]):
            print("\n⤵  High Correlations (>|0.9|):")
            for col1, col2 in res_qc["high_correlations"]:
                print(f"  {col1} and {col2}")

        if "vif" in res_qc:
            print("\n⤵  Features with High VIF (>|5|):")
            display(res_qc["vif"].style.background_gradient(cmap="coolwarm", axis=0))

        if any(res_qc["high_cardinality_categoricals"]):
            print("\n⤵  High Cardinality Categorical Columns (>|50 unique|):")
            print(res_qc["high_cardinality_categoricals"])
        if any(res_qc["inconsistent_types"]):
            print("\n⤵  Inconsistent Data Types:")
            display(res_qc["inconsistent_types"])
        if any(res_qc["text_length_analysis"]):
            print("\n⤵  Text Length Analysis:")
            for col, stats in res_qc["text_length_analysis"].items():
                print(
                    f"{col}: Avg Length={round(stats['avg_length'],1)}, Length Variance={round(stats['length_variance'],1)}"
                )

        if res_qc["warnings"]:
            print("\nWarnings:")
            for warning in res_qc["warnings"]:
                print("  -", warning)

    pd.reset_option("display.max_seq_items")
    if plot_:
        df_qc_plots(
            data=data, res_qc=res_qc, max_cols=max_cols, hue=hue, dir_save=dir_save
        )
    if output or not plot_:
        return res_qc
    return None


def df_qc_plots(
    data: pd.DataFrame,
    columns=None,
    res_qc: dict = None,
    max_cols=20,
    hue=None,
    dir_save=None,
):
    import matplotlib.pyplot as plt
    import seaborn as sns
    from .plot import subplot, figsets, get_color
    from datetime import datetime

    now_ = datetime.now().strftime("%y%m%d_%H%M%S")

    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]
    len_total = len(res_qc)
    n_row, n_col = int((len_total + 10)), 3
    nexttile = subplot(n_row, n_col, figsize=[5 * n_col, 5 * n_row], verbose=False)

    missing_data = res_qc["missing_values"][res_qc["missing_values"] > 0].sort_values(
        ascending=False
    )
    if len(missing_data) > max_cols:
        missing_data = missing_data[:max_cols]
    ax_missing_data = sns.barplot(
        y=missing_data.index,
        x=missing_data.values,
        hue=missing_data.index,
        palette=get_color(len(missing_data), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        title="Missing (#)",
        xlabel="#",
        ax=ax_missing_data,
        ylabel=None,
        fontsize=8 if len(missing_data) <= 20 else 6,
    )

    outlier_num = res_qc["outlier_num"].sort_values(ascending=False)
    if len(outlier_num) > max_cols:
        outlier_num = outlier_num[:max_cols]
    ax_outlier_num = sns.barplot(
        y=outlier_num.index,
        x=outlier_num.values,
        hue=outlier_num.index,
        palette=get_color(len(outlier_num), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        ax=ax_outlier_num,
        title="Outliers (#)",
        xlabel="#",
        ylabel=None,
        fontsize=8 if len(outlier_num) <= 20 else 6,
    )

    #!
    try:
        for col in data.select_dtypes(include="category").columns:
            sns.countplot(
                y=data[col],
                palette=get_color(
                    data.select_dtypes(include="category").shape[1], cmap="coolwarm"
                )[::-1],
                ax=nexttile(),
            )
            figsets(title=f"Count Plot: {col}", xlabel="Count", ylabel=col)
    except Exception as e:
        pass

    # Skewness and Kurtosis Plots
    skewness = res_qc["skewness"].sort_values(ascending=False)
    kurtosis = res_qc["kurtosis"].sort_values(ascending=False)
    if not skewness.empty:
        ax_skewness = sns.barplot(
            y=skewness.index,
            x=skewness.values,
            hue=skewness.index,
            palette=get_color(len(skewness), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Highly Skewed Numeric Columns (Skewness > 1)",
            xlabel="Skewness",
            ylabel=None,
            ax=ax_skewness,
            fontsize=8 if len(skewness) <= 20 else 6,
        )
    if not kurtosis.empty:
        ax_kurtosis = sns.barplot(
            y=kurtosis.index,
            x=kurtosis.values,
            hue=kurtosis.index,
            palette=get_color(len(kurtosis), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Highly Kurtotic Numeric Columns (Kurtosis > 3)",
            xlabel="Kurtosis",
            ylabel=None,
            ax=ax_kurtosis,
            fontsize=8 if len(kurtosis) <= 20 else 6,
        )

    # Entropy for Categorical Variables
    entropy_data = pd.Series(res_qc["entropy_categoricals"]).sort_values(
        ascending=False
    )
    ax_entropy_data = sns.barplot(
        y=entropy_data.index,
        x=entropy_data.values,
        hue=entropy_data.index,
        palette=get_color(len(entropy_data), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        ylabel="Categorical Columns",
        title="Entropy of Categorical Variables",
        xlabel="Entropy (bits)",
        ax=ax_entropy_data,
        fontsize=8 if len(entropy_data) <= 20 else 6,
    )

    # unique counts
    unique_counts = res_qc["unique_counts"].sort_values(ascending=False)
    ax_unique_counts_ = sns.barplot(
        y=unique_counts.index,
        x=unique_counts.values,
        hue=unique_counts.index,
        palette=get_color(len(unique_counts), cmap="coolwarm")[::-1],
        ax=nexttile(),
    )
    figsets(
        title="Unique Counts",
        ylabel=None,
        xlabel="#",
        ax=ax_unique_counts_,
        fontsize=8 if len(unique_counts) <= 20 else 6,
    )
    # Binary Checking
    ax_unique_counts = sns.barplot(
        y=unique_counts[unique_counts < 8].index,
        x=unique_counts[unique_counts < 8].values,
        hue=unique_counts[unique_counts < 8].index,
        palette=get_color(len(unique_counts[unique_counts < 8].index), cmap="coolwarm")[
            ::-1
        ],
        ax=nexttile(),
    )
    plt.axvline(x=2, color="r", linestyle="--", lw=2)
    figsets(
        ylabel=None,
        title="Binary Checking",
        xlabel="#",
        ax=ax_unique_counts,
        fontsize=8 if len(unique_counts[unique_counts < 10].index) <= 20 else 6,
    )

    # dtypes counts
    dtype_counts = res_qc["dtype_counts"]
    txt = []
    for tp in dtype_counts.index:
        txt.append(list(data.select_dtypes(include=tp).columns))

    ax_dtype_counts = sns.barplot(
        x=dtype_counts.index,
        y=dtype_counts.values,
        color="#F3C8B2",
        ax=nexttile(),
    )
    max_columns_per_row = 1  # Maximum number of columns per row
    for i, tp in enumerate(dtype_counts.index):
        if i <= 20:
            column_names = txt[i]
            # Split the column names into multiple lines if too long
            column_name_str = ", ".join(column_names)
            if len(column_name_str) > 40:  # If column names are too long, split them
                column_name_str = "\n".join(
                    [
                        ", ".join(column_names[j : j + max_columns_per_row])
                        for j in range(0, len(column_names), max_columns_per_row)
                    ]
                )
            # Place text annotation with line breaks and rotate the text if needed
            ax_dtype_counts.text(
                i,
                dtype_counts.values[i],
                f"{column_name_str}",
                ha="center",
                va="top",
                c="k",
                fontsize=8 if len(dtype_counts.index) <= 20 else 6,
                rotation=0,
            )
    figsets(
        xlabel=None,
        title="Dtypes",
        ylabel="#",
        ax=ax_dtype_counts,
        fontsize=8 if len(dtype_counts.index) <= 20 else 6,
    )
    # from .plot import pie
    # pie()

    # High cardinality: Show top categorical columns by unique value count
    high_cardinality = res_qc["high_cardinality_categoricals"]
    if high_cardinality and len(high_cardinality) > max_cols:
        high_cardinality = dict(
            sorted(high_cardinality.items(), key=lambda x: x[1], reverse=True)[
                :max_cols
            ]
        )

    if high_cardinality:
        ax_high_cardinality = sns.barplot(
            y=list(high_cardinality.keys()),
            x=list(high_cardinality.values()),
            hue=list(high_cardinality.keys()),
            palette=get_color(len(list(high_cardinality.keys())), cmap="coolwarm")[
                ::-1
            ],
            ax=nexttile(),
        )
        figsets(
            title="High Cardinality Categorical Columns",
            xlabel="Unique Value Count",
            ax=ax_high_cardinality,
            fontsize=8 if len(list(high_cardinality.keys())) <= 20 else 6,
        )
    if res_qc["low_variance_features"]:
        low_variance_data = data[res_qc["low_variance_features"]].copy()
        for col in low_variance_data.columns:
            ax_low_variance_features = sns.histplot(
                low_variance_data[col], bins=20, kde=True, color="coral", ax=nexttile()
            )
            figsets(
                title=f"Low Variance Feature: {col}",
                ax=ax_low_variance_features,
                fontsize=8 if len(low_variance_data[col]) <= 20 else 6,
            )

    # VIF plot for multicollinearity detection
    if "vif" in res_qc and not res_qc["vif"].empty:
        vif_data = res_qc["vif"].sort_values(by="VIF", ascending=False)
        if len(vif_data) > max_cols:
            vif_data = vif_data[:max_cols]
        ax_vif = sns.barplot(
            data=vif_data,
            x="VIF",
            y="feature",
            hue="VIF",
            palette=get_color(len(vif_data), cmap="coolwarm")[::-1],
            ax=nexttile(),
        )
        figsets(
            title="Variance Inflation Factor(VIF)",
            xlabel="VIF",
            ylabel="Features",
            legend=None,
            ax=ax_vif,
            fontsize=8 if len(vif_data) <= 20 else 6,
        )

    # Correlation heatmap for numeric columns with high correlation pairs
    if any(data.dtypes.apply(pd.api.types.is_numeric_dtype)):
        corr = data.select_dtypes(include=[np.number]).corr()
        if corr.shape[1] <= 33:
            mask = np.triu(np.ones_like(corr, dtype=bool))
            num_columns = corr.shape[1]
            fontsize = max(
                6, min(12, 12 - (num_columns - 10) * 0.2)
            )  # Scale between 8 and 12

            ax_heatmap = sns.heatmap(
                corr,
                mask=mask,
                annot=True,
                cmap="coolwarm",
                center=0,
                fmt=".1f",
                linewidths=0.5,
                vmin=-1,
                vmax=1,
                ax=nexttile(2, 2),
                cbar_kws=dict(shrink=0.2, ticks=np.arange(-1, 2, 1)),
                annot_kws={"size": fontsize},
            )

            figsets(xangle=45, title="Correlation Heatmap", ax=ax_heatmap)
    # # save figure
    # if dir_save:
    #     figsave(dir_save,f"qc_plot_{now_}.pdf")

    if columns is not None:
        if isinstance(columns, (list, pd.core.indexes.base.Index)):
            data = data[columns]

    # len_total = len(res_qc)
    # n_row, n_col = int((len_total + 10) / 3), 3
    # nexttile = subplot(n_row, n_col, figsize=[5 * n_col, 5 * n_row],verbose=False)
    #! check distribution
    data_num = data.select_dtypes(include=np.number)
    if len(data_num) > max_cols:
        data_num = data_num.iloc[:, :max_cols]

    data_num = df_scaler(data=data_num, method="standard")

    import scipy.stats as stats

    for column in data_num.columns:
        # * Shapiro-Wilk test for normality
        stat, p_value = stats.shapiro(data_num[column])
        normality = "norm" if p_value > 0.05 else "not_norm"
        # * Plot histogram
        ax_hist = sns.histplot(data_num[column], kde=True, ax=nexttile())
        x_min, x_max = ax_hist.get_xlim()
        y_min, y_max = ax_hist.get_ylim()
        ax_hist.text(
            x_min + (x_max - x_min) * 0.5,
            y_min + (y_max - y_min) * 0.75,
            f"p(Shapiro-Wilk)={p_value:.3f}\n{normality}",
            ha="center",
            va="top",
        )
        figsets(title=column, ax=ax_hist)
        ax_twin = ax_hist.twinx()
        # * Q-Q plot
        stats.probplot(data_num[column], dist="norm", plot=ax_twin)
        figsets(ylabel=f"Q-Q Plot:{column}", title=None)
    # save figure
    if dir_save:
        figsave(dir_save, f"qc_plot_{now_}.pdf")


def df_corr(df: pd.DataFrame, method="pearson"):
    """
    Compute correlation coefficients and p-values for a DataFrame.

    Parameters:
    - df (pd.DataFrame): Input DataFrame with numeric data.
    - method (str): Correlation method ("pearson", "spearman", "kendall").

    Returns:
    - corr_matrix (pd.DataFrame): Correlation coefficient matrix.
    - pval_matrix (pd.DataFrame): P-value matrix.
    """
    from scipy.stats import pearsonr, spearmanr, kendalltau

    methods = ["pearson", "spearman", "kendall"]
    method = strcmp(method, methods)[0]
    methods_dict = {"pearson": pearsonr, "spearman": spearmanr, "kendall": kendalltau}

    cols = df.columns
    corr_matrix = pd.DataFrame(index=cols, columns=cols, dtype=float)
    pval_matrix = pd.DataFrame(index=cols, columns=cols, dtype=float)
    correlation_func = methods_dict[method]

    for col1 in cols:
        for col2 in cols:
            if col1 == col2:
                corr_matrix.loc[col1, col2] = 1.0
                pval_matrix.loc[col1, col2] = 0.0
            else:
                corr, pval = correlation_func(df[col1], df[col2])
                corr_matrix.loc[col1, col2] = corr
                pval_matrix.loc[col1, col2] = pval

    return corr_matrix, pval_matrix


def use_pd(
    func_name="excel",
    verbose=True,
    dir_json="./data/usages_pd.json",
):
    try:
        default_settings = fload(dir_json, output="json")
        valid_kinds = list(default_settings.keys())
        kind = strcmp(func_name, valid_kinds)[0]
        usage = default_settings[kind]
        if verbose:
            for i, i_ in enumerate(ssplit(usage, by=",")):
                i_ = i_.replace("=", "\t= ") + ","
                print(i_) if i == 0 else print("\t", i_)
        else:
            print(usage)
    except Exception as e:
        if verbose:
            print(e)


def get_phone(phone_number: str, region: str = None, verbose=True):
    """
    usage:
        info = get_phone(15237654321, "DE")
        preview(info)

    Extremely advanced phone number analysis function.

    Args:
        phone_number (str): The phone number to analyze.
        region (str): None (Default). Tries to work with international numbers including country codes; otherwise, uses the specified region.

    Returns:
        dict: Comprehensive information about the phone number.
    """
    import phonenumbers
    from phonenumbers import geocoder, carrier, timezone, number_type
    from datetime import datetime
    import pytz
    from tzlocal import get_localzone

    if not isinstance(phone_number, str):
        phone_number = str(phone_number)
    if isinstance(region, str):
        region = region.upper()

    try:
        # Parse the phone number
        parsed_number = phonenumbers.parse(phone_number, region)

        # Validate the phone number
        valid = phonenumbers.is_valid_number(parsed_number)
        possible = phonenumbers.is_possible_number(parsed_number)

        if not valid:
            suggested_fix = phonenumbers.example_number(region) if region else "Unknown"
            return {
                "valid": False,
                "error": "Invalid phone number",
                "suggested_fix": suggested_fix,
            }

        # Basic details
        formatted_international = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.INTERNATIONAL
        )
        formatted_national = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.NATIONAL
        )
        formatted_e164 = phonenumbers.format_number(
            parsed_number, phonenumbers.PhoneNumberFormat.E164
        )
        country_code = parsed_number.country_code
        region_code = geocoder.region_code_for_number(parsed_number)
        country_name = geocoder.country_name_for_number(parsed_number, "en")

        location = geocoder.description_for_number(parsed_number, "en")
        carrier_name = carrier.name_for_number(parsed_number, "en") or "Unknown Carrier"
        time_zones = timezone.time_zones_for_number(parsed_number)[0]
        current_times = datetime.now(pytz.timezone(time_zones)).strftime(
            "%Y-%m-%d %H:%M:%S %Z"
        )
        number_type_str = {
            phonenumbers.PhoneNumberType.FIXED_LINE: "Fixed Line",
            phonenumbers.PhoneNumberType.MOBILE: "Mobile",
            phonenumbers.PhoneNumberType.FIXED_LINE_OR_MOBILE: "Fixed Line or Mobile",
            phonenumbers.PhoneNumberType.TOLL_FREE: "Toll Free",
            phonenumbers.PhoneNumberType.PREMIUM_RATE: "Premium Rate",
            phonenumbers.PhoneNumberType.SHARED_COST: "Shared Cost",
            phonenumbers.PhoneNumberType.VOIP: "VOIP",
            phonenumbers.PhoneNumberType.PERSONAL_NUMBER: "Personal Number",
            phonenumbers.PhoneNumberType.PAGER: "Pager",
            phonenumbers.PhoneNumberType.UAN: "UAN",
            phonenumbers.PhoneNumberType.UNKNOWN: "Unknown",
        }.get(number_type(parsed_number), "Unknown")

        # Advanced Features
        is_toll_free = (
            number_type(parsed_number) == phonenumbers.PhoneNumberType.TOLL_FREE
        )
        is_premium_rate = (
            number_type(parsed_number) == phonenumbers.PhoneNumberType.PREMIUM_RATE
        )

        # Dialing Information
        dialing_instructions = f"Dial {formatted_national} within {country_name}. Dial {formatted_e164} from abroad."

        # Advanced Timezone Handling
        gmt_offsets = (
            pytz.timezone(time_zones).utcoffset(datetime.now()).total_seconds() / 3600
        )
        # Get the local timezone (current computer's time)
        local_timezone = get_localzone()
        # local_timezone = pytz.timezone(pytz.country_timezones[region_code][0])
        local_offset = local_timezone.utcoffset(datetime.now()).total_seconds() / 3600
        offset_diff = local_offset - gmt_offsets
        head_time = "earlier" if offset_diff < 0 else "later" if offset_diff > 0 else ""
        res = {
            "valid": True,
            "possible": possible,
            "formatted": {
                "international": formatted_international,
                "national": formatted_national,
                "e164": formatted_e164,
            },
            "country_code": country_code,
            "country_name": country_name,
            "region_code": region_code,
            "location": location if location else "Unknown",
            "carrier": carrier_name,
            "time_zone": time_zones,
            "current_times": current_times,
            "local_offset": f"{local_offset} utcoffset",
            "time_zone_diff": f"{head_time} {int(np.abs(offset_diff))} h",
            "number_type": number_type_str,
            "is_toll_free": is_toll_free,
            "is_premium_rate": is_premium_rate,
            "dialing_instructions": dialing_instructions,
            "suggested_fix": None,  # Use phonenumbers.example_number if invalid
            "logs": {
                "number_analysis_completed": datetime.now().strftime(
                    "%Y-%m-%d %H:%M:%S"
                ),
                "raw_input": phone_number,
                "parsed_number": str(parsed_number),
            },
        }

    except phonenumbers.NumberParseException as e:
        res = {"valid": False, "error": str(e)}
    if verbose:
        preview(res)
    return res


def decode_pluscode(
    pluscode: str, reference: tuple = (52.5200, 13.4050), return_bbox: bool = False
):
    """
    Decodes a Plus Code into latitude and longitude (and optionally returns a bounding box).

    Parameters:
        pluscode (str): The Plus Code to decode. Can be full or short.
        reference (tuple, optional): Reference latitude and longitude for decoding short Plus Codes.
                                     Default is None, required if Plus Code is short.
        return_bbox (bool): If True, returns the bounding box coordinates (latitude/longitude bounds).
                            Default is False.

    Returns:
        tuple: (latitude, longitude) if `return_bbox` is False.
               (latitude, longitude, bbox) if `return_bbox` is True.
               bbox = (latitudeLo, latitudeHi, longitudeLo, longitudeHi)
    Raises:
        ValueError: If the Plus Code is invalid or reference is missing for a short code.

    Usage:
    lat, lon = decode_pluscode("7FG6+89")
    print(f"Decoded Short Plus Code: Latitude: {lat}, Longitude: {lon}, Bounding Box: {bbox}")

    lat, lon = decode_pluscode("9F4M7FG6+89")
    print(f"Decoded Full Plus Code: Latitude: {lat}, Longitude: {lon}")
    """
    from openlocationcode import openlocationcode as olc

    # Validate Plus Code
    if not olc.isValid(pluscode):
        raise ValueError(f"Invalid Plus Code: {pluscode}")

    # Handle Short Plus Codes
    if olc.isShort(pluscode):
        if reference is None:
            raise ValueError(
                "Reference location (latitude, longitude) is required for decoding short Plus Codes."
            )
        # Recover the full Plus Code using the reference location
        pluscode = olc.recoverNearest(pluscode, reference[0], reference[1])

    # Decode the Plus Code
    decoded = olc.decode(pluscode)

    # Calculate the center point of the bounding box
    latitude = (decoded.latitudeLo + decoded.latitudeHi) / 2
    longitude = (decoded.longitudeLo + decoded.longitudeHi) / 2

    if return_bbox:
        bbox = (
            decoded.latitudeLo,
            decoded.latitudeHi,
            decoded.longitudeLo,
            decoded.longitudeHi,
        )
        return latitude, longitude, bbox

    return latitude, longitude


def get_loc(input_data, user_agent="0413@mygmail.com)", verbose=True):
    """
        Determine if the input is a city name, lat/lon, or DMS and perform geocoding or reverse geocoding.
    Usage:
        get_loc("Berlin, Germany")  # Example city
        # get_loc((48.8566, 2.3522))  # Example latitude and longitude
        # get_loc("48 51 24.3 N")  # Example DMS input
    """
    from geopy.geocoders import Nominatim
    import re

    def dms_to_decimal(dms):
        """
        Convert DMS (Degrees, Minutes, Seconds) to Decimal format.
        Input should be in the format of "DD MM SS" or "D M S".
        """
        # Regex pattern for DMS input
        pattern = r"(\d{1,3})[^\d]*?(\d{1,2})[^\d]*?(\d{1,2})"
        match = re.match(pattern, dms)

        if match:
            degrees, minutes, seconds = map(float, match.groups())
            decimal = degrees + (minutes / 60) + (seconds / 3600)
            return decimal
        else:
            raise ValueError("Invalid DMS format")

    geolocator = Nominatim(user_agent="0413@mygmail.com)")
    # Case 1: Input is a city name (string)
    if isinstance(input_data, str) and not re.match(r"^\d+(\.\d+)?$", input_data):
        location = geolocator.geocode(input_data)
        try:
            if verbose:
                print(
                    f"Latitude and Longitude for {input_data}: {location.latitude}, {location.longitude}"
                )
            else:
                print(f"Could not find {input_data}.")
            return location
        except Exception as e:
            print(f"Error: {e}")
            return

    # Case 2: Input is latitude and longitude (float or tuple)
    elif isinstance(input_data, (float, tuple)):
        if isinstance(input_data, tuple) and len(input_data) == 2:
            latitude, longitude = input_data
        elif isinstance(input_data, float):
            latitude = input_data
            longitude = None  # No longitude provided for a single float

        # Reverse geocoding
        location_reversed = geolocator.reverse(
            (latitude, longitude) if longitude else latitude
        )
        if verbose:
            print(
                f"Address from coordinates ({latitude}, {longitude if longitude else ''}): {location_reversed.address}"
            )
        else:
            print("Could not reverse geocode the coordinates.")
        return location_reversed

    # Case 3: Input is a DMS string
    elif isinstance(input_data, str):
        try:
            decimal_lat = dms_to_decimal(input_data)
            print(f"Converted DMS to decimal latitude: {decimal_lat}")

            location_reversed = geolocator.reverse(decimal_lat)
            if verbose:
                print(f"Address from coordinates: {location_reversed.address}")
            else:
                print("Could not reverse geocode the coordinates.")
            return location_reversed
        except ValueError:
            print(
                "Invalid input format. Please provide a city name, latitude/longitude, or DMS string."
            )


def enpass(code: str, method: str = "AES", key: str = None):
    """
    usage: enpass("admin")
    Master encryption function that supports multiple methods: AES, RSA, and SHA256.
    :param code: The input data to encrypt or hash.
    :param method: The encryption or hashing method ('AES', 'RSA', or 'SHA256').
    :param key: The key to use for encryption. For AES and RSA, it can be a password or key in PEM format.
    :return: The encrypted data or hashed value.
    """
    import hashlib

    # AES Encryption (Advanced)
    def aes_encrypt(data: str, key: str):
        """
        Encrypts data using AES algorithm in CBC mode.
        :param data: The data to encrypt.
        :param key: The key to use for AES encryption.
        :return: The encrypted data, base64 encoded.
        """
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.backends import default_backend
        from cryptography.hazmat.primitives import padding
        import base64
        import os

        # Generate a 256-bit key from the provided password
        key = hashlib.sha256(key.encode()).digest()

        # Generate a random initialization vector (IV)
        iv = os.urandom(16)  # 16 bytes for AES block size

        # Pad the data to be a multiple of 16 bytes using PKCS7
        padder = padding.PKCS7(128).padder()
        padded_data = padder.update(data.encode()) + padder.finalize()

        # Create AES cipher object using CBC mode
        cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
        encryptor = cipher.encryptor()
        encrypted_data = encryptor.update(padded_data) + encryptor.finalize()

        # Return the base64 encoded result (IV + encrypted data)
        return base64.b64encode(iv + encrypted_data).decode()

    # RSA Encryption (Advanced)
    def rsa_encrypt(data: str, public_key: str):
        """
        Encrypts data using RSA encryption with OAEP padding.
        :param data: The data to encrypt.
        :param public_key: The public key in PEM format.
        :return: The encrypted data, base64 encoded.
        """
        import base64
        from Crypto.PublicKey import RSA
        from Crypto.Cipher import PKCS1_OAEP

        public_key_obj = RSA.import_key(public_key)
        cipher_rsa = PKCS1_OAEP.new(public_key_obj)
        encrypted_data = cipher_rsa.encrypt(data.encode())
        return base64.b64encode(encrypted_data).decode()

    # SHA256 Hashing (Non-reversible)
    def sha256_hash(data: str):
        """
        Generates a SHA256 hash of the data.
        :param data: The data to hash.
        :return: The hashed value (hex string).
        """
        return hashlib.sha256(data.encode()).hexdigest()
    def md5_hash(data:str):
        return hashlib.md5(data.encode()).hexdigest()

    if key is None:
        key = "worldpeace"
    method = strcmp(method, ["AES", "RSA", "SHA256","md5","hash"])[0]
    if method.upper() == "AES":
        return aes_encrypt(code, key)
    elif method.upper() == "RSA":
        return rsa_encrypt(code, key)
    elif method.upper() == "SHA256":
        return sha256_hash(code)
    elif method.lower() == "md5" or method.lower() == "hash":
        return md5_hash(code)
    else:
        raise ValueError("Unsupported encryption method")


# Master Decryption Function (Supports AES, RSA)
def depass(encrypted_code: str, method: str = "AES", key: str = None):
    """
    Master decryption function that supports multiple methods: AES and RSA.
    :param encrypted_code: The encrypted data to decrypt.
    :param method: The encryption method ('AES' or 'RSA').
    :param key: The key to use for decryption. For AES and RSA, it can be a password or key in PEM format.
    :return: The decrypted data.
    """
    import hashlib

    def aes_decrypt(encrypted_data: str, key: str):
        """
        Decrypts data encrypted using AES in CBC mode.
        :param encrypted_data: The encrypted data, base64 encoded.
        :param key: The key to use for AES decryption.
        :return: The decrypted data (string).
        """
        from cryptography.hazmat.primitives.ciphers import Cipher, algorithms, modes
        from cryptography.hazmat.backends import default_backend
        from cryptography.hazmat.primitives import padding
        import base64

        # Generate the same 256-bit key from the password
        key = hashlib.sha256(key.encode()).digest()

        # Decode the encrypted data from base64
        encrypted_data = base64.b64decode(encrypted_data)

        # Extract the IV and the actual encrypted data
        iv = encrypted_data[:16]  # First 16 bytes are the IV
        encrypted_data = encrypted_data[16:]  # Remaining data is the encrypted message

        # Create AES cipher object using CBC mode
        cipher = Cipher(algorithms.AES(key), modes.CBC(iv), backend=default_backend())
        decryptor = cipher.decryptor()
        decrypted_data = decryptor.update(encrypted_data) + decryptor.finalize()

        # Unpad the decrypted data using PKCS7
        unpadder = padding.PKCS7(128).unpadder()
        unpadded_data = unpadder.update(decrypted_data) + unpadder.finalize()

        return unpadded_data.decode()

    def rsa_decrypt(encrypted_data: str, private_key: str):
        """
        Decrypts RSA-encrypted data using the private key.
        :param encrypted_data: The encrypted data, base64 encoded.
        :param private_key: The private key in PEM format.
        :return: The decrypted data (string).
        """
        from Crypto.PublicKey import RSA
        from Crypto.Cipher import PKCS1_OAEP
        import base64

        encrypted_data = base64.b64decode(encrypted_data)
        private_key_obj = RSA.import_key(private_key)
        cipher_rsa = PKCS1_OAEP.new(private_key_obj)
        decrypted_data = cipher_rsa.decrypt(encrypted_data)
        return decrypted_data.decode()

    if key is None:
        key = "worldpeace"
    method = strcmp(method, ["AES", "RSA", "SHA256","md5",'hash'])[0]
    if method == "AES":
        return aes_decrypt(encrypted_code, key)
    elif method == "RSA":
        return rsa_decrypt(encrypted_code, key)
    elif method == "SHA256":
        raise ValueError("SHA256 is a hash function and cannot be decrypted.")
    else:
        raise ValueError("Unsupported decryption method")


def get_clip(dir_save=None):
    """
    Master function to extract content from the clipboard (text, URL, or image).

    Parameters:
        dir_save (str, optional): If an image is found, save it to this path.

    Returns:
        dict: A dictionary with extracted content:
              {
                  "type": "text" | "url" | "image" | "none",
                  "content": <str|Image|None>,
                  "saved_to": <str|None>  # Path if an image is saved
              }
    """
    result = {"type": "none", "content": None, "saved_to": None}

    try:
        import pyperclip
        from PIL import ImageGrab, Image
        import validators

        # 1. Check for text in the clipboard
        clipboard_content = pyperclip.paste()
        if clipboard_content:
            if validators.url(clipboard_content.strip()):
                result["type"] = "url"
                result["content"] = clipboard_content.strip()

            else:
                result["type"] = "text"
                result["content"] = clipboard_content.strip()
            return clipboard_content.strip()

        # 2. Check for image in the clipboard
        image = ImageGrab.grabclipboard()
        if isinstance(image, Image.Image):
            result["type"] = "image"
            result["content"] = image
            if dir_save:
                image.save(dir_save)
                result["saved_to"] = dir_save
                print(f"Image saved to {dir_save}.")
            else:
                print("Image detected in clipboard but not saved.")
            return image
        print("No valid text, URL, or image found in clipboard.")
        return result

    except Exception as e:
        print(f"An error occurred: {e}")
        return result


def keyboard(*args, action="press", n_click=1, interval=0, verbose=False, **kwargs):
    """
    Simulates keyboard input using pyautogui.

    Parameters:
        input_key (str): The key to simulate. Check the list of supported keys with verbose=True.
        action (str): The action to perform. Options are 'press', 'keyDown', or 'keyUp'.
        n_click (int): Number of times to press the key (only for 'press' action).
        interval (float): Time interval between key presses for 'press' action.
        verbose (bool): Print detailed output, including supported keys and debug info.
        kwargs: Additional arguments (reserved for future extensions).

    keyboard("command", "d", action="shorcut")
    """
    import pyautogui

    input_key = args

    actions = ["press", "keyDown", "keyUp", "hold", "release", "hotkey", "shortcut"]
    action = strcmp(action, actions)[0]
    keyboard_keys_ = [
        "\t",
        "\n",
        "\r",
        " ",
        "!",
        '"',
        "#",
        "$",
        "%",
        "&",
        "'",
        "(",
        ")",
        "*",
        "+",
        ",",
        "-",
        ".",
        "/",
        "0",
        "1",
        "2",
        "3",
        "4",
        "5",
        "6",
        "7",
        "8",
        "9",
        ":",
        ";",
        "<",
        "=",
        ">",
        "?",
        "@",
        "[",
        "\\",
        "]",
        "^",
        "_",
        "`",
        "a",
        "b",
        "c",
        "d",
        "e",
        "f",
        "g",
        "h",
        "i",
        "j",
        "k",
        "l",
        "m",
        "n",
        "o",
        "p",
        "q",
        "r",
        "s",
        "t",
        "u",
        "v",
        "w",
        "x",
        "y",
        "z",
        "{",
        "|",
        "}",
        "~",
        "accept",
        "add",
        "alt",
        "altleft",
        "altright",
        "apps",
        "backspace",
        "browserback",
        "browserfavorites",
        "browserforward",
        "browserhome",
        "browserrefresh",
        "browsersearch",
        "browserstop",
        "capslock",
        "clear",
        "convert",
        "ctrl",
        "ctrlleft",
        "ctrlright",
        "decimal",
        "del",
        "delete",
        "divide",
        "down",
        "end",
        "enter",
        "esc",
        "escape",
        "execute",
        "f1",
        "f10",
        "f11",
        "f12",
        "f13",
        "f14",
        "f15",
        "f16",
        "f17",
        "f18",
        "f19",
        "f2",
        "f20",
        "f21",
        "f22",
        "f23",
        "f24",
        "f3",
        "f4",
        "f5",
        "f6",
        "f7",
        "f8",
        "f9",
        "final",
        "fn",
        "hanguel",
        "hangul",
        "hanja",
        "help",
        "home",
        "insert",
        "junja",
        "kana",
        "kanji",
        "launchapp1",
        "launchapp2",
        "launchmail",
        "launchmediaselect",
        "left",
        "modechange",
        "multiply",
        "nexttrack",
        "nonconvert",
        "num0",
        "num1",
        "num2",
        "num3",
        "num4",
        "num5",
        "num6",
        "num7",
        "num8",
        "num9",
        "numlock",
        "pagedown",
        "pageup",
        "pause",
        "pgdn",
        "pgup",
        "playpause",
        "prevtrack",
        "print",
        "printscreen",
        "prntscrn",
        "prtsc",
        "prtscr",
        "return",
        "right",
        "scrolllock",
        "select",
        "separator",
        "shift",
        "shiftleft",
        "shiftright",
        "sleep",
        "space",
        "stop",
        "subtract",
        "tab",
        "up",
        "volumedown",
        "volumemute",
        "volumeup",
        "win",
        "winleft",
        "winright",
        "yen",
        "command",
        "option",
        "optionleft",
        "optionright",
    ]
    if verbose:
        print(f"supported keys: {keyboard_keys_}")

    if action not in ["hotkey", "shortcut"]:
        if not isinstance(input_key, list):
            input_key = list(input_key)
        input_key = [strcmp(i, keyboard_keys_)[0] for i in input_key]

    # correct action
    cmd_keys = [
        "command",
        "option",
        "optionleft",
        "optionright",
        "win",
        "winleft",
        "winright",
        "ctrl",
        "ctrlleft",
        "ctrlright",
    ]
    try:
        if any([i in cmd_keys for i in input_key]):
            action = "hotkey"
    except:
        pass

    print(f"\n{action}: {input_key}")
    # keyboard
    if action in ["press"]:
        # pyautogui.press(input_key, presses=n_click,interval=interval)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.press(key)
                pyautogui.sleep(interval)
    elif action in ["keyDown", "hold"]:
        # pyautogui.keyDown(input_key)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.keyDown(key)
                pyautogui.sleep(interval)

    elif action in ["keyUp", "release"]:
        # pyautogui.keyUp(input_key)
        for _ in range(n_click):
            for key in input_key:
                pyautogui.keyUp(key)
                pyautogui.sleep(interval)

    elif action in ["hotkey", "shortcut"]:
        pyautogui.hotkey(input_key)


def mouse(
    *args,  # loc
    action: str = "move",
    duration: float = 0.5,
    loc_type: str = "absolute",  # 'absolute', 'relative'
    region: tuple = None,  # (tuple, optional): A region (x, y, width, height) to search for the image.
    image_path: str = None,
    wait: float = 0,
    text: str = None,
    confidence: float = 0.8,
    button: str = "left",
    n_click: int = 1,  # number of clicks
    interval: float = 0.25,  # time between clicks
    scroll_amount: int = -500,
    fail_safe: bool = True,
    grayscale: bool = False,
    n_try: int = 10,
    verbose: bool = True,
    **kwargs,
):
    """
    Master function to handle pyautogui actions.

    Parameters:
        action (str): The action to perform ('click', 'double_click', 'type', 'drag', 'scroll', 'move', 'locate', etc.).
        image_path (str, optional): Path to the image for 'locate' or 'click' actions.
        text (str, optional): Text to type for 'type' action.
        confidence (float, optional): Confidence level for image recognition (default 0.8).
        duration (float, optional): Duration for smooth movements in seconds (default 0.5).
        region (tuple, optional): A region (x, y, width, height) to search for the image.
        button (str, optional): Mouse button to use ('left', 'right', 'middle').
        n_click (int, optional): Number of times to click for 'click' actions.
        interval (float, optional): Interval between clicks for 'click' actions.
        offset (tuple, optional): Horizontal offset from the located image. y_offset (int, optional): Vertical offset from the located image.
        scroll_amount (int, optional): Amount to scroll (positive for up, negative for down).
        fail_safe (bool, optional): Enable/disable pyautogui's fail-safe feature.
        grayscale (bool, optional): Search for the image in grayscale mode.

    Returns:
        tuple or None: Returns coordinates for 'locate' actions, otherwise None.
    """
    import pyautogui
    import time

    # import logging
    # logging.basicConfig(level=logging.DEBUG, filename="debug.log")

    pyautogui.FAILSAFE = fail_safe  # Enable/disable fail-safe
    loc_type = "absolute" if "abs" in loc_type else "relative"
    if len(args) == 1:
        if isinstance(args[0], str):
            image_path = args[0]
            x_offset, y_offset = None, None
        else:
            x_offset, y_offset = args
    elif len(args) == 2:
        x_offset, y_offset = args
    elif len(args) == 3:
        x_offset, y_offset, action = args
    elif len(args) == 4:
        x_offset, y_offset, action, duration = args
    else:
        x_offset, y_offset = None, None

    what_action = [
        "locate",
        "click",
        "double_click",
        "triple_click",
        "input",
        "write",
        "type",
        "drag",
        "move",
        "scroll",
        "down",
        "up",
        "hold",
        "press",
        "release",
    ]
    action = strcmp(action, what_action)[0]
    # get the locations
    location = None
    if any([x_offset is None, y_offset is None]):
        if region is None:
            w, h = pyautogui.size()
            region = (0, 0, w, h)
        retries = 0
        while location is None and retries <= n_try:
            try:
                confidence_ = round(float(confidence - 0.05 * retries), 2)
                location = pyautogui.locateOnScreen(
                    image_path,
                    confidence=confidence_,
                    region=region,
                    grayscale=grayscale,
                )
            except Exception as e:
                if verbose:
                    print(f"confidence={confidence_},{e}")
                location = None
            retries += 1

    # try:
    if location:
        x, y = pyautogui.center(location)
        x += x_offset if x_offset else 0
        if x_offset is not None:
            x += x_offset
        if y_offset is not None:
            y += y_offset
        x_offset, y_offset = x, y
    print(action) if verbose else None
    if action in ["locate"]:
        x, y = pyautogui.position()
    elif action in ["click", "double_click", "triple_click"]:
        if action == "click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.click(
                x=x_offset, y=y_offset, clicks=n_click, interval=interval, button=button
            )
        elif action == "double_click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.doubleClick(
                x=x_offset, y=y_offset, interval=interval, button=button
            )
        elif action == "triple_click":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
            time.sleep(wait)
            pyautogui.tripleClick(
                x=x_offset, y=y_offset, interval=interval, button=button
            )

    elif action in ["type", "write", "input"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        if text is not None:
            pyautogui.typewrite(text, interval=interval)
        else:
            print("Text must be provided for the 'type' action.") if verbose else None

    elif action == "drag":
        if loc_type == "absolute":
            pyautogui.dragTo(x_offset, y_offset, duration=duration, button=button)
        else:
            pyautogui.dragRel(x_offset, y_offset, duration=duration, button=button)

    elif action in ["move"]:
        if loc_type == "absolute":
            pyautogui.moveTo(x_offset, y_offset, duration=duration)
        else:
            pyautogui.moveRel(x_offset, y_offset, duration=duration)

    elif action == "scroll":
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.scroll(scroll_amount)

    elif action in ["down", "hold", "press"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.mouseDown(x_offset, y_offset, button=button, duration=duration)

    elif action in ["up", "release"]:
        pyautogui.moveTo(x_offset, y_offset, duration=duration)
        time.sleep(wait)
        pyautogui.mouseUp(x_offset, y_offset, button=button, duration=duration)

    else:
        raise ValueError(f"Unsupported action: {action}")


def py2installer(
    script_path: str = None,
    flatform: str = "mingw64",
    output_dir: str = "dist",
    icon_path: str = None,
    include_data: list = None,
    include_import: list = None,
    exclude_import: list = None,
    plugins: list = None,
    use_nuitka: bool = True,
    console: bool = True,
    clean_build: bool = False,
    additional_args: list = None,
    verbose: bool = True,
    standalone: bool = True,
    onefile: bool = False,
    use_docker: bool = False,
    docker_image: str = "python:3.12-slim",
):
    """
    to package Python scripts into standalone application.

    script_path (str): Path to the Python script to package.
    output_dir (str): Directory where the executable will be stored.
    icon_path (str): Path to the .ico file for the executable icon.
    include_data (list): List of additional data files or directories in "source:dest" format.
    exclude_import (list): List of hidden imports to include.
    plugins (list): List of plugins imports to include.e.g., 'tk-inter'
    use_nuitka (bool): Whether to use Nuitka instead of PyInstaller.
    console (bool): If False, hides the console window (GUI mode).
    clean_build (bool): If True, cleans previous build and dist directories.
    additional_args (list): Additional arguments for PyInstaller/Nuitka.
    verbose (bool): If True, provides detailed logs.
    use_docker (bool): If True, uses Docker to package the script.
    docker_image (str): Docker image to use for packaging.

    """
    import glob
    from pathlib import Path

    if run_once_within():
        usage_str = """
            # build locally
            py2installer(
                script_path="update_tab.py",
                output_dir="dist",
                icon_path="icon4app.ico",
                include_data=["dat/*.xlsx:dat"],
                exclude_import=["msoffcrypto", "tkinter", "pandas", "numpy"],
                onefile=True,
                console=False,
                clean_build=True,
                verbose=True,
            )
            # build via docker
            py2installer(
                "my_script.py",
                output_dir="dist",
                onefile=True,
                clean_build=True,
                use_docker=True,
                docker_image="python:3.12-slim"
            )
            # 尽量不要使用--include-package,这可能导致冲突
            py2installer(
                script_path="update_tab.py",
                # flatform=None,
                output_dir="dist_simp_subprocess",
                icon_path="icon4app.ico",
                standalone=True,
                onefile=False,
                include_data=["dat/*.xlsx=dat"],
                plugins=[
                    "tk-inter",
                ],
                use_nuitka=True,
                console=True,
                clean_build=False,
                verbose=0,
            )
            # 最终文件大小对比
            900 MB: nuitka --mingw64 --standalone --windows-console-mode=attach --show-progress --output-dir=dist --macos-create-app-bundle --macos-app-icon=icon4app.ico --nofollow-import-to=timm,paddle,torch,torchmetrics,torchvision,tensorflow,tensorboard,tensorboardx,tensorboard-data-server,textblob,PIL,sklearn,scienceplots,scikit-image,scikit-learn,scikit-surprise,scipy,spikeinterface,spike-sort-lfpy,stanza,statsmodels,streamlit,streamlit-autorefresh,streamlit-folium,pkg2ls,plotly --include-package=msoffcrypto,tkinter,datetime,pandas,numpy --enable-plugin=tk-inter update_tab.py;
            470 MB: nuitka --mingw64 --standalone --windows-console-mode=attach --show-progress --output-dir=dist_direct_nuitka --macos-create-app-bundle --macos-app-icon=icon4app.ico --enable-plugin=tk-inter update_taby ;
        """
        print(usage_str)
        if verbose:
            return
        else:
            pass
    # Check if the script path exists
    script_path = Path(script_path)
    if not script_path.exists():
        raise FileNotFoundError(f"Script '{script_path}' not found.")

    # Clean build and dist directories if requested
    if clean_build:
        for folder in ["build", "dist"]:
            folder_path = Path(folder)
            if folder_path.exists():
                shutil.rmtree(folder_path, ignore_errors=True)
        # Recreate the folders
        for folder in ["build", "dist"]:
            folder_path = Path(folder)
            folder_path.mkdir(parents=True, exist_ok=True)

    if use_docker:
        # Ensure Docker is installed
        try:
            subprocess.run(
                ["docker", "--version"], check=True, capture_output=True, text=True
            )
        except FileNotFoundError:
            raise EnvironmentError("Docker is not installed or not in the PATH.")

        # Prepare Docker volume mappings
        script_dir = script_path.parent.resolve()
        dist_path = Path(output_dir).resolve()
        volumes = [
            f"{script_dir}:/app:rw",
            f"{dist_path}:/output:rw",
        ]
        docker_cmd = [
            "docker",
            "run",
            "--rm",
            "-v",
            volumes[0],
            "-v",
            volumes[1],
            docker_image,
            "bash",
            "-c",
        ]

        # Build the packaging command inside the container
        cmd = ["nuitka"] if use_nuitka else ["pyinstaller"]
        if onefile:
            cmd.append("--onefile")
        if not console:
            cmd.append("--windowed")
        cmd.extend(["--distpath", "/output"])
        if icon_path:
            cmd.extend(["--icon", f"/app/{Path(icon_path).name}"])
        if include_data:
            for data in include_data:
                cmd.extend(["--add-data", f"/app/{data}"])
        if exclude_import:
            for hidden in exclude_import:
                cmd.extend(["--hidden-import", hidden])
        if additional_args:
            cmd.extend(additional_args)
        cmd.append(f"/app/{script_path.name}")

        # Full command to execute inside the container
        docker_cmd.append(" ".join(cmd))

        if verbose:
            print(f"Running Docker command: {' '.join(docker_cmd)}")

        # Run Docker command
        try:
            subprocess.run(
                docker_cmd,
                capture_output=not verbose,
                text=True,
                check=True,
            )
        except subprocess.CalledProcessError as e:
            print(f"Error during Docker packaging:\n{e.stderr}", file=sys.stderr)
            raise
    else:
        # Handle local packaging (native build)
        cmd = ["nuitka"] if use_nuitka else ["pyinstaller"]
        if "min" in flatform.lower() and use_nuitka:
            cmd.append("--mingw64")
        cmd.append("--standalone") if use_nuitka and standalone else None
        cmd.append("--onefile") if onefile else None
        if not console:
            cmd.append("--windows-console-mode=disable")
        else:
            cmd.append("--windows-console-mode=attach")

        cmd.extend(["--show-progress", f"--output-dir={output_dir}"])
        if icon_path:
            icon_path = Path(icon_path)
            if not icon_path.exists():
                raise FileNotFoundError(f"Icon file '{icon_path}' not found.")
            if sys.platform == "darwin":  # macOS platform
                cmd.extend(
                    ["--macos-create-app-bundle", f"--macos-app-icon={icon_path}"]
                )
            elif sys.platform == "win32":  # Windows platform
                cmd.extend([f"--windows-icon-from-ico={icon_path}"])
            elif sys.platform == "linux":  # Linux platform
                cmd.append("--linux-onefile")

        if include_data:
            for data in include_data:
                if "*" in data:
                    matches = glob.glob(data.split(":")[0])
                    for match in matches:
                        dest = data.split(":")[1]
                        cmd.extend(
                            [
                                "--include-data-file=" if use_nuitka else "--add-data",
                                f"{match}:{dest}",
                            ]
                        )
                else:
                    cmd.extend(
                        ["--include-data-file=" if use_nuitka else "--add-data", data]
                    )
        if exclude_import is not None:
            if any(exclude_import):
                cmd.extend([f"--nofollow-import-to={','.join(exclude_import)}"])
        if include_import is not None:
            if any(
                include_import
            ):  # are included in the final build. Some packages may require manual inclusion.
                cmd.extend([f"--include-package={','.join(include_import)}"])
        if plugins:
            for plugin in plugins:
                # Adds support for tkinter, ensuring it works correctly in the standalone build.
                cmd.extend([f"--enable-plugin={plugin}"])

        if additional_args:
            cmd.extend(additional_args)

        # # clean
        # cmd.extend(
        #     [   "--noinclude-numba-mode=nofollow", #Prevents the inclusion of the numba library and its dependencies, reducing the executable size.
        #         "--noinclude-dask-mode=nofollow",#Excludes the dask library
        #         "--noinclude-IPython-mode=nofollow",#Excludes the IPython library and its dependencies.
        #         "--noinclude-unittest-mode=nofollow",#Excludes the unittest module (used for testing) from the build
        #         "--noinclude-pytest-mode=nofollow",#Excludes the pytest library (used for testing) from the build.
        #         "--noinclude-setuptools-mode=nofollow",#Excludes setuptools, which is not needed for the standalone executable.
        #         "--lto=no",#Disables Link-Time Optimization (LTO), which reduces the compilation time but may slightly increase the size of the output.
        #     ]
        # )

        if clean_build:
            cmd.append(
                "--remove-output"
            )  # Removes intermediate files created during the build process, keeping only the final executable.
        # Add the script path (final positional argument)
        cmd.append(str(script_path))
        # Ensure Windows shell compatibility
        shell_flag = sys.platform.startswith("win")
        print(f"Running command: ⤵ \n{' '.join(cmd)}\n")
        try:
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                shell=shell_flag,
                check=True,
            )
            if verbose:
                print(result.stdout)
        except subprocess.CalledProcessError as e:
            print(f"Error during packaging:\n{e.stderr}", file=sys.stderr)
            print(" ".join(cmd))
            raise

    print("\nPackaging complete. Check the output directory for the executable.")


def set_theme(
    context="paper",
    style="whitegrid",
    palette="deep",
    font="sans-serif",
    font_scale=1.0,
    color_codes=True,
    grid_alpha=0.5,
    grid_linewidth=0.8,
    grid_linestyle="--",
    tick_direction="out",
    # tick_length=4,
    spine_visibility=False,
    # figsize=(8, 6),
    # linewidth=2,
    dpi=100,
    rc=None,
):
    """
    to configure Seaborn theme with maximum flexibility.

    # Example Usage
    set_sns_theme(font_scale=1.2, grid_alpha=0.8, tick_direction="in", dpi=150)

        Parameters:
        - context: Plotting context ('notebook', 'paper', 'talk', 'poster')
        - style: Style of the plot ('darkgrid', 'whitegrid', 'dark', 'white', 'ticks')
        - palette: Color palette (string or list of colors)
        - font: Font family ('sans-serif', 'serif', etc.)
        - font_scale: Scaling factor for fonts
        - color_codes: Boolean, whether to use seaborn color codes
        - grid_alpha: Opacity of the grid lines
        - grid_linewidth: Thickness of grid lines
        - grid_linestyle: Style of grid lines ('-', '--', '-.', ':')
        - tick_direction: Direction of ticks ('in', 'out', 'inout')
        - tick_length: Length of ticks
        - spine_visibility: Whether to show plot spines (True/False)
        - figsize: Default figure size as tuple (width, height)
        - linewidth: Default line width for plots
        - dpi: Resolution of the figure
        - rc: Dictionary of additional rc settings
    """
    import seaborn as sns
    import matplotlib.pyplot as plt

    # Define additional rc parameters for fine-tuning
    rc_params = {
        # "axes.grid": True,
        "grid.alpha": grid_alpha,
        "grid.linewidth": grid_linewidth,
        "grid.linestyle": grid_linestyle,
        "xtick.direction": tick_direction,
        "ytick.direction": tick_direction,
        # "xtick.major.size": tick_length,
        # "ytick.major.size": tick_length,
        # "axes.linewidth": linewidth,
        # "figure.figsize": figsize,
        "figure.dpi": dpi,
        "axes.spines.top": spine_visibility,
        "axes.spines.right": spine_visibility,
        "axes.spines.bottom": spine_visibility,
        "axes.spines.left": spine_visibility,
    }

    # Merge user-provided rc settings
    if rc:
        rc_params.update(rc)

    # Apply the theme settings
    sns.set_theme(
        context=context,
        style=style,
        palette=palette,
        font=font,
        font_scale=font_scale,
        color_codes=color_codes,
        rc=rc_params,
    )



def df_wide_long(df):
    rows, columns = df.shape
    if columns > rows:
        return "Wide"
    elif rows > columns:
        return "Long"

def df2array(data: pd.DataFrame, x=None, y=None, hue=None, sort=False):
    
    def sort_rows_move_nan(arr, sort=False):
        # Handle edge cases where all values are NaN
        if np.all(np.isnan(arr)):
            return arr  # Return unchanged if the entire array is NaN

        if sort:
            # Replace NaNs with a temporary large value for sorting
            temp_value = (
                np.nanmax(arr[np.isfinite(arr)]) + 1 if np.any(np.isfinite(arr)) else np.inf
            )
            arr_no_nan = np.where(np.isnan(arr), temp_value, arr)

            # Sort each row
            sorted_arr = np.sort(arr_no_nan, axis=1)

            # Move NaNs to the end
            result_arr = np.where(sorted_arr == temp_value, np.nan, sorted_arr)
        else:
            result_rows = []
            for row in arr:
                # Separate non-NaN and NaN values
                non_nan_values = row[~np.isnan(row)]
                nan_count = np.isnan(row).sum()
                # Create a new row with non-NaN values followed by NaNs
                new_row = np.concatenate([non_nan_values, [np.nan] * nan_count])
                result_rows.append(new_row)
            # Convert the list of rows back into a 2D NumPy array
            result_arr = np.array(result_rows)

        # Remove rows/columns that contain only NaNs
        clean_arr = result_arr[~np.isnan(result_arr).all(axis=1)]
        clean_arr_ = clean_arr[:, ~np.isnan(clean_arr).all(axis=0)]

        return clean_arr_
    # data = data.copy()
    # data[y] = pd.to_numeric(data[y], errors="coerce")
    # data = data.dropna(subset=[y])
    if hue is None:
        a = []
        if sort:
            cat_x = np.sort(data[x].unique().tolist()).tolist()
        else:
            cat_x = data[x].unique().tolist()
        for i, x_ in enumerate(cat_x):
            new_ = data.loc[data[x] == x_, y].to_list()
            a = padcat(a, new_, axis=0)
        return sort_rows_move_nan(a).T
    else:
        a = []
        if sort:
            cat_x = np.sort(data[x].unique().tolist()).tolist()
            cat_hue = np.sort(data[hue].unique().tolist()).tolist()
        else:
            cat_x = data[x].unique().tolist()
            cat_hue = data[hue].unique().tolist()
        for i, x_ in enumerate(cat_x):
            for j, hue_ in enumerate(cat_hue):
                new_ = data.loc[(data[x] == x_) & (data[hue] == hue_), y].to_list()
                a = padcat(a, new_, axis=0)
        return sort_rows_move_nan(a).T


def array2df(data: np.ndarray):
    df = pd.DataFrame()
    df["group"] = (
        np.tile(
            ["group" + str(i) for i in range(1, data.shape[1] + 1)], [data.shape[0], 1]
        )
        .reshape(-1, 1, order="F")[:, 0]
        .tolist()
    )
    df["value"] = data.reshape(-1, 1, order="F")
    return df


def padcat(*args, fill_value=np.nan, axis=1, order="row"):
    """
    Concatenate vectors with padding.

    Parameters:
    *args : variable number of list or 1D arrays
        Input arrays to concatenate.
    fill_value : scalar, optional
        The value to use for padding the shorter lists (default is np.nan).
    axis : int, optional
        The axis along which to concatenate (0 for rows, 1 for columns, default is 1).
    order : str, optional
        The order for flattening when required: "row" or "column" (default is "row").

    Returns:
    np.ndarray
        A 2D array with the input arrays concatenated along the specified axis,
        padded with fill_value where necessary.


# Example usage:
a = [1, np.nan]
b = [1, 3, 4, np.nan, 2, np.nan]
c = [1, 2, 3, 4, 5, 6, 7, 8, 10]
d = padcat(a, b)
result1 = padcat(d, c)
result2 = padcat(a, b, c)
print("Result of padcat(d, c):\n", result1)
print("Result of padcat(a, b, c):\n", result2)
    """
    # Set the order for processing
    if "ro" in order.lower():
        order = "C"  # row-major order
    else:
        order = "F"  # column-major order

    # Process input arrays based on their dimensions
    processed_arrays = []
    for arg in args:
        arr = np.asarray(arg)
        if arr.ndim == 1:
            processed_arrays.append(arr)  # Keep 1D arrays as is
        elif arr.ndim == 2:
            if axis == 0:
                # If concatenating along rows, split 2D arrays into 1D arrays row-wise
                processed_arrays.extend(arr)
            elif axis == 1:
                # If concatenating along columns, split 2D arrays into 1D arrays column-wise
                processed_arrays.extend(arr.T)
            else:
                raise ValueError("axis must be 0 or 1")
        else:
            raise ValueError("Input arrays must be 1D or 2D")

    if axis == 0:
        # Concatenate along rows
        max_len = max(arr.size for arr in processed_arrays)
        result = np.full((len(processed_arrays), max_len), fill_value)
        for i, arr in enumerate(processed_arrays):
            result[i, : arr.size] = arr
    elif axis == 1:
        # Concatenate along columns
        max_len = max(arr.size for arr in processed_arrays)
        result = np.full((max_len, len(processed_arrays)), fill_value)
        for i, arr in enumerate(processed_arrays):
            result[: arr.size, i] = arr
    else:
        raise ValueError("axis must be 0 or 1")

    return result

 
# ========== memory cleaner ========== 
import gc 
import psutil 
import weakref
import time
import inspect 
import tracemalloc
from collections import defaultdict

class MemoryOptimizer:
    def __init__(self, 
                 verbose: bool = True, 
                 aggressive_mode: bool = True,
                 track_leaks: bool = False,
                 max_history: int = 100):
        self.verbose = verbose
        self.aggressive_mode = aggressive_mode
        self.track_leaks = track_leaks
        self.max_history = max_history
        self.system = platform.system()
        self.process = psutil.Process(os.getpid())
        self.start_time = time.time()
        self.memory_history = []
        self.leak_tracker = None

        if track_leaks:
            self._setup_leak_tracking()
    
    def _setup_leak_tracking(self):
        self.leak_tracker = {
            'snapshots': [],
            'diff_stats': [],
            'object_types': defaultdict(int),
            'suspected_leaks': []
        }
        tracemalloc.start(25)
    
    def log(self, msg: str, level: str = "INFO"):
        if self.verbose:
            rss = self.process.memory_info().rss / (1024 ** 2)
            elapsed = time.time() - self.start_time
            caller = inspect.currentframe().f_back.f_code.co_name
            print(f"[{level}][{elapsed:.2f}s][{rss:.1f}MB][{caller}] {msg}")
    
    def collect_garbage(self, generations: List[int] = None) -> Dict[str, Any]:
        self.log("Starting deep garbage collection...")
        stats = {
            'collected': defaultdict(int),
            'garbage_cleared': 0,
            'freed_mb': 0
        }

        before_mem = self.process.memory_info().rss

        if self.aggressive_mode:
            gc.set_threshold(1, 1, 1)
            gc.set_debug(gc.DEBUG_SAVEALL)

        gens = generations if generations is not None else [2, 1, 0]
        for gen in gens:
            collected = gc.collect(gen)
            stats['collected'][f'gen_{gen}'] = collected
            self.log(f"GC Gen {gen}: Collected {collected} objects")

        stats['garbage_cleared'] = len(gc.garbage)
        gc.garbage.clear()

        self._clear_weakref_caches()

        after_mem = self.process.memory_info().rss
        stats['freed_mb'] = (before_mem - after_mem) / (1024 ** 2)

        return stats

    def _clear_weakref_caches(self):
        self.log("Clearing weak reference caches...")
        try:
            for obj in gc.get_objects():
                if isinstance(obj, weakref.WeakValueDictionary):
                    obj.clear()
        except Exception as e:
            self.log(f"Failed to clear weakref caches: {e}", "WARNING")
    
    def clear_frameworks(self) -> Dict[str, Any]:
        result = {}

        # PyTorch
        try:
            import torch
            if torch.cuda.is_available():
                self.log("Clearing PyTorch CUDA cache...")
                torch.cuda.empty_cache()
                torch.cuda.ipc_collect()
                result['pytorch'] = {
                    'cuda_cache_cleared': True,
                    'allocated_mb': torch.cuda.memory_allocated() / (1024 ** 2),
                    'cached_mb': torch.cuda.memory_reserved() / (1024 ** 2)
                }
        except Exception as e:
            self.log(f"PyTorch skipped: {e}", "WARNING")
            result['pytorch'] = {'error': str(e)}
        
        # TensorFlow
        try:
            import tensorflow as tf
            self.log("Clearing TensorFlow session...")
            tf.keras.backend.clear_session()
            result['tensorflow'] = {'session_cleared': True}
        except Exception as e:
            self.log(f"TensorFlow skipped: {e}", "WARNING")
            result['tensorflow'] = {'error': str(e)}
        
        # OpenCV
        try:
            import cv2
            self.log("Closing OpenCV windows...")
            cv2.destroyAllWindows()
            result['opencv'] = {'windows_closed': True}
        except Exception as e:
            self.log(f"OpenCV skipped: {e}", "WARNING")
            result['opencv'] = {'error': str(e)}

        # Matplotlib
        try:
            import matplotlib.pyplot as plt
            self.log("Closing matplotlib figures...")
            plt.close('all')
            result['matplotlib'] = {'figures_closed': True}
        except Exception as e:
            self.log(f"Matplotlib skipped: {e}", "WARNING")
            result['matplotlib'] = {'error': str(e)}

        # IPython
        try:
            from IPython import get_ipython
            ipython = get_ipython()
            if ipython is not None:
                self.log("Clearing IPython outputs...")
                ipython.run_line_magic('reset', '-f')
                result['ipython'] = {'outputs_cleared': True}
        except Exception as e:
            self.log(f"IPython skipped: {e}", "WARNING")
            result['ipython'] = {'error': str(e)}

        return result
 

    def profile(self, deep: bool = False) -> Dict[str, Any]:
        mem = self.process.memory_info()
        vm = psutil.virtual_memory()
        swap = psutil.swap_memory()

        profile = {
            'timestamp': time.time(),
            'process': {
                'rss_mb': mem.rss / (1024 ** 2),
                'vms_mb': mem.vms / (1024 ** 2),
            },
            'system': {
                'used_gb': vm.used / (1024 ** 3),
                'available_gb': vm.available / (1024 ** 3),
                'percent': vm.percent,
                'swap_used_gb': swap.used / (1024 ** 3),
                'swap_free_gb': swap.free / (1024 ** 3),
            },
            'gc': {
                'objects': len(gc.get_objects()),
                'garbage': len(gc.garbage),
                'thresholds': gc.get_threshold(),
            }
        }

        if deep:
            profile['deep'] = self._deep_memory_analysis()

        self.memory_history.append(profile)
        if len(self.memory_history) > self.max_history:
            self.memory_history.pop(0)

        return profile

    def _deep_memory_analysis(self) -> Dict[str, Any]:
        self.log("Performing deep memory analysis...")
        type_sizes = defaultdict(int)
        for obj in gc.get_objects():
            try:
                obj_type = type(obj).__name__
                type_sizes[obj_type] += sys.getsizeof(obj)
            except Exception:
                continue

        top_types = sorted(type_sizes.items(), key=lambda x: x[1], reverse=True)[:10]
        return {'top_object_types': top_types}
 
    
    def detect_leaks(self, min_growth_mb: float = 5.0) -> Optional[Dict[str, Any]]:
        """
        Detect potential memory leaks by comparing snapshots.
        
        Args:
            min_growth_mb: Minimum growth in MB to consider a leak
            
        Returns:
            Leak detection report or None if no leaks detected
        """
        if not self.track_leaks or len(self.memory_history) < 2:
            return None
        
        current = self.memory_history[-1]
        previous = self.memory_history[-2]
        
        growth_mb = current['process']['rss_mb'] - previous['process']['rss_mb']
        if growth_mb < min_growth_mb:
            return None
        
        leak_report = {
            'growth_mb': growth_mb,
            'time_elapsed': current['timestamp'] - previous['timestamp'],
            'suspected_causes': [],
        }
        
        # Try to identify potential causes
        if 'deep' in current and 'deep' in previous:
            current_counts = current['deep']['object_counts']
            previous_counts = previous['deep']['object_counts']
            
            for obj_type, count in current_counts.items():
                prev_count = previous_counts.get(obj_type, 0)
                if count > prev_count * 1.5 and count - prev_count > 100:
                    leak_report['suspected_causes'].append({
                        'type': obj_type,
                        'growth': count - prev_count,
                        'percent_growth': ((count - prev_count) / prev_count * 100) if prev_count else float('inf')
                    })
        
        if leak_report['suspected_causes']:
            self.leak_tracker['suspected_leaks'].append(leak_report)
            return leak_report
        
        return None
    
    def optimize(self, full: bool = True) -> Dict[str, Any]:
        """
        Perform comprehensive memory optimization.
        
        Args:
            full: Whether to perform all optimization steps
            
        Returns:
            Dictionary with optimization results
        """
        result = {
            'timestamp': time.time(),
            'before': self.profile(deep=self.track_leaks),
            'steps': {}
        }
        
        # Step 1: Garbage collection
        result['steps']['gc'] = self.collect_garbage()
        
        # Step 2: Framework-specific memory clearing
        result['steps']['frameworks'] = self.clear_frameworks()
        
        # # Step 3: System-level cache clearing
        # if full:
        #     result['steps']['system'] = self.clear_system_caches()
        
        # Step 4: Additional aggressive measures
        if self.aggressive_mode and full:
            result['steps']['aggressive'] = self._aggressive_optimizations()
        
        # Final profile and results
        result['after'] = self.profile(deep=self.track_leaks)
        
        # Calculate savings
        saved_mb = result['before']['process']['rss_mb'] - result['after']['process']['rss_mb']
        result['saved_mb'] = saved_mb
        result['saved_percent'] = (saved_mb / result['before']['process']['rss_mb']) * 100
        
        # Check for leaks if tracking enabled
        if self.track_leaks:
            leak_report = self.detect_leaks()
            if leak_report:
                result['leak_detected'] = leak_report
        
        self.log(
            f"Optimization complete: Saved {saved_mb:.2f} MB "
            f"({result['saved_percent']:.1f}% reduction)",
            "SUCCESS"
        )
        
        return result
     
    def _aggressive_optimizations(self):
        self.log("Aggressively clearing known caches...")

        errors = {}
        try:
            gc.collect()
            self.log("Basic garbage collection done.")
        except Exception as e:
            errors['gc_collect'] = str(e)

        try:
            import numpy as np
            _ = np.empty(0)  # trigger allocation to finalize previous arrays
        except Exception as e:
            errors['numpy'] = str(e)

        try:
            import pandas as pd
            _ = pd.DataFrame()  # no effect but helps ensure cleanup
        except Exception as e:
            errors['pandas'] = str(e)

        return {'status': 'done', 'errors': errors}

    def memory_report(self, detailed: bool = False) -> str:
        """Generate a comprehensive memory usage report."""
        current = self.profile(deep=detailed)
        report = [
            "="*80,
            f"Memory Report (PID: {os.getpid()})",
            "="*80,
            f"Process RSS: {current['process']['rss_mb']:.1f} MB",
            f"Process VMS: {current['process']['vms_mb']:.1f} MB",
            f"System Memory Used: {current['system']['used_gb']:.1f} GB ({current['system']['percent']}%)",
            f"Available Memory: {current['system']['available_gb']:.1f} GB",
            f"Swap Used: {current['system']['swap_used_gb']:.1f} GB",
            f"GC Objects: {current['gc']['objects']:,}",
            f"GC Garbage: {current['gc']['garbage']:,}",
        ]
        
        if detailed and 'deep' in current:
            report.append("\n[Object Type Breakdown (Top 10)]")
            sorted_types = sorted(
                current['deep']['object_counts'].items(),
                key=lambda x: x[1],
                reverse=True
            )[:10]
            
            for obj_type, count in sorted_types:
                size_mb = current['deep']['estimated_sizes'].get(obj_type, 0)
                report.append(f"{obj_type}: {count:,} objects ({size_mb:.2f} MB)")
        
        if self.track_leaks and self.leak_tracker['suspected_leaks']:
            report.append("\n[Potential Memory Leaks]")
            for i, leak in enumerate(self.leak_tracker['suspected_leaks'], 1):
                report.append(
                    f"Leak {i}: +{leak['growth_mb']:.1f}MB in "
                    f"{leak['time_elapsed']:.1f}s"
                )
                for cause in leak['suspected_causes']:
                    report.append(
                        f"  - {cause['type']}: +{cause['growth']:,} "
                        f"({cause['percent_growth']:.1f}%)"
                    )
        
        return "\n".join(report)


def cleaner(
    verbose: bool = True, 
    aggressive: bool = True,
    track_leaks: bool = False,
    full_clean: bool = True,
    return_output:bool=False
) -> Dict[str, Any]:
    """
    Ultimate memory cleaning function with all optimizations.
    
    Args:
        verbose: Print detailed progress information
        aggressive: Use aggressive memory clearing techniques
        track_leaks: Enable memory leak detection
        full_clean: Perform all cleaning steps (including system-level)
    
    Returns:
        Dictionary with optimization results
    """
    optimizer = MemoryOptimizer(
        verbose=verbose,
        aggressive_mode=aggressive,
        track_leaks=track_leaks
    )
    output=optimizer.optimize(full=full_clean)
    return output if return_output else None