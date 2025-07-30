"""
Generic File System Service

This service provides a unified interface for file operations across local and S3 storage systems.
It handles various file types and returns appropriate Python objects for each format.

Features:
- Unified read/write interface for local and S3 storage
- Automatic file type detection and handling
- Support for images, documents, audio, text, and structured data
- Uses polars for CSV/dataframe operations
- Extensible file type handlers
"""

import os
import io
import tempfile
from pathlib import Path
from typing import Union, Any, Dict, BinaryIO, Optional, List, Literal, Callable
from abc import ABC, abstractmethod
import typing
import mimetypes
from datetime import datetime, timezone

# Core dependencies
import polars as pl
from PIL import Image

# Optional dependencies - gracefully handle imports
try:
    import PyPDF2
    import fitz  # PyMuPDF
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import librosa
    import soundfile as sf
    HAS_AUDIO = True
except ImportError:
    HAS_AUDIO = False

try:
    import openpyxl
    import xlrd
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import markdown
    HAS_MARKDOWN = True
except ImportError:
    HAS_MARKDOWN = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from percolate.services.S3Service import S3Service
from percolate.utils import logger


class FileSystemProvider(ABC):
    """Abstract base class for file system providers"""
    
    @abstractmethod
    def exists(self, path: str) -> bool:
        pass
    
    @abstractmethod
    def read_bytes(self, path: str) -> bytes:
        pass
    
    @abstractmethod
    def write_bytes(self, path: str, data: bytes) -> None:
        pass
    
    @abstractmethod
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        pass
    
    @abstractmethod
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        pass


class LocalFileSystemProvider(FileSystemProvider):
    """Local file system provider"""
    
    def _normalize_path(self, path: str) -> str:
        """Normalize path by removing file:// prefix if present"""
        if path.startswith('file://'):
            return path[7:]  # Remove 'file://' prefix
        return path
    
    def exists(self, path: str) -> bool:
        return Path(self._normalize_path(path)).exists()
    
    def read_bytes(self, path: str) -> bytes:
        normalized_path = self._normalize_path(path)
        with open(normalized_path, 'rb') as f:
            return f.read()
    
    def write_bytes(self, path: str, data: bytes) -> None:
        normalized_path = self._normalize_path(path)
        Path(normalized_path).parent.mkdir(parents=True, exist_ok=True)
        with open(normalized_path, 'wb') as f:
            f.write(data)
    
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        normalized_path = self._normalize_path(path)
        with open(normalized_path, 'r', encoding=encoding) as f:
            return f.read()
    
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        normalized_path = self._normalize_path(path)
        Path(normalized_path).parent.mkdir(parents=True, exist_ok=True)
        with open(normalized_path, 'w', encoding=encoding) as f:
            f.write(text)


class S3FileSystemProvider(FileSystemProvider):
    """S3 file system provider using existing S3Service"""
    
    def __init__(self, s3_service: Optional[S3Service] = None):
        self.s3_service = s3_service or S3Service()
    
    def exists(self, path: str) -> bool:
        try:
            parsed = self.s3_service.parse_s3_uri(path)
            bucket_name = parsed["bucket"]
            object_key = parsed["key"]
            self.s3_service.s3_client.head_object(Bucket=bucket_name, Key=object_key)
            return True
        except:
            return False
    
    def read_bytes(self, path: str) -> bytes:
        result = self.s3_service.download_file_from_uri(path)
        return result["content"]
    
    def write_bytes(self, path: str, data: bytes) -> None:
        self.s3_service.upload_filebytes_to_uri(path, data)
    
    def read_text(self, path: str, encoding: str = 'utf-8') -> str:
        data = self.read_bytes(path)
        return data.decode(encoding)
    
    def write_text(self, path: str, text: str, encoding: str = 'utf-8') -> None:
        data = text.encode(encoding)
        self.write_bytes(path, data)


class FileTypeHandler(ABC):
    """Abstract base class for file type handlers"""
    
    @abstractmethod
    def can_handle(self, file_path: str) -> bool:
        pass
    
    @abstractmethod
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Any:
        pass
    
    @abstractmethod
    def write(self, provider: FileSystemProvider, file_path: str, data: Any, **kwargs) -> None:
        pass


class ImageHandler(FileTypeHandler):
    """Handler for image files (PNG, JPG, JPEG, TIFF, etc.)"""
    
    SUPPORTED_FORMATS = {'.png', '.jpg', '.jpeg', '.tiff', '.tif', '.bmp', '.gif', '.webp'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Image.Image:
        data = provider.read_bytes(file_path)
        return Image.open(io.BytesIO(data))
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Image.Image, **kwargs) -> None:
        buffer = io.BytesIO()
        format_type = kwargs.get('format', Path(file_path).suffix[1:].upper())
        if format_type.upper() == 'JPG':
            format_type = 'JPEG'
        data.save(buffer, format=format_type, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class TextHandler(FileTypeHandler):
    """Handler for text files (TXT, MD, HTML, etc.)"""
    
    SUPPORTED_FORMATS = {'.txt', '.md', '.html', '.htm', '.css', '.js', '.py', '.json', '.xml', '.yaml', '.yml'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> str:
        encoding = kwargs.get('encoding', 'utf-8')
        return provider.read_text(file_path, encoding)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: str, **kwargs) -> None:
        encoding = kwargs.get('encoding', 'utf-8')
        provider.write_text(file_path, data, encoding)


class CSVHandler(FileTypeHandler):
    """Handler for CSV files using Polars"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.csv'
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> pl.DataFrame:
        data = provider.read_bytes(file_path)
        return pl.read_csv(io.BytesIO(data), **kwargs)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: pl.DataFrame, **kwargs) -> None:
        buffer = io.BytesIO()
        data.write_csv(buffer, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class ParquetHandler(FileTypeHandler):
    """Handler for Parquet files using Polars"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.parquet'
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> pl.DataFrame:
        data = provider.read_bytes(file_path)
        return pl.read_parquet(io.BytesIO(data), **kwargs)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: pl.DataFrame, **kwargs) -> None:
        buffer = io.BytesIO()
        data.write_parquet(buffer, **kwargs)
        provider.write_bytes(file_path, buffer.getvalue())


class PDFHandler(FileTypeHandler):
    """Handler for PDF files"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.pdf' and HAS_PDF
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Read PDF and return a dictionary with text content and metadata.
        Enhanced version based on the existing PDF parser.
        """
        if not HAS_PDF:
            raise ImportError("PDF support requires PyPDF2 and PyMuPDF: pip install PyPDF2 PyMuPDF")
        
        data = provider.read_bytes(file_path)
        pdf_stream = io.BytesIO(data)
        
        # Extract text using PyPDF2
        pdf_reader = PyPDF2.PdfReader(stream=pdf_stream)
        text_pages = []
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text = page.extract_text()
            text_pages.append(text.replace('\n \n', ' '))  # Clean text
        
        # Extract images using PyMuPDF
        pdf_stream.seek(0)  # Reset stream
        images = []
        image_info = []
        
        with fitz.open(stream=pdf_stream) as pdf_document:
            for page_num in range(pdf_document.page_count):
                page = pdf_document.load_page(page_num)
                page_images = []
                
                for img in page.get_images(full=True):
                    xref = img[0]
                    base_image = pdf_document.extract_image(xref)
                    image_bytes = base_image["image"]
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    # Filter out small images (likely logos/decorations)
                    min_size = kwargs.get('min_image_size', (300, 300))
                    if image.size[0] >= min_size[0] and image.size[1] >= min_size[1]:
                        page_images.append(image)
                
                images.append(page_images)
                image_info.append(page.get_image_info())
        
        return {
            'text_pages': text_pages,
            'images': images,
            'image_info': image_info,
            'num_pages': len(text_pages),
            'metadata': {
                'title': pdf_reader.metadata.get('/Title', '') if pdf_reader.metadata else '',
                'author': pdf_reader.metadata.get('/Author', '') if pdf_reader.metadata else '',
                'subject': pdf_reader.metadata.get('/Subject', '') if pdf_reader.metadata else '',
            }
        }
    
    def write(self, provider: FileSystemProvider, file_path: str, data: bytes, **kwargs) -> None:
        """Write PDF bytes to file"""
        provider.write_bytes(file_path, data)


class AudioHandler(FileTypeHandler):
    """Handler for audio files (WAV, MP3, etc.)"""
    
    SUPPORTED_FORMATS = {'.wav', '.mp3', '.flac', '.ogg', '.m4a'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS and HAS_AUDIO
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read audio file and return audio data with metadata"""
        if not HAS_AUDIO:
            raise ImportError("Audio support requires librosa and soundfile: pip install librosa soundfile")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file for librosa (it needs file path)
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            # Load audio data
            sr = kwargs.get('sr', None)  # Sample rate
            audio_data, sample_rate = librosa.load(tmp_path, sr=sr)
            
            # Get metadata
            info = sf.info(tmp_path)
            
            return {
                'audio_data': audio_data,
                'sample_rate': sample_rate,
                'duration': info.duration,
                'channels': info.channels,
                'format': info.format,
                'subtype': info.subtype
            }
        finally:
            os.unlink(tmp_path)  # Clean up temp file
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Dict[str, Any], **kwargs) -> None:
        """Write audio data to file"""
        if not HAS_AUDIO:
            raise ImportError("Audio support requires soundfile: pip install soundfile")
        
        audio_data = data['audio_data']
        sample_rate = data['sample_rate']
        
        # Write to temporary file first
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            sf.write(tmp_path, audio_data, sample_rate, **kwargs)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class ExcelHandler(FileTypeHandler):
    """Handler for Excel files (XLS, XLSX) using Polars"""
    
    SUPPORTED_FORMATS = {'.xls', '.xlsx'}
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED_FORMATS and HAS_EXCEL
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, pl.DataFrame]:
        """Read Excel file and return dictionary of sheet name -> DataFrame"""
        if not HAS_EXCEL:
            raise ImportError("Excel support requires openpyxl: pip install openpyxl")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=Path(file_path).suffix, delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            import pandas as pd
            
            # Filter out kwargs that aren't for pandas.read_excel
            pandas_kwargs = {k: v for k, v in kwargs.items() 
                           if k not in ['mode']}  # Remove 'mode' and other non-pandas args
            
            # Try multiple engines in order of preference for performance and reliability
            engines_to_try = []
            
            # Check if calamine is available (fastest, but newer)
            try:
                import python_calamine
                engines_to_try.append('calamine')
            except ImportError:
                pass
            
            # Always have openpyxl as fallback
            engines_to_try.append('openpyxl')
            
            excel_data = None
            last_error = None
            
            for engine in engines_to_try:
                try:
                    logger.info(f"Attempting to read Excel file with {engine} engine")
                    excel_data = pd.read_excel(tmp_path, sheet_name=None, engine=engine, **pandas_kwargs)
                    logger.info(f"Successfully read Excel file with {engine} engine")
                    break
                except Exception as e:
                    logger.warning(f"Failed to read Excel with {engine} engine: {e}")
                    last_error = e
                    continue
            
            if excel_data is None:
                raise last_error or Exception("Failed to read Excel file with any available engine")
            
            # Convert to Polars DataFrames with robust error handling
            result = {}
            for sheet_name, df in excel_data.items():
                try:
                    # First, try direct conversion
                    result[sheet_name] = pl.from_pandas(df)
                except Exception as e:
                    logger.warning(f"Direct Polars conversion failed for sheet '{sheet_name}': {e}")
                    try:
                        # Fallback 1: Handle NaN values by filling with empty strings
                        df_filled = df.fillna('')
                        result[sheet_name] = pl.from_pandas(df_filled)
                    except Exception as e2:
                        logger.warning(f"NaN-filled conversion failed for sheet '{sheet_name}': {e2}")
                        try:
                            # Fallback 2: Convert all to string to avoid PyArrow type issues
                            df_str = df.astype(str)
                            result[sheet_name] = pl.from_pandas(df_str)
                        except Exception as e3:
                            logger.error(f"All conversion methods failed for sheet '{sheet_name}': {e3}")
                            # Return the original pandas DataFrame if Polars conversion fails completely
                            result[sheet_name] = df
            
            return result
        finally:
            os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Dict[str, pl.DataFrame], **kwargs) -> None:
        """Write dictionary of DataFrames to Excel sheets"""
        import tempfile
        
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            # Convert to pandas and write
            import pandas as pd
            with pd.ExcelWriter(tmp_path, engine='openpyxl') as writer:
                for sheet_name, df in data.items():
                    pandas_df = df.to_pandas()
                    pandas_df.to_excel(writer, sheet_name=sheet_name, index=False, **kwargs)
            
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class DocxHandler(FileTypeHandler):
    """Handler for DOCX files"""
    
    def can_handle(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() == '.docx' and HAS_DOCX
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """Read DOCX file and extract text content"""
        if not HAS_DOCX:
            raise ImportError("DOCX support requires python-docx: pip install python-docx")
        
        data = provider.read_bytes(file_path)
        
        # Save to temporary file
        import tempfile
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
            tmp_file.write(data)
            tmp_path = tmp_file.name
        
        try:
            doc = Document(tmp_path)
            
            # Extract paragraphs
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            return {
                'paragraphs': paragraphs,
                'tables': tables,
                'full_text': '\n'.join(paragraphs)
            }
        finally:
            os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: str, **kwargs) -> None:
        """Create a simple DOCX file with the provided text"""
        if not HAS_DOCX:
            raise ImportError("DOCX support requires python-docx: pip install python-docx")
        
        import tempfile
        
        doc = Document()
        doc.add_paragraph(data)
        
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            doc.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class PPTXHandler(FileTypeHandler):
    """Handler for PowerPoint presentations (.pptx, .ppt)"""
    
    def can_handle(self, file_path: str) -> bool:
        """Check if this handler can process the file"""
        ext = Path(file_path).suffix.lower()
        return ext in ['.pptx', '.ppt']
    
    def read(self, provider: FileSystemProvider, file_path: str, **kwargs) -> Dict[str, Any]:
        """
        Read PPTX file and return structured data with text and optionally images.
        
        Args:
            provider: File system provider
            file_path: Path to the PPTX file
            mode: 'simple' for text only, 'enriched' for text + LLM image analysis
            max_images: Maximum number of images to analyze (default: 50)
            
        Returns:
            Dict with slides, text_content, and optionally image_analysis
        """
        if not HAS_PPTX:
            raise ImportError("PPTX support requires python-pptx: pip install python-pptx")
        
        # Get enhanced PPTX provider from our parsing utilities
        try:
            from percolate.utils.parsing.providers import PPTXContentProvider
            pptx_provider = PPTXContentProvider()
            
            # Use temporary file for PPTX processing
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                tmp_file.write(provider.read_bytes(file_path))
            
            try:
                mode = kwargs.get('mode', 'simple')
                max_images = kwargs.get('max_images', 50)
                enriched = (mode == 'enriched')
                
                # Use our enhanced provider
                text_content = pptx_provider.extract_text(tmp_path, enriched=enriched, max_images=max_images)
                
                # Also get slide-by-slide breakdown
                prs = Presentation(tmp_path)
                slides = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slides.append({
                        'slide_number': slide_num,
                        'text': slide_text,
                        'combined_text': '\n'.join(slide_text)
                    })
                
                result = {
                    'slides': slides,
                    'text_content': text_content,
                    'total_slides': len(slides),
                    'mode': mode,
                    'file_path': file_path
                }
                
                # Add metadata if enriched mode and images were analyzed
                if enriched and '=== SLIDE IMAGE ANALYSIS ===' in text_content:
                    result['has_image_analysis'] = True
                    result['max_images_processed'] = max_images
                else:
                    result['has_image_analysis'] = False
                
                return result
                
            finally:
                os.unlink(tmp_path)
                
        except ImportError:
            # Fallback to basic PPTX handling without our enhanced provider
            logger.warning("Enhanced PPTX provider not available, using basic extraction")
            
            import tempfile
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_path = tmp_file.name
                tmp_file.write(provider.read_bytes(file_path))
            
            try:
                prs = Presentation(tmp_path)
                slides = []
                all_text = []
                
                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    slides.append({
                        'slide_number': slide_num,
                        'text': slide_text,
                        'combined_text': '\n'.join(slide_text)
                    })
                    all_text.extend(slide_text)
                
                return {
                    'slides': slides,
                    'text_content': '\n\n'.join(all_text),
                    'total_slides': len(slides),
                    'mode': 'simple',
                    'has_image_analysis': False,
                    'file_path': file_path
                }
                
            finally:
                os.unlink(tmp_path)
    
    def write(self, provider: FileSystemProvider, file_path: str, data: Any, **kwargs) -> None:
        """
        Write PPTX file (limited functionality - creates basic presentation)
        
        Args:
            provider: File system provider
            file_path: Path where to write the PPTX file
            data: Either string (single slide) or list of strings (multiple slides)
        """
        if not HAS_PPTX:
            raise ImportError("PPTX support requires python-pptx: pip install python-pptx")
        
        import tempfile
        
        prs = Presentation()
        
        if isinstance(data, str):
            # Single slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
            slide.shapes.title.text = "Generated Slide"
            slide.shapes.placeholders[1].text = data
        elif isinstance(data, list):
            # Multiple slides
            for i, slide_content in enumerate(data):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Slide {i + 1}"
                slide.shapes.placeholders[1].text = str(slide_content)
        else:
            raise ValueError(f"Unsupported data type for PPTX: {type(data)}")
        
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            tmp_path = tmp_file.name
        
        try:
            prs.save(tmp_path)
            with open(tmp_path, 'rb') as f:
                provider.write_bytes(file_path, f.read())
        finally:
            os.unlink(tmp_path)


class FileSystemService:
    """
    Unified file system service that provides a single interface for file operations
    across local and S3 storage with automatic file type detection and handling.
    """
    
    def __init__(self, s3_service: Optional[S3Service] = None):
        self.s3_service = s3_service
        self._providers = {}
        self._handlers = []
        
        # Register default handlers
        self.register_handler(ImageHandler())
        self.register_handler(TextHandler())
        self.register_handler(CSVHandler())
        self.register_handler(ParquetHandler())
        
        if HAS_PDF:
            self.register_handler(PDFHandler())
        if HAS_AUDIO:
            self.register_handler(AudioHandler())
        if HAS_EXCEL:
            self.register_handler(ExcelHandler())
        if HAS_DOCX:
            self.register_handler(DocxHandler())
        if HAS_PPTX:
            self.register_handler(PPTXHandler())
    
    def register_handler(self, handler: FileTypeHandler):
        """Register a new file type handler"""
        self._handlers.append(handler)
    
    def _get_provider(self, path: str) -> FileSystemProvider:
        """Get the appropriate provider for the given path"""
        if path.startswith('s3://'):
            if 's3' not in self._providers:
                self._providers['s3'] = S3FileSystemProvider(self.s3_service)
            return self._providers['s3']
        else:
            if 'local' not in self._providers:
                self._providers['local'] = LocalFileSystemProvider()
            return self._providers['local']
    
    def _get_handler(self, file_path: str) -> Optional[FileTypeHandler]:
        """Get the appropriate handler for the given file path"""
        for handler in self._handlers:
            if handler.can_handle(file_path):
                return handler
        return None
    
    def exists(self, path: str) -> bool:
        """Check if a file exists"""
        provider = self._get_provider(path)
        return provider.exists(path)
    
    def read_bytes(self, path: str) -> bytes:
        """
        Read a file as raw bytes, bypassing any handlers.
        Useful for copying files without processing them.
        """
        provider = self._get_provider(path)
        return provider.read_bytes(path)
    
    def read(self, path: str, **kwargs) -> Any:
        """
        Read a file and return the appropriate Python object based on file type.
        
        Args:
            path: File path (local or s3://)
            **kwargs: Additional arguments passed to the specific handler
            
        Returns:
            Appropriate Python object:
            - Images: PIL.Image
            - Text files: str
            - CSV/Parquet: polars.DataFrame
            - PDF: dict with text_pages, images, metadata
            - Audio: dict with audio_data, sample_rate, metadata
            - Excel: dict of sheet_name -> polars.DataFrame
            - DOCX: dict with paragraphs, tables, full_text
            - Unknown types: bytes
        """
        logger.info(f"Reading file: {path}")
        
        provider = self._get_provider(path)
        handler = self._get_handler(path)
        
        if handler:
            try:
                result = handler.read(provider, path, **kwargs)
                
                # Handle extended mode for PDFs
                if isinstance(handler, PDFHandler):
                    mode = kwargs.get('mode', 'simple')
                    if mode == 'extended':
                        # Add raw bytes for page conversion
                        if isinstance(result, dict):
                            result['raw_bytes'] = provider.read_bytes(path)
                        
                        file_name = Path(path).name
                        extended_content = self._extract_extended_pdf_content(result, file_name, path)
                        
                        # Return enhanced result with extended content
                        if isinstance(result, dict):
                            result['content'] = extended_content
                            result['mode'] = 'extended'
                        
                        return result
                    else:
                        # Simple mode - add basic content field
                        if isinstance(result, dict):
                            text_content = '\n'.join(result.get('text_pages', []))
                            result['content'] = text_content
                            result['mode'] = 'simple'
                        
                        return result
                
                return result
            except Exception as e:
                logger.warning(f"Handler failed for {path}: {e}. Falling back to raw bytes.")
                return provider.read_bytes(path)
        else:
            logger.warning(f"No specific handler for {path}. Returning raw bytes.")
            return provider.read_bytes(path)
    
    def write(self, path: str, data: Any, **kwargs) -> None:
        """
        Write data to a file using the appropriate handler based on file type.
        
        Args:
            path: File path (local or s3://)
            data: Data to write (type depends on file format)
            **kwargs: Additional arguments passed to the specific handler
        """
        logger.info(f"Writing file: {path}")
        
        provider = self._get_provider(path)
        handler = self._get_handler(path)
        
        # Special case: if we're trying to write parsed PDF content,
        # we need to read the original bytes instead
        if (isinstance(data, dict) and 'text_pages' in data and 
            Path(path).suffix.lower() == '.pdf'):
            logger.warning(f"Cannot write parsed PDF content to {path}. PDF writing requires original bytes.")
            raise ValueError("Cannot write parsed PDF content. PDF files must be written as bytes.")
        
        # If we have raw bytes and a handler that expects processed data,
        # just write the bytes directly to preserve file integrity
        if isinstance(data, bytes):
            provider.write_bytes(path, data)
        elif handler:
            handler.write(provider, path, data, **kwargs)
        else:
            # Fallback: treat as text
            if isinstance(data, str):
                provider.write_text(path, data)
            else:
                raise ValueError(f"No handler for file type {Path(path).suffix} and data type {type(data)}")
    
    def copy(self, source_path: str, dest_path: str, **kwargs) -> None:
        """Copy a file from source to destination"""
        logger.info(f"Copying file: {source_path} -> {dest_path}")
        
        # Read from source and write to destination
        data = self.read(source_path, **kwargs)
        self.write(dest_path, data, **kwargs)
    
    def get_file_info(self, path: str) -> Dict[str, Any]:
        """Get information about a file"""
        provider = self._get_provider(path)
        
        if not provider.exists(path):
            return {'exists': False}
        
        # Get basic info
        info = {
            'exists': True,
            'path': path,
            'extension': Path(path).suffix.lower(),
            'name': Path(path).name,
            'storage_type': 's3' if path.startswith('s3://') else 'local'
        }
        
        # Try to get file size
        try:
            if path.startswith('s3://'):
                # For S3, we'd need to get object metadata
                pass  # Could implement S3 head_object here
            else:
                info['size'] = Path(path).stat().st_size
        except:
            pass
        
        # Check if we have a handler for this file type
        handler = self._get_handler(path)
        info['has_handler'] = handler is not None
        info['handler_type'] = type(handler).__name__ if handler else None
        
        # Get MIME type
        mime_type, _ = mimetypes.guess_type(path)
        info['mime_type'] = mime_type
        
        return info
    
    def apply(self, uri: str, fn: typing.Callable, **kwargs) -> typing.Any:
        """
        Apply a function to a file that requires a local file path.
        
        This method handles the case where a library function needs a file path
        (not bytes or file-like object) but we want to work with our unified
        file system interface including S3 files.
        
        Args:
            uri: File URI (local file://, S3 s3://, or HTTP/HTTPS URL)
            fn: Function that takes a file path as its first argument
            **kwargs: Additional arguments to pass to the function
            
        Returns:
            Whatever the function returns
            
        Example:
            from pdf2image import convert_from_path
            images = fs.apply("s3://bucket/document.pdf", convert_from_path, dpi=200)
        """
        import tempfile
        
        # For local files, we can pass the path directly if it's a simple file:// URI
        if uri.startswith('file://'):
            local_path = uri[7:]  # Remove 'file://' prefix
            if os.path.exists(local_path):
                return fn(local_path, **kwargs)
        elif not uri.startswith(('s3://', 'http://', 'https://')):
            # Assume it's a local path
            if os.path.exists(uri):
                return fn(uri, **kwargs)
        
        # For S3 or remote files, download to temp file
        provider = self._get_provider(uri)
        file_data = provider.read_bytes(uri)
        
        # Get file extension from URI
        file_extension = Path(uri).suffix
        
        # Create temporary file with same extension
        with tempfile.NamedTemporaryFile(
            suffix=file_extension, 
            prefix="fs_apply_", 
            delete=False
        ) as temp_file:
            temp_file.write(file_data)
            temp_file.flush()
            temp_path = temp_file.name
        
        try:
            # Apply the function to the temporary file
            result = fn(temp_path, **kwargs)
            return result
        finally:
            # Clean up temporary file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def _extract_extended_pdf_content(self, pdf_data: Dict[str, Any], file_name: str, uri: str = None) -> str:
        """Extract content from PDF using LLM vision analysis of page images."""
        try:
            # Get image interpreter service
            from percolate.services.llm.ImageInterpreter import get_image_interpreter
            interpreter = get_image_interpreter()
            
            if not interpreter.is_available():
                logger.warning("Image interpreter not available, falling back to simple PDF parsing")
                return self._extract_simple_content(pdf_data, 'pdf')
            
            # Convert PDF pages to images for LLM analysis
            page_images = self._convert_pdf_pages_to_images(pdf_data, file_name, uri)
            
            if not page_images:
                logger.warning("No page images generated, falling back to simple PDF parsing")
                return self._extract_simple_content(pdf_data, 'pdf')
            
            logger.info(f"Analyzing {len(page_images)} PDF pages with LLM vision")
            
            # Analyze each page with LLM
            analyzed_pages = []
            for i, page_image in enumerate(page_images):
                try:
                    prompt = """
                    Analyze this PDF page image in detail. Extract and describe:
                    1. All text content (transcribe exactly what you see)
                    2. Any tables, charts, or structured data
                    3. Images, diagrams, or visual elements 
                    4. Layout and formatting structure
                    5. Any important visual information not captured in plain text
                    
                    Provide a comprehensive description that captures both the textual content and visual elements.
                    """
                    
                    result = interpreter.describe_images(
                        images=page_image,
                        prompt=prompt,
                        context=f"PDF page {i+1} from document '{file_name}'",
                        max_tokens=2000
                    )
                    
                    if result["success"]:
                        page_content = f"=== PAGE {i+1} ===\n{result['content']}\n"
                        analyzed_pages.append(page_content)
                        logger.info(f"Successfully analyzed page {i+1}")
                    else:
                        logger.warning(f"Failed to analyze page {i+1}: {result.get('error', 'Unknown error')}")
                        # Fallback to simple text for this page
                        if i < len(pdf_data.get('text_pages', [])):
                            simple_text = pdf_data['text_pages'][i]
                            page_content = f"=== PAGE {i+1} (TEXT ONLY) ===\n{simple_text}\n"
                            analyzed_pages.append(page_content)
                
                except Exception as e:
                    logger.error(f"Error analyzing page {i+1}: {str(e)}")
                    # Fallback to simple text for this page
                    if i < len(pdf_data.get('text_pages', [])):
                        simple_text = pdf_data['text_pages'][i]
                        page_content = f"=== PAGE {i+1} (TEXT ONLY) ===\n{simple_text}\n"
                        analyzed_pages.append(page_content)
            
            # Combine all analyzed pages
            full_content = "\n".join(analyzed_pages)
            
            # Add summary information
            summary = f"""
DOCUMENT ANALYSIS SUMMARY:
- Document: {file_name}
- Total Pages: {len(page_images)}
- Analysis Method: LLM Vision + Text Extraction
- Pages Successfully Analyzed: {len([p for p in analyzed_pages if 'TEXT ONLY' not in p])}

FULL CONTENT:
{full_content}
"""
            
            logger.info(f"Extended PDF analysis complete: {len(full_content)} characters")
            return summary
            
        except Exception as e:
            logger.error(f"Error in extended PDF processing: {str(e)}")
            logger.info("Falling back to simple PDF parsing")
            return self._extract_simple_content(pdf_data, 'pdf')
    
    def _convert_pdf_pages_to_images(self, pdf_data: Dict[str, Any], file_name: str, uri: str = None) -> List[Image.Image]:
        """Convert PDF pages to PIL Images for LLM analysis"""
        pdf2image_failed = False
        fitz_failed = False
        
        try:
            # Try pdf2image first (requires poppler)
            # Prefer convert_from_path if we have a URI, as it's more reliable
            if uri:
                from pdf2image import convert_from_path
                logger.info(f"Converting PDF pages to images using pdf2image convert_from_path for {file_name}")
                images = self.apply(uri, convert_from_path)
                logger.info(f"Successfully converted {len(images)} pages to images using convert_from_path")
                return images
            else:
                # Fallback to convert_from_bytes if no URI provided
                from pdf2image import convert_from_bytes
                raw_bytes = pdf_data.get('raw_bytes')
                if raw_bytes:
                    logger.info(f"Converting PDF pages to images using pdf2image convert_from_bytes for {file_name}")
                    images = convert_from_bytes(raw_bytes)
                    logger.info(f"Successfully converted {len(images)} pages to images using convert_from_bytes")
                    return images
        except ImportError:
            logger.warning("pdf2image not available, falling back to fitz rendering")
            pdf2image_failed = True
        except Exception as e:
            logger.warning(f"pdf2image conversion failed: {e}, falling back to fitz rendering")
            pdf2image_failed = True
        
        # Fallback to fitz rendering
        try:
            import fitz
            
            raw_bytes = pdf_data.get('raw_bytes')
            if raw_bytes:
                logger.info(f"Converting PDF pages to images using fitz for {file_name}")
                pdf_document = fitz.open(stream=raw_bytes, filetype="pdf")
                images = []
                
                for page_num in range(pdf_document.page_count):
                    page = pdf_document.load_page(page_num)
                    # Render page as image (default DPI is 72, increase for better quality)
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scale for better quality
                    img_data = pix.tobytes("png")
                    image = Image.open(io.BytesIO(img_data))
                    images.append(image)
                
                pdf_document.close()
                logger.info(f"Successfully converted {len(images)} pages to images using fitz")
                return images
        except ImportError:
            logger.error("fitz (PyMuPDF) not available for PDF page rendering")
            fitz_failed = True
        except Exception as e:
            logger.error(f"Fitz PDF page conversion failed: {e}")
            fitz_failed = True
        
        # If both methods failed, raise an exception in extended mode
        if pdf2image_failed and fitz_failed:
            raise Exception(
                f"Failed to convert PDF pages to images for '{file_name}'. "
                "Both pdf2image and fitz (PyMuPDF) conversion methods failed. "
                "Extended PDF processing requires successful page-to-image conversion. "
                "Please ensure poppler-utils and/or PyMuPDF are properly installed."
            )
        
        # This should not be reached, but just in case
        return []
    
    def _extract_simple_content(self, file_data: Any, file_type: str) -> str:
        """Extract simple text content from file data."""
        if file_type == 'pdf':
            # PDF data - combine text pages
            text_pages = file_data.get('text_pages', [])
            return '\n'.join(text_pages)
        elif file_type == 'text':
            # Already a string
            return file_data if isinstance(file_data, str) else str(file_data)
        elif file_type == 'csv':
            # CSV data - convert to simple text format
            if isinstance(file_data, dict) and 'data' in file_data:
                rows = file_data['data']
                if not rows:
                    return ""
                # Create simple text representation
                lines = []
                for row in rows:
                    line = ', '.join(str(cell) for cell in row.values())
                    lines.append(line)
                return '\n'.join(lines)
            else:
                # Polars DataFrame
                return str(file_data)
        elif file_type in ['image']:
            # For images, we can't extract simple text
            return f"[Image file: {file_type}]"
        else:
            # Try to convert to string
            if isinstance(file_data, dict):
                # Other structured data - convert to simple string representation
                import json
                return json.dumps(file_data, indent=2, default=str)
        
        # Unknown format - convert to string
        return str(file_data)
    
    def read_chunks(self, path: str, mode: str = 'simple', target_model=None, **kwargs):
        """
        Read a file and yield chunked Resources using the integrated ResourceChunker.
        
        This is a convenient method that combines file reading and chunking in one call.
        
        Args:
            path: File path (local or s3://)
            mode: Processing mode - 'simple' (fast, text-based) or 'extended' (LLM-enhanced, expensive)
            target_model: Custom target model class for the chunks (default: Resources)
            **kwargs: Additional arguments passed to the chunker:
                - chunk_size: Maximum size of each chunk (default: 1000)
                - chunk_overlap: Overlap between chunks (default: 200) 
                - max_chunks: Maximum number of chunks to create (default: None)
                - save_to_db: Whether to save chunks to database (default: False)
                - For audio files:
                    - max_file_size_mb: Maximum file size for processing (default: 250)
                    - chunk_duration_minutes: Audio chunk duration (default: 10)
                
        Yields:
            Resources: Individual chunked Resources ready for use
            
        Example:
            # Simple chunking (fast) - iterate over chunks
            for chunk in fs.read_chunks("document.pdf"):
                print(f"Chunk: {chunk.content[:100]}...")
            
            # Extended chunking with LLM analysis (expensive but comprehensive)  
            chunks = list(fs.read_chunks("document.pdf", mode='extended', chunk_size=500))
            
            # Audio chunking with custom settings
            for chunk in fs.read_chunks("audio.wav", chunk_duration_minutes=5):
                process_audio_chunk(chunk)
            
            # Custom model with database saving
            for chunk in fs.read_chunks("spreadsheet.csv", target_model=MyCustomModel, save_to_db=True):
                handle_chunk(chunk)
        """
        logger.info(f"Reading and chunking file: {path} (mode: {mode})")
        
        # Import here to avoid circular imports
        from percolate.models.p8.types import Resources
        
        # Use custom model or default to Resources
        model_class = target_model or Resources
        
        # Use the integrated ResourceChunker
        chunker = ResourceChunker(self)
        
        # Create chunked resources using the integrated ResourceChunker to avoid circular calls
        try:
            # Always use ResourceChunker directly to avoid circular dependency with Resources.chunked_resource
            chunker = ResourceChunker(self)
            
            # Create chunks directly using the ResourceChunker
            chunks = chunker.chunk_resource_from_uri(
                uri=path,
                parsing_mode=mode,
                chunk_size=kwargs.get('chunk_size', 1000),
                chunk_overlap=kwargs.get('chunk_overlap', 200),
                user_id=kwargs.get('userid'),
                metadata=kwargs.get('metadata')
            )
            
            # Override any additional properties if provided and the model supports them
            if chunks:
                for chunk in chunks:
                    # Override category if provided
                    if kwargs.get('category'):
                        chunk.category = kwargs['category']
                    
                    # Override name if provided
                    if kwargs.get('name'):
                        chunk_index = chunk.metadata.get('chunk_index', 0) if chunk.metadata else 0
                        if chunk_index > 0:
                            chunk.name = f"{kwargs['name']} (chunk {chunk_index + 1})"
                        else:
                            chunk.name = kwargs['name']
                    
                    # If using a custom model class, convert the chunk
                    if model_class and model_class != type(chunk):
                        # Try to create an instance of the custom model with the chunk data
                        try:
                            chunk_dict = chunk.model_dump() if hasattr(chunk, 'model_dump') else chunk.__dict__
                            # Filter to only include fields that the custom model accepts
                            filtered_dict = {k: v for k, v in chunk_dict.items() if hasattr(model_class, k)}
                            custom_chunk = model_class(**filtered_dict)
                            chunks[chunks.index(chunk)] = custom_chunk
                        except Exception as e:
                            logger.warning(f"Failed to convert chunk to {model_class.__name__}: {e}")
                            # Keep the original chunk if conversion fails
            
            logger.info(f"Successfully created {len(chunks)} chunks from {path}")
            
            # Yield each chunk
            for chunk in chunks:
                yield chunk
                
        except Exception as e:
            logger.error(f"Error chunking file {path}: {str(e)}")
            raise


class ResourceChunker:
    """
    Service for creating chunked resources from files using FileSystemService.
    Supports both simple and extended parsing modes.
    """
    
    def __init__(self, fs: Optional[FileSystemService] = None):
        """Initialize the resource chunker."""
        self.fs = fs or FileSystemService()
        self._transcription_service = None
        
    def _get_transcription_service(self):
        """Lazy load transcription service to avoid circular imports."""
        if self._transcription_service is None:
            from percolate.services.llm.TranscriptionService import get_transcription_service
            self._transcription_service = get_transcription_service()
        return self._transcription_service
        
    def chunk_resource_from_uri(
        self,
        uri: str,
        parsing_mode: Literal["simple", "extended"] = "simple",
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        user_id: Optional[str] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> List["Resources"]:
        """
        Create chunked resources from a file URI.
        
        Args:
            uri: File URI (local file://, S3 s3://, or HTTP/HTTPS URL)
            parsing_mode: "simple" for basic text extraction, "extended" for LLM-enhanced parsing
            chunk_size: Number of characters per chunk
            chunk_overlap: Number of characters to overlap between chunks
            user_id: Optional user ID to associate with resources
            metadata: Optional metadata to include with resources
            
        Returns:
            List of Resources representing the chunks
            
        Raises:
            ValueError: If file type is not supported or transcription is required
            Exception: If parsing fails
        """
        from percolate.utils import make_uuid
        from percolate.models.p8.types import Resources
        
        logger.info(f"Chunking resource from URI: {uri} (mode: {parsing_mode})")
        
        # Extract file info from URI
        file_info = self._extract_file_info(uri)
        file_type = file_info['type']
        file_name = file_info['name']
        
        # Check if this is audio/video and handle accordingly
        if file_type in ['audio', 'video']:
            return self._chunk_media_resource(
                uri, file_type, parsing_mode, user_id, metadata
            )
        
        # For other file types, use FileSystemService to read content
        try:
            # Read the file using FileSystemService
            file_data = self.fs.read(uri)
            
            if parsing_mode == "simple":
                content = self._extract_simple_content(file_data, file_type)
            else:  # extended
                content = self._extract_extended_content(file_data, file_type, file_name, uri)
            
            # Create chunks from the content
            chunks = self._create_text_chunks(
                content, chunk_size, chunk_overlap
            )
            
            # Create Resource objects for each chunk
            resources = []
            for i, chunk_text in enumerate(chunks):
                resource_id = make_uuid(f"{uri}_chunk_{i}")
                
                resource = Resources(
                    id=resource_id,
                    name=f"{file_name}_chunk_{i+1}",
                    category=f"{file_type}_chunk",
                    content=chunk_text,
                    uri=uri,
                    metadata={
                        **(metadata or {}),
                        "source_file": file_name,
                        "parsing_mode": parsing_mode,
                        "chunk_index": i,
                        "total_chunks": len(chunks),
                        "chunk_size": chunk_size,
                        "chunk_overlap": chunk_overlap,
                        "file_type": file_type,
                        "original_uri": uri
                    },
                    userid=user_id,
                    resource_timestamp=datetime.now(timezone.utc)
                )
                resources.append(resource)
            
            logger.info(f"Created {len(resources)} chunks from {file_name}")
            return resources
            
        except Exception as e:
            logger.error(f"Error chunking resource from {uri}: {str(e)}")
            raise
    
    def _extract_file_info(self, uri: str) -> Dict[str, str]:
        """Extract file information from URI."""
        # Get filename from URI
        if uri.startswith('http'):
            file_name = uri.split('/')[-1].split('?')[0]  # Remove query params
        else:
            file_name = os.path.basename(uri.replace('file://', '').replace('s3://', ''))
        
        # Determine file type from extension
        ext = Path(file_name).suffix.lower()
        
        if ext in ['.txt', '.md', '.markdown']:
            file_type = 'text'
        elif ext in ['.pdf']:
            file_type = 'pdf'
        elif ext in ['.docx', '.doc']:
            file_type = 'docx'
        elif ext in ['.pptx', '.ppt']:
            file_type = 'pptx'
        elif ext in ['.csv']:
            file_type = 'csv'
        elif ext in ['.xlsx', '.xls']:
            file_type = 'xlsx'
        elif ext in ['.json']:
            file_type = 'json'
        elif ext in ['.html', '.htm']:
            file_type = 'html'
        elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff']:
            file_type = 'image'
        elif ext in ['.wav', '.mp3', '.m4a', '.flac', '.ogg']:
            file_type = 'audio'
        elif ext in ['.mp4', '.mov', '.avi', '.mkv']:
            file_type = 'video'
        else:
            file_type = 'unknown'
        
        return {
            'name': file_name,
            'type': file_type,
            'extension': ext
        }
    
    def _chunk_media_resource(
        self,
        uri: str,
        file_type: str,
        parsing_mode: str,
        user_id: Optional[str],
        metadata: Optional[Dict[str, Any]]
    ) -> List["Resources"]:
        """Handle audio/video files that require transcription with intelligent chunking."""
        if parsing_mode == "simple":
            raise ValueError(
                f"Simple parsing mode not supported for {file_type} files. "
                "Transcription is required. Use extended mode or process through audio pipeline."
            )
        
        # For extended mode, we need transcription
        transcription_service = self._get_transcription_service()
        if not transcription_service.is_available():
            raise ValueError(
                "OpenAI API key not available for transcription. "
                "Cannot process audio/video files in extended mode."
            )
        
        # Parameters for audio chunking and transcription
        max_transcription_size = 25 * 1024 * 1024  # 25MB limit for OpenAI API
        chunk_duration_seconds = 10 * 60  # 10 minutes per chunk (safe for most audio)
        
        logger.info(f"Processing {file_type} file: {uri}")
        
        # Download and analyze the file
        with tempfile.NamedTemporaryFile(delete=False, suffix=Path(uri).suffix) as temp_file:
            temp_path = temp_file.name
            
        try:
            # Use FileSystemService to download the file
            file_data = self.fs.read_bytes(uri)  # Get raw bytes for audio files
            
            with open(temp_path, 'wb') as f:
                f.write(file_data)
            
            file_size = len(file_data)
            logger.info(f"Audio file size: {file_size / (1024*1024):.1f}MB")
            
            # Check if file type is supported by transcription service
            if not transcription_service.supports_file_type(temp_path):
                raise ValueError(f"File type {Path(uri).suffix} not supported for transcription")
            
            # Determine if we need to chunk the audio file
            if file_size <= max_transcription_size:
                logger.info("File size within transcription limits, processing as single file")
                audio_chunks = [(temp_path, 0, None)]  # (path, start_time, end_time)
            else:
                logger.info(f"File size ({file_size / (1024*1024):.1f}MB) exceeds limit, chunking audio")
                audio_chunks = self._chunk_large_audio_file(
                    temp_path, max_transcription_size, chunk_duration_seconds
                )
            
            # Transcribe all audio chunks
            all_transcriptions = []
            total_duration = 0
            
            for i, (chunk_path, start_time, end_time) in enumerate(audio_chunks):
                logger.info(f"Transcribing audio chunk {i+1}/{len(audio_chunks)}")
                
                try:
                    # TranscriptionService.transcribe_file is now synchronous
                    transcription, confidence = transcription_service.transcribe_file(chunk_path)
                    
                    # Create timestamped transcription entry
                    if start_time is not None and end_time is not None:
                        duration = end_time - start_time
                        timestamp_text = f"[{start_time:.1f}s - {end_time:.1f}s]: {transcription}"
                        total_duration = max(total_duration, end_time)
                    else:
                        timestamp_text = transcription
                        duration = None
                    
                    all_transcriptions.append({
                        'text': transcription,
                        'timestamped_text': timestamp_text,
                        'start_time': start_time,
                        'end_time': end_time,
                        'duration': duration,
                        'confidence': confidence,
                        'chunk_index': i
                    })
                    
                    logger.info(f"Chunk {i+1} transcribed: {len(transcription)} characters")
                    
                except Exception as e:
                    logger.error(f"Failed to transcribe chunk {i+1}: {e}")
                    # Continue with other chunks rather than failing completely
                    continue
                finally:
                    # Clean up temporary chunk file (if different from original)
                    if chunk_path != temp_path and os.path.exists(chunk_path):
                        os.unlink(chunk_path)
            
            # Combine all transcriptions into full text
            full_transcription = "\n\n".join([t['timestamped_text'] for t in all_transcriptions])
            average_confidence = sum(t['confidence'] for t in all_transcriptions) / len(all_transcriptions)
            
            logger.info(f"Complete transcription: {len(full_transcription)} characters from {len(audio_chunks)} chunks")
            
            # Now chunk the transcription into resources based on text chunk size
            # This handles cases where transcription might be very long
            from percolate.utils import make_uuid
            from percolate.models.p8.types import Resources
            
            # Use the regular text chunking for the transcription
            text_chunks = self._create_text_chunks(
                full_transcription,
                chunk_size=1000,  # Default text chunk size
                chunk_overlap=200  # Default overlap
            )
            
            file_info = self._extract_file_info(uri)
            resources = []
            
            for i, chunk_text in enumerate(text_chunks):
                resource_id = make_uuid(f"{uri}_transcription_chunk_{i}")
                
                resource = Resources(
                    id=resource_id,
                    name=f"{file_info['name']}_transcription_chunk_{i+1}",
                    category=f"{file_type}_transcription",
                    content=chunk_text,
                    uri=uri,
                    metadata={
                        **(metadata or {}),
                        "source_file": file_info['name'],
                        "parsing_mode": parsing_mode,
                        "file_type": file_type,
                        "transcription_confidence": average_confidence,
                        "original_uri": uri,
                        "transcription_service": "openai_whisper",
                        "audio_chunks_processed": len(audio_chunks),
                        "total_audio_duration_seconds": total_duration,
                        "chunk_index": i,
                        "total_chunks": len(text_chunks),
                        "audio_file_size_mb": file_size / (1024*1024),
                        "transcription_chunks": [
                            {
                                "start_time": t["start_time"],
                                "end_time": t["end_time"], 
                                "confidence": t["confidence"]
                            } for t in all_transcriptions
                        ]
                    },
                    userid=user_id,
                    resource_timestamp=datetime.now(timezone.utc)
                )
                resources.append(resource)
            
            logger.info(f"Created {len(resources)} transcription resources for {file_info['name']}")
            return resources
            
        finally:
            # Clean up main temp file
            if os.path.exists(temp_path):
                os.unlink(temp_path)
    
    def _chunk_large_audio_file(
        self,
        audio_path: str,
        max_size_bytes: int,
        chunk_duration_seconds: int
    ) -> List[tuple]:
        """
        Chunk a large audio file into smaller pieces for transcription.
        
        Args:
            audio_path: Path to the audio file
            max_size_bytes: Maximum size per chunk in bytes
            chunk_duration_seconds: Maximum duration per chunk in seconds
            
        Returns:
            List of (chunk_path, start_time, end_time) tuples
        """
        try:
            # Try to use pydub for audio processing
            from pydub import AudioSegment
            
            logger.info(f"Loading audio file for chunking: {audio_path}")
            audio = AudioSegment.from_file(audio_path)
            
            # Get audio properties
            duration_seconds = len(audio) / 1000.0
            file_size = os.path.getsize(audio_path)
            
            logger.info(f"Audio duration: {duration_seconds:.1f}s, file size: {file_size / (1024*1024):.1f}MB")
            
            # Calculate number of chunks needed
            # Use the more restrictive of size or time limits
            chunks_by_size = max(1, file_size // max_size_bytes + (1 if file_size % max_size_bytes else 0))
            chunks_by_duration = max(1, int(duration_seconds // chunk_duration_seconds) + (1 if duration_seconds % chunk_duration_seconds else 0))
            
            num_chunks = max(chunks_by_size, chunks_by_duration)
            chunk_duration_ms = len(audio) // num_chunks
            
            logger.info(f"Splitting into {num_chunks} chunks of ~{chunk_duration_ms/1000:.1f}s each")
            
            chunks = []
            chunk_dir = tempfile.mkdtemp(prefix="audio_chunks_")
            
            for i in range(num_chunks):
                start_ms = i * chunk_duration_ms
                end_ms = min((i + 1) * chunk_duration_ms, len(audio))
                
                # Extract chunk
                chunk_audio = audio[start_ms:end_ms]
                
                # Save chunk to temporary file
                chunk_filename = f"chunk_{i+1}.wav"
                chunk_path = os.path.join(chunk_dir, chunk_filename)
                
                # Export as WAV for best transcription compatibility
                chunk_audio.export(chunk_path, format="wav")
                
                start_time = start_ms / 1000.0
                end_time = end_ms / 1000.0
                
                chunks.append((chunk_path, start_time, end_time))
                
                logger.info(f"Created chunk {i+1}: {start_time:.1f}s - {end_time:.1f}s ({os.path.getsize(chunk_path) / (1024*1024):.1f}MB)")
            
            return chunks
            
        except ImportError:
            logger.error("pydub is required for audio chunking but not available")
            raise ImportError("pydub library is required for large audio file processing")
        
        except Exception as e:
            logger.error(f"Error chunking audio file: {e}")
            raise Exception(f"Failed to chunk audio file: {e}")
    
    def _extract_simple_content(self, file_data: Any, file_type: str) -> str:
        """Extract content using simple parsing (no LLM)."""
        if isinstance(file_data, str):
            # Already parsed as text
            return file_data
        elif isinstance(file_data, dict):
            # Structured data (JSON, CSV as dict, Excel sheets, etc.)
            if file_type == 'csv' and 'data' in file_data:
                # CSV data - convert to simple text format
                rows = file_data['data']
                if not rows:
                    return ""
                
                # Create simple text representation
                lines = []
                if len(rows) > 0:
                    # Add header if available
                    headers = list(rows[0].keys())
                    lines.append(" | ".join(headers))
                    lines.append("-" * len(" | ".join(headers)))
                    
                    # Add data rows
                    for row in rows:
                        values = [str(row.get(h, "")) for h in headers]
                        lines.append(" | ".join(values))
                
                return "\n".join(lines)
            elif file_type in ['xlsx', 'xls'] or any(k for k in file_data.keys() if hasattr(file_data[k], 'shape')):
                # Excel data - convert sheets to text format
                lines = []
                for sheet_name, df in file_data.items():
                    lines.append(f"=== SHEET: {sheet_name} ===")
                    try:
                        if hasattr(df, 'to_pandas'):
                            # Polars DataFrame
                            pandas_df = df.to_pandas()
                            lines.append(pandas_df.to_string(index=False, max_rows=100))
                        elif hasattr(df, 'to_string'):
                            # Pandas DataFrame
                            lines.append(df.to_string(index=False, max_rows=100))
                        elif hasattr(df, 'shape'):
                            # Other DataFrame-like object
                            lines.append(str(df))
                        else:
                            lines.append(str(df))
                    except Exception as e:
                        logger.warning(f"Error converting sheet '{sheet_name}' to string: {e}")
                        lines.append(f"[Error displaying sheet data: {e}]")
                    lines.append("")  # Add blank line between sheets
                
                return "\n".join(lines)
            else:
                # Other structured data - convert to simple string representation
                import json
                return json.dumps(file_data, indent=2, default=str)
        else:
            # Unknown format - convert to string
            return str(file_data)
    
    def _extract_extended_content(self, file_data: Any, file_type: str, file_name: str, uri: str = None) -> str:
        """Extract content using extended parsing (with LLM if needed)."""
        
        if file_type == 'pdf':
            return self._extract_extended_pdf_content(file_data, file_name, uri)
        elif file_type == 'image':
            return self._extract_extended_image_content(file_data, file_name)
        elif file_type == 'pptx':
            return self._extract_extended_pptx_content(file_data, file_name, uri)
        else:
            # For other file types, use simple content extraction for now
            simple_content = self._extract_simple_content(file_data, file_type)
            logger.info(f"Extended parsing not yet implemented for {file_type}, using simple parsing")
            return simple_content
    
    def _extract_extended_pdf_content(self, pdf_data: Dict[str, Any], file_name: str, uri: str = None) -> str:
        """Extract content from PDF using LLM vision analysis of page images."""
        try:
            # Get image interpreter service
            from percolate.services.llm.ImageInterpreter import get_image_interpreter
            interpreter = get_image_interpreter()
            
            if not interpreter.is_available():
                logger.warning("Image interpreter not available, falling back to simple PDF parsing")
                return self._extract_simple_content(pdf_data, 'pdf')
            
            # Convert PDF pages to images for LLM analysis
            page_images = self.fs._convert_pdf_pages_to_images(pdf_data, file_name, uri)
            
            if not page_images:
                logger.warning("No page images generated, falling back to simple PDF parsing")
                return self._extract_simple_content(pdf_data, 'pdf')
            
            logger.info(f"Analyzing {len(page_images)} PDF pages with LLM vision")
            
            # Analyze each page with LLM
            analyzed_pages = []
            for i, page_image in enumerate(page_images):
                try:
                    prompt = """
                    Analyze this PDF page image in detail. Extract and describe:
                    1. All text content (transcribe exactly what you see)
                    2. Any tables, charts, or structured data
                    3. Images, diagrams, or visual elements 
                    4. Layout and formatting structure
                    5. Any important visual information not captured in plain text
                    
                    Provide a comprehensive description that captures both the textual content and visual elements.
                    """
                    
                    result = interpreter.describe_images(
                        images=page_image,
                        prompt=prompt,
                        context=f"PDF page {i+1} from document '{file_name}'",
                        max_tokens=2000
                    )
                    
                    if result["success"]:
                        page_content = f"=== PAGE {i+1} ===\n{result['content']}\n"
                        analyzed_pages.append(page_content)
                        logger.info(f"Successfully analyzed page {i+1}")
                    else:
                        logger.warning(f"Failed to analyze page {i+1}: {result.get('error', 'Unknown error')}")
                        # Fallback to simple text for this page
                        if i < len(pdf_data.get('text_pages', [])):
                            simple_text = pdf_data['text_pages'][i]
                            page_content = f"=== PAGE {i+1} (TEXT ONLY) ===\n{simple_text}\n"
                            analyzed_pages.append(page_content)
                
                except Exception as e:
                    logger.error(f"Error analyzing page {i+1}: {str(e)}")
                    # Fallback to simple text for this page
                    if i < len(pdf_data.get('text_pages', [])):
                        simple_text = pdf_data['text_pages'][i]
                        page_content = f"=== PAGE {i+1} (TEXT ONLY) ===\n{simple_text}\n"
                        analyzed_pages.append(page_content)
            
            # Combine all analyzed pages
            full_content = "\n".join(analyzed_pages)
            
            # Add summary information
            summary = f"""
DOCUMENT ANALYSIS SUMMARY:
- Document: {file_name}
- Total Pages: {len(page_images)}
- Analysis Method: LLM Vision + Text Extraction
- Pages Successfully Analyzed: {len([p for p in analyzed_pages if 'TEXT ONLY' not in p])}

FULL CONTENT:
{full_content}
"""
            
            logger.info(f"Extended PDF analysis complete: {len(full_content)} characters")
            return summary
            
        except Exception as e:
            logger.error(f"Error in extended PDF processing: {str(e)}")
            logger.info("Falling back to simple PDF parsing")
            return self._extract_simple_content(pdf_data, 'pdf')
    
    def _extract_extended_image_content(self, image_data: Image.Image, file_name: str) -> str:
        """Extract content from images using LLM vision analysis."""
        try:
            from percolate.services.llm.ImageInterpreter import get_image_interpreter
            interpreter = get_image_interpreter()
            
            if not interpreter.is_available():
                logger.warning("Image interpreter not available for image analysis")
                return f"Image file: {file_name} (analysis not available)"
            
            prompt = """
            Analyze this image in detail and provide:
            1. A comprehensive description of what you see
            2. Any text content visible in the image (OCR)
            3. Objects, people, scenes, or subjects present
            4. Colors, composition, and visual style
            5. Any technical diagrams, charts, or structured information
            6. Context clues about the purpose or meaning of the image
            
            Provide a thorough analysis that would be useful for document processing and search.
            """
            
            result = interpreter.describe_images(
                images=image_data,
                prompt=prompt,
                context=f"Image file: {file_name}",
                max_tokens=1500
            )
            
            if result["success"]:
                content = f"""
IMAGE ANALYSIS: {file_name}

{result['content']}

Analysis provided by: {result['provider']} ({result.get('model', 'unknown model')})
"""
                logger.info(f"Successfully analyzed image {file_name}")
                return content
            else:
                logger.warning(f"Failed to analyze image {file_name}: {result.get('error', 'Unknown error')}")
                return f"Image file: {file_name} (analysis failed: {result.get('error', 'Unknown error')})"
        
        except Exception as e:
            logger.error(f"Error in extended image processing: {str(e)}")
            return f"Image file: {file_name} (analysis error: {str(e)})"
    
    def _extract_extended_pptx_content(self, pptx_data: Dict[str, Any], file_name: str, uri: str = None) -> str:
        """Extract content from PPTX using enhanced parsing with LLM image analysis."""
        try:
            # The pptx_data already contains the results from our PPTXHandler
            # which includes image analysis if available
            
            if isinstance(pptx_data, dict) and 'text_content' in pptx_data:
                # Check if our enhanced parsing already included image analysis
                if pptx_data.get('has_image_analysis', False):
                    logger.info(f"PPTX {file_name} already includes image analysis from enhanced provider")
                    return pptx_data['text_content']
                else:
                    # No image analysis was done, but we have the structured data
                    # Let's try to trigger image analysis if the LLM is available
                    logger.info(f"Attempting additional image analysis for PPTX {file_name}")
                    
                    # Re-read the file with enriched mode to get image analysis
                    if uri:
                        try:
                            # Use the FileSystemService to re-read with enriched mode
                            temp_fs = FileSystemService(self.s3_service if hasattr(self, 's3_service') else None)
                            enhanced_result = temp_fs.read(uri, mode='enriched', max_images=50)
                            if isinstance(enhanced_result, dict) and enhanced_result.get('has_image_analysis', False):
                                return enhanced_result['text_content']
                        except Exception as e:
                            logger.warning(f"Failed to get enhanced PPTX analysis: {e}")
                    
                    # Fallback to existing content
                    return pptx_data['text_content']
            else:
                # Fallback for non-dict data
                return str(pptx_data)
                
        except Exception as e:
            logger.error(f"Error in extended PPTX processing: {str(e)}")
            return f"PPTX file: {file_name} (extended analysis error: {str(e)})"
    
    def _create_text_chunks(
        self,
        text: str,
        chunk_size: int,
        chunk_overlap: int
    ) -> List[str]:
        """Create text chunks with overlap."""
        if not text or len(text) <= chunk_size:
            return [text] if text else []
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # If this isn't the last chunk, try to break at word boundaries
            if end < len(text):
                # Look for the last space within the chunk
                last_space = text.rfind(' ', start, end)
                if last_space > start:
                    end = last_space
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            # Move start position considering overlap
            if end >= len(text):
                break
            start = end - chunk_overlap
            
            # Ensure we don't go backwards
            if len(chunks) > 0 and start <= len(''.join(chunks)) - len(chunks[-1]):
                start = len(''.join(chunks)) - chunk_overlap
        
        return chunks
    
    def save_chunks_to_database(self, resources: List["Resources"]) -> bool:
        """Save chunked resources to the database."""
        try:
            if not resources:
                logger.warning("No resources to save")
                return True
            
            logger.info(f"Saving {len(resources)} chunked resources to database")
            import percolate as p8
            p8.repository(type(resources[0])).update_records(resources)
            logger.info("Successfully saved all chunked resources")
            return True
            
        except Exception as e:
            logger.error(f"Error saving chunked resources: {str(e)}")
            return False


# Global resource chunker instance
_resource_chunker = None

def get_resource_chunker() -> ResourceChunker:
    """Get a global resource chunker instance."""
    global _resource_chunker
    if _resource_chunker is None:
        _resource_chunker = ResourceChunker()
    return _resource_chunker