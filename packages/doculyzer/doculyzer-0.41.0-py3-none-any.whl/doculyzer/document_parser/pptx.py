"""
PPTX document parser module for the document pointer system.

This module parses PowerPoint (PPTX) documents into structured elements
with comprehensive date extraction and temporal analysis.
"""

import json
import logging
import os
import re
from typing import Dict, Any, List, Optional, Union

from ..relationships import RelationshipType
from ..storage import ElementType

try:
    import pptx
    # noinspection PyUnresolvedReferences
    from pptx import Presentation
    from pptx.shapes.autoshape import Shape
    # noinspection PyUnresolvedReferences
    from pptx.shapes.group import GroupShape
    # noinspection PyUnresolvedReferences
    from pptx.shapes.picture import Picture
    # noinspection PyUnresolvedReferences
    from pptx.slide import Slide, SlideLayout
    # noinspection PyUnresolvedReferences
    from pptx.text.text import TextFrame

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. Install with 'pip install python-pptx' to use PPTX parser")

from .base import DocumentParser
from .extract_dates import DateExtractor

logger = logging.getLogger(__name__)


class PptxParser(DocumentParser):
    """Parser for PowerPoint (PPTX) documents with enhanced date extraction."""

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """Initialize the PPTX parser."""
        super().__init__(config)

        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx is required for PPTX parsing")

        # Configuration options
        self.config = config or {}
        self.extract_notes = self.config.get("extract_notes", True)
        self.extract_hidden_slides = self.config.get("extract_hidden_slides", False)
        self.extract_comments = self.config.get("extract_comments", True)
        self.extract_shapes = self.config.get("extract_shapes", True)
        self.extract_images = self.config.get("extract_images", True)
        self.extract_tables = self.config.get("extract_tables", True)
        self.extract_text_boxes = self.config.get("extract_text_boxes", True)
        self.extract_charts = self.config.get("extract_charts", True)
        self.extract_masters = self.config.get("extract_masters", False)
        self.extract_templates = self.config.get("extract_templates", False)
        self.temp_dir = self.config.get("temp_dir", os.path.join(os.path.dirname(__file__), 'temp'))

        # Date extraction configuration
        self.extract_dates = self.config.get("extract_dates", True)
        self.date_context_chars = self.config.get("date_context_chars", 50)  # Small context window
        self.min_year = self.config.get("min_year", 1900)
        self.max_year = self.config.get("max_year", 2100)
        self.fiscal_year_start_month = self.config.get("fiscal_year_start_month", 10)
        self.default_locale = self.config.get("default_locale", "US")

        # Initialize date extractor if enabled
        self.date_extractor = None
        if self.extract_dates:
            try:
                self.date_extractor = DateExtractor(
                    context_chars=self.date_context_chars,
                    min_year=self.min_year,
                    max_year=self.max_year,
                    fiscal_year_start_month=self.fiscal_year_start_month,
                    default_locale=self.default_locale
                )
                logger.debug("Date extraction enabled with comprehensive temporal analysis")
            except ImportError as e:
                logger.warning(f"Date extraction disabled: {e}")
                self.extract_dates = False

    def _resolve_element_text(self, location_data: Dict[str, Any], source_content: Optional[Union[str, bytes]] = None) -> str:
        """
        Resolve the plain text representation of a PPTX element.

        Args:
            location_data: Content location data
            source_content: Optional preloaded source content

        Returns:
            Plain text representation of the element
        """
        # Get the content using the improved _resolve_element_content method
        content = self._resolve_element_content(location_data, source_content)
        element_type = location_data.get("type", "")

        # Handle specific element types
        if element_type == ElementType.PRESENTATION_BODY.value:
            return content.strip()

        elif element_type == ElementType.SLIDE.value:
            return content.strip()

        elif element_type == ElementType.TEXT_BOX.value or element_type == ElementType.PARAGRAPH.value:
            return content.strip()

        elif element_type == ElementType.TABLE.value or element_type == ElementType.TABLE_CELL.value:
            # The improved _resolve_element_content already formats tables properly
            return content.strip()

        elif element_type == ElementType.SLIDE_NOTES.value:
            return content.strip()

        elif element_type == ElementType.COMMENT.value:
            # For comments, extract just the comment text without metadata
            if ": " in content:
                return content.split(": ", 1)[1].strip()
            return content.strip()

        elif element_type == ElementType.IMAGE.value:
            if "Alt text: " in content:
                return content.split("Alt text: ", 1)[1].strip()
            return "Image"

        elif element_type == ElementType.CHART.value:
            if "Chart: " in content and "\n" in content:
                return content.split("\n")[0].replace("Chart: ", "").strip()
            return content.strip()

        # Default
        return content.strip()

    def _resolve_element_content(self, location_data: Dict[str, Any],
                                 source_content: Optional[Union[str, bytes]] = None) -> str:
        """
        Resolve content for specific PPTX element types.

        Args:
            location_data: Content location data
            source_content: Optional pre-loaded source content

        Returns:
            Resolved content string
        """
        source = location_data.get("source", "")
        element_type = location_data.get("type", "")
        slide_index = location_data.get("slide_index", 0)

        # Load the document if source content is not provided
        presentation = None
        temp_file = None
        try:
            if source_content is None:
                # Check if source is a file path
                if os.path.exists(source):
                    try:
                        presentation = Presentation(source)
                    except Exception as e:
                        raise ValueError(f"Error loading PPTX document: {str(e)}")
                else:
                    raise ValueError(f"Source file not found: {source}")
            else:
                # Save content to a temporary file
                if not os.path.exists(self.temp_dir):
                    os.makedirs(self.temp_dir, exist_ok=True)

                import uuid
                temp_file = os.path.join(self.temp_dir, f"temp_{uuid.uuid4().hex}.pptx")
                with open(temp_file, 'wb') as f:
                    if isinstance(source_content, str):
                        f.write(source_content.encode('utf-8'))
                    else:
                        f.write(source_content)

                # Load the document
                try:
                    presentation = Presentation(temp_file)
                except Exception as e:
                    raise ValueError(f"Error loading PPTX document: {str(e)}")

            # Handle different element types
            if element_type == ElementType.PRESENTATION_BODY.value:
                # Return basic presentation information
                slide_count = len(presentation.slides)
                return f"Presentation with {slide_count} slides"

            elif element_type == ElementType.SLIDE.value:
                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}. Presentation has {len(presentation.slides)} slides."

                # Get the slide
                slide = presentation.slides[slide_index]

                # Extract all text from the slide
                all_text = []
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        text = shape.text_frame.text
                        if text:
                            all_text.append(text)

                return "\n\n".join(all_text)

            elif element_type == ElementType.TEXT_BOX.value:
                # Extract text from a text box shape
                shape_path = location_data.get("shape_path", "")
                if not shape_path:
                    return "No shape path specified"

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                    return f"Text shape not found at path: {shape_path}"

                return shape.text_frame.text

            elif element_type == ElementType.PARAGRAPH.value:
                # Extract specific paragraph from a text shape
                shape_path = location_data.get("shape_path", "")
                paragraph_index = location_data.get("paragraph_index", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                    return f"Text shape not found at path: {shape_path}"

                # Check if paragraph index is valid
                if not hasattr(shape.text_frame, 'paragraphs') or paragraph_index >= len(shape.text_frame.paragraphs):
                    return f"Invalid paragraph index: {paragraph_index}"

                return shape.text_frame.paragraphs[paragraph_index].text

            elif element_type == ElementType.TABLE.value:
                # Extract table content with proper formatting
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_table') or not shape.has_table:
                    return f"Table shape not found at path: {shape_path}"

                # Get table object
                table = shape.table

                # Format table with consistent structure
                rows_text = []

                # Process each row
                for row in table.rows:
                    cells_text = []
                    for cell in row.cells:
                        cell_text = cell.text_frame.text.strip() if hasattr(cell, 'text_frame') else ""
                        cells_text.append(cell_text)

                    # Join cells with pipe separator
                    rows_text.append(" | ".join(cells_text))

                # Return formatted table
                return "\n".join(rows_text)

            elif element_type == ElementType.TABLE_CELL.value:
                # Extract cell content from a table
                shape_path = location_data.get("shape_path", "")
                row = location_data.get("row", 0)
                col = location_data.get("col", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_table') or not shape.has_table:
                    return f"Table shape not found at path: {shape_path}"

                # Get table object
                table = shape.table

                # Check if row and column indices are valid
                if row < 0 or row >= len(table.rows) or col < 0 or col >= len(table.columns):
                    return f"Invalid cell coordinates: row={row}, col={col}"

                # Get cell text
                cell = table.cell(row, col)
                return cell.text_frame.text if hasattr(cell, 'text_frame') else ""

            elif element_type == ElementType.SLIDE_NOTES.value:
                # Extract notes from a slide
                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Check if slide has notes
                if not hasattr(slide, 'notes_slide') or not slide.notes_slide:
                    return "No notes for this slide"

                # Return notes text
                return slide.notes_slide.notes_text_frame.text

            elif element_type == ElementType.COMMENT.value:
                # Extract a specific comment
                comment_index = location_data.get("comment_index", 0)

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Check if slide has comments
                if not hasattr(slide, 'comments') or not slide.comments:
                    return "No comments for this slide"

                # Check if comment index is valid
                if comment_index < 0 or comment_index >= len(slide.comments):
                    return f"Invalid comment index: {comment_index}"

                # Get comment
                comment = slide.comments[comment_index]

                # Format comment details
                author = comment.author if hasattr(comment, 'author') else "Unknown"
                text = comment.text if hasattr(comment, 'text') else ""
                date = comment.date if hasattr(comment, 'date') else None

                if date:
                    return f"Comment by {author} on {date}: {text}"
                else:
                    return f"Comment by {author}: {text}"

            elif element_type == ElementType.IMAGE.value:
                # Return information about an image
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not isinstance(shape, Picture):
                    return f"Image shape not found at path: {shape_path}"

                # Get image details
                image_name = shape.image.filename if hasattr(shape, 'image') and hasattr(shape.image,
                                                                                         'filename') else "Unknown"
                alt_text = shape.alt_text if hasattr(shape, 'alt_text') else ""

                return f"Image: {image_name}\nAlt text: {alt_text}"

            elif element_type == ElementType.CHART.value:
                # Return information about a chart
                shape_path = location_data.get("shape_path", "")

                # Check if slide index is valid
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    return f"Invalid slide index: {slide_index}"

                # Get the slide
                slide = presentation.slides[slide_index]

                # Parse shape path to locate the shape
                shape_indices = shape_path.split('/')
                shape = self._find_shape_by_path(slide.shapes, shape_indices)

                if not shape or not hasattr(shape, 'has_chart') or not shape.has_chart:
                    return f"Chart shape not found at path: {shape_path}"

                # Get chart details
                chart = shape.chart
                chart_type = str(chart.chart_type) if hasattr(chart, 'chart_type') else "Unknown"
                chart_title = chart.chart_title.text_frame.text if hasattr(chart, 'chart_title') and hasattr(
                    chart.chart_title, 'text_frame') else "Untitled Chart"

                # Get categories and series if available
                categories = []
                series_names = []

                if hasattr(chart, 'plots') and chart.plots:
                    plot = chart.plots[0]

                    if hasattr(plot, 'categories'):
                        for category in plot.categories:
                            if category:
                                categories.append(str(category))

                    if hasattr(plot, 'series'):
                        for series in plot.series:
                            if hasattr(series, 'name') and series.name:
                                series_names.append(str(series.name))

                # Format chart description
                description = f"Chart: {chart_title}\nType: {chart_type}"

                if categories:
                    description += f"\nCategories: {', '.join(categories)}"

                if series_names:
                    description += f"\nSeries: {', '.join(series_names)}"

                return description

            else:
                # For other element types or if no specific handler,
                # return basic information about the presentation
                return f"PowerPoint presentation with {len(presentation.slides)} slides"

        finally:
            # Clean up temporary file
            if temp_file and os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    logger.warning(f"Failed to delete temporary file {temp_file}: {str(e)}")

    def supports_location(self, content_location: Dict[str, Any]) -> bool:
        """
        Check if this parser supports resolving the given location.

        Args:
            content_location: Content location pointer

        Returns:
            True if supported, False otherwise
        """
        try:
            location_data = content_location
            source = location_data.get("source", "")

            # Check if source exists and is a file
            if not os.path.exists(source) or not os.path.isfile(source):
                return False

            # Check file extension for PPTX
            _, ext = os.path.splitext(source.lower())
            return ext in ['.pptx', '.pptm']

        except (json.JSONDecodeError, TypeError):
            return False

    def _extract_links(self, content: str, element_id: str) -> List[Dict[str, Any]]:
        """
        Base method for extracting links from content.

        Args:
            content: Text content
            element_id: ID of the element containing the links

        Returns:
            List of extracted link dictionaries
        """
        links = []

        # Extract URLs using a regular expression
        url_pattern = r'https?://[^\s<>)"\']+|www\.[^\s<>)"\']+|ftp://[^\s<>)"\']+|file://[^\s<>)"\']+|mailto:[^\s<>)"\']+|[^\s<>)"\']+\.(?:com|org|net|edu|gov|io|ai|app)[^\s<>)"\']*'
        urls = re.findall(url_pattern, content)

        # Create link entries for each URL found
        for url in urls:
            # Clean up URL
            if url.startswith('www.'):
                url = 'http://' + url

            links.append({
                "source_id": element_id,
                "link_text": url,
                "link_target": url,
                "link_type": "url"
            })

        # Look for slide references (e.g., "See slide 5")
        slide_refs = re.findall(r'slide\s+(\d+)', content, re.IGNORECASE)

        for ref in slide_refs:
            try:
                slide_num = int(ref)

                links.append({
                    "source_id": element_id,
                    "link_text": f"Slide {slide_num}",
                    "link_target": f"slide_{slide_num}",
                    "link_type": "slide_reference"
                })
            except ValueError:
                pass

        return links

    def _extract_dates_from_text(self, text: str, element_id: str, element_dates: Dict[str, List[Dict[str, Any]]]):
        """
        Extract dates from text content and add to element_dates.

        Args:
            text: Text content to extract dates from
            element_id: ID of the element containing the text
            element_dates: Dictionary to store extracted dates
        """
        if not self.extract_dates or not self.date_extractor or not text.strip():
            return

        try:
            dates = self.date_extractor.extract_dates_as_dicts(text)
            if dates:
                element_dates[element_id] = dates
                logger.debug(f"Extracted {len(dates)} dates from element {element_id}")
        except Exception as e:
            logger.warning(f"Error extracting dates from element {element_id}: {e}")

    @staticmethod
    def _extract_document_links(presentation, elements: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Helper method to extract hyperlinks from PowerPoint presentation.
        This is called during the parsing phase.

        Args:
            presentation: The PPTX presentation
            elements: List of extracted elements

        Returns:
            List of hyperlink dictionaries
        """
        links = []

        # Map element IDs to elements for quick lookup
        # element_map = {elem["element_id"]: elem for elem in elements}

        # Extract hyperlinks from text shapes
        for element in elements:
            if element["element_type"] in [ElementType.TEXT_BOX.value, ElementType.PARAGRAPH.value,
                                           ElementType.TABLE_CELL.value, ElementType.SLIDE_NOTES.value,
                                           ElementType.COMMENT.value]:
                element_id = element["element_id"]
                text = element.get("metadata", {}).get("text", "")

                # Look for hyperlink patterns in text
                url_pattern = r'https?://[^\s<>)"\']+|www\.[^\s<>)"\']+|ftp://[^\s<>)"\']+|file://[^\s<>)"\']+|mailto:[^\s<>)"\']+|[^\s<>)"\']+\.(?:com|org|net|edu|gov|io|ai|app)[^\s<>)"\']*'
                urls = re.findall(url_pattern, text)

                for url in urls:
                    # Clean up URL
                    if url.startswith('www.'):
                        url = 'http://' + url

                    # Add link
                    links.append({
                        "source_id": element_id,
                        "link_text": url,
                        "link_target": url,
                        "link_type": "url"
                    })

                # Look for slide references (e.g., "See slide 5")
                slide_refs = re.findall(r'slide\s+(\d+)', text, re.IGNORECASE)

                for ref in slide_refs:
                    try:
                        slide_num = int(ref)
                        # Adjust for 0-based indexing
                        # slide_idx = slide_num - 1

                        # Find target slide element
                        target_slide = None
                        for slide_elem in elements:
                            if (slide_elem["element_type"] == ElementType.SLIDE.value and
                                    slide_elem.get("metadata", {}).get("number") == slide_num):
                                target_slide = slide_elem
                                break

                        if target_slide:
                            # Add link
                            links.append({
                                "source_id": element_id,
                                "link_text": f"Slide {slide_num}",
                                "link_target": target_slide["element_id"],
                                "link_type": "slide_reference"
                            })
                    except (ValueError, IndexError):
                        pass

        return links

    def parse(self, doc_content: Dict[str, Any]) -> Dict[str, Any]:
        """
        Parse a PPTX document into structured elements with comprehensive date extraction.

        Args:
            doc_content: Document content and metadata

        Returns:
            Dictionary with document metadata, elements, relationships, extracted links, and dates
        """
        # Extract metadata from doc_content
        source_id = doc_content["id"]
        metadata = doc_content.get("metadata", {}).copy()  # Make a copy to avoid modifying original

        # Check if we have a binary path or content
        binary_path = doc_content.get("binary_path")
        binary_content = doc_content.get("content")

        # If we have content but no path, save it to a temp file
        if not binary_path and binary_content:
            if not os.path.exists(self.temp_dir):
                os.makedirs(self.temp_dir, exist_ok=True)

            import uuid
            temp_file_path = os.path.join(self.temp_dir, f"temp_{uuid.uuid4().hex}.pptx")
            with open(temp_file_path, 'wb') as f:
                if isinstance(binary_content, str):
                    f.write(binary_content.encode('utf-8'))
                else:
                    f.write(binary_content)

            binary_path = temp_file_path

        if not binary_path:
            raise ValueError("PPTX parser requires either binary_path or content to process the presentation")

        # Generate document ID if not present
        doc_id = metadata.get("doc_id", self._generate_id("doc_"))

        # Load PPTX document
        try:
            presentation = Presentation(binary_path)
        except Exception as e:
            logger.error(f"Error loading PPTX document: {str(e)}")
            raise

        # Create document record with metadata
        document = {
            "doc_id": doc_id,
            "doc_type": "pptx",
            "source": source_id,
            "metadata": self._extract_document_metadata(presentation, metadata),
            "content_hash": doc_content.get("content_hash", "")
        }

        # Create root element
        elements = [self._create_root_element(doc_id, source_id)]
        root_id = elements[0]["element_id"]

        # Initialize relationships list and element_dates dictionary
        relationships = []
        element_dates = {}

        # Parse document elements and create relationships
        new_elements = self._parse_presentation(presentation, doc_id, root_id, source_id, relationships, element_dates)
        elements.extend(new_elements)

        # Extract links from the document using the helper method
        links = self._extract_document_links(presentation, elements)

        # Add date statistics to document metadata
        if element_dates:
            total_dates = sum(len(dates) for dates in element_dates.values())
            document["metadata"]["date_extraction"] = {
                "total_dates_found": total_dates,
                "elements_with_dates": len(element_dates),
                "extraction_enabled": True
            }
        else:
            document["metadata"]["date_extraction"] = {
                "total_dates_found": 0,
                "elements_with_dates": 0,
                "extraction_enabled": self.extract_dates
            }

        # Clean up temporary file if we created one
        if binary_path != doc_content.get("binary_path") and os.path.exists(binary_path):
            try:
                os.remove(binary_path)
                logger.debug(f"Deleted temporary file: {binary_path}")
            except Exception as e:
                logger.warning(f"Failed to delete temporary file {binary_path}: {str(e)}")

        # Return the parsed document with extracted links, relationships, and dates
        result = {
            "document": document,
            "elements": elements,
            "links": links,
            "relationships": relationships
        }

        # Add dates if any were extracted
        if element_dates:
            result["element_dates"] = element_dates

        return result

    @staticmethod
    def _extract_document_metadata(presentation: Presentation, base_metadata: Dict[str, Any]) -> Dict[str, Any]:
        """
        Extract metadata from PPTX document.

        Args:
            presentation: The PPTX presentation
            base_metadata: Base metadata from content source

        Returns:
            Dictionary of document metadata
        """
        # Combine base metadata with document properties
        metadata = base_metadata.copy()

        try:
            # Add core properties
            core_props = presentation.core_properties
            if core_props:
                # Add core properties to metadata
                if core_props.title:
                    metadata["title"] = core_props.title
                if core_props.author:
                    metadata["author"] = core_props.author
                if core_props.created:
                    metadata["created"] = core_props.created.timestamp() if hasattr(core_props.created,
                                                                                    'timestamp') else str(
                        core_props.created)
                if core_props.modified:
                    metadata["modified"] = core_props.modified.timestamp() if hasattr(core_props.modified,
                                                                                      'timestamp') else str(
                        core_props.modified)
                if core_props.last_modified_by:
                    metadata["last_modified_by"] = core_props.last_modified_by
                if core_props.keywords:
                    metadata["keywords"] = core_props.keywords
                if core_props.subject:
                    metadata["subject"] = core_props.subject
                if core_props.comments:
                    metadata["comments"] = core_props.comments
                if core_props.category:
                    metadata["category"] = core_props.category

            # Add presentation statistics
            metadata["slide_count"] = len(presentation.slides)
            metadata["slide_width"] = presentation.slide_width
            metadata["slide_height"] = presentation.slide_height

        except Exception as e:
            logger.warning(f"Error extracting document metadata: {str(e)}")

        return metadata

    def _parse_presentation(self, presentation: Presentation, doc_id: str, parent_id: str, source_id: str,
                            relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Parse PowerPoint presentation into structured elements and create relationships.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of parsed elements
        """
        elements: List = []

        # Create presentation body element
        body_id = self._generate_id("body_")
        body_element = {
            "element_id": body_id,
            "doc_id": doc_id,
            "element_type": ElementType.PRESENTATION_BODY.value,
            "parent_id": parent_id,
            "content_preview": "Presentation body",
            "content_location": json.dumps({
                "source": source_id,
                "type": ElementType.PRESENTATION_BODY.value
            }),
            "content_hash": "",
            "metadata": {
                "slide_count": len(presentation.slides)
            }
        }
        elements.append(body_element)

        # Create relationship between document root and presentation body
        contains_relationship = {
            "relationship_id": self._generate_id("rel_"),
            "source_id": parent_id,
            "target_id": body_id,
            "relationship_type": RelationshipType.CONTAINS.value,
            "metadata": {
                "confidence": 1.0
            }
        }
        relationships.append(contains_relationship)

        # Create inverse relationship
        contained_by_relationship = {
            "relationship_id": self._generate_id("rel_"),
            "source_id": body_id,
            "target_id": parent_id,
            "relationship_type": RelationshipType.CONTAINED_BY.value,
            "metadata": {
                "confidence": 1.0
            }
        }
        relationships.append(contained_by_relationship)

        # Process slides
        for slide_idx, slide in enumerate(presentation.slides):
            # Skip hidden slides if not configured to extract them
            if not self.extract_hidden_slides and hasattr(slide, 'hidden') and slide.hidden:
                continue

            # Process this slide
            slide_elements = self._process_slide(slide, slide_idx, doc_id, body_id, source_id, relationships, element_dates)
            elements.extend(slide_elements)

        # Extract masters if configured
        if self.extract_masters and hasattr(presentation, 'slide_masters'):
            master_elements = self._extract_slide_masters(presentation, doc_id, parent_id, source_id, relationships, element_dates)
            elements.extend(master_elements)

        # Extract templates if configured
        if self.extract_templates and hasattr(presentation, 'slide_layouts'):
            template_elements = self._extract_slide_templates(presentation, doc_id, parent_id, source_id, relationships, element_dates)
            elements.extend(template_elements)

        return elements

    def _extract_slide_masters(self, presentation, doc_id: str, parent_id: str, source_id: str,
                               relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Extract slide masters from presentation and create relationships.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of slide master elements
        """
        elements = list()

        try:
            if not hasattr(presentation, 'slide_masters'):
                return elements

            # Create masters container element
            masters_id = self._generate_id("masters_")

            masters_element = {
                "element_id": masters_id,
                "doc_id": doc_id,
                "element_type": ElementType.SLIDE_MASTERS.value,
                "parent_id": parent_id,
                "content_preview": "Slide Masters",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.SLIDE_MASTERS.value
                }),
                "content_hash": "",
                "metadata": {
                    "master_count": len(presentation.slide_masters)
                }
            }
            elements.append(masters_element)

            # Create relationship from parent to masters container
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": masters_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": masters_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process individual masters
            for master_idx, master in enumerate(presentation.slide_masters):
                # Generate master ID
                master_id = self._generate_id(f"master_{master_idx}_")

                # Create master element
                master_element = {
                    "element_id": master_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.SLIDE_MASTER.value,
                    "parent_id": masters_id,
                    "content_preview": f"Slide Master {master_idx + 1}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.SLIDE_MASTER.value,
                        "index": master_idx
                    }),
                    "content_hash": self._generate_hash(f"master_{master_idx}"),
                    "metadata": {
                        "index": master_idx,
                        "layout_count": len(master.slide_layouts) if hasattr(master, 'slide_layouts') else 0
                    }
                }
                elements.append(master_element)

                # Create relationship from masters container to master
                contains_master_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": masters_id,
                    "target_id": master_id,
                    "relationship_type": RelationshipType.CONTAINS.value,
                    "metadata": {
                        "confidence": 1.0,
                        "index": master_idx
                    }
                }
                relationships.append(contains_master_relationship)

                # Create inverse relationship
                master_contained_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": master_id,
                    "target_id": masters_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(master_contained_relationship)

                # Process master shapes if desired
                if self.extract_shapes and hasattr(master, 'shapes'):
                    shape_elements = self._process_shapes(master.shapes, doc_id, master_id, source_id, -1,
                                                          relationships, f"master_{master_idx}", element_dates)
                    elements.extend(shape_elements)

        except Exception as e:
            logger.warning(f"Error extracting slide masters: {str(e)}")

        return elements

    def _extract_slide_templates(self, presentation, doc_id: str, parent_id: str, source_id: str,
                                 relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Extract slide templates (layouts) from presentation and create relationships.

        Args:
            presentation: The PPTX presentation
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of slide template elements
        """
        elements = []

        try:
            # Collect all slide layouts from all masters
            layouts = []
            layout_names = set()

            if hasattr(presentation, 'slide_masters'):
                for master in presentation.slide_masters:
                    if hasattr(master, 'slide_layouts'):
                        for layout in master.slide_layouts:
                            # Avoid duplicates by name
                            layout_name = layout.name if hasattr(layout, 'name') else ""
                            if layout_name not in layout_names:
                                layouts.append(layout)
                                layout_names.add(layout_name)

            if not layouts:
                return elements

            # Create templates container element
            templates_id = self._generate_id("templates_")

            templates_element = {
                "element_id": templates_id,
                "doc_id": doc_id,
                "element_type": ElementType.SLIDE_TEMPLATES.value,
                "parent_id": parent_id,
                "content_preview": "Slide Templates",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.SLIDE_TEMPLATES.value
                }),
                "content_hash": "",
                "metadata": {
                    "template_count": len(layouts)
                }
            }
            elements.append(templates_element)

            # Create relationship from parent to templates container
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": templates_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": templates_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process individual templates
            for layout_idx, layout in enumerate(layouts):
                # Generate layout ID
                layout_id = self._generate_id(f"layout_{layout_idx}_")

                # Get layout name
                layout_name = layout.name if hasattr(layout, 'name') else f"Layout {layout_idx + 1}"

                # Create layout element
                layout_element = {
                    "element_id": layout_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.SLIDE_LAYOUT.value,
                    "parent_id": templates_id,
                    "content_preview": layout_name,
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.SLIDE_LAYOUT.value,
                        "index": layout_idx
                    }),
                    "content_hash": self._generate_hash(f"layout_{layout_idx}"),
                    "metadata": {
                        "index": layout_idx,
                        "name": layout_name
                    }
                }
                elements.append(layout_element)

                # Create relationship from templates container to layout
                contains_layout_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": templates_id,
                    "target_id": layout_id,
                    "relationship_type": RelationshipType.CONTAINS.value,
                    "metadata": {
                        "confidence": 1.0,
                        "index": layout_idx
                    }
                }
                relationships.append(contains_layout_relationship)

                # Create inverse relationship
                layout_contained_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": layout_id,
                    "target_id": templates_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(layout_contained_relationship)

                # Process layout shapes if desired
                if self.extract_shapes and hasattr(layout, 'shapes'):
                    shape_elements = self._process_shapes(layout.shapes, doc_id, layout_id, source_id, -1,
                                                          relationships, f"layout_{layout_idx}", element_dates)
                    elements.extend(shape_elements)

        except Exception as e:
            logger.warning(f"Error extracting slide templates: {str(e)}")

        return elements

    def _find_shape_by_path(self, shapes, shape_indices):
        """
        Find a shape by following the shape path indices.

        Args:
            shapes: Collection of shapes to search in
            shape_indices: List of indices to follow

        Returns:
            The shape if found, None otherwise
        """
        try:
            if not shape_indices:
                return None

            # Get the first index
            current_idx = int(shape_indices[0])

            # Check if index is valid
            if current_idx < 0 or current_idx >= len(shapes):
                return None

            # Get the shape at this index
            shape = shapes[current_idx]

            # If this is the last index, return the shape
            if len(shape_indices) == 1:
                return shape

            # If this is a group shape, recurse into it
            if isinstance(shape, GroupShape) and hasattr(shape, 'shapes'):
                return self._find_shape_by_path(shape.shapes, shape_indices[1:])

            # If we get here, the path is invalid
            return None

        except (ValueError, IndexError, TypeError):
            return None

    def _process_slide(self, slide: Slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str,
                       relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process a PowerPoint slide into structured elements and create relationships.

        Args:
            slide: The PPTX slide
            slide_idx: Slide index (0-based)
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of slide-related elements
        """
        elements = []

        # Generate slide ID
        slide_id = self._generate_id(f"slide_{slide_idx}_")

        # Get slide title if available
        slide_title = self._get_slide_title(slide)

        # Extract all text from the slide for date extraction
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text = shape.text_frame.text
                if text:
                    slide_text += text + "\n"

        # Create slide element
        slide_element = {
            "element_id": slide_id,
            "doc_id": doc_id,
            "element_type": ElementType.SLIDE.value,
            "parent_id": parent_id,
            "content_preview": f"Slide {slide_idx + 1}: {slide_title}" if slide_title else f"Slide {slide_idx + 1}",
            "content_location": json.dumps({
                "source": source_id,
                "type": ElementType.SLIDE.value,
                "slide_index": slide_idx
            }),
            "content_hash": self._generate_hash(f"slide_{slide_idx}"),
            "metadata": {
                "index": slide_idx,
                "number": slide_idx + 1,
                "title": slide_title,
                "layout": self._get_slide_layout_name(slide),
                "has_notes": bool(slide.notes_slide and slide.notes_slide.notes_text_frame.text),
                "shape_count": len(slide.shapes)
            }
        }
        elements.append(slide_element)

        # Extract dates from slide text
        self._extract_dates_from_text(slide_text, slide_id, element_dates)

        # Create relationship from parent to slide
        contains_relationship = {
            "relationship_id": self._generate_id("rel_"),
            "source_id": parent_id,
            "target_id": slide_id,
            "relationship_type": RelationshipType.CONTAINS.value,
            "metadata": {
                "confidence": 1.0,
                "order": slide_idx
            }
        }
        relationships.append(contains_relationship)

        # Create inverse relationship
        contained_by_relationship = {
            "relationship_id": self._generate_id("rel_"),
            "source_id": slide_id,
            "target_id": parent_id,
            "relationship_type": RelationshipType.CONTAINED_BY.value,
            "metadata": {
                "confidence": 1.0
            }
        }
        relationships.append(contained_by_relationship)

        # Process slide shapes
        if self.extract_shapes:
            shape_elements = self._process_shapes(slide.shapes, doc_id, slide_id, source_id, slide_idx, relationships, "", element_dates)
            elements.extend(shape_elements)

        # Process slide notes
        if self.extract_notes and slide.notes_slide and slide.notes_slide.notes_text_frame.text:
            notes_elements = self._process_notes(slide.notes_slide, slide_idx, doc_id, slide_id, source_id,
                                                 relationships, element_dates)
            elements.extend(notes_elements)

        # Process slide comments if available and configured
        if self.extract_comments and hasattr(slide, 'comments'):
            comment_elements = self._process_comments(slide, slide_idx, doc_id, slide_id, source_id, relationships, element_dates)
            elements.extend(comment_elements)

        return elements

    def _process_shapes(self, shapes, doc_id: str, parent_id: str, source_id: str, slide_idx: int,
                        relationships: List[Dict[str, Any]], shape_path: str = "", element_dates: Dict[str, List[Dict[str, Any]]] = None) -> List[Dict[str, Any]]:
        """
        Process PowerPoint shapes into structured elements and create relationships.

        Args:
            shapes: Collection of shapes
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            relationships: List to add relationships to
            shape_path: Path to the shape (for nested shapes)
            element_dates: Dictionary to store extracted dates

        Returns:
            List of shape-related elements
        """
        if element_dates is None:
            element_dates = {}

        elements = []

        for shape_idx, shape in enumerate(shapes):
            # Generate current shape path
            current_shape_path = f"{shape_path}/{shape_idx}" if shape_path else f"{shape_idx}"

            # Process shape based on type
            if isinstance(shape, GroupShape):
                # Process group shape and its children
                group_id = self._generate_id(f"group_{current_shape_path}_")

                # Create group element
                group_element = {
                    "element_id": group_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.SHAPE_GROUP.value,
                    "parent_id": parent_id,
                    "content_preview": f"Shape Group {current_shape_path}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.SHAPE_GROUP.value,
                        "slide_index": slide_idx,
                        "shape_path": current_shape_path
                    }),
                    "content_hash": self._generate_hash(f"group_{current_shape_path}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_path": current_shape_path,
                        "shape_type": "group",
                        "shape_count": len(shape.shapes) if hasattr(shape, 'shapes') else 0
                    }
                }
                elements.append(group_element)

                # Create relationship from parent to group
                contains_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": parent_id,
                    "target_id": group_id,
                    "relationship_type": RelationshipType.CONTAINS.value,
                    "metadata": {
                        "confidence": 1.0,
                        "index": shape_idx
                    }
                }
                relationships.append(contains_relationship)

                # Create inverse relationship
                contained_by_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": group_id,
                    "target_id": parent_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(contained_by_relationship)

                # Process child shapes
                child_elements = self._process_shapes(shape.shapes, doc_id, group_id, source_id, slide_idx,
                                                      relationships, current_shape_path, element_dates)
                elements.extend(child_elements)

            elif hasattr(shape, 'has_table') and shape.has_table and self.extract_tables:
                # Process table shape
                table_elements = self._process_table_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                           relationships, current_shape_path, element_dates)
                elements.extend(table_elements)

            elif hasattr(shape, 'has_chart') and shape.has_chart and self.extract_charts:
                # Process chart shape
                chart_elements = self._process_chart_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                           relationships, current_shape_path, element_dates)
                elements.extend(chart_elements)

            elif isinstance(shape, Picture) and self.extract_images:
                # Process picture shape
                picture_elements = self._process_picture_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                               relationships, current_shape_path, element_dates)
                elements.extend(picture_elements)

            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame and self.extract_text_boxes:
                # Process text shape
                text_elements = self._process_text_shape(shape, doc_id, parent_id, source_id, slide_idx,
                                                         relationships, current_shape_path, element_dates)
                elements.extend(text_elements)

            elif self.extract_shapes:
                # Process generic shape
                shape_id = self._generate_id(f"shape_{current_shape_path}_")

                # Get shape name and type info
                shape_name = shape.name if hasattr(shape, 'name') else ""
                shape_type = self._get_shape_type(shape)

                # Create shape element
                shape_element = {
                    "element_id": shape_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.SHAPE.value,
                    "parent_id": parent_id,
                    "content_preview": f"Shape: {shape_name}" if shape_name else f"Shape {current_shape_path}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.SHAPE.value,
                        "slide_index": slide_idx,
                        "shape_path": current_shape_path
                    }),
                    "content_hash": self._generate_hash(f"shape_{current_shape_path}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_index": shape_idx,
                        "shape_path": current_shape_path,
                        "shape_type": shape_type,
                        "shape_name": shape_name
                    }
                }
                elements.append(shape_element)

                # Create relationship from parent to shape
                contains_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": parent_id,
                    "target_id": shape_id,
                    "relationship_type": RelationshipType.CONTAINS.value,
                    "metadata": {
                        "confidence": 1.0,
                        "index": shape_idx
                    }
                }
                relationships.append(contains_relationship)

                # Create inverse relationship
                contained_by_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": shape_id,
                    "target_id": parent_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(contained_by_relationship)

        return elements

    def _process_text_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                            slide_idx: int, relationships: List[Dict[str, Any]],
                            shape_path: str, element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process a text shape into structured elements and create relationships.

        Args:
            shape: The text shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            relationships: List to add relationships to
            shape_path: Path to the shape
            element_dates: Dictionary to store extracted dates

        Returns:
            List of text-related elements
        """
        elements = []

        try:
            if not shape.has_text_frame or not hasattr(shape, 'text_frame'):
                return elements

            text_frame = shape.text_frame
            text = text_frame.text

            if not text:
                return elements

            # Generate text element ID
            text_id = self._generate_id(f"text_{shape_path}_")

            # Get shape name if available
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Create text element
            text_element = {
                "element_id": text_id,
                "doc_id": doc_id,
                "element_type": ElementType.TEXT_BOX.value,
                "parent_id": parent_id,
                "content_preview": text[:100] + ("..." if len(text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.TEXT_BOX.value,
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(text),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "text": text,
                    "is_title": self._is_title_shape(shape),
                    "level": self._get_paragraph_level(text_frame) if hasattr(text_frame, 'paragraphs') else 0
                }
            }
            elements.append(text_element)

            # Extract dates from text
            self._extract_dates_from_text(text, text_id, element_dates)

            # Create relationship from parent to text box
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": text_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": text_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process paragraphs if detailed paragraph extraction is desired
            if hasattr(text_frame, 'paragraphs') and len(text_frame.paragraphs) > 1:
                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    para_text = paragraph.text
                    if not para_text:
                        continue

                    para_id = self._generate_id(f"para_{shape_path}_{para_idx}_")

                    para_element = {
                        "element_id": para_id,
                        "doc_id": doc_id,
                        "element_type": ElementType.PARAGRAPH.value,
                        "parent_id": text_id,
                        "content_preview": para_text[:100] + ("..." if len(para_text) > 100 else ""),
                        "content_location": json.dumps({
                            "source": source_id,
                            "type": ElementType.PARAGRAPH.value,
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "paragraph_index": para_idx
                        }),
                        "content_hash": self._generate_hash(para_text),
                        "metadata": {
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "paragraph_index": para_idx,
                            "text": para_text,
                            "level": paragraph.level if hasattr(paragraph, 'level') else 0
                        }
                    }
                    elements.append(para_element)

                    # Extract dates from paragraph text
                    self._extract_dates_from_text(para_text, para_id, element_dates)

                    # Create relationship from text box to paragraph
                    contains_para_relationship = {
                        "relationship_id": self._generate_id("rel_"),
                        "source_id": text_id,
                        "target_id": para_id,
                        "relationship_type": RelationshipType.CONTAINS_TEXT.value,
                        "metadata": {
                            "confidence": 1.0,
                            "index": para_idx
                        }
                    }
                    relationships.append(contains_para_relationship)

                    # Create inverse relationship
                    para_contained_relationship = {
                        "relationship_id": self._generate_id("rel_"),
                        "source_id": para_id,
                        "target_id": text_id,
                        "relationship_type": RelationshipType.CONTAINED_BY.value,
                        "metadata": {
                            "confidence": 1.0
                        }
                    }
                    relationships.append(para_contained_relationship)

        except Exception as e:
            logger.warning(f"Error processing text shape: {str(e)}")

        return elements

    def _process_picture_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                               slide_idx: int, relationships: List[Dict[str, Any]],
                               shape_path: str, element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process a picture shape into structured elements and create relationships.

        Args:
            shape: The picture shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            relationships: List to add relationships to
            shape_path: Path to the shape
            element_dates: Dictionary to store extracted dates

        Returns:
            List of picture-related elements
        """
        elements = []

        try:
            # Generate image element ID
            image_id = self._generate_id(f"image_{shape_path}_")

            # Get shape name and image info
            shape_name = shape.name if hasattr(shape, 'name') else ""
            image_name = ""

            if hasattr(shape, 'image') and hasattr(shape.image, 'filename'):
                image_name = shape.image.filename

            # Get alt text for date extraction
            alt_text = shape.alt_text if hasattr(shape, 'alt_text') else ""

            # Create image element
            image_element = {
                "element_id": image_id,
                "doc_id": doc_id,
                "element_type": ElementType.IMAGE.value,
                "parent_id": parent_id,
                "content_preview": f"Image: {image_name}" if image_name else f"Image {shape_path}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.IMAGE.value,
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(f"image_{shape_path}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "image_name": image_name,
                    "alt_text": alt_text
                }
            }
            elements.append(image_element)

            # Extract dates from alt text and image name
            if alt_text:
                self._extract_dates_from_text(alt_text, image_id, element_dates)
            if image_name:
                self._extract_dates_from_text(image_name, image_id, element_dates)

            # Create relationship from parent to image
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": image_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": image_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process image caption (text) if available
            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                caption_elements = self._process_text_shape(shape, doc_id, image_id, source_id, slide_idx,
                                                            relationships, f"{shape_path}_caption", element_dates)
                elements.extend(caption_elements)

        except Exception as e:
            logger.warning(f"Error processing picture shape: {str(e)}")

        return elements

    def _process_table_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                             slide_idx: int, relationships: List[Dict[str, Any]],
                             shape_path: str, element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process a table shape into structured elements and create relationships.

        Args:
            shape: The table shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            relationships: List to add relationships to
            shape_path: Path to the shape
            element_dates: Dictionary to store extracted dates

        Returns:
            List of table-related elements
        """
        elements = list()

        try:
            if not shape.has_table or not hasattr(shape, 'table'):
                return elements

            table = shape.table
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Generate table element ID
            table_id = self._generate_id(f"table_{shape_path}_")

            # Extract basic table content for preview and date extraction
            table_text = ""
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame.text:
                        table_text += cell.text_frame.text + " | "
                table_text += "\n"

            # Create table element
            table_element = {
                "element_id": table_id,
                "doc_id": doc_id,
                "element_type": ElementType.TABLE.value,
                "parent_id": parent_id,
                "content_preview": f"Table: {table_text[:100]}" + ("..." if len(table_text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.TABLE.value,
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(table_text),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "rows": len(table.rows),
                    "columns": len(table.columns)
                }
            }
            elements.append(table_element)

            # Extract dates from table text
            self._extract_dates_from_text(table_text, table_id, element_dates)

            # Create relationship from parent to table
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": table_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": table_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process table rows and cells
            for row_idx, row in enumerate(table.rows):
                # Generate row element ID
                row_id = self._generate_id(f"row_{shape_path}_{row_idx}_")

                # Create row element
                row_element = {
                    "element_id": row_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.TABLE_ROW.value,
                    "parent_id": table_id,
                    "content_preview": f"Row {row_idx + 1}",
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.TABLE_ROW.value,
                        "slide_index": slide_idx,
                        "shape_path": shape_path,
                        "row": row_idx
                    }),
                    "content_hash": self._generate_hash(f"row_{shape_path}_{row_idx}"),
                    "metadata": {
                        "slide_index": slide_idx,
                        "shape_path": shape_path,
                        "row": row_idx
                    }
                }
                elements.append(row_element)

                # Create relationship from table to row
                contains_row_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": table_id,
                    "target_id": row_id,
                    "relationship_type": RelationshipType.CONTAINS_TABLE_ROW.value,
                    "metadata": {
                        "confidence": 1.0,
                        "row_index": row_idx
                    }
                }
                relationships.append(contains_row_relationship)

                # Create inverse relationship
                row_contained_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": row_id,
                    "target_id": table_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(row_contained_relationship)

                # Process cells in this row
                for col_idx, cell in enumerate(row.cells):
                    cell_text = cell.text_frame.text if hasattr(cell, 'text_frame') else ""

                    if not cell_text:
                        continue

                    # Generate cell element ID
                    cell_id = self._generate_id(f"cell_{shape_path}_{row_idx}_{col_idx}_")

                    # Create cell element
                    cell_element = {
                        "element_id": cell_id,
                        "doc_id": doc_id,
                        "element_type": ElementType.TABLE_CELL.value,
                        "parent_id": row_id,
                        "content_preview": cell_text[:100] + ("..." if len(cell_text) > 100 else ""),
                        "content_location": json.dumps({
                            "source": source_id,
                            "type": ElementType.TABLE_CELL.value,
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "row": row_idx,
                            "col": col_idx
                        }),
                        "content_hash": self._generate_hash(cell_text),
                        "metadata": {
                            "slide_index": slide_idx,
                            "shape_path": shape_path,
                            "row": row_idx,
                            "col": col_idx,
                            "text": cell_text
                        }
                    }
                    elements.append(cell_element)

                    # Extract dates from cell text
                    self._extract_dates_from_text(cell_text, cell_id, element_dates)

                    # Create relationship from row to cell
                    contains_cell_relationship = {
                        "relationship_id": self._generate_id("rel_"),
                        "source_id": row_id,
                        "target_id": cell_id,
                        "relationship_type": RelationshipType.CONTAINS_TABLE_CELL.value,
                        "metadata": {
                            "confidence": 1.0,
                            "col_index": col_idx
                        }
                    }
                    relationships.append(contains_cell_relationship)

                    # Create inverse relationship
                    cell_contained_relationship = {
                        "relationship_id": self._generate_id("rel_"),
                        "source_id": cell_id,
                        "target_id": row_id,
                        "relationship_type": RelationshipType.CONTAINED_BY.value,
                        "metadata": {
                            "confidence": 1.0
                        }
                    }
                    relationships.append(cell_contained_relationship)

        except Exception as e:
            logger.warning(f"Error processing table shape: {str(e)}")

        return elements

    def _process_chart_shape(self, shape, doc_id: str, parent_id: str, source_id: str,
                             slide_idx: int, relationships: List[Dict[str, Any]],
                             shape_path: str, element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process a chart shape into structured elements and create relationships.

        Args:
            shape: The chart shape
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            slide_idx: Slide index
            relationships: List to add relationships to
            shape_path: Path to the shape
            element_dates: Dictionary to store extracted dates

        Returns:
            List of chart-related elements
        """
        elements = []

        try:
            if not shape.has_chart or not hasattr(shape, 'chart'):
                return elements

            chart = shape.chart
            shape_name = shape.name if hasattr(shape, 'name') else ""

            # Generate chart element ID
            chart_id = self._generate_id(f"chart_{shape_path}_")

            # Get chart type
            chart_type = "unknown"
            if hasattr(chart, 'chart_type'):
                chart_type = str(chart.chart_type)

            # Get chart title
            chart_title = ""
            if hasattr(chart, 'chart_title') and hasattr(chart.chart_title, 'text_frame'):
                chart_title = chart.chart_title.text_frame.text

            # Extract chart data for date extraction
            chart_text = chart_title
            categories = []
            series_names = []

            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]

                if hasattr(plot, 'categories'):
                    for category in plot.categories:
                        if category:
                            categories.append(str(category))
                            chart_text += " " + str(category)

                if hasattr(plot, 'series'):
                    for series in plot.series:
                        if hasattr(series, 'name') and series.name:
                            series_names.append(str(series.name))
                            chart_text += " " + str(series.name)

            # Create chart element
            chart_element = {
                "element_id": chart_id,
                "doc_id": doc_id,
                "element_type": ElementType.CHART.value,
                "parent_id": parent_id,
                "content_preview": f"Chart: {chart_title}" if chart_title else f"Chart {shape_path}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.CHART.value,
                    "slide_index": slide_idx,
                    "shape_path": shape_path
                }),
                "content_hash": self._generate_hash(f"chart_{shape_path}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "shape_path": shape_path,
                    "shape_name": shape_name,
                    "chart_type": chart_type,
                    "chart_title": chart_title
                }
            }

            if categories:
                chart_element["metadata"]["categories"] = categories

            if series_names:
                chart_element["metadata"]["series"] = series_names

            elements.append(chart_element)

            # Extract dates from chart text
            self._extract_dates_from_text(chart_text, chart_id, element_dates)

            # Create relationship from parent to chart
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": chart_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": chart_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

        except Exception as e:
            logger.warning(f"Error processing chart shape: {str(e)}")

        return elements

    def _process_notes(self, notes_slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str,
                       relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process slide notes into structured elements and create relationships.

        Args:
            notes_slide: The notes slide
            slide_idx: Slide index
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of notes-related elements
        """
        elements = []

        try:
            if not notes_slide or not hasattr(notes_slide, 'notes_text_frame'):
                return elements

            notes_text = notes_slide.notes_text_frame.text

            if not notes_text:
                return elements

            # Generate notes element ID
            notes_id = self._generate_id(f"notes_{slide_idx}_")

            # Create notes element
            notes_element = {
                "element_id": notes_id,
                "doc_id": doc_id,
                "element_type": ElementType.SLIDE_NOTES.value,
                "parent_id": parent_id,
                "content_preview": notes_text[:100] + ("..." if len(notes_text) > 100 else ""),
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.SLIDE_NOTES.value,
                    "slide_index": slide_idx
                }),
                "content_hash": self._generate_hash(notes_text),
                "metadata": {
                    "slide_index": slide_idx,
                    "text": notes_text
                }
            }
            elements.append(notes_element)

            # Extract dates from notes text
            self._extract_dates_from_text(notes_text, notes_id, element_dates)

            # Create relationship from parent to notes
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": notes_id,
                "relationship_type": RelationshipType.CONTAINS_NOTES.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": notes_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

        except Exception as e:
            logger.warning(f"Error processing slide notes: {str(e)}")

        return elements

    def _process_comments(self, slide, slide_idx: int, doc_id: str, parent_id: str, source_id: str,
                          relationships: List[Dict[str, Any]], element_dates: Dict[str, List[Dict[str, Any]]]) -> List[Dict[str, Any]]:
        """
        Process slide comments into structured elements and create relationships.

        Args:
            slide: The slide
            slide_idx: Slide index
            doc_id: Document ID
            parent_id: Parent element ID
            source_id: Source identifier
            relationships: List to add relationships to
            element_dates: Dictionary to store extracted dates

        Returns:
            List of comment-related elements
        """
        elements = []

        try:
            # Check if slide has comments
            if not hasattr(slide, 'comments') or not slide.comments:
                return elements

            # Create comments container element
            comments_id = self._generate_id(f"comments_{slide_idx}_")

            comments_element = {
                "element_id": comments_id,
                "doc_id": doc_id,
                "element_type": ElementType.COMMENTS_CONTAINER.value,
                "parent_id": parent_id,
                "content_preview": f"Comments for Slide {slide_idx + 1}",
                "content_location": json.dumps({
                    "source": source_id,
                    "type": ElementType.COMMENTS_CONTAINER.value,
                    "slide_index": slide_idx
                }),
                "content_hash": self._generate_hash(f"comments_{slide_idx}"),
                "metadata": {
                    "slide_index": slide_idx,
                    "comment_count": len(slide.comments)
                }
            }
            elements.append(comments_element)

            # Create relationship from parent to comments container
            contains_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": parent_id,
                "target_id": comments_id,
                "relationship_type": RelationshipType.CONTAINS.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contains_relationship)

            # Create inverse relationship
            contained_by_relationship = {
                "relationship_id": self._generate_id("rel_"),
                "source_id": comments_id,
                "target_id": parent_id,
                "relationship_type": RelationshipType.CONTAINED_BY.value,
                "metadata": {
                    "confidence": 1.0
                }
            }
            relationships.append(contained_by_relationship)

            # Process individual comments
            for comment_idx, comment in enumerate(slide.comments):
                comment_text = comment.text if hasattr(comment, 'text') else ""

                if not comment_text:
                    continue

                # Generate comment element ID
                comment_id = self._generate_id(f"comment_{slide_idx}_{comment_idx}_")

                # Get comment author and date if available
                author = comment.author if hasattr(comment, 'author') else "Unknown"
                date = comment.date if hasattr(comment, 'date') else None

                # Create comment element
                comment_element = {
                    "element_id": comment_id,
                    "doc_id": doc_id,
                    "element_type": ElementType.COMMENT.value,
                    "parent_id": comments_id,
                    "content_preview": comment_text[:100] + ("..." if len(comment_text) > 100 else ""),
                    "content_location": json.dumps({
                        "source": source_id,
                        "type": ElementType.COMMENT.value,
                        "slide_index": slide_idx,
                        "comment_index": comment_idx
                    }),
                    "content_hash": self._generate_hash(comment_text),
                    "metadata": {
                        "slide_index": slide_idx,
                        "comment_index": comment_idx,
                        "author": author,
                        "date": date.timestamp() if date and hasattr(date, 'timestamp') else None,
                        "text": comment_text
                    }
                }
                elements.append(comment_element)

                # Extract dates from comment text
                self._extract_dates_from_text(comment_text, comment_id, element_dates)

                # Create relationship from comments container to comment
                contains_comment_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": comments_id,
                    "target_id": comment_id,
                    "relationship_type": RelationshipType.CONTAINS.value,
                    "metadata": {
                        "confidence": 1.0,
                        "index": comment_idx
                    }
                }
                relationships.append(contains_comment_relationship)

                # Create inverse relationship
                comment_contained_relationship = {
                    "relationship_id": self._generate_id("rel_"),
                    "source_id": comment_id,
                    "target_id": comments_id,
                    "relationship_type": RelationshipType.CONTAINED_BY.value,
                    "metadata": {
                        "confidence": 1.0
                    }
                }
                relationships.append(comment_contained_relationship)

        except Exception as e:
            logger.warning(f"Error processing slide comments: {str(e)}")

        return elements

    """
    Add the missing methods to the PowerPoint parser.
    """

    # Add these methods to the PptxParser class

    @staticmethod
    def _get_slide_title(slide):
        """
        Extract the title from a slide.

        Args:
            slide: The slide object

        Returns:
            Title text or empty string if no title found
        """
        try:
            # Look for title placeholder
            for shape in slide.shapes:
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title placeholder
                        if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                            return shape.text_frame.text

                # Also check if any shape with text has a name indicating it's a title
                if hasattr(shape, 'name') and hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                    name = shape.name.lower()
                    if 'title' in name:
                        return shape.text_frame.text

                # If there's a text shape with large text at the top of the slide, it's likely a title
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text:
                    # Check if this is likely a title by position (top of slide) and text length
                    if shape.top < (slide.slide_layout.slide_height * 0.25) and len(shape.text_frame.text) < 100:
                        return shape.text_frame.text

            # Fallback: Look for the first text shape with text
            for shape in slide.shapes:
                if hasattr(shape, 'has_text_frame') and shape.has_text_frame and shape.text_frame.text:
                    text = shape.text_frame.text.strip()
                    # If it's reasonable length for a title, use it
                    if len(text) < 100:
                        return text
                    else:
                        # If it's long, try to use just the first line
                        lines = text.split('\n')
                        if lines and len(lines[0]) < 100:
                            return lines[0]

            return ""
        except Exception as e:
            logger.warning(f"Error extracting slide title: {str(e)}")
            return ""

    @staticmethod
    def _get_slide_layout_name(slide):
        """
        Get the name of the slide layout.

        Args:
            slide: The slide object

        Returns:
            Layout name or empty string if not available
        """
        try:
            if hasattr(slide, 'slide_layout') and hasattr(slide.slide_layout, 'name'):
                return slide.slide_layout.name
            return ""
        except Exception as e:
            logger.warning(f"Error getting slide layout name: {str(e)}")
            return ""

    @staticmethod
    def _is_title_shape(shape):
        """
        Determine if a shape is likely a title shape.

        Args:
            shape: The shape object

        Returns:
            True if shape is likely a title, False otherwise
        """
        try:
            # Check if it's explicitly a title placeholder
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title placeholder
                    return True

            # Check the shape name
            if hasattr(shape, 'name'):
                name = shape.name.lower()
                if 'title' in name:
                    return True

            return False
        except Exception as e:
            logger.warning(f"Error checking if shape is title: {str(e)}")
            return False

    @staticmethod
    def _get_paragraph_level(text_frame):
        """
        Get the indentation level of paragraphs in a text frame.

        Args:
            text_frame: The text frame

        Returns:
            Indentation level (0 for top level)
        """
        try:
            if not hasattr(text_frame, 'paragraphs') or not text_frame.paragraphs:
                return 0

            # Get the level of the first paragraph
            if hasattr(text_frame.paragraphs[0], 'level'):
                return text_frame.paragraphs[0].level

            return 0
        except Exception as e:
            logger.warning(f"Error getting paragraph level: {str(e)}")
            return 0

    @staticmethod
    def _get_shape_type(shape):
        """
        Determine the type of a shape.

        Args:
            shape: The shape object

        Returns:
            String describing the shape type
        """
        try:
            # Try to get shape type from various attributes
            if hasattr(shape, 'shape_type'):
                return str(shape.shape_type)

            # Check for common shape characteristics
            if isinstance(shape, GroupShape):
                return "group"
            elif hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                return "text"
            elif hasattr(shape, 'has_table') and shape.has_table:
                return "table"
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                return "chart"
            elif isinstance(shape, Picture):
                return "picture"

            # Generic fallback
            return "shape"
        except Exception as e:
            logger.warning(f"Error determining shape type: {str(e)}")
            return "shape"

    def _create_root_element(self, doc_id, source_id):
        """
        Create the root element for the document.

        Args:
            doc_id: Document ID
            source_id: Source identifier

        Returns:
            Root element dictionary
        """
        return {
            "element_id": self._generate_id("root_"),
            "doc_id": doc_id,
            "element_type": "root",
            "parent_id": None,
            "content_preview": "Document Root",
            "content_location": json.dumps({
                "source": source_id,
                "type": "root"
            }),
            "content_hash": "",
            "metadata": {}
        }

    def _generate_id(self, prefix="id_"):
        """
        Generate a unique ID.

        Args:
            prefix: Prefix for the ID

        Returns:
            Generated ID string
        """
        import uuid
        return f"{prefix}{uuid.uuid4().hex}"

    def _generate_hash(self, content):
        """
        Generate a hash for content.

        Args:
            content: Content to hash

        Returns:
            Hash string
        """
        import hashlib

        if isinstance(content, str):
            content_bytes = content.encode('utf-8')
        elif isinstance(content, bytes):
            content_bytes = content
        else:
            content_bytes = str(content).encode('utf-8')

        return hashlib.sha256(content_bytes).hexdigest()
